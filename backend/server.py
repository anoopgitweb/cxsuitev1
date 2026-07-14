from __future__ import annotations

import base64
import cgi
import datetime as _dt
import hashlib
import hmac
import secrets
import json
import mimetypes
import os
import pickle
import socket
import subprocess
import sys
import tempfile
import time
import traceback
import uuid
import webbrowser
from dataclasses import dataclass, field
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib.parse import parse_qs, unquote, urlparse
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen
import threading
import math
import random
import re
import shutil
import importlib.util

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

def _app_root() -> Path:
    if getattr(sys, "frozen", False):
        exe_dir = Path(sys.executable).resolve().parent
        for candidate in (exe_dir, exe_dir.parent, exe_dir.parent.parent):
            if (candidate / "frontend").exists() and (candidate / "models").exists():
                return candidate
        return exe_dir
    return Path(__file__).resolve().parents[1]


ROOT = _app_root()
FRONTEND = ROOT / "frontend"
APPS = ROOT / "apps"
MODELS = ROOT / "models"
PYTHON_PACKAGES = ROOT / "python_packages"
OUTPUTS = ROOT / "outputs"
SECURITY_DIR = ROOT / "security"
USERS_FILE = SECURITY_DIR / "users.json"
AUDIT_DIR = ROOT / "logs" / "audit"
USER_LOG_DIR = ROOT / "logs" / "user_logs"
ANALYSIS_LOG_DIR = ROOT / "logs" / "analysis_runs"
SESSION_LOG_DIR = ROOT / "logs" / "server_sessions"
SESSION_LOCK = threading.Lock()
SESSIONS: dict[str, dict[str, Any]] = {}
SERVER_SESSION_ID = uuid.uuid4().hex
SERVER_SESSION_STARTED = time.time()
SERVER_SESSION_LOG: Path | None = None
SERVER_SESSION_LOG_LOCK = threading.Lock()
PASSWORD_ITERATIONS = 220000
THEME_TRAINING_MODULE: Any | None = None
sys.path.insert(0, str(Path(__file__).resolve().parent))
sys.path.insert(0, str(APPS / "nps-analyzer" / "backend"))
if PYTHON_PACKAGES.exists():
    sys.path.insert(0, str(PYTHON_PACKAGES))

from analysis_engine import (  # noqa: E402
    add_theme_acpt_resolution_outputs,
    add_reporting_calendar_columns,
    build_analysis,
    build_analysis_with_local_model,
    build_summaries,
    executive_snapshot_insights,
    export_workbook,
    nps_composition_counts,
    nps_summary,
    normalize_fiscal_start_month,
    resolve_roberta_model_path,
    resolve_theme_acpt_resolution_model_path,
    sentiment_summary,
    week_period_start,
)


@dataclass
class AppState:
    base_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    lookup_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    analyzed_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    weekly_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    agent_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    manager_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    reason_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    complaints_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    passive_df: pd.DataFrame = field(default_factory=pd.DataFrame)
    base_column_profile: dict[str, dict[str, Any]] = field(default_factory=dict)
    lookup_column_profile: dict[str, dict[str, Any]] = field(default_factory=dict)
    analyzed_column_profile: dict[str, dict[str, Any]] = field(default_factory=dict)
    files: dict[str, str] = field(default_factory=dict)
    file_sizes: dict[str, int] = field(default_factory=dict)
    status: str = "Ready"
    progress: float = 0.0
    analysis_running: bool = False
    analysis_error: str = ""
    analysis_id: str = ""
    analysis_log_file: str = ""
    analysis_stage: str = ""
    analysis_rows_processed: int = 0
    analysis_total_rows: int = 0
    date_filter: dict[str, str] = field(default_factory=lambda: {"mode": "All Time", "start": "", "end": ""})
    dynamic_dimensions: list[str] = field(default_factory=list)
    analysis_engines: dict[str, str] = field(default_factory=lambda: {"sentiment": "local", "theme": "local"})
    model_paths: dict[str, str] = field(default_factory=dict)
    calendar_settings: dict[str, Any] = field(default_factory=lambda: {"weekStart": "Sun", "fiscalYearStartMonth": 1})
    last_run_config: dict[str, Any] = field(default_factory=dict)
    analysis_started_at: float = 0.0
    analysis_completed_at: float = 0.0


STATE = AppState()
STATE_LOCK = threading.Lock()
UPLOAD_PROGRESS: dict[str, dict[str, Any]] = {}
UPLOAD_PROGRESS_LOCK = threading.Lock()
MODULE_PROGRESS: dict[str, dict[str, Any]] = {}
MODULE_PROGRESS_LOCK = threading.Lock()
AUDIT_PROCESS: subprocess.Popen | None = None
KESTRELIQ_PROCESS: subprocess.Popen | None = None
SERVER_LAUNCHER_PROCESS_ID: int | None = None
SPARROW_TRAINING_LOCK = threading.Lock()
SPARROW_TRAINING_JOB: dict[str, Any] = {
    "running": False,
    "status": "Idle",
    "progress": 0,
    "logs": [],
    "telemetry": [],
    "result": None,
    "error": "",
}
ANALYSIS_RUN_LOGS: dict[str, dict[str, Any]] = {}
ANALYSIS_RUN_LOG_LOCK = threading.Lock()


def _calendar_settings(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    raw = payload.get("calendar") if isinstance(payload, dict) and isinstance(payload.get("calendar"), dict) else {}
    week_start = str(raw.get("weekStart") or raw.get("week_start") or "Sun").strip().title()[:3] or "Sun"
    if week_start not in {"Sun", "Mon", "Tue", "Wed", "Thu", "Fri", "Sat"}:
        week_start = "Sun"
    fiscal_start = normalize_fiscal_start_month(
        raw.get("fiscalYearStartMonth") or raw.get("fiscal_year_start_month") or raw.get("fiscalStartMonth") or 1
    )
    return {"weekStart": week_start, "fiscalYearStartMonth": fiscal_start}


def _state_calendar_settings() -> dict[str, Any]:
    with STATE_LOCK:
        return _calendar_settings({"calendar": dict(STATE.calendar_settings)})


def _summaries_for_calendar(df: pd.DataFrame, calendar: dict[str, Any] | None = None) -> dict[str, pd.DataFrame]:
    settings = _calendar_settings({"calendar": calendar or {}})
    return build_summaries(
        df,
        week_start=settings["weekStart"],
        fiscal_year_start_month=settings["fiscalYearStartMonth"],
    )


def _apply_reporting_calendar(df: pd.DataFrame, calendar: dict[str, Any] | None = None) -> pd.DataFrame:
    settings = _calendar_settings({"calendar": calendar or {}})
    return add_reporting_calendar_columns(
        df,
        week_start=settings["weekStart"],
        fiscal_year_start_month=settings["fiscalYearStartMonth"],
    )



def _audit_timestamp() -> str:
    return time.strftime("%Y-%m-%d %H:%M:%S")


def _start_server_session_log() -> Path:
    global SERVER_SESSION_LOG
    SESSION_LOG_DIR.mkdir(parents=True, exist_ok=True)
    stamp = time.strftime("%Y-%m-%d_%H-%M-%S")
    path = SESSION_LOG_DIR / f"CXSuite_Server_Session_{stamp}_{SERVER_SESSION_ID[:8]}.log"
    header = {
        "timestamp": _audit_timestamp(),
        "level": "INFO",
        "sessionId": SERVER_SESSION_ID,
        "event": "SESSION_LOG_CREATED",
        "action": "Continuous server session logging started.",
        "details": {"root": str(ROOT), "logFile": str(path)},
    }
    path.write_text(json.dumps(header, ensure_ascii=False) + "\n", encoding="utf-8")
    SERVER_SESSION_LOG = path
    return path


def _write_server_session_record(record: dict[str, Any]) -> None:
    path = SERVER_SESSION_LOG
    if path is None:
        return
    entry = dict(record)
    entry.setdefault("sessionId", SERVER_SESSION_ID)
    try:
        with SERVER_SESSION_LOG_LOCK:
            with path.open("a", encoding="utf-8", buffering=1) as handle:
                handle.write(json.dumps(entry, ensure_ascii=False) + "\n")
                handle.flush()
    except Exception:
        traceback.print_exc()


def _analysis_log_entry(status: str, component: str, message: str) -> str:
    return f"[{_audit_timestamp()}] [{status}] [{component}] {message}"


def _analysis_log_write(analysis_id: str, status: str, component: str, message: str) -> None:
    with ANALYSIS_RUN_LOG_LOCK:
        context = ANALYSIS_RUN_LOGS.get(analysis_id)
        if not context:
            return
        path = Path(context["path"])
        try:
            with path.open("a", encoding="utf-8", buffering=1) as handle:
                handle.write(_analysis_log_entry(status, component, message) + "\n")
                handle.flush()
                os.fsync(handle.fileno())
        except Exception:
            traceback.print_exc()


def _analysis_log_section(analysis_id: str, name: str) -> None:
    with ANALYSIS_RUN_LOG_LOCK:
        context = ANALYSIS_RUN_LOGS.get(analysis_id)
        if not context or name in context["sections"]:
            return
        context["sections"].add(name)
        path = Path(context["path"])
        with path.open("a", encoding="utf-8", buffering=1) as handle:
            handle.write(f"\n---\n\n## {name}\n\n")
            handle.flush()
            os.fsync(handle.fileno())


def _analysis_log_component(message: str) -> tuple[str, str]:
    text = str(message or "").lower()
    if "acpt" in text or "agent, customer, process, or technology" in text:
        return "ACPT CLASSIFICATION", "ACPT"
    if "owl" in text or "theme" in text:
        return "THEME ANALYSIS", "THEMES"
    if "sparrow" in text or "sentiment" in text or "local rules" in text:
        return "SPARROW SENTIMENT ENGINE", "SPARROW"
    if "summar" in text or "weekly" in text or "insight" in text:
        return "INSIGHT GENERATION", "INSIGHTS"
    if "dashboard" in text or "profil" in text or "publish" in text or "export" in text:
        return "OUTPUT GENERATION", "OUTPUT"
    return "DATA VALIDATION", "VALIDATION"


def _parse_progress_rows(message: str, default_total: int = 0) -> tuple[int, int]:
    text = str(message or "")
    patterns = (
        r"(?P<done>[\d,]+)\s*/\s*(?P<total>[\d,]+)\s*rows?",
        r"rows?\s*(?P<done>[\d,]+)\s*/\s*(?P<total>[\d,]+)",
        r"row\s*(?P<done>[\d,]+)\s+of\s+(?P<total>[\d,]+)",
        r"(?P<done>[\d,]+)\s+of\s+(?P<total>[\d,]+)\s*rows?",
    )
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if match:
            return int(match.group("done").replace(",", "")), int(match.group("total").replace(",", ""))
    return 0, max(0, int(default_total or 0))


def _analysis_log_progress(analysis_id: str, progress: float, message: str) -> None:
    if not analysis_id or analysis_id not in ANALYSIS_RUN_LOGS:
        return
    section, component = _analysis_log_component(message)
    context = ANALYSIS_RUN_LOGS.get(analysis_id)
    if context:
        now = time.time()
        context.setdefault("stage_first", {}).setdefault(section, now)
        context.setdefault("stage_last", {})[section] = now
        rows_done, rows_total = _parse_progress_rows(message, int(context.get("rows", 0)))
        previous_stage = str(context.get("last_stage", ""))
        previous_rows = int(context.get("last_rows_logged", 0) or 0)
        context["last_stage"] = section
        context["last_component"] = component
        context["last_progress"] = float(progress)
        context["last_message"] = str(message)
        if rows_done:
            context["rows_done"] = rows_done
            context["rows_total"] = rows_total
        should_checkpoint = (
            section != previous_stage
            or progress >= 100
            or (rows_done and (rows_done - previous_rows >= 500 or rows_done >= rows_total))
        )
        if should_checkpoint and rows_done:
            context["last_rows_logged"] = rows_done
    _analysis_log_section(analysis_id, section)
    _analysis_log_write(analysis_id, "RUNNING" if progress < 100 else "SUCCESS", component, f"{progress:.1f}% - {message}")
    if context and should_checkpoint:
        rows_text = f"; rows={int(context.get('rows_done', 0)):,}/{int(context.get('rows_total', context.get('rows', 0))):,}"
        _analysis_log_write(analysis_id, "CHECKPOINT", component, f"Stage={section}; progress={progress:.1f}%{rows_text}; status={message}")


def _start_analysis_log(payload: dict[str, Any], analysis_id: str) -> Path:
    ANALYSIS_LOG_DIR.mkdir(parents=True, exist_ok=True)
    file_name = str(STATE.files.get("base") or "InputFile")
    safe_stem = "".join(ch if ch.isalnum() or ch in {"-", "_"} else "_" for ch in Path(file_name).stem).strip("_") or "InputFile"
    stamp = time.strftime("%Y-%m-%d_%H-%M-%S")
    path = ANALYSIS_LOG_DIR / f"CXSuite_Analysis_Log_{safe_stem}_{stamp}.txt"
    mapping = payload.get("mapping") if isinstance(payload.get("mapping"), dict) else {}
    engines = payload.get("engines") if isinstance(payload.get("engines"), dict) else {}
    model_paths = payload.get("modelPaths") if isinstance(payload.get("modelPaths"), dict) else {}
    base = STATE.base_df
    feedback = str(mapping.get("feedback") or "")
    empty_comments = int(base[feedback].fillna("").astype(str).str.strip().eq("").sum()) if feedback in base.columns else len(base)
    missing_values = int(base.isna().sum().sum()) if not base.empty else 0
    total_cells = max(int(base.shape[0] * base.shape[1]), 1)
    health_score = max(0.0, 100.0 - (missing_values / total_cells * 100.0))
    context = {
        "path": str(path), "start": time.time(), "sections": set(), "payload": payload,
        "file_name": file_name, "rows": len(base), "columns": len(base.columns),
        "empty_comments": empty_comments, "duplicates": int(base.duplicated().sum()),
        "missing_values": missing_values, "health_score": health_score,
        "stage_started": time.time(), "stage_times": {},
        "last_stage": "SESSION INFORMATION", "last_component": "SYSTEM",
        "last_progress": 0.0, "last_message": "Analysis log created",
        "rows_done": 0, "rows_total": len(base), "last_rows_logged": 0,
    }
    with ANALYSIS_RUN_LOG_LOCK:
        ANALYSIS_RUN_LOGS[analysis_id] = context
    header = (
        "==================================================================\n"
        "CXSuite_Analysis_Log ANALYSIS EXECUTION LOG\n"
        "==================================================================\n\n"
        "## SESSION INFORMATION\n\n"
    )
    path.write_text(header, encoding="utf-8")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", f"Analysis ID: {analysis_id}")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", "Application Version: Feedback Intelligence Suite 19.06")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", f"Analysis Start Time: {_audit_timestamp()}")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", "Analysis End Time: Pending")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", "Total Duration: Running")
    _analysis_log_write(analysis_id, "RUNNING", "SYSTEM", "Analysis Status: Started")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", f"User Selected Analysis Options: {json.dumps({'mode': payload.get('mode'), 'engines': engines, 'mapping': mapping}, ensure_ascii=True)}")
    _analysis_log_section(analysis_id, "INPUT DATASET")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Input File Name: {file_name}")
    _analysis_log_write(analysis_id, "INFO", "FILE", "File Location: Browser upload; the browser does not expose the source filesystem path")
    input_size = int(STATE.file_sizes.get("base", 0))
    _analysis_log_write(analysis_id, "INFO", "FILE", f"File Size: {input_size:,} bytes ({input_size / 1024:.1f} KB)" if input_size else "File Size: Not available")
    selected_sheet = str(STATE.files.get("base_sheet") or "First worksheet loaded by the local analyzer")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Sheet Name(s): {selected_sheet}")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Total Rows: {len(base):,}")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Total Columns: {len(base.columns):,}")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Selected Comment Column: {feedback or 'Not selected'}")
    _analysis_log_write(analysis_id, "INFO", "FILE", "Encoding: XLSX binary workbook / UTF-8 internal text")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Empty Comments: {empty_comments:,}")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Duplicate Records: {context['duplicates']:,}")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Missing Values: {missing_values:,}")
    _analysis_log_write(analysis_id, "INFO", "FILE", f"Dataset Health Score: {health_score:.1f}%")
    _analysis_log_section(analysis_id, "DATA VALIDATION")
    _analysis_log_write(analysis_id, "RUNNING", "VALIDATION", "Validation Started")
    _analysis_log_write(analysis_id, "INFO", "VALIDATION", f"Required Columns Found: feedback={bool(feedback)}, score={bool(mapping.get('score'))}, satisfaction={bool(mapping.get('satisfaction'))}")
    _analysis_log_write(analysis_id, "INFO", "SPARROW", f"Selected Sentiment Engine: {engines.get('sentiment', 'local')}; Model Location: {model_paths.get('sparrow') or _default_model_path('sparrow')}")
    _analysis_log_write(analysis_id, "INFO", "THEMES", f"Selected Theme Engine: {engines.get('theme', 'local')}; Model Location: {model_paths.get('theme') or _default_model_path('theme')}")
    return path


def _analysis_log_failure(analysis_id: str, error: BaseException) -> dict[str, Any]:
    context = ANALYSIS_RUN_LOGS.get(analysis_id) or {}
    details = {
        "stage": str(context.get("last_stage", "Unknown")),
        "component": str(context.get("last_component", "SYSTEM")),
        "progress": float(context.get("last_progress", 0.0) or 0.0),
        "status": str(context.get("last_message", "No status recorded")),
        "rowsProcessed": int(context.get("rows_done", 0) or 0),
        "totalRows": int(context.get("rows_total", context.get("rows", 0)) or 0),
        "logFile": str(context.get("path", "")),
    }
    _analysis_log_section(analysis_id, "FAILURE DETAIL")
    _analysis_log_write(analysis_id, "ERROR", details["component"], f"Failure Stage: {details['stage']}")
    _analysis_log_write(analysis_id, "ERROR", details["component"], f"Last Progress: {details['progress']:.1f}%")
    _analysis_log_write(analysis_id, "ERROR", details["component"], f"Last Row Checkpoint: {details['rowsProcessed']:,}/{details['totalRows']:,}")
    _analysis_log_write(analysis_id, "ERROR", details["component"], f"Last Status: {details['status']}")
    _analysis_log_write(analysis_id, "ERROR", details["component"], f"Exception: {type(error).__name__}: {error}")
    _analysis_log_write(analysis_id, "ERROR", details["component"], "Traceback:\n" + "".join(traceback.format_exception(type(error), error, error.__traceback__)).strip())
    return details


def _finalize_analysis_log(analysis_id: str, analyzed: pd.DataFrame | None, error: BaseException | None = None) -> None:
    context = ANALYSIS_RUN_LOGS.get(analysis_id)
    if not context:
        return
    elapsed = max(0.0, time.time() - float(context["start"]))
    result = analyzed if isinstance(analyzed, pd.DataFrame) else pd.DataFrame()
    total = len(result)
    sentiment_col = next((name for name in ("Sentiment", "Sparrow Sentiment") if name in result.columns), "")
    sentiment_counts = result[sentiment_col].astype(str).str.lower().value_counts().to_dict() if sentiment_col else {}
    confidence_col = next((name for name in result.columns if "confidence" in str(name).lower() and "theme" not in str(name).lower()), "")
    confidence = pd.to_numeric(result[confidence_col], errors="coerce").dropna() if confidence_col else pd.Series(dtype=float)
    theme_col = next((name for name in ("Owl Primary Driver", "Primary Reason", "Bucket Category") if name in result.columns), "")
    theme_count = int(result[theme_col].dropna().astype(str).str.strip().replace("", pd.NA).dropna().nunique()) if theme_col else 0
    engines = context["payload"].get("engines") if isinstance(context["payload"].get("engines"), dict) else {}
    _analysis_log_section(analysis_id, "INSIGHT GENERATION")
    _analysis_log_write(analysis_id, "SUCCESS" if error is None else "ERROR", "INSIGHTS", "Executive Summary Completed" if error is None else "Insight generation did not complete")
    _analysis_log_write(analysis_id, "INFO", "INSIGHTS", "Leadership, client, and agent insight structures prepared from the analyzed dataset")
    _analysis_log_section(analysis_id, "OUTPUT GENERATION")
    _analysis_log_write(analysis_id, "INFO", "OUTPUT", "Dashboard response and export-ready tables prepared")
    _analysis_log_write(analysis_id, "INFO", "OUTPUT", "Excel output is generated only when the user selects Download Excel")
    _analysis_log_write(analysis_id, "INFO", "OUTPUT", f"Output File Location: {OUTPUTS}")
    _analysis_log_section(analysis_id, "PERFORMANCE")
    _analysis_log_write(analysis_id, "INFO", "PERFORMANCE", f"Total Processing Time: {elapsed:.2f} seconds")
    _analysis_log_write(analysis_id, "INFO", "PERFORMANCE", f"Records Per Second: {(total / elapsed if elapsed else 0):.2f}")
    _analysis_log_write(analysis_id, "INFO", "PERFORMANCE", "Peak/Average Memory, CPU, and GPU utilization: Not captured by the current portable runtime")
    for stage, started in context.get("stage_first", {}).items():
        finished = context.get("stage_last", {}).get(stage, started)
        _analysis_log_write(analysis_id, "INFO", "PERFORMANCE", f"Time Taken - {stage}: {max(0.0, finished - started):.2f} seconds of recorded stage activity")
    _analysis_log_section(analysis_id, "WARNINGS")
    _analysis_log_write(analysis_id, "WARNING" if context["empty_comments"] else "INFO", "WARNINGS", f"Missing Comments: {context['empty_comments']:,}")
    _analysis_log_write(analysis_id, "WARNING" if context["duplicates"] else "INFO", "WARNINGS", f"Duplicate Records: {context['duplicates']:,}")
    _analysis_log_write(analysis_id, "INFO", "WARNINGS", "Engine warnings and recovery actions are recorded in the live stage entries above")
    _analysis_log_section(analysis_id, "ERRORS")
    if error is None:
        _analysis_log_write(analysis_id, "SUCCESS", "ERRORS", "No unhandled analysis errors")
    else:
        _analysis_log_write(analysis_id, "ERROR", "ERRORS", f"Exception Message: {error}")
        _analysis_log_write(analysis_id, "ERROR", "ERRORS", f"Full Stack Trace: {''.join(traceback.format_exception(type(error), error, error.__traceback__)).strip()}")
        _analysis_log_write(analysis_id, "ERROR", "ERRORS", "Final Outcome: Analysis failed; partial execution history has been preserved")
    _analysis_log_section(analysis_id, "FINAL SUMMARY")
    _analysis_log_write(analysis_id, "SUCCESS" if error is None else "ERROR", "SYSTEM", f"Analysis Completed Successfully: {'Yes' if error is None else 'No'}")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", f"Analysis Status: {'Completed' if error is None else 'Failed'}")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", f"Analysis End Time: {_audit_timestamp()}")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", f"Total Records Analysed: {total:,}")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", f"Total Records Skipped: {max(0, int(context['rows']) - total):,}")
    _analysis_log_write(analysis_id, "INFO", "SPARROW", f"Sparrow Engine Used: {'Yes' if str(engines.get('sentiment', 'local')).lower() == 'sparrow' else 'No - Local Rules'}")
    _analysis_log_write(analysis_id, "INFO", "THEMES", f"Theme Engine Used: {engines.get('theme', 'local')}")
    _analysis_log_write(analysis_id, "INFO", "SPARROW", f"Records Successfully Processed: {total:,}; Failed Records: {max(0, int(context['rows']) - total):,}")
    _analysis_log_write(analysis_id, "INFO", "SPARROW", f"Positive Count: {int(sentiment_counts.get('positive', 0)):,}; Neutral Count: {int(sentiment_counts.get('neutral', 0)):,}; Negative Count: {int(sentiment_counts.get('negative', 0)):,}")
    _analysis_log_write(analysis_id, "INFO", "SPARROW", f"Overall Average Confidence: {confidence.mean():.4f}" if not confidence.empty else "Overall Average Confidence: Not available")
    _analysis_log_write(analysis_id, "INFO", "SPARROW", f"Highest Confidence: {confidence.max():.4f}; Lowest Confidence: {confidence.min():.4f}; Low Confidence Records (<0.60): {int((confidence < 0.60).sum()):,}" if not confidence.empty else "Highest/Lowest Confidence and Low Confidence Records: Not available")
    _analysis_log_write(analysis_id, "INFO", "THEMES", f"Themes Generated: {theme_count:,} unique values in {theme_col or 'no theme column'}")
    _analysis_log_write(analysis_id, "INFO", "SYSTEM", f"Total Processing Time: {elapsed:.2f} seconds")
    _analysis_log_write(analysis_id, "SUCCESS", "SYSTEM", f"Log File Saved Successfully: {context['path']}")
    with Path(context["path"]).open("a", encoding="utf-8", buffering=1) as handle:
        handle.write("\n==================================================================\nEND OF ANALYSIS\n==================================================================\n")
        handle.flush()
        os.fsync(handle.fileno())
    with ANALYSIS_RUN_LOG_LOCK:
        ANALYSIS_RUN_LOGS.pop(analysis_id, None)


def _hash_password(password: str, salt_hex: str, iterations: int = PASSWORD_ITERATIONS) -> str:
    salt = bytes.fromhex(salt_hex)
    digest = hashlib.pbkdf2_hmac("sha256", str(password or "").encode("utf-8"), salt, int(iterations))
    return digest.hex()


def _write_audit_record(username: str, event: str, action: str, details: dict[str, Any] | None = None, level: str = "INFO") -> None:
    try:
        AUDIT_DIR.mkdir(parents=True, exist_ok=True)
        USER_LOG_DIR.mkdir(parents=True, exist_ok=True)
        record = {
            "timestamp": _audit_timestamp(),
            "level": level,
            "user": username or "unknown",
            "event": event,
            "action": action,
            "details": details or {},
        }
        line = json.dumps(record, ensure_ascii=False)
        day = time.strftime("%Y-%m-%d")
        (AUDIT_DIR / f"audit_{day}.log").open("a", encoding="utf-8").write(line + "\n")
        safe_user = "".join(ch for ch in str(username or "unknown") if ch.isalnum() or ch in {"-", "_"}) or "unknown"
        (USER_LOG_DIR / f"{safe_user}_{day}.log").open("a", encoding="utf-8").write(line + "\n")
        _write_server_session_record(record)
    except Exception:
        pass


def _ensure_security_files() -> None:
    SECURITY_DIR.mkdir(parents=True, exist_ok=True)
    AUDIT_DIR.mkdir(parents=True, exist_ok=True)
    USER_LOG_DIR.mkdir(parents=True, exist_ok=True)
    if not USERS_FILE.exists():
        salt = secrets.token_hex(16)
        payload = {
            "version": 1,
            "createdAt": _audit_timestamp(),
            "users": [
                {
                    "username": "admin",
                    "displayName": "Suite Administrator",
                    "role": "admin",
                    "active": True,
                    "salt": salt,
                    "passwordHash": _hash_password("admin123", salt),
                    "iterations": PASSWORD_ITERATIONS,
                    "createdAt": _audit_timestamp(),
                    "updatedAt": _audit_timestamp(),
                    "mustChangePassword": True,
                }
            ],
        }
        USERS_FILE.write_text(json.dumps(payload, indent=2), encoding="utf-8")
        _write_audit_record("SYSTEM", "USER_STORE_CREATED", "Default admin user created. Change password before UAT.", {"file": str(USERS_FILE)})


def _verify_password(password: str, user: dict[str, Any]) -> bool:
    try:
        expected = str(user.get("passwordHash") or "")
        salt = str(user.get("salt") or "")
        iterations = int(user.get("iterations") or PASSWORD_ITERATIONS)
        actual = _hash_password(password, salt, iterations)
        return hmac.compare_digest(actual, expected)
    except Exception:
        return False


def _load_user_store() -> dict[str, Any]:
    _ensure_security_files()
    try:
        data = json.loads(USERS_FILE.read_text(encoding="utf-8"))
        if not isinstance(data, dict):
            raise ValueError("Invalid users file")
        data.setdefault("users", [])
        return data
    except Exception:
        backup = USERS_FILE.with_suffix(f".corrupt-{int(time.time())}.json")
        if USERS_FILE.exists():
            shutil.copy2(USERS_FILE, backup)
        payload = {"version": 1, "createdAt": _audit_timestamp(), "users": []}
        USERS_FILE.write_text(json.dumps(payload, indent=2), encoding="utf-8")
        _write_audit_record("SYSTEM", "USER_STORE_RESET", "Users file was unreadable and has been reset.", {"backup": str(backup)})
        return payload


def _save_user_store(store: dict[str, Any]) -> None:
    _ensure_security_files()
    store["updatedAt"] = _audit_timestamp()
    USERS_FILE.write_text(json.dumps(store, indent=2), encoding="utf-8")


def _public_user(user: dict[str, Any]) -> dict[str, Any]:
    return {
        "username": user.get("username", ""),
        "displayName": user.get("displayName", ""),
        "role": user.get("role", "user"),
        "active": bool(user.get("active", True)),
        "createdAt": user.get("createdAt", ""),
        "updatedAt": user.get("updatedAt", ""),
        "lastLoginAt": user.get("lastLoginAt", ""),
        "mustChangePassword": bool(user.get("mustChangePassword", False)),
    }


def _audit_from_handler(handler: BaseHTTPRequestHandler, username: str, event: str, action: str, details: dict[str, Any] | None = None, level: str = "INFO") -> None:
    enriched = dict(details or {})
    enriched.setdefault("path", getattr(handler, "path", ""))
    enriched.setdefault("client", getattr(handler, "client_address", ["", ""])[0])
    _write_audit_record(username, event, action, enriched, level)


def _session_from_token(token: str) -> dict[str, Any] | None:
    if not token:
        return None
    with SESSION_LOCK:
        session = SESSIONS.get(token)
        if session:
            session["lastSeenAt"] = _audit_timestamp()
        return dict(session) if session else None


def _session_from_handler(handler: BaseHTTPRequestHandler) -> dict[str, Any] | None:
    token = handler.headers.get("X-Session-Token", "")
    if not token:
        parsed = urlparse(getattr(handler, "path", ""))
        token = parse_qs(parsed.query).get("token", [""])[0]
    return _session_from_token(token)


def _require_admin(handler: BaseHTTPRequestHandler) -> tuple[bool, dict[str, Any] | None]:
    session = _session_from_handler(handler)
    if not session or session.get("role") != "admin":
        _json_response(handler, {"ok": False, "error": "Admin access required."}, 403)
        return False, session
    return True, session


def _audit_log_entries(limit: int = 300) -> list[dict[str, Any]]:
    _ensure_security_files()
    entries: list[dict[str, Any]] = []
    for path in sorted(AUDIT_DIR.glob("audit_*.log"), reverse=True):
        try:
            lines = path.read_text(encoding="utf-8").splitlines()
            for line in reversed(lines):
                if not line.strip():
                    continue
                try:
                    entries.append(json.loads(line))
                except Exception:
                    entries.append({"timestamp": "", "level": "INFO", "user": "", "event": "RAW", "action": line, "details": {}})
                if len(entries) >= limit:
                    return entries
        except Exception:
            continue
    return entries

def _training_log(message: str, progress: int | None = None) -> None:
    with SPARROW_TRAINING_LOCK:
        SPARROW_TRAINING_JOB["status"] = message
        if progress is not None:
            SPARROW_TRAINING_JOB["progress"] = max(0, min(100, int(progress)))
        logs = SPARROW_TRAINING_JOB.setdefault("logs", [])
        logs.append(f"[{time.strftime('%H:%M:%S')}] {message}")
        del logs[:-250]


def _training_snapshot() -> dict[str, Any]:
    with SPARROW_TRAINING_LOCK:
        return {
            "ok": True,
            "running": bool(SPARROW_TRAINING_JOB.get("running")),
            "status": SPARROW_TRAINING_JOB.get("status", "Idle"),
            "progress": int(SPARROW_TRAINING_JOB.get("progress", 0)),
            "logs": list(SPARROW_TRAINING_JOB.get("logs", [])),
            "telemetry": list(SPARROW_TRAINING_JOB.get("telemetry", [])),
            "result": SPARROW_TRAINING_JOB.get("result"),
            "error": SPARROW_TRAINING_JOB.get("error", ""),
        }


def _training_metric(metric: dict[str, Any]) -> None:
    with SPARROW_TRAINING_LOCK:
        telemetry = SPARROW_TRAINING_JOB.setdefault("telemetry", [])
        telemetry.append(metric)
        del telemetry[:-500]


def _resolve_output_path(raw_path: str) -> Path:
    cleaned = str(raw_path or "").strip().strip('"')
    if not cleaned:
        cleaned = "models/sparrow_cnx_sentimentmodel_new"
    path = Path(cleaned)
    if not path.is_absolute():
        path = ROOT / cleaned
    return path.resolve()


def _normalize_training_label(value: Any) -> str:
    text = str(value or "").strip().lower()
    if text in {"0", "neg", "negative", "detractor"} or "negative" in text:
        return "negative"
    if text in {"1", "neu", "neutral", "passive"} or "neutral" in text:
        return "neutral"
    if text in {"2", "pos", "positive", "promoter"} or "positive" in text:
        return "positive"
    return ""


def _macro_f1_score(actual: list[int], predicted: list[int], label_count: int = 3) -> float:
    scores = []
    for label in range(label_count):
        tp = sum(1 for a, p in zip(actual, predicted) if a == label and p == label)
        fp = sum(1 for a, p in zip(actual, predicted) if a != label and p == label)
        fn = sum(1 for a, p in zip(actual, predicted) if a == label and p != label)
        precision = tp / (tp + fp) if (tp + fp) else 0.0
        recall = tp / (tp + fn) if (tp + fn) else 0.0
        scores.append((2 * precision * recall / (precision + recall)) if (precision + recall) else 0.0)
    return sum(scores) / len(scores)


def _weighted_f1_score(actual: list[int], predicted: list[int], label_count: int = 3) -> float:
    total = max(1, len(actual))
    weighted_total = 0.0
    for label in range(label_count):
        support = sum(1 for value in actual if value == label)
        tp = sum(1 for a, p in zip(actual, predicted) if a == label and p == label)
        fp = sum(1 for a, p in zip(actual, predicted) if a != label and p == label)
        fn = sum(1 for a, p in zip(actual, predicted) if a == label and p != label)
        precision = tp / (tp + fp) if (tp + fp) else 0.0
        recall = tp / (tp + fn) if (tp + fn) else 0.0
        f1 = (2 * precision * recall / (precision + recall)) if (precision + recall) else 0.0
        weighted_total += f1 * support
    return weighted_total / total


class _TextClassificationDataset:
    def __init__(self, encodings: dict[str, Any], labels: list[int]) -> None:
        self.encodings = encodings
        self.labels = labels

    def __len__(self) -> int:
        return len(self.labels)

    def __getitem__(self, idx: int) -> dict[str, Any]:
        import torch

        item = {key: value[idx].clone().detach() for key, value in self.encodings.items()}
        item["labels"] = torch.tensor(self.labels[idx], dtype=torch.long)
        return item


def _run_sparrow_training(payload: dict[str, Any]) -> None:
    with SPARROW_TRAINING_LOCK:
        SPARROW_TRAINING_JOB.update({"running": True, "status": "Starting", "progress": 1, "logs": [], "telemetry": [], "result": None, "error": ""})
    try:
        started_at = time.time()
        rows = payload.get("rows") if isinstance(payload.get("rows"), list) else []
        text_column = str(payload.get("textColumn") or "").strip()
        label_column = str(payload.get("labelColumn") or "").strip()
        output_path = _resolve_output_path(str(payload.get("outputPath") or ""))
        base_model = _default_model_path("sparrow")
        config = payload.get("config") if isinstance(payload.get("config"), dict) else {}
        epochs = max(1, min(20, int(config.get("epochs") or 3)))
        batch_size = max(1, min(64, int(config.get("batchSize") or 8)))
        learning_rate = float(config.get("learningRate") or 2e-5)
        max_length = max(32, min(512, int(config.get("maxLength") or 192)))
        validation_split_raw = str(config.get("validationSplit") or "20%").strip().replace("%", "")
        validation_split = max(0.05, min(0.4, float(validation_split_raw) / 100 if float(validation_split_raw) > 1 else float(validation_split_raw)))
        seed = int(config.get("seed") or 42)
        label_to_id = {"negative": 0, "neutral": 1, "positive": 2}
        id_to_label = {0: "negative", 1: "neutral", 2: "positive"}

        _training_log("Preparing labelled rows: reading browser payload, normalizing labels, and excluding blank text.", 4)
        examples: list[tuple[str, int]] = []
        for row in rows:
            if not isinstance(row, dict):
                continue
            text = str(row.get(text_column) or "").strip()
            label = _normalize_training_label(row.get(label_column))
            if text and label:
                examples.append((text, label_to_id[label]))
        if len(examples) < 12:
            raise ValueError("At least 12 valid labelled rows are required for a fine-tuning run.")
        if len({label for _, label in examples}) < 2:
            raise ValueError("Training data must contain at least two sentiment classes.")

        class_counts = {label_name: sum(1 for _, label_id in examples if label_id == label_to_id[label_name]) for label_name in label_to_id}
        _training_log(
            f"Dataset accepted: {len(examples):,} usable rows. "
            f"Class balance negative={class_counts['negative']:,}, neutral={class_counts['neutral']:,}, positive={class_counts['positive']:,}.",
            6,
        )
        random.Random(seed).shuffle(examples)
        eval_size = max(1, min(len(examples) - 1, math.ceil(len(examples) * validation_split)))
        eval_examples = examples[:eval_size]
        train_examples = examples[eval_size:]
        if not train_examples:
            raise ValueError("Not enough rows after validation split.")
        _training_log(
            f"Split complete: train={len(train_examples):,}, validation={len(eval_examples):,}, "
            f"validation_split={validation_split:.2f}, seed={seed}.",
            7,
        )

        _training_log("Loading PyTorch runtime from bundled portable Python. First run can be slow while Windows scans native DLL files.", 8)
        import torch
        _training_log("PyTorch imported. Loading Torch data utilities.", 9)
        from torch.utils.data import DataLoader
        _training_log("Torch data utilities ready. Loading Transformers runtime from bundled packages.", 10)
        from transformers import AutoConfig, AutoModelForSequenceClassification, AutoTokenizer

        device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        _training_log(f"Using device: {device}. Reading base Sparrow configuration from {base_model}.", 12)
        hf_config = AutoConfig.from_pretrained(str(base_model))
        hf_config.id2label = {str(key): value for key, value in id_to_label.items()}
        hf_config.label2id = label_to_id
        _training_log("Base configuration ready. Loading Sparrow tokenizer from local files.", 13)
        tokenizer = AutoTokenizer.from_pretrained(str(base_model), local_files_only=True)
        _training_log("Tokenizer loaded. Loading Sparrow model weights from local files.", 14)
        model = AutoModelForSequenceClassification.from_pretrained(str(base_model), config=hf_config, local_files_only=True)
        _training_log(f"Model weights loaded. Moving model to {device}.", 15)
        model.to(device)

        train_texts = [text for text, _ in train_examples]
        train_labels = [label for _, label in train_examples]
        eval_texts = [text for text, _ in eval_examples]
        eval_labels = [label for _, label in eval_examples]
        _training_log(f"Tokenizing train and validation rows with max_length={max_length}.", 18)
        train_encodings = tokenizer(train_texts, truncation=True, padding=True, max_length=max_length, return_tensors="pt")
        eval_encodings = tokenizer(eval_texts, truncation=True, padding=True, max_length=max_length, return_tensors="pt")
        train_loader = DataLoader(_TextClassificationDataset(train_encodings, train_labels), batch_size=batch_size, shuffle=True)
        eval_loader = DataLoader(_TextClassificationDataset(eval_encodings, eval_labels), batch_size=batch_size)
        optimizer = torch.optim.AdamW(model.parameters(), lr=learning_rate)

        _training_log(
            f"Fine-tuning Sparrow model: epochs={epochs}, batch_size={batch_size}, "
            f"learning_rate={learning_rate:g}, optimizer=AdamW, total_steps={len(train_loader) * epochs}.",
            24,
        )
        model.train()
        total_steps = max(1, len(train_loader) * epochs)
        completed_steps = 0
        loss_value = 0.0
        rolling_losses: list[float] = []
        epoch_metrics: list[dict[str, Any]] = []
        for epoch in range(epochs):
            epoch_loss = 0.0
            epoch_started_at = time.time()
            for batch_index, batch in enumerate(train_loader, start=1):
                optimizer.zero_grad(set_to_none=True)
                batch = {key: value.to(device) for key, value in batch.items()}
                output = model(**batch)
                loss = output.loss
                loss.backward()
                grad_norm = 0.0
                grad_terms = 0
                for parameter in model.parameters():
                    if parameter.grad is not None:
                        grad_norm += float(parameter.grad.detach().data.norm(2).cpu().item())
                        grad_terms += 1
                avg_grad_norm = grad_norm / max(grad_terms, 1)
                optimizer.step()
                completed_steps += 1
                loss_value = float(loss.detach().cpu().item())
                epoch_loss += loss_value
                rolling_losses.append(loss_value)
                rolling_losses = rolling_losses[-20:]
                rolling_loss = sum(rolling_losses) / len(rolling_losses)
                progress = 24 + int(52 * completed_steps / total_steps)
                elapsed = time.time() - started_at
                processed_rows = min(completed_steps * batch_size, len(train_examples) * epochs)
                _training_metric(
                    {
                        "phase": "train",
                        "epoch": epoch + 1,
                        "batch": batch_index,
                        "batches": len(train_loader),
                        "step": completed_steps,
                        "total_steps": total_steps,
                        "loss": loss_value,
                        "rolling_loss": rolling_loss,
                        "learning_rate": learning_rate,
                        "grad_norm": avg_grad_norm,
                        "progress": progress,
                        "elapsed_seconds": round(elapsed, 2),
                    }
                )
                _training_log(
                    f"Epoch {epoch + 1}/{epochs} | batch {batch_index}/{len(train_loader)} | "
                    f"global_step {completed_steps}/{total_steps} | rows_seen~{processed_rows:,} | "
                    f"loss={loss_value:.4f} | rolling_loss_20={rolling_loss:.4f} | "
                    f"avg_grad_norm={avg_grad_norm:.4f} | lr={learning_rate:g} | elapsed={elapsed:.1f}s.",
                    progress,
                )
            epoch_average_loss = epoch_loss / max(1, len(train_loader))
            epoch_metrics.append({"epoch": epoch + 1, "average_loss": epoch_average_loss, "elapsed_seconds": round(time.time() - epoch_started_at, 2)})
            _training_log(
                f"Epoch {epoch + 1} complete: average_loss={epoch_average_loss:.4f}, "
                f"epoch_elapsed={time.time() - epoch_started_at:.1f}s, cumulative_elapsed={time.time() - started_at:.1f}s.",
                24 + int(52 * completed_steps / total_steps),
            )

        _training_log("Evaluating fine-tuned model.", 80)
        model.eval()
        predictions: list[int] = []
        confidences: list[float] = []
        confusion = [[0, 0, 0], [0, 0, 0], [0, 0, 0]]
        with torch.no_grad():
            eval_batch_count = len(eval_loader)
            for eval_batch_index, batch in enumerate(eval_loader, start=1):
                labels = batch.pop("labels")
                batch = {key: value.to(device) for key, value in batch.items()}
                output = model(**batch)
                probabilities = torch.softmax(output.logits, dim=-1).detach().cpu()
                batch_predictions = probabilities.argmax(dim=-1).tolist()
                predictions.extend(batch_predictions)
                confidences.extend(probabilities.max(dim=-1).values.tolist())
                for actual, pred in zip(labels.tolist(), batch_predictions):
                    confusion[int(actual)][int(pred)] += 1
                _training_log(
                    f"Evaluation batch {eval_batch_index}/{eval_batch_count}: "
                    f"rows={len(batch_predictions)}, running_predictions={len(predictions)}/{len(eval_labels)}.",
                    80 + int(6 * eval_batch_index / max(eval_batch_count, 1)),
                )
        accuracy = sum(1 for actual, pred in zip(eval_labels, predictions) if actual == pred) / len(eval_labels)
        macro_f1 = _macro_f1_score(eval_labels, predictions)
        weighted_f1 = _weighted_f1_score(eval_labels, predictions)
        _training_log(f"Evaluation complete: accuracy={accuracy:.4f}, macro_f1={macro_f1:.4f}, weighted_f1={weighted_f1:.4f}.", 86)

        def prediction_rows(examples_for_audit: list[tuple[str, int]], split_name: str, cap: int = 5000) -> list[dict[str, Any]]:
            audit_examples = examples_for_audit[:cap]
            if not audit_examples:
                return []
            audit_texts = [text for text, _ in audit_examples]
            audit_labels = [label for _, label in audit_examples]
            audit_rows: list[dict[str, Any]] = []
            audit_encodings = tokenizer(audit_texts, truncation=True, padding=True, max_length=max_length, return_tensors="pt")
            audit_loader = DataLoader(_TextClassificationDataset(audit_encodings, audit_labels), batch_size=batch_size)
            row_number = 1
            with torch.no_grad():
                for audit_batch in audit_loader:
                    labels_tensor = audit_batch.pop("labels")
                    audit_batch = {key: value.to(device) for key, value in audit_batch.items()}
                    output = model(**audit_batch)
                    probabilities = torch.softmax(output.logits, dim=-1).detach().cpu()
                    batch_predictions = probabilities.argmax(dim=-1).tolist()
                    batch_confidences = probabilities.max(dim=-1).values.tolist()
                    for actual, pred, confidence in zip(labels_tensor.tolist(), batch_predictions, batch_confidences):
                        text = audit_texts[row_number - 1]
                        audit_rows.append(
                            {
                                "index": row_number,
                                "split": split_name,
                                "text": text,
                                "actual": id_to_label[int(actual)],
                                "predicted": id_to_label[int(pred)],
                                "confidence": float(confidence),
                            }
                        )
                        row_number += 1
            return audit_rows

        _training_log("Building train and validation verbatim audit tables with predictions and confidence scores.", 87)
        train_samples = prediction_rows(train_examples, "train")
        test_samples = [
            {
                "index": index,
                "split": "evaluation",
                "text": text,
                "actual": id_to_label[int(actual)],
                "predicted": id_to_label[int(pred)],
                "confidence": float(confidence),
            }
            for index, ((text, actual), pred, confidence) in enumerate(zip(eval_examples, predictions, confidences), start=1)
        ]

        if output_path == base_model:
            raise ValueError("Choose a new output folder. Do not overwrite the bundled production Sparrow model directly.")
        if output_path.exists():
            backup_path = output_path.with_name(f"{output_path.name}_backup_{time.strftime('%Y%m%d_%H%M%S')}")
            shutil.move(str(output_path), str(backup_path))
            _training_log(f"Existing output folder moved to {backup_path}.", 84)
        output_path.mkdir(parents=True, exist_ok=True)
        _training_log(f"Saving fine-tuned model to {output_path}.", 88)
        model.save_pretrained(str(output_path), safe_serialization=True)
        tokenizer.save_pretrained(str(output_path))
        metadata = {
            "base_model_path": str(base_model),
            "output_path": str(output_path),
            "created_at": time.strftime("%Y-%m-%d %H:%M:%S"),
            "text_column": text_column,
            "label_column": label_column,
            "train_rows": len(train_examples),
            "validation_rows": len(eval_examples),
            "epochs": epochs,
            "batch_size": batch_size,
            "learning_rate": learning_rate,
            "max_length": max_length,
            "validation_split": validation_split,
            "seed": seed,
            "eval_accuracy": accuracy,
            "macro_f1": macro_f1,
            "weighted_f1": weighted_f1,
            "last_loss": loss_value,
            "epoch_metrics": epoch_metrics,
            "telemetry_points": completed_steps,
            "confusion_matrix": confusion,
            "label_mapping": label_to_id,
        }
        (output_path / "training_metadata.json").write_text(json.dumps(metadata, indent=2), encoding="utf-8")
        failed_rows = [row for row in test_samples if row["actual"] != row["predicted"]]
        (output_path / "failed_samples.json").write_text(json.dumps(failed_rows, indent=2), encoding="utf-8")
        artifact_validation = {
            "config.json": (output_path / "config.json").exists(),
            "model.safetensors": (output_path / "model.safetensors").exists(),
            "tokenizer.json": (output_path / "tokenizer.json").exists(),
            "tokenizer_config.json": (output_path / "tokenizer_config.json").exists(),
            "training_metadata.json": (output_path / "training_metadata.json").exists(),
        }
        result = {
            "ok": True,
            "output_path": str(output_path),
            "accuracy": accuracy,
            "macro_f1": macro_f1,
            "weighted_f1": weighted_f1,
            "train_rows": len(train_examples),
            "validation_rows": len(eval_examples),
            "metadata": metadata,
            "telemetry": _training_snapshot().get("telemetry", []),
            "epoch_metrics": epoch_metrics,
            "confusion_matrix": confusion,
            "train_samples": train_samples,
            "test_samples": test_samples,
            "artifact_validation": artifact_validation,
        }
        with SPARROW_TRAINING_LOCK:
            SPARROW_TRAINING_JOB.update({"running": False, "status": "Training complete.", "progress": 100, "result": result, "error": ""})
        _training_log("Training complete. New Sparrow model folder is ready.", 100)
    except Exception as exc:
        traceback.print_exc()
        with SPARROW_TRAINING_LOCK:
            SPARROW_TRAINING_JOB.update({"running": False, "status": "Training failed.", "error": str(exc)})
        _training_log(f"Training failed: {exc}", 100)


def _start_sparrow_training(payload: dict[str, Any]) -> dict[str, Any]:
    with SPARROW_TRAINING_LOCK:
        if SPARROW_TRAINING_JOB.get("running"):
            return {"ok": False, "error": "A Sparrow training job is already running."}
    thread = threading.Thread(target=_run_sparrow_training, args=(payload,), daemon=True)
    thread.start()
    return {"ok": True, "message": "Sparrow fine-tuning started."}


def _is_port_open(host: str = "127.0.0.1", port: int = 5000) -> bool:
    try:
        with socket.create_connection((host, port), timeout=0.6):
            return True
    except OSError:
        return False


def _model_status() -> dict[str, Any]:
    sparrow = MODELS / "sparrow_cnx_sentimentmodel"
    theme = MODELS / "theme_acpt_resolution_model"
    return {
        "sparrow": {
            "ready": (sparrow / "config.json").exists()
            and (sparrow / "model.safetensors").exists()
            and ((sparrow / "tokenizer.json").exists() or (sparrow / "vocab.json").exists()),
            "path": str(sparrow),
        },
        "theme": {
            "ready": (theme / "theme_classifier.joblib").exists() or any(theme.glob("*.joblib")) if theme.exists() else False,
            "path": str(theme),
        },
    }


def _default_model_path(kind: str) -> Path:
    if kind == "sparrow":
        return MODELS / "sparrow_cnx_sentimentmodel"
    if kind == "theme":
        return MODELS / "theme_acpt_resolution_model"
    if kind == "owl":
        return MODELS / "theme_acpt_resolution_model"
    raise ValueError(f"Unknown model kind: {kind}")


def _validate_model_path(kind: str, raw_path: str = "") -> dict[str, Any]:
    path = Path(str(raw_path or "").strip().strip('"')) if str(raw_path or "").strip() else _default_model_path(kind)
    try:
        if kind == "sparrow":
            resolved = resolve_roberta_model_path(str(path))
        elif kind in {"theme", "owl"}:
            resolved = resolve_theme_acpt_resolution_model_path(str(path))
        else:
            raise ValueError(f"Unknown model kind: {kind}")
        display_kind = "Owl" if kind in {"theme", "owl"} else kind.title()
        return {"ok": True, "kind": kind, "path": str(path), "resolved_path": str(resolved), "message": f"{display_kind} model is valid."}
    except Exception as exc:
        return {"ok": False, "kind": kind, "path": str(path), "error": str(exc)}


def _choose_folder(initial_path: str = "") -> dict[str, Any]:
    try:
        import tkinter as tk
        from tkinter import filedialog

        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        selected = filedialog.askdirectory(initialdir=initial_path or str(ROOT), title="Select model folder")
        root.destroy()
        return {"ok": True, "path": selected or ""}
    except Exception as exc:
        return {"ok": False, "error": f"Folder picker unavailable: {exc}"}


def _list_owl_models() -> dict[str, Any]:
    manifest_path = ROOT / "models" / "theme_acpt_resolution_model" / "models.json"
    models: list[dict[str, Any]] = []
    if manifest_path.exists():
        try:
            loaded = json.loads(manifest_path.read_text(encoding="utf-8"))
            if isinstance(loaded, list):
                models.extend(row for row in loaded if isinstance(row, dict))
        except Exception:
            models = []
    bundle_dir = ROOT / "models" / "theme_acpt_resolution_model"
    for model_file in sorted(bundle_dir.glob("*.joblib"), key=lambda path: path.name.lower()):
        resolved = str(model_file.resolve())
        if any(str(item.get("path", "")).lower() == resolved.lower() for item in models):
            continue
        models.append({
            "id": model_file.stem,
            "name": model_file.stem.replace("-", " ").title(),
            "path": resolved,
            "trainedRows": "",
            "labels": [],
            "outputs": ["Theme", "ACPT", "Resolution Status"],
        })
    return {"ok": True, "models": models}


def _launch_training_tool(kind: str) -> dict[str, Any]:
    if kind == "sparrow":
        return {
            "ok": True,
            "url": "/apps/sparrow-training/index.html?login=1",
            "message": "Opening Sparrow sentiment training workspace.",
        }
    elif kind in {"theme", "owl"}:
        return {
            "ok": True,
            "url": "/apps/theme-model-training/index.html",
            "healthUrl": "/api/health",
            "message": "Opening Owl training workspace.",
        }
    else:
        return {"ok": False, "error": "Unknown training tool."}


def _theme_training_module() -> Any:
    global THEME_TRAINING_MODULE
    if THEME_TRAINING_MODULE is not None:
        return THEME_TRAINING_MODULE
    module_path = APPS / "theme-model-training" / "app.py"
    if not module_path.exists():
        raise FileNotFoundError("Owl training app was not found.")
    spec = importlib.util.spec_from_file_location("owl_theme_training_app", module_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("Could not load Owl training app.")
    module = importlib.util.module_from_spec(spec)
    sys.modules["owl_theme_training_app"] = module
    spec.loader.exec_module(module)
    THEME_TRAINING_MODULE = module
    return module


def _send_training_workbook(handler: BaseHTTPRequestHandler, payload: dict[str, Any]) -> None:
    module = _theme_training_module()
    metrics = payload.get("metrics") or {}
    if not metrics:
        _json_response(handler, {"ok": False, "error": "Train a model before exporting."}, 400)
        return
    filename = f"theme_training_{module.model_slug(str(metrics.get('modelName', 'model')))}_{time.strftime('%Y%m%d_%H%M%S')}.xlsx"
    data = module.build_training_workbook(metrics).getvalue()
    handler.send_response(200)
    handler.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    handler.send_header("Content-Disposition", f'attachment; filename="{filename}"')
    handler.send_header("Content-Length", str(len(data)))
    handler.end_headers()
    handler.wfile.write(data)


def _handle_theme_training_api(handler: BaseHTTPRequestHandler) -> bool:
    path = urlparse(handler.path).path
    if path not in {"/api/list-models", "/api/test-models", "/api/inspect", "/api/train", "/api/export-training", "/api/bulk-inspect", "/api/bulk-predict"}:
        return False
    module = _theme_training_module()
    try:
        if path == "/api/list-models":
            payload = _read_json(handler)
            folder = payload.get("folder") or str(module.TRAINED_DIR)
            _json_response(handler, {"ok": True, "folder": str(Path(folder).expanduser()), "models": module.discover_model_files(folder)})
            return True
        if path == "/api/test-models":
            payload = _read_json(handler)
            text = module.clean_text(payload.get("text"))
            model_paths = [module.clean_text(item) for item in payload.get("models", []) if module.clean_text(item)]
            if not model_paths:
                _json_response(handler, {"ok": False, "error": "Select at least one model."}, 400)
                return True
            if len(model_paths) > 3:
                _json_response(handler, {"ok": False, "error": "Select up to three models."}, 400)
                return True
            _json_response(handler, {"ok": True, "text": text, "results": [module.predict_one_verbatim(model_path, text) for model_path in model_paths]})
            return True
        if path == "/api/export-training":
            _send_training_workbook(handler, _read_json(handler))
            return True

        form = cgi.FieldStorage(fp=handler.rfile, headers=handler.headers, environ={"REQUEST_METHOD": "POST"})
        upload = form["file"] if "file" in form else None
        if upload is None or not getattr(upload, "filename", ""):
            _json_response(handler, {"ok": False, "error": "Please upload a CSV or Excel file."}, 400)
            return True
        filename = Path(upload.filename).name
        file_bytes = upload.file.read()
        df = module.records_from_upload(filename, file_bytes)
        df.columns = [str(column) for column in df.columns]
        df = df.where(pd.notna(df), "")
        columns = list(df.columns)
        if path in {"/api/inspect", "/api/bulk-inspect"}:
            _json_response(handler, {"ok": True, "rows": len(df), "columns": columns, "feedbackColumn": module.infer_feedback_column(columns)})
            return True
        if path == "/api/bulk-predict":
            feedback_col = form.getfirst("feedbackColumn", "") or module.infer_feedback_column(columns)
            model_path = form.getfirst("modelPath", "") or ""
            save_folder = form.getfirst("saveFolder", "") or ""
            output_name = form.getfirst("outputName", "") or ""
            result = module.bulk_predict_file(filename, file_bytes, feedback_col, model_path, save_folder, output_name)
            _json_response(handler, result)
            return True
        feedback_col = form.getfirst("feedbackColumn", "") or module.infer_feedback_column(columns)
        label_col = form.getfirst("labelColumn", "") or ""
        acpt_col = form.getfirst("acptColumn", "") or ""
        sentiment_col = form.getfirst("sentimentColumn", "") or ""
        resolution_col = form.getfirst("resolutionColumn", "") or ""
        max_rows = int(form.getfirst("maxRows", "5000") or 5000)
        model_name = form.getfirst("modelName", "") or ""
        save_folder = form.getfirst("saveFolder", "") or ""
        if feedback_col not in df.columns:
            _json_response(handler, {"ok": False, "error": "Select a valid feedback column."}, 400)
            return True
        if label_col not in df.columns:
            _json_response(handler, {"ok": False, "error": "Select a valid human label column."}, 400)
            return True
        if acpt_col and acpt_col not in df.columns:
            _json_response(handler, {"ok": False, "error": "Select a valid ACPT column."}, 400)
            return True
        if sentiment_col and sentiment_col not in df.columns:
            _json_response(handler, {"ok": False, "error": "Select a valid sentiment column."}, 400)
            return True
        if resolution_col and resolution_col not in df.columns:
            _json_response(handler, {"ok": False, "error": "Select a valid resolution status column."}, 400)
            return True
        metrics = module.train_model(df, feedback_col, label_col, max_rows, model_name, acpt_col, resolution_col, save_folder, sentiment_col)
        _json_response(handler, {"ok": True, "metrics": metrics})
        return True
    except Exception as exc:
        _json_response(handler, {"ok": False, "error": str(exc)}, 400)
        return True


def _sparrow_runtime_enabled() -> bool:
    return os.environ.get("NPSHTML_DISABLE_SPARROW_MODEL", "").strip().lower() not in {"1", "true", "yes", "on"}


def _theme_model_runtime_enabled() -> bool:
    return os.environ.get("NPSHTML_DISABLE_THEME_MODEL", "").strip().lower() not in {"1", "true", "yes", "on"}


def _isolated_worker_enabled() -> bool:
    return os.environ.get("NPSHTML_USE_ISOLATED_WORKER", "").strip().lower() in {"1", "true", "yes", "on"}


def _json_response(handler: BaseHTTPRequestHandler, payload: dict[str, Any], status: int = 200) -> None:
    body = json.dumps(payload, default=str).encode("utf-8")
    handler.send_response(status)
    handler.send_header("Content-Type", "application/json; charset=utf-8")
    handler.send_header("Content-Length", str(len(body)))
    handler.end_headers()
    handler.wfile.write(body)


def _ollama_category_preview(payload: dict[str, Any]) -> dict[str, Any]:
    text = str(payload.get("text") or "").strip()
    model = str(payload.get("model") or "qwen2.5:0.5b").strip()
    schema = payload.get("schema") if isinstance(payload.get("schema"), list) else []
    examples = payload.get("examples") if isinstance(payload.get("examples"), list) else []
    if not text:
        return {"ok": False, "error": "Text is required."}
    slots = []
    for item in schema[:20]:
        if not isinstance(item, dict):
            continue
        slot = str(item.get("slot") or "").strip()
        meaning = str(item.get("meaning") or "").strip()
        if slot:
            slots.append({"slot": slot, "meaning": meaning})
    if not slots:
        slots = [{"slot": f"Category{i}", "meaning": ""} for i in range(1, 21)]
    clean_examples = []
    for item in examples[:50]:
        if not isinstance(item, dict):
            continue
        slot = str(item.get("slot") or "").strip()
        label = str(item.get("label") or "").strip()
        example_text = str(item.get("text") or "").strip()
        if slot and label and example_text:
            clean_examples.append({"slot": slot, "label": label, "text": example_text[:500]})
    active_slots = [item for item in slots if str(item.get("meaning", "")).strip()]
    system_prompt = (
        "You are a strict customer feedback multi-label classifier. Return only JSON. "
        "Each key must be the provided Category slot. Each value must be exactly Yes, No, or NA. "
        "Do not overtag. Use Yes only when the specific category meaning is actually discussed. "
        "Use No when the text discusses other topics but not that category. "
        "Use NA only for blank, no-comment, or impossible-to-understand text. "
        "Generic words such as issue, process, concern, customer, request, and service do not prove a match."
    )
    user_prompt = (
        f"Text:\n{text}\n\n"
        f"Categories:\n{json.dumps(active_slots, ensure_ascii=False)}\n\n"
        "Return one JSON object only. Include every provided category slot exactly once."
    )
    if clean_examples:
        user_prompt += f"\n\nTraining examples:\n{json.dumps(clean_examples[:30], ensure_ascii=False)}"
    request_body = json.dumps(
        {
            "model": model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            "stream": False,
            "format": "json",
            "options": {"temperature": 0.0, "num_predict": 900},
        },
        ensure_ascii=False,
    ).encode("utf-8")
    request = Request(
        "http://127.0.0.1:11434/api/chat",
        data=request_body,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urlopen(request, timeout=180) as response:
            raw = json.loads(response.read().decode("utf-8"))
    except Exception as exc:
        return {"ok": False, "error": f"Ollama request failed: {exc}"}
    message = raw.get("message") if isinstance(raw.get("message"), dict) else {}
    content = str(message.get("content") or raw.get("response") or "").strip()
    start = content.find("{")
    end = content.rfind("}")
    parsed: dict[str, Any] = {}
    if start >= 0 and end > start:
        try:
            parsed = json.loads(content[start : end + 1])
        except Exception:
            parsed = {}
    normalized = {}
    raw_outputs = {"batch": content}
    for item in slots:
        slot = item["slot"]
        if not str(item.get("meaning", "")).strip():
            normalized[slot] = "NA"
            continue
        value = str(parsed.get(slot) or parsed.get(slot.lower()) or "").strip().upper()
        if value in {"Y", "YES", "TRUE", "1"}:
            normalized[slot] = "Yes"
        elif value in {"N", "NO", "FALSE", "0"}:
            normalized[slot] = "No"
        elif value in {"NA", "N/A", "UNKNOWN"}:
            normalized[slot] = "NA"
        else:
            normalized[slot] = "No"
    yes_count = sum(1 for value in normalized.values() if value == "Yes")
    na_count = sum(1 for value in normalized.values() if value == "NA")
    if len(slots) >= 10 and yes_count >= max(8, int(len(slots) * 0.75)):
        return {
            "ok": False,
            "error": "Model output failed validation: too many categories were marked Yes. Choose a stronger model or reduce category scope.",
            "partial_result": normalized,
            "raw": raw_outputs,
        }
    if len(slots) >= 10 and na_count == len(slots) and len(text.split()) > 5:
        return {
            "ok": False,
            "error": "Model output failed validation: every category came back NA for meaningful text. Choose a stronger model.",
            "partial_result": normalized,
            "raw": raw_outputs,
        }
    return {"ok": True, "model": model, "result": normalized, "raw": raw_outputs}


def _ollama_models() -> list[str]:
    try:
        request = Request("http://127.0.0.1:11434/api/tags", method="GET")
        with urlopen(request, timeout=10) as response:
            raw = json.loads(response.read().decode("utf-8"))
        models = raw.get("models") if isinstance(raw.get("models"), list) else []
        return [str(item.get("name") or "").strip() for item in models if isinstance(item, dict) and str(item.get("name") or "").strip()]
    except Exception:
        return []


def _resolve_ollama_model(model: str) -> str:
    requested = str(model or "qwen2.5:7b-instruct").strip() or "qwen2.5:7b-instruct"
    available = _ollama_models()
    if not available or requested in available:
        return requested
    aliases = {
        "qwen 7b instruct": ["qwen2.5:7b-instruct", "qwen2.5:7b", "qwen3:8b"],
        "qwen2.5:7b-instruct": ["qwen2.5:7b-instruct", "qwen2.5:7b"],
        "qwen2.5:14b-instruct": ["qwen2.5:14b-instruct", "qwen2.5:14b", "qwen2.5:7b-instruct", "qwen2.5:7b"],
        "llama3.2:latest": ["llama3.2:latest", "llama3.1:8b"],
    }
    for candidate in aliases.get(requested.lower(), []):
        if candidate in available:
            return candidate
    qwen = next((name for name in available if name.startswith("qwen2.5:7b")), None)
    if qwen:
        return qwen
    return requested


def _provider_chat_json(provider: str, api_key: str, ollama_model: str, system_prompt: str, user_prompt: str, timeout: int = 180, max_tokens: int = 3200) -> dict[str, Any]:
    provider = provider.strip().lower()
    if provider == "ollama":
        resolved_model = _resolve_ollama_model(ollama_model)
        request_body = json.dumps(
            {
                "model": resolved_model,
                "messages": [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
                "stream": False,
                "format": "json",
                "options": {"temperature": 0.0, "num_predict": int(max_tokens), "top_p": 0.9},
            },
            ensure_ascii=False,
        ).encode("utf-8")
        request = Request("http://127.0.0.1:11434/api/chat", data=request_body, headers={"Content-Type": "application/json"}, method="POST")
        try:
            with urlopen(request, timeout=timeout) as response:
                raw = json.loads(response.read().decode("utf-8"))
        except URLError as exc:
            raise RuntimeError("Ollama is not reachable. Please start Ollama and try Qwen again.") from exc
        except HTTPError as exc:
            detail = exc.read().decode("utf-8", errors="ignore")
            raise RuntimeError(f"Ollama rejected model '{resolved_model}'. Installed models: {', '.join(_ollama_models()) or 'none'}. Details: {detail[:220]}") from exc
        message = raw.get("message") if isinstance(raw.get("message"), dict) else {}
        content = str(message.get("content") or raw.get("response") or "")
        if not content.strip():
            raise RuntimeError(f"Ollama model '{resolved_model}' returned an empty response.")
        return {"model": raw.get("model") or resolved_model, "content": content}
    if provider == "openai":
        if not api_key:
            raise ValueError("OpenAI API key is required.")
        request_body = json.dumps(
            {
                "model": "gpt-4o-mini",
                "messages": [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
                "temperature": 0,
                "response_format": {"type": "json_object"},
            },
            ensure_ascii=False,
        ).encode("utf-8")
        request = Request("https://api.openai.com/v1/chat/completions", data=request_body, headers={"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}, method="POST")
        with urlopen(request, timeout=timeout) as response:
            raw = json.loads(response.read().decode("utf-8"))
        choice = (raw.get("choices") or [{}])[0]
        message = choice.get("message") if isinstance(choice.get("message"), dict) else {}
        return {"model": raw.get("model") or "OpenAI API", "content": str(message.get("content") or "")}
    if provider == "claude":
        if not api_key:
            raise ValueError("Claude API key is required.")
        request_body = json.dumps(
            {
                "model": "claude-3-5-haiku-latest",
                "max_tokens": 1600,
                "temperature": 0,
                "system": system_prompt,
                "messages": [{"role": "user", "content": user_prompt}],
            },
            ensure_ascii=False,
        ).encode("utf-8")
        request = Request("https://api.anthropic.com/v1/messages", data=request_body, headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"}, method="POST")
        with urlopen(request, timeout=timeout) as response:
            raw = json.loads(response.read().decode("utf-8"))
        parts = raw.get("content") if isinstance(raw.get("content"), list) else []
        content = "\n".join(str(part.get("text") or "") for part in parts if isinstance(part, dict))
        return {"model": raw.get("model") or "Claude API", "content": content}
    raise ValueError(f"Unsupported provider: {provider}")


def _json_from_model_content(content: str) -> dict[str, Any]:
    text = str(content or "").strip()
    if text.startswith("```"):
        text = text.strip("`").strip()
        if text.lower().startswith("json"):
            text = text[4:].strip()
    start = text.find("{")
    end = text.rfind("}")
    if start >= 0 and end > start:
        return json.loads(text[start : end + 1])
    return json.loads(text)


def _provider_chat_json_parsed(provider: str, api_key: str, model: str, system_prompt: str, user_prompt: str, timeout: int = 180, max_tokens: int = 3200) -> tuple[dict[str, Any], dict[str, Any]]:
    answer = _provider_chat_json(provider, api_key, model, system_prompt, user_prompt, timeout=timeout, max_tokens=max_tokens)
    try:
        return answer, _json_from_model_content(answer["content"])
    except Exception:
        repair_prompt = (
            "The previous response was not valid JSON. Convert it into strict JSON only. "
            "Do not add markdown, explanation, or extra text.\n\n"
            f"Required response:\n{user_prompt}\n\nPrevious response:\n{answer.get('content', '')[:4000]}"
        )
        repaired = _provider_chat_json(provider, api_key, model, "Return only valid JSON.", repair_prompt, timeout=timeout, max_tokens=max_tokens)
        return repaired, _json_from_model_content(repaired["content"])


def _executive_lens_benchmark(payload: dict[str, Any]) -> dict[str, Any]:
    api_key = str(payload.get("apiKey") or "").strip()
    provider_raw = str(payload.get("provider") or "openai").strip().lower()
    provider = "claude" if provider_raw in {"claude", "anthropic", "claude api"} else "openai"
    question = str(payload.get("question") or "").strip()
    if not api_key:
        return {"ok": False, "error": f"{'Claude' if provider == 'claude' else 'OpenAI'} API key is required."}
    if not question:
        return {"ok": False, "error": "Enter an industry benchmark question."}
    if len(question) > 500:
        return {"ok": False, "error": "Keep the benchmark question under 500 characters."}
    system_prompt = (
        "You are an executive customer-experience benchmark assistant. "
        "Answer only from general industry knowledge and public benchmark context. "
        "Do not ask for or infer private company data, customer verbatims, names, internal scores, uploaded files, or row-level data. "
        "Return concise JSON with keys: summary, benchmark_range, interpretation, cautions, leadership_actions. "
        "leadership_actions must be an array of 3 to 5 short actions."
    )
    user_prompt = (
        "Leader typed benchmark question only. No internal NPS data is included.\n\n"
        f"Question: {question}"
    )
    try:
        response = _provider_chat_json(provider, api_key, "", system_prompt, user_prompt, timeout=60)
        try:
            parsed = _json_from_model_content(response.get("content", ""))
        except Exception:
            parsed = {
                "summary": response.get("content", "").strip(),
                "benchmark_range": "",
                "interpretation": "",
                "cautions": "Treat external benchmark ranges as directional because NPS varies by industry, geography, channel, and customer segment.",
                "leadership_actions": [],
            }
        return {
            "ok": True,
            "provider": provider,
            "model": response.get("model") or ("Claude API" if provider == "claude" else "OpenAI API"),
            "question": question,
            "result": {
                "summary": str(parsed.get("summary") or ""),
                "benchmark_range": str(parsed.get("benchmark_range") or ""),
                "interpretation": str(parsed.get("interpretation") or ""),
                "cautions": str(parsed.get("cautions") or ""),
                "leadership_actions": parsed.get("leadership_actions") if isinstance(parsed.get("leadership_actions"), list) else [],
            },
        }
    except Exception as exc:
        return {"ok": False, "error": f"Benchmark lookup failed: {exc}"}


THEME_DISCOVERY_STOP_WORDS = {
    "the", "and", "for", "with", "that", "this", "was", "were", "are", "but", "not", "you", "your", "our",
    "they", "them", "have", "has", "had", "from", "get", "got", "can", "could", "would", "should", "there",
    "their", "about", "into", "just", "very", "more", "less", "than", "then", "when", "what", "why", "how",
    "because", "been", "being", "also", "still", "after", "before", "again", "all", "any", "some", "such",
    "customer", "customers", "feedback", "comment", "comments", "issue", "issues", "problem", "problems",
    "service", "agent", "call", "team", "support", "experience", "good", "bad", "nice", "great", "poor",
    "helpful", "unhelpful", "time", "thing", "things", "everything", "nothing", "yes", "no", "none",
}


def _clean_theme_token(value: Any) -> str:
    return "".join(ch.lower() if ch.isalnum() else " " for ch in str(value or "")).strip()


def _is_theme_candidate_name(name: str) -> bool:
    cleaned = _clean_theme_token(name)
    tokens = [token for token in cleaned.split() if len(token) >= 3]
    meaningful = [token for token in tokens if token not in THEME_DISCOVERY_STOP_WORDS]
    if not meaningful:
        return False
    if len(meaningful) == 1 and meaningful[0] in THEME_DISCOVERY_STOP_WORDS:
        return False
    return True


def _filter_theme_keywords(keywords: list[str]) -> list[str]:
    filtered: list[str] = []
    for keyword in keywords:
        cleaned = _clean_theme_token(keyword)
        if not cleaned:
            continue
        tokens = [token for token in cleaned.split() if token not in THEME_DISCOVERY_STOP_WORDS and len(token) >= 3]
        if not tokens:
            continue
        filtered.append(keyword)
    return filtered


def _theme_builder_discover(payload: dict[str, Any]) -> dict[str, Any]:
    provider_raw = str(payload.get("provider") or "ollama").strip().lower()
    provider = "ollama" if provider_raw in {"llama", "ollama"} else provider_raw
    api_key = str(payload.get("apiKey") or "").strip()
    model = str(payload.get("ollamaModel") or "qwen2.5:7b-instruct").strip()
    comments = payload.get("comments") if isinstance(payload.get("comments"), list) else []
    max_themes = max(3, min(int(payload.get("maxThemes") or 12), 30))
    clean_comments = [
        {"row": int(item.get("row") or idx + 1), "text": str(item.get("text") or "")[:700]}
        for idx, item in enumerate(comments[:200])
        if isinstance(item, dict) and str(item.get("text") or "").strip()
    ]
    if not clean_comments:
        return {"ok": False, "error": "No comments were supplied for theme discovery."}
    system_prompt = (
        "You are a customer feedback theme discovery engine. Return only JSON. "
        "Create clean business issue categories, not sentiment labels. Group similar issues together. "
        "Do not return stop words, adjectives, single filler words, or generic labels like Service, Issue, Customer, Feedback, Good, Helpful, Time, or Problem. "
        "Use business-readable categories such as Long Wait, Billing Dispute, Resolution Delay, Agent Knowledge, or App Login Issue."
    )
    user_prompt = (
        f"Discover up to {max_themes} theme categories from these NPS verbatims.\n"
        "Return JSON exactly like {\"themes\":[{\"name\":\"Long Wait\",\"keywords\":[\"wait\",\"hold\"]}]}.\n\n"
        f"Verbatims:\n{json.dumps(clean_comments, ensure_ascii=False)}"
    )
    try:
        answer, parsed = _provider_chat_json_parsed(provider, api_key, model, system_prompt, user_prompt, timeout=240, max_tokens=3000)
    except Exception as exc:
        return {"ok": False, "error": f"Theme discovery failed: {exc}"}
    themes = []
    raw_themes = parsed.get("themes") if isinstance(parsed.get("themes"), list) else []
    for item in raw_themes:
        if not isinstance(item, dict):
            continue
        name = str(item.get("name") or "").strip()
        if not name or not _is_theme_candidate_name(name):
            continue
        keywords = [str(word).strip().lower() for word in item.get("keywords", []) if str(word).strip()] if isinstance(item.get("keywords"), list) else []
        keywords = _filter_theme_keywords(keywords)
        themes.append({"name": name[:80], "keywords": keywords[:16], "source": provider})
    return {"ok": True, "model": answer.get("model") or model, "themes": themes[:max_themes]}


def _theme_builder_classify(payload: dict[str, Any]) -> dict[str, Any]:
    provider_raw = str(payload.get("provider") or "ollama").strip().lower()
    provider = "ollama" if provider_raw in {"llama", "ollama"} else provider_raw
    api_key = str(payload.get("apiKey") or "").strip()
    model = str(payload.get("ollamaModel") or "qwen2.5:7b-instruct").strip()
    comments = payload.get("comments") if isinstance(payload.get("comments"), list) else []
    themes = payload.get("themes") if isinstance(payload.get("themes"), list) else []
    clean_comments = [
        {"row": int(item.get("row") or idx + 1), "text": str(item.get("text") or "")[:900]}
        for idx, item in enumerate(comments[:40])
        if isinstance(item, dict)
    ]
    clean_themes = [
        {"name": str(item.get("name") or "").strip()[:80], "keywords": [str(word).strip().lower() for word in item.get("keywords", [])[:16] if str(word).strip()] if isinstance(item.get("keywords"), list) else []}
        for item in themes[:30]
        if isinstance(item, dict) and str(item.get("name") or "").strip()
    ]
    if not clean_comments or not clean_themes:
        return {"ok": False, "error": "Comments and themes are required for classification."}
    system_prompt = (
        "You are a strict NPS feedback multi-label classifier. Return only JSON. "
        "For each row and each theme, label exactly Yes, No, or NA. "
        "Use Yes only when the theme is clearly discussed. Use No for other understandable feedback. "
        "Use NA only for blank, no-comment, or impossible-to-understand text."
    )
    user_prompt = (
        "Classify each verbatim against every theme.\n"
        "Return JSON exactly like {\"assignments\":[{\"row\":1,\"categories\":{\"Long Wait\":\"Yes\",\"Billing\":\"No\"}}]}.\n\n"
        f"Themes:\n{json.dumps(clean_themes, ensure_ascii=False)}\n\n"
        f"Verbatims:\n{json.dumps(clean_comments, ensure_ascii=False)}"
    )
    try:
        token_budget = max(3200, min(9000, len(clean_comments) * max(len(clean_themes), 1) * 18 + 1200))
        answer, parsed = _provider_chat_json_parsed(provider, api_key, model, system_prompt, user_prompt, timeout=300, max_tokens=token_budget)
    except Exception as exc:
        return {"ok": False, "error": f"Theme classification failed: {exc}"}
    theme_names = [item["name"] for item in clean_themes]
    assignments = []
    raw_assignments = parsed.get("assignments") if isinstance(parsed.get("assignments"), list) else []
    for item in raw_assignments:
        if not isinstance(item, dict):
            continue
        row = int(item.get("row") or 0)
        cats = item.get("categories") if isinstance(item.get("categories"), dict) else {}
        normalized = {}
        for name in theme_names:
            raw = str(cats.get(name) or cats.get(name.lower()) or "").strip().upper()
            normalized[name] = "Yes" if raw in {"Y", "YES", "TRUE", "1"} else ("NA" if raw in {"NA", "N/A", "UNKNOWN"} else "No")
        assignments.append({"row": row, "categories": normalized})
    return {"ok": True, "model": answer.get("model") or model, "assignments": assignments}


def _read_json(handler: BaseHTTPRequestHandler) -> dict[str, Any]:
    length = int(handler.headers.get("Content-Length", "0"))
    raw = handler.rfile.read(length)
    if not raw:
        return {}
    return json.loads(raw.decode("utf-8"))


def _weekly_trend_dashboard(payload: dict[str, Any]) -> dict[str, Any]:
    """Aggregate the completed analysis into a configurable weekly scorecard."""
    with STATE_LOCK:
        source = STATE.analyzed_df.copy()
        calendar = dict(STATE.calendar_settings)
    if source.empty:
        return {"ok": False, "error": "Run an analysis before building a weekly trend dashboard."}

    work = _apply_date_filter(source)
    if work.empty:
        return {"ok": False, "error": "No analyzed rows are available for the current date filter."}
    if "Week" not in work.columns:
        if "Feedback Date" not in work.columns:
            return {"ok": False, "error": "Weekly trends require a mapped feedback date."}
        work["Week"] = week_period_start(work["Feedback Date"], _calendar_settings({"calendar": calendar})["weekStart"])

    mode = str(payload.get("mode") or "nps").strip().lower()
    mode = "csat" if mode == "csat" else "nps"
    dimension = str(payload.get("dimension") or "__overall__").strip()
    dimension_value = str(payload.get("dimensionValue") or "__all__").strip()
    if dimension != "__overall__":
        if dimension not in work.columns:
            return {"ok": False, "error": f"Dimension '{dimension}' is not available in this analysis."}
        if dimension_value != "__all__":
            work = work[work[dimension].fillna("").astype(str).str.strip().eq(dimension_value)].copy()
    if work.empty:
        return {"ok": False, "error": "No analyzed rows match the selected dimension value."}

    work["__week"] = pd.to_datetime(work["Week"], errors="coerce").dt.normalize()
    work = work.dropna(subset=["__week"])
    if work.empty:
        return {"ok": False, "error": "No valid week values were found in the analyzed data."}
    all_weeks = sorted(work["__week"].drop_duplicates().tolist())
    try:
        requested_weeks = max(0, int(payload.get("weeks") or 0))
    except (TypeError, ValueError):
        requested_weeks = 12
    selected_weeks = all_weeks[-requested_weeks:] if requested_weeks else all_weeks
    work = work[work["__week"].isin(selected_weeks)]

    metric_set = str(payload.get("metricSet") or "full").strip().lower()
    sentiment_col = next((name for name in ["Sentiment", "Sentiment Label", "Overall Sentiment"] if name in work.columns), "")
    type_candidates = ["CSAT Type", "CSAT Segment", "Satisfaction Segment"] if mode == "csat" else ["NPS Type", "NPS Segment"]
    type_col = next((name for name in type_candidates if name in work.columns), "")

    metrics: list[dict[str, Any]] = []
    volume_values: list[int] = []
    score_values: list[float | None] = []
    if metric_set == "sentiment":
        definitions = [("Positive", "positive"), ("Neutral", "neutral"), ("Negative", "negative")]
        count_rows = {label: [] for label, _ in definitions}
        pct_rows = {label: [] for label, _ in definitions}
        for week in selected_weeks:
            frame = work[work["__week"].eq(week)]
            total = int(len(frame)); volume_values.append(total)
            values = frame[sentiment_col].fillna("").astype(str).str.strip().str.lower() if sentiment_col else pd.Series([], dtype=str)
            for label, key in definitions:
                count = int(values.eq(key).sum())
                count_rows[label].append(count)
                pct_rows[label].append(round((count / total * 100.0), 2) if total else 0.0)
        metrics.append({"label": "Survey Volume", "format": "integer", "values": volume_values})
        for label, _ in definitions:
            metrics.append({"label": label, "format": "integer", "values": count_rows[label]})
        for label, _ in definitions:
            metrics.append({"label": f"{label} %", "format": "percent", "values": pct_rows[label]})
        score_values = pct_rows["Positive"]
        score_name = "Positive %"
    else:
        if mode == "csat":
            definitions = [("Satisfied", {"satisfied", "promoter"}), ("Neutral", {"neutral", "passive"}), ("Dissatisfied", {"dissatisfied", "detractor"})]
            score_name = "CSAT"
        else:
            definitions = [("Promoters", {"promoter"}), ("Passives", {"passive"}), ("Detractors", {"detractor"})]
            score_name = "NPS"
        count_rows = {label: [] for label, _ in definitions}
        pct_rows = {label: [] for label, _ in definitions}
        for week in selected_weeks:
            frame = work[work["__week"].eq(week)]
            total = int(len(frame)); volume_values.append(total)
            values = frame[type_col].fillna("").astype(str).str.strip().str.lower() if type_col else pd.Series([], dtype=str)
            for label, accepted in definitions:
                count = int(values.isin(accepted).sum())
                count_rows[label].append(count)
                pct_rows[label].append(round((count / total * 100.0), 2) if total else 0.0)
        if mode == "csat":
            score_values = list(pct_rows["Satisfied"])
        else:
            score_values = [round(p - d, 2) for p, d in zip(pct_rows["Promoters"], pct_rows["Detractors"])]
        if metric_set in {"full", "score"}:
            metrics.append({"label": "Survey Volume", "format": "integer", "values": volume_values})
        if metric_set in {"full", "segments"}:
            for label, _ in definitions:
                metrics.append({"label": label, "format": "integer", "values": count_rows[label]})
            for label, _ in definitions:
                metrics.append({"label": f"{label} %", "format": "percent", "values": pct_rows[label]})
        if metric_set in {"full", "score"}:
            metrics.append({"label": score_name, "format": "score", "values": score_values})

    comparison_count = 0
    comparison_omitted = 0
    if dimension != "__overall__" and dimension_value == "__all__":
        clean_dimension = work[dimension].fillna("").astype(str).str.strip()
        ranked_values = clean_dimension[clean_dimension.ne("")].value_counts()
        comparison_count = min(20, int(len(ranked_values)))
        comparison_omitted = max(0, int(len(ranked_values)) - comparison_count)
        comparison_values = [str(value) for value in ranked_values.head(20).index.tolist()]
        comparison_metrics: list[dict[str, Any]] = []
        segment_definitions = (
            [("Satisfied %", {"satisfied", "promoter"}), ("Neutral %", {"neutral", "passive"}), ("Dissatisfied %", {"dissatisfied", "detractor"})]
            if mode == "csat"
            else [("Promoters %", {"promoter"}), ("Passives %", {"passive"}), ("Detractors %", {"detractor"})]
        )
        for value in comparison_values:
            value_frame = work[clean_dimension.eq(value)]
            group_volume: list[int] = []
            group_score: list[float] = []
            group_segments = {label: [] for label, _ in segment_definitions}
            group_sentiments = {label: [] for label in ["Positive %", "Neutral %", "Negative %"]}
            for week in selected_weeks:
                frame = value_frame[value_frame["__week"].eq(week)]
                total = int(len(frame))
                group_volume.append(total)
                if metric_set == "sentiment":
                    labels = frame[sentiment_col].fillna("").astype(str).str.strip().str.lower() if sentiment_col else pd.Series([], dtype=str)
                    for label, key in [("Positive %", "positive"), ("Neutral %", "neutral"), ("Negative %", "negative")]:
                        group_sentiments[label].append(round(float(labels.eq(key).sum()) / total * 100.0, 2) if total else 0.0)
                else:
                    labels = frame[type_col].fillna("").astype(str).str.strip().str.lower() if type_col else pd.Series([], dtype=str)
                    for label, accepted in segment_definitions:
                        group_segments[label].append(round(float(labels.isin(accepted).sum()) / total * 100.0, 2) if total else 0.0)
                    if mode == "csat":
                        group_score.append(group_segments["Satisfied %"][-1])
                    else:
                        group_score.append(round(group_segments["Promoters %"][-1] - group_segments["Detractors %"][-1], 2))
            if metric_set == "sentiment":
                for label in ["Positive %", "Neutral %", "Negative %"]:
                    comparison_metrics.append({
                        "label": f"{value} · {label}", "format": "percent-volume",
                        "values": [{"value": metric_value, "volume": volume} for metric_value, volume in zip(group_sentiments[label], group_volume)],
                        "groupStart": label == "Positive %",
                    })
            elif metric_set == "segments":
                for label, _ in segment_definitions:
                    comparison_metrics.append({
                        "label": f"{value} · {label}", "format": "percent-volume",
                        "values": [{"value": metric_value, "volume": volume} for metric_value, volume in zip(group_segments[label], group_volume)],
                        "groupStart": label == segment_definitions[0][0],
                    })
            else:
                comparison_metrics.append({
                    "label": f"{value} · {score_name}", "format": "score-volume",
                    "values": [{"value": score, "volume": volume} for score, volume in zip(group_score, group_volume)],
                    "emphasis": True, "groupStart": True,
                })
        metrics = comparison_metrics

    chronological_score = list(score_values)
    chronological_volume = list(volume_values)
    week_items = [{"key": week.strftime("%Y-%m-%d"), "label": f"WE {(week + pd.Timedelta(days=6)).strftime('%-d %b') if os.name != 'nt' else str((week + pd.Timedelta(days=6)).day) + ' ' + (week + pd.Timedelta(days=6)).strftime('%b')}"} for week in selected_weeks]
    if str(payload.get("order") or "newest").lower() == "newest":
        week_items.reverse()
        for metric in metrics:
            metric["values"] = list(reversed(metric["values"]))
    latest = chronological_score[-1] if chronological_score else None
    previous = chronological_score[-2] if len(chronological_score) > 1 else None
    filter_label = "Overall"
    if dimension != "__overall__":
        filter_label = f"{dimension}: {dimension_value if dimension_value != '__all__' else f'Comparing {comparison_count} values'}"
    return {
        "ok": True, "mode": mode, "scoreName": score_name, "filterLabel": filter_label,
        "weeks": week_items, "metrics": metrics, "rowsUsed": int(len(work)),
        "dimensionComparison": dimension != "__overall__" and dimension_value == "__all__",
        "comparisonCount": comparison_count, "comparisonOmitted": comparison_omitted,
        "summary": {
            "latestScore": latest, "previousScore": previous,
            "movement": round(latest - previous, 2) if latest is not None and previous is not None else None,
            "latestVolume": chronological_volume[-1] if chronological_volume else 0,
            "averageVolume": round(sum(chronological_volume) / len(chronological_volume), 1) if chronological_volume else 0,
        },
    }


def _weekly_trend_matrix(payload: dict[str, Any]) -> dict[str, Any]:
    """Build a one-week matrix with dimension values across the columns."""
    with STATE_LOCK:
        source = STATE.analyzed_df.copy()
        calendar = dict(STATE.calendar_settings)
    if source.empty:
        return {"ok": False, "error": "Run an analysis before building a weekly matrix dashboard."}
    work = _apply_date_filter(source)
    if "Week" not in work.columns:
        if "Feedback Date" not in work.columns:
            return {"ok": False, "error": "Weekly dashboards require a mapped feedback date."}
        work["Week"] = week_period_start(work["Feedback Date"], _calendar_settings({"calendar": calendar})["weekStart"])
    work["__week"] = pd.to_datetime(work["Week"], errors="coerce").dt.normalize()
    work = work.dropna(subset=["__week"])
    if work.empty:
        return {"ok": False, "error": "No valid reporting weeks were found."}

    all_weeks = sorted(work["__week"].drop_duplicates().tolist())
    week_items = []
    for week in reversed(all_weeks):
        week_end = week + pd.Timedelta(days=6)
        week_items.append({"key": week.strftime("%Y-%m-%d"), "label": f"WE {week_end.day} {week_end.strftime('%b')}"})
    requested_week = pd.to_datetime(str(payload.get("week") or ""), errors="coerce")
    selected_week = requested_week.normalize() if not pd.isna(requested_week) else all_weeks[-1]
    if selected_week not in all_weeks:
        selected_week = all_weeks[-1]
    week_frame = work[work["__week"].eq(selected_week)].copy()

    dimension = str(payload.get("dimension") or "").strip()
    if not dimension or dimension not in week_frame.columns:
        return {"ok": False, "error": "Select an available dimension field."}
    dimension_aliases: list[str] = [dimension]
    dimension_key = dimension.lower()
    if any(token in dimension_key for token in ["agent", "advisor", "employee"]):
        dimension_aliases.extend(["Agent Name", "Agent", "AgentName", "Advisor Name", "Employee Name", "Associate"])
    elif any(token in dimension_key for token in ["manager", "supervisor", "team lead", "tl"]):
        dimension_aliases.extend(["Manager/TL", "Manager", "Manager Name", "TL Name", "Team Lead", "Supervisor"])
    resolved_dimension = dimension
    clean_dimension = week_frame[dimension].fillna("").astype(str).str.strip()
    for candidate in dict.fromkeys(dimension_aliases):
        if candidate not in week_frame.columns:
            continue
        candidate_values = week_frame[candidate].fillna("").astype(str).str.strip()
        if candidate_values.ne("").any():
            resolved_dimension = candidate
            clean_dimension = candidate_values
            break
    week_frame = week_frame[clean_dimension.ne("")].copy()
    if week_frame.empty:
        return {"ok": False, "error": f"No values are available for {dimension} in the selected week."}
    week_frame["__dimension"] = week_frame[resolved_dimension].fillna("").astype(str).str.strip()
    volumes = week_frame["__dimension"].value_counts()
    sort_mode = str(payload.get("sort") or "volume").strip().lower()
    dimension_values = sorted(volumes.index.tolist(), key=lambda value: str(value).lower()) if sort_mode == "name" else volumes.index.tolist()
    try:
        max_columns = max(0, min(50, int(payload.get("maxColumns") or 15)))
    except (TypeError, ValueError):
        max_columns = 15
    if max_columns:
        dimension_values = dimension_values[:max_columns]

    mode = "csat" if str(payload.get("mode") or "nps").strip().lower() == "csat" else "nps"
    type_candidates = ["CSAT Type", "CSAT Segment", "Satisfaction Segment"] if mode == "csat" else ["NPS Type", "NPS Segment"]
    type_col = next((name for name in type_candidates if name in week_frame.columns), "")
    sentiment_col = next((name for name in ["Sentiment", "Sentiment Label", "Overall Sentiment"] if name in week_frame.columns), "")
    requested_metrics = payload.get("metrics") if isinstance(payload.get("metrics"), list) else []
    default_metrics = ["volume", "positiveCount", "middleCount", "negativeCount", "positivePct", "middlePct", "negativePct", "score"]
    metric_keys = [str(item) for item in requested_metrics] or default_metrics
    allowed_metrics = {"volume", "positiveCount", "middleCount", "negativeCount", "positivePct", "middlePct", "negativePct", "score", "sentimentPositivePct", "sentimentNeutralPct", "sentimentNegativePct"}
    metric_keys = [key for key in metric_keys if key in allowed_metrics]
    if not metric_keys:
        return {"ok": False, "error": "Select at least one metric row."}

    if mode == "csat":
        positive_label, middle_label, negative_label, score_label = "Satisfied", "Neutral", "Dissatisfied", "CSAT"
        positive_values, middle_values, negative_values = {"satisfied", "promoter"}, {"neutral", "passive"}, {"dissatisfied", "detractor"}
    else:
        positive_label, middle_label, negative_label, score_label = "Promoters", "Passives", "Detractors", "NPS"
        positive_values, middle_values, negative_values = {"promoter"}, {"passive"}, {"detractor"}
    row_definitions = {
        "volume": ("Survey Volume", "integer"),
        "positiveCount": (positive_label, "integer"), "middleCount": (middle_label, "integer"), "negativeCount": (negative_label, "integer"),
        "positivePct": (f"{positive_label} %", "percent"), "middlePct": (f"{middle_label} %", "percent"), "negativePct": (f"{negative_label} %", "percent"),
        "score": (score_label, "score"),
        "sentimentPositivePct": ("Positive Sentiment %", "percent"), "sentimentNeutralPct": ("Neutral Sentiment %", "percent"), "sentimentNegativePct": ("Negative Sentiment %", "percent"),
    }
    rows = [{"key": key, "label": row_definitions[key][0], "format": row_definitions[key][1], "values": []} for key in metric_keys]
    for value in dimension_values:
        frame = week_frame[week_frame["__dimension"].eq(value)]
        total = int(len(frame))
        types = frame[type_col].fillna("").astype(str).str.strip().str.lower() if type_col else pd.Series([], dtype=str)
        positive = int(types.isin(positive_values).sum()); middle = int(types.isin(middle_values).sum()); negative = int(types.isin(negative_values).sum())
        positive_pct = round(positive / total * 100.0, 2) if total else 0.0
        middle_pct = round(middle / total * 100.0, 2) if total else 0.0
        negative_pct = round(negative / total * 100.0, 2) if total else 0.0
        score = positive_pct if mode == "csat" else round(positive_pct - negative_pct, 2)
        sentiments = frame[sentiment_col].fillna("").astype(str).str.strip().str.lower() if sentiment_col else pd.Series([], dtype=str)
        values = {
            "volume": total, "positiveCount": positive, "middleCount": middle, "negativeCount": negative,
            "positivePct": positive_pct, "middlePct": middle_pct, "negativePct": negative_pct, "score": score,
            "sentimentPositivePct": round(float(sentiments.eq("positive").sum()) / total * 100.0, 2) if total else 0.0,
            "sentimentNeutralPct": round(float(sentiments.eq("neutral").sum()) / total * 100.0, 2) if total else 0.0,
            "sentimentNegativePct": round(float(sentiments.eq("negative").sum()) / total * 100.0, 2) if total else 0.0,
        }
        for row in rows:
            row["values"].append(values[row["key"]])

    value_field = str(payload.get("valueField") or "").strip()
    aggregation = str(payload.get("aggregation") or "average").strip().lower()
    if value_field and value_field in week_frame.columns:
        custom_values = []
        for value in dimension_values:
            numeric = pd.to_numeric(week_frame.loc[week_frame["__dimension"].eq(value), value_field], errors="coerce").dropna()
            result = float(numeric.sum()) if aggregation == "sum" else (float(numeric.mean()) if not numeric.empty else None)
            custom_values.append(round(result, 2) if result is not None else None)
        rows.append({"key": "customValue", "label": f"{'Sum' if aggregation == 'sum' else 'Average'} {value_field}", "format": "number", "values": custom_values})

    selected_week_end = selected_week + pd.Timedelta(days=6)
    return {
        "ok": True, "mode": mode, "scoreName": score_label, "dimension": dimension,
        "week": selected_week.strftime("%Y-%m-%d"), "weekLabel": f"WE {selected_week_end.day} {selected_week_end.strftime('%b')}", "availableWeeks": week_items,
        "columns": [{"key": str(value), "label": str(value), "volume": int(volumes.get(value, 0))} for value in dimension_values],
        "metrics": rows, "rowsUsed": int(len(week_frame)), "totalDimensionValues": int(len(volumes)),
    }


def _set_upload_progress(upload_id: str, percent: float, stage: str, message: str, complete: bool = False) -> None:
    if not upload_id:
        return
    with UPLOAD_PROGRESS_LOCK:
        UPLOAD_PROGRESS[upload_id] = {
            "ok": True,
            "uploadId": upload_id,
            "percent": max(0, min(100, round(float(percent), 1))),
            "stage": stage,
            "message": message,
            "complete": complete,
            "updatedAt": time.time(),
        }


def _decode_excel(payload: dict[str, Any], upload_id: str = "") -> tuple[pd.DataFrame, list[str], str]:
    raw = payload.get("data", "")
    if "," in raw:
        raw = raw.split(",", 1)[1]
    _set_upload_progress(upload_id, 66, "Decode", "Transfer complete. Decoding the workbook payload in local memory.")
    content = base64.b64decode(raw)
    _set_upload_progress(upload_id, 74, "Workbook", "Opening the Excel workbook and checking available sheets. Large workbooks can remain in this stage for several minutes.")
    workbook = pd.ExcelFile(BytesIO(content))
    sheet_names = [str(name) for name in workbook.sheet_names]
    requested_sheet = str(payload.get("sheetName") or "").strip()
    if len(sheet_names) > 1 and not requested_sheet:
        _set_upload_progress(upload_id, 82, "Sheet selection", f"Workbook has {len(sheet_names):,} sheets. Waiting for the user to choose one.", True)
        return pd.DataFrame(), sheet_names, ""
    selected_sheet = requested_sheet or (sheet_names[0] if sheet_names else "")
    if requested_sheet and requested_sheet not in sheet_names:
        raise ValueError(f"Selected sheet '{requested_sheet}' was not found in this workbook.")
    frame = pd.read_excel(workbook, sheet_name=selected_sheet)
    _set_upload_progress(upload_id, 84, "Rows", f"Workbook opened on sheet '{selected_sheet}'. Found {len(frame):,} rows and {len(frame.columns):,} columns.")
    return frame, sheet_names, selected_sheet


def _analysis_file_preview(payload: dict[str, Any]) -> dict[str, Any]:
    base_key = str(payload.get("baseKey") or "").strip()
    lookup_key = str(payload.get("lookupKey") or "").strip()
    with STATE_LOCK:
        base = STATE.base_df.copy()
        lookup = STATE.lookup_df.copy()
        files = dict(STATE.files)
    if base.empty:
        return {"ok": False, "error": "Upload a Base File before previewing the analysis file."}
    merged = _merge_lookup(base, lookup, base_key, lookup_key) if base_key and lookup_key and not lookup.empty else base
    return {
        "ok": True,
        "rows": int(len(merged)),
        "columns": list(merged.columns),
        "columnCount": int(len(merged.columns)),
        "previewRows": _safe_records(merged, 8),
        "baseSheet": files.get("base_sheet", ""),
        "lookupSheet": files.get("lookup_sheet", ""),
        "usedLookup": bool(base_key and lookup_key and not lookup.empty),
    }

def _excel_sheet_preview(payload: dict[str, Any]) -> dict[str, Any]:
    raw = payload.get("data", "")
    if "," in raw:
        raw = raw.split(",", 1)[1]
    content = base64.b64decode(raw)
    sheet_name = str(payload.get("sheetName") or "").strip()
    if not sheet_name:
        return {"ok": False, "error": "Select a worksheet to preview."}
    workbook = pd.ExcelFile(BytesIO(content))
    sheet_names = [str(name) for name in workbook.sheet_names]
    if sheet_name not in sheet_names:
        return {"ok": False, "error": f"Sheet '{sheet_name}' was not found in this workbook.", "sheetNames": sheet_names}
    frame = pd.read_excel(workbook, sheet_name=sheet_name, nrows=8)
    total_rows = 0
    total_columns = len(frame.columns)
    try:
        worksheet = workbook.book[sheet_name]
        total_rows = max(int(getattr(worksheet, "max_row", 0) or 0) - 1, 0)
        total_columns = int(getattr(worksheet, "max_column", total_columns) or total_columns)
    except Exception:
        total_rows = len(frame)
    total_rows = max(total_rows, len(frame))
    return {
        "ok": True,
        "sheetName": sheet_name,
        "sheetNames": sheet_names,
        "rows": total_rows,
        "columns": list(frame.columns),
        "columnCount": total_columns,
        "previewRows": _safe_records(frame, 8),
    }

def _uploaded_file_preview(payload: dict[str, Any]) -> dict[str, Any]:
    kind = str(payload.get("kind") or "base").strip().lower()
    page = max(1, int(_safe_float(payload.get("page"), 1)))
    page_size = int(_safe_float(payload.get("pageSize"), 100))
    page_size = min(max(page_size, 25), 500)
    search = str(payload.get("search") or "").strip().lower()
    with STATE_LOCK:
        frame = STATE.lookup_df.copy() if kind == "lookup" else STATE.base_df.copy()
        files = dict(STATE.files)
    if frame.empty:
        return {"ok": False, "error": f"Upload a {'Lookup' if kind == 'lookup' else 'Base'} File before previewing it."}
    filtered = frame
    if search:
        text = frame.astype(str).agg(" ".join, axis=1).str.lower()
        filtered = frame[text.str.contains(search, na=False)]
    total_rows = int(len(filtered))
    start = (page - 1) * page_size
    end = start + page_size
    visible = filtered.iloc[start:end]
    return {
        "ok": True,
        "kind": kind,
        "fileName": files.get(kind, ""),
        "sheetName": files.get(f"{kind}_sheet", ""),
        "rows": int(len(frame)),
        "filteredRows": total_rows,
        "page": page,
        "pageSize": page_size,
        "pageCount": max(1, math.ceil(total_rows / page_size)) if total_rows else 1,
        "columns": [str(column) for column in frame.columns],
        "columnCount": int(len(frame.columns)),
        "previewRows": _safe_records(visible, max(len(visible), 1)),
    }

def _decode_tabular(payload: dict[str, Any]) -> pd.DataFrame:
    raw = payload.get("data", "")
    if "," in raw:
        raw = raw.split(",", 1)[1]
    content = base64.b64decode(raw)
    name = str(payload.get("name") or "").lower()
    if name.endswith(".csv"):
        return pd.read_csv(BytesIO(content))
    if name.endswith(".tsv") or name.endswith(".txt"):
        return pd.read_csv(BytesIO(content), sep="\t")
    return pd.read_excel(BytesIO(content))


def _generic_guess_columns(columns: list[str], mode: str = "csat") -> dict[str, str]:
    lowered = {str(column).lower(): column for column in columns}

    def pick(*needles: str) -> str:
        for needle in needles:
            for lower, original in lowered.items():
                if needle in lower:
                    return original
        return ""

    return {
        "feedback": pick("comment", "verbatim", "feedback", "review", "response", "text"),
        "score": pick("csat", "satisfaction", "rating", "score") if mode == "csat" else "",
        "agent": pick("agent", "advisor", "employee", "associate"),
        "manager": pick("manager/tl", "tl name", "team manager", "manager", "supervisor"),
        "lob": pick("lob", "line of business", "business unit", "business"),
        "date": pick("date", "survey date", "response date", "created"),
    }


def _infer_column_type(series: pd.Series, column_name: str = "") -> str:
    non_blank = series.dropna()
    non_blank = non_blank[non_blank.astype(str).str.strip() != ""]
    if non_blank.empty:
        return "Blank"
    numeric = pd.to_numeric(non_blank, errors="coerce")
    if numeric.notna().mean() >= 0.9:
        return "Number"
    name_lower = str(column_name).lower()
    sample_text = non_blank.astype(str).str.strip()
    looks_date_named = any(token in name_lower for token in ["date", "time", "day", "month", "year"])
    looks_date_values = sample_text.str.contains(r"[-/.:]|jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec", case=False, regex=True).mean() >= 0.6
    dates = pd.to_datetime(non_blank, errors="coerce") if looks_date_named or looks_date_values else pd.Series(dtype="datetime64[ns]")
    if (looks_date_named or looks_date_values) and dates.notna().mean() >= 0.9:
        return "Date/Time"
    unique_values = non_blank.astype(str).nunique(dropna=True)
    if unique_values <= max(20, len(non_blank) * 0.2):
        return "Category"
    return "Text"


def _column_profile(df: pd.DataFrame) -> dict[str, dict[str, Any]]:
    profiles: dict[str, dict[str, Any]] = {}
    total_rows = len(df)
    for column in df.columns:
        series = df[column]
        text_series = series.astype(str).str.strip()
        blanks = int(series.isna().sum() + (text_series == "").sum())
        unique_series = series.dropna().astype(str).str.strip().replace("", pd.NA).dropna()
        unique = int(unique_series.nunique())
        unique_values = []
        if unique and unique < 8:
            unique_values = sorted(unique_series.drop_duplicates().astype(str).tolist(), key=lambda value: value.lower())
        profiles[str(column)] = {
            "totalEntries": total_rows,
            "totalBlanks": blanks,
            "unique": unique,
            "uniqueLabel": "All unique" if total_rows > 0 and unique == (total_rows - blanks) and blanks == 0 else f"{unique:,}",
            "uniqueValues": unique_values,
            "dataType": _infer_column_type(series, str(column)),
            "sampleValue": str(unique_series.iloc[0]) if not unique_series.empty else "",
        }
    return profiles


def _safe_number(value: Any) -> float | None:
    try:
        number = float(value)
        if pd.isna(number):
            return None
        return number
    except Exception:
        return None


def _csat_label(value: Any) -> str:
    number = _safe_number(value)
    if number is None:
        text = str(value or "").strip().lower()
        if any(token in text for token in ["satisfied", "good", "excellent", "happy", "positive"]):
            return "Satisfied"
        if any(token in text for token in ["neutral", "average", "ok"]):
            return "Neutral"
        if any(token in text for token in ["dissatisfied", "bad", "poor", "negative", "unhappy"]):
            return "Dissatisfied"
        return "Unknown"
    if number <= 5:
        if number >= 4:
            return "Satisfied"
        if number >= 3:
            return "Neutral"
        return "Dissatisfied"
    if number >= 8:
        return "Satisfied"
    if number >= 6:
        return "Neutral"
    return "Dissatisfied"


def _word_tokens(value: Any) -> set[str]:
    text = str(value or "").lower()
    cleaned = "".join(char if char.isalnum() else " " for char in text)
    return {token for token in cleaned.split() if len(token) > 2}


def _custom_category_for_text(text: Any, categories: list[str]) -> tuple[str, float]:
    clean_categories = [str(item).strip() for item in categories if str(item).strip()]
    if not clean_categories:
        return "Uncategorized", 0.0
    feedback_tokens = _word_tokens(text)
    keyword_bank = {
        "wait": {"wait", "waiting", "delay", "delayed", "slow", "queue", "hold", "response"},
        "time": {"time", "delay", "delayed", "quick", "fast", "slow", "waiting"},
        "agent": {"agent", "advisor", "representative", "staff", "executive", "associate", "rude", "helpful"},
        "behavior": {"rude", "polite", "helpful", "attitude", "empathy", "professional", "courteous"},
        "resolution": {"resolved", "resolution", "fix", "fixed", "solve", "solved", "unresolved", "issue"},
        "knowledge": {"knowledge", "know", "understand", "trained", "expert", "information"},
        "billing": {"bill", "billing", "charge", "charged", "payment", "refund", "invoice", "fee"},
        "technical": {"technical", "tech", "system", "app", "website", "login", "error", "crash"},
        "process": {"process", "policy", "procedure", "steps", "documents", "approval"},
        "communication": {"communication", "communicate", "explained", "explain", "update", "email", "call"},
        "follow": {"follow", "callback", "called", "respond", "response", "update"},
        "quality": {"quality", "accuracy", "accurate", "wrong", "incorrect", "mistake"},
    }
    best = clean_categories[0]
    best_score = -1
    for category in clean_categories:
        category_tokens = _word_tokens(category)
        expanded = set(category_tokens)
        for token in list(category_tokens):
            expanded |= keyword_bank.get(token, set())
        score = len(feedback_tokens & expanded)
        if score > best_score:
            best = category
            best_score = score
    if best_score <= 0:
        for category in clean_categories:
            if category.lower() == "other":
                return category, 0.25
        return clean_categories[-1], 0.2
    return best, round(min(0.98, 0.45 + (best_score * 0.12)), 2)


def _build_verbatim_insights(analyzed: pd.DataFrame, category_col: str = "Custom Category") -> list[dict[str, Any]]:
    if analyzed.empty:
        return []
    insights: list[dict[str, Any]] = []
    total = len(analyzed)
    sentiment_counts = analyzed.get("Sentiment", pd.Series(dtype=str)).fillna("Unknown").astype(str).value_counts()
    negative_count = int(sentiment_counts.get("Negative", 0))
    positive_count = int(sentiment_counts.get("Positive", 0))
    if negative_count:
        insights.append({
            "Title": "Negative feedback concentration",
            "Insight": f"{round((negative_count / max(total, 1)) * 100, 1)}% of verbatims are negative. Review the top negative custom categories first.",
        })
    if positive_count:
        insights.append({
            "Title": "Positive feedback signal",
            "Insight": f"{round((positive_count / max(total, 1)) * 100, 1)}% of verbatims are positive. Use these comments to identify repeatable good behaviors.",
        })
    if category_col in analyzed.columns:
        category_counts = analyzed[category_col].fillna("Uncategorized").astype(str).value_counts()
        if not category_counts.empty:
            top_category = str(category_counts.index[0])
            top_count = int(category_counts.iloc[0])
            insights.append({
                "Title": "Largest custom category",
                "Insight": f"{top_category} has {top_count:,} verbatims, making it the biggest category in this run.",
            })
        if "Sentiment" in analyzed.columns:
            negative_by_category = analyzed[analyzed["Sentiment"].astype(str) == "Negative"][category_col].fillna("Uncategorized").astype(str).value_counts()
            if not negative_by_category.empty:
                insights.append({
                    "Title": "Highest negative category",
                    "Insight": f"{negative_by_category.index[0]} has the most negative verbatims and should be reviewed for coaching or process action.",
                })
    return insights[:6]


def _emotion_for_text(text: Any, sentiment: str = "") -> str:
    lowered = str(text or "").lower()
    emotion_keywords = {
        "Frustration": ["frustrat", "annoy", "angry", "upset", "irritat", "fed up", "terrible", "worst"],
        "Urgency": ["urgent", "asap", "immediately", "critical", "escalat", "deadline", "emergency"],
        "Confusion": ["confus", "unclear", "not sure", "don't understand", "no idea", "misleading"],
        "Effort": ["again", "multiple", "repeat", "repeated", "chase", "follow up", "call back", "long time"],
        "Trust Risk": ["cancel", "complaint", "legal", "refund", "promise", "lied", "trust", "unacceptable"],
        "Appreciation": ["thank", "thanks", "appreciate", "helpful", "great", "excellent", "amazing", "polite"],
    }
    scores = {
        emotion: sum(lowered.count(keyword) for keyword in keywords)
        for emotion, keywords in emotion_keywords.items()
    }
    best, score = max(scores.items(), key=lambda item: item[1])
    if score > 0:
        return best
    if str(sentiment).lower() == "positive":
        return "Appreciation"
    if str(sentiment).lower() == "negative":
        return "Frustration"
    return "Neutral"


def _top_terms_for_frame(df: pd.DataFrame, limit: int = 5) -> list[tuple[str, int]]:
    if df.empty or "Verbatim Feedback" not in df.columns:
        return []
    stop = {
        "the", "and", "for", "with", "that", "this", "was", "were", "are", "but", "not", "you", "your", "our",
        "they", "them", "have", "has", "had", "from", "get", "got", "can", "could", "would", "should", "there",
        "their", "about", "into", "just", "very", "more", "less", "than", "then", "when", "what", "why", "how",
        "because", "been", "being", "also", "still", "after", "before", "again",
        "customer", "service", "agent", "issue", "problem", "call", "chat", "email", "support", "team",
    }
    counts: dict[str, int] = {}
    for text in df["Verbatim Feedback"].fillna("").astype(str):
        for token in _word_tokens(text):
            if len(token) < 4 or token in stop:
                continue
            counts[token] = counts.get(token, 0) + 1
    return sorted(counts.items(), key=lambda item: item[1], reverse=True)[:limit]


def _build_holistic_sentiment_intelligence(analyzed: pd.DataFrame, category_col: str = "Custom Category") -> list[dict[str, Any]]:
    if analyzed.empty:
        return []
    total = len(analyzed)
    working = analyzed.copy()
    working["__sentiment"] = working.get("Sentiment", pd.Series(["Unknown"] * total)).fillna("Unknown").astype(str)
    working["__emotion"] = [
        _emotion_for_text(text, sentiment)
        for text, sentiment in zip(
            working.get("Verbatim Feedback", pd.Series([""] * total)).fillna("").astype(str),
            working["__sentiment"],
        )
    ]
    emotion_counts = working["__emotion"].value_counts()
    top_emotions = emotion_counts.head(3)
    negative_df = working[working["__sentiment"].str.lower().eq("negative")]
    positive_df = working[working["__sentiment"].str.lower().eq("positive")]
    neutral_df = working[working["__sentiment"].str.lower().eq("neutral")]
    category_counts = working[category_col].fillna("Uncategorized").astype(str).value_counts() if category_col in working.columns else pd.Series(dtype=int)
    negative_categories = negative_df[category_col].fillna("Uncategorized").astype(str).value_counts() if category_col in negative_df.columns else pd.Series(dtype=int)
    top_negative_terms = _top_terms_for_frame(negative_df, 4)
    top_all_terms = _top_terms_for_frame(working, 6)

    def pct(count: int) -> str:
        return f"{round((count / max(total, 1)) * 100, 1)}%"

    cards: list[dict[str, Any]] = []
    cards.append({
        "Title": "Emotion Mix",
        "Metric": ", ".join(f"{emotion} {pct(int(count))}" for emotion, count in top_emotions.items()) or "No emotion signal",
        "Insight": "Holistic emotion mix across all verbatims, inferred from sentiment plus recurring language patterns.",
        "Evidence": [f"{emotion}: {int(count):,} verbatims" for emotion, count in top_emotions.items()],
    })

    if not negative_categories.empty:
        top_driver = str(negative_categories.index[0])
        driver_count = int(negative_categories.iloc[0])
        term_text = ", ".join(term for term, _ in top_negative_terms) or "negative language"
        cards.append({
            "Title": "Sentiment Driver",
            "Metric": f"{top_driver} drives {driver_count:,} negative verbatims",
            "Insight": f"The strongest negative driver is {top_driver}. Common language around this driver includes {term_text}.",
            "Evidence": [f"{label}: {int(count):,} negative verbatims" for label, count in negative_categories.head(4).items()],
        })
    else:
        cards.append({
            "Title": "Sentiment Driver",
            "Metric": "No dominant negative driver",
            "Insight": "No single category is currently concentrating negative sentiment.",
            "Evidence": [],
        })

    pos_count, neg_count, neu_count = len(positive_df), len(negative_df), len(neutral_df)
    leading_category = str(category_counts.index[0]) if not category_counts.empty else "Uncategorized"
    cards.append({
        "Title": "Verbatim Intelligence Summary",
        "Metric": f"{pct(pos_count)} positive / {pct(neg_count)} negative",
        "Insight": f"Across {total:,} verbatims, the dominant discussion area is {leading_category}. The overall tone is shaped by {pos_count:,} positive, {neu_count:,} neutral, and {neg_count:,} negative comments.",
        "Evidence": [f"{label}: {int(count):,} verbatims" for label, count in category_counts.head(4).items()],
    })

    common_theme_labels = [str(label) for label in category_counts.head(5).index] if not category_counts.empty else []
    common_terms = [term for term, _ in top_all_terms[:5]]
    theme_metric = ", ".join(common_theme_labels[:3]) if common_theme_labels else "No common themes yet"
    theme_language = ", ".join(common_terms) if common_terms else "no repeated language yet"
    cards.append({
        "Title": "Common Customer Themes",
        "Metric": theme_metric,
        "Insight": f"Customers are most often talking about {theme_metric}. The repeated language across comments includes {theme_language}.",
        "Evidence": [f"{label}: {int(count):,} verbatims" for label, count in category_counts.head(5).items()],
    })

    if top_all_terms:
        emerging_terms = [term for term, count in top_all_terms if count >= max(2, int(total * 0.03))][:4]
    else:
        emerging_terms = []
    emerging_label = ", ".join(emerging_terms) if emerging_terms else (str(category_counts.index[0]) if not category_counts.empty else "No emerging issue")
    cards.append({
        "Title": "Emerging Issues",
        "Metric": emerging_label,
        "Insight": "Emerging issues are identified from recurring language across the full verbatim set. Add a date column later to convert this into a true time-based trend.",
        "Evidence": [f"{term}: {count:,} mentions" for term, count in top_all_terms[:5]],
    })

    repeat_source = negative_categories if not negative_categories.empty else category_counts
    repeat_label = str(repeat_source.index[0]) if not repeat_source.empty else "No repeat pain point"
    repeat_count = int(repeat_source.iloc[0]) if not repeat_source.empty else 0
    cards.append({
        "Title": "Repeat Pain Points",
        "Metric": f"{repeat_label} appears {repeat_count:,} times",
        "Insight": f"{repeat_label} is the most repeated pain point across the uploaded verbatims. Prioritize it when building coaching, process, or knowledge-base actions.",
        "Evidence": [f"{label}: {int(count):,} verbatims" for label, count in repeat_source.head(5).items()],
    })
    return cards


def _generic_dimension_rows(df: pd.DataFrame, column: str, mode: str) -> list[dict[str, Any]]:
    if df.empty or not column or column not in df.columns:
        return []
    rows: list[dict[str, Any]] = []
    for value, group in df.groupby(df[column].fillna("Unknown").astype(str), dropna=False):
        total = len(group)
        if mode == "csat":
            satisfied = int((group["CSAT Segment"] == "Satisfied").sum()) if "CSAT Segment" in group.columns else 0
            neutral = int((group["CSAT Segment"] == "Neutral").sum()) if "CSAT Segment" in group.columns else 0
            dissatisfied = int((group["CSAT Segment"] == "Dissatisfied").sum()) if "CSAT Segment" in group.columns else 0
            score_values = pd.to_numeric(group.get("CSAT Score", pd.Series(dtype=float)), errors="coerce").dropna()
            rows.append({
                "Segment": value,
                "Responses": total,
                "Satisfied": satisfied,
                "Neutral": neutral,
                "Dissatisfied": dissatisfied,
                "CSAT %": round((satisfied / max(total, 1)) * 100, 1),
                "Avg CSAT": round(float(score_values.mean()), 2) if not score_values.empty else 0,
                "Negative Sentiment %": round((int((group.get("Sentiment", pd.Series(dtype=str)) == "Negative").sum()) / max(total, 1)) * 100, 1),
            })
        else:
            positive = int((group.get("Sentiment", pd.Series(dtype=str)) == "Positive").sum())
            neutral = int((group.get("Sentiment", pd.Series(dtype=str)) == "Neutral").sum())
            negative = int((group.get("Sentiment", pd.Series(dtype=str)) == "Negative").sum())
            rows.append({
                "Segment": value,
                "Responses": total,
                "Positive": positive,
                "Neutral": neutral,
                "Negative": negative,
                "Positive %": round((positive / max(total, 1)) * 100, 1),
                "Negative %": round((negative / max(total, 1)) * 100, 1),
            })
    return sorted(rows, key=lambda row: (row.get("CSAT %", row.get("Positive %", 0)), row["Responses"]), reverse=True)


def _generic_analysis_payload(payload: dict[str, Any]) -> dict[str, Any]:
    mode = str(payload.get("mode") or "csat").strip().lower()
    calendar_settings = _calendar_settings(payload)
    analysis_id = str(payload.get("analysisId") or "").strip()
    mapping = payload.get("mapping") if isinstance(payload.get("mapping"), dict) else {}
    engines = payload.get("engines") if isinstance(payload.get("engines"), dict) else {}
    model_paths = payload.get("modelPaths") if isinstance(payload.get("modelPaths"), dict) else {}
    dimensions = [str(item).strip() for item in payload.get("dimensions", []) if str(item).strip()]
    custom_categories = [str(item).strip() for item in payload.get("customCategories", []) if str(item).strip()]
    df = _decode_tabular(payload)
    feedback = str(mapping.get("feedback") or "").strip()
    score = str(mapping.get("score") or "").strip()
    agent = str(mapping.get("agent") or "").strip()
    manager = str(mapping.get("manager") or "").strip()
    lob = str(mapping.get("lob") or "").strip()
    date = str(mapping.get("date") or "").strip()
    if not feedback or feedback not in df.columns:
        raise ValueError("Map a valid feedback/comment column before analysis.")
    if mode == "csat" and (not score or score not in df.columns):
        raise ValueError("Map a valid CSAT score/rating column before CSAT analysis.")

    total_rows = int(len(df))

    def set_module_progress(
        percent: float,
        message: str,
        done: int = 0,
        total: int | None = None,
        current_row: int | None = None,
    ) -> None:
        if not analysis_id:
            return
        row_total = int(total if total is not None else total_rows)
        row_done = max(0, min(int(done or 0), max(row_total, 0)))
        with MODULE_PROGRESS_LOCK:
            MODULE_PROGRESS[analysis_id] = {
                "ok": True,
                "id": analysis_id,
                "mode": mode,
                "percent": round(max(0.0, min(float(percent), 100.0)), 1),
                "message": message,
                "done": row_done,
                "total": row_total,
                "currentRow": int(current_row if current_row is not None else row_done),
                "complete": False,
                "updatedAt": time.time(),
            }

    set_module_progress(1, f"Validating upload and mapping for {total_rows:,} rows...", 0)

    sentiment_engine = str(engines.get("sentiment") or "sparrow").strip().lower()
    theme_engine = str(engines.get("theme") or "local").strip().lower()
    sparrow_model_path = str(model_paths.get("sparrow") or _default_model_path("sparrow")).strip()
    theme_model_path = str(model_paths.get("theme") or model_paths.get("owl") or _default_model_path("theme")).strip()

    analysis_score = score if mode == "csat" else None
    if sentiment_engine in {"local", "rules", "local rules", "openai", "open ai", "openai api", "claude", "claude api"}:
        def local_progress(done: int, total: int, message: str | None = None) -> None:
            pct = 5 + (min(max(done, 0), max(total, 1)) / max(total, 1)) * 35
            estimated_row = int((pct / 100) * total_rows)
            set_module_progress(pct, message or f"Preparing local sentiment step {done}/{total}...", estimated_row)

        analyzed = build_analysis(df, feedback, analysis_score, agent or None, date or None, progress_callback=local_progress)
        if sentiment_engine in {"openai", "open ai", "openai api"}:
            analyzed["Analysis Source"] = "Local Rules: OpenAI API connector placeholder"
        elif sentiment_engine in {"claude", "claude api"}:
            analyzed["Analysis Source"] = "Local Rules: Claude API connector placeholder"
        else:
            analyzed["Analysis Source"] = "Local Rules"
    else:
        try:
            def sparrow_progress(done: int, total: int, message: str | None = None) -> None:
                if total and total == total_rows and done > 0:
                    pct = 18 + (min(done, total) / max(total, 1)) * 66
                    label = message or f"Sparrow is analyzing row {done:,} of {total:,}."
                    set_module_progress(pct, label, done, total, done)
                else:
                    pct = 5 + (min(max(done, 0), max(total, 1)) / max(total, 1)) * 12
                    set_module_progress(pct, message or "Preparing Sparrow sentiment engine...", 0, total_rows, 0)

            analyzed = build_analysis_with_local_model(
                df,
                feedback,
                analysis_score,
                agent or None,
                date or None,
                model_path=sparrow_model_path,
                progress_callback=sparrow_progress,
            )
        except Exception as exc:
            set_module_progress(18, "Sparrow unavailable. Continuing with local sentiment rules...", 0)

            def fallback_progress(done: int, total: int, message: str | None = None) -> None:
                pct = 18 + (min(max(done, 0), max(total, 1)) / max(total, 1)) * 50
                estimated_row = int((pct / 100) * total_rows)
                set_module_progress(pct, message or f"Local fallback step {done}/{total}...", estimated_row)

            analyzed = build_analysis(df, feedback, analysis_score, agent or None, date or None, progress_callback=fallback_progress)
            analyzed["Analysis Source"] = f"Local Fallback: Sparrow unavailable: {exc}"

    set_module_progress(86, "Mapping selected fields into the analyzed output...", total_rows)

    if manager and manager in df.columns:
        analyzed["Manager/TL"] = df[manager].reset_index(drop=True).reindex(range(len(analyzed))).fillna("Unknown").astype(str).to_numpy()
    if lob and lob in df.columns:
        analyzed["LOB"] = df[lob].reset_index(drop=True).reindex(range(len(analyzed))).fillna("Unknown").astype(str).to_numpy()
    if mode == "csat":
        analyzed["CSAT Score"] = pd.to_numeric(df[score].reset_index(drop=True).reindex(range(len(analyzed))), errors="coerce")
        analyzed["CSAT Segment"] = df[score].reset_index(drop=True).reindex(range(len(analyzed))).apply(_csat_label).astype(str)
    analyzed = _neutralize_blank_feedback_sentiment(analyzed)
    analyzed = _apply_reporting_calendar(analyzed, calendar_settings)
    if mode == "sentiment":
        if not custom_categories:
            custom_categories = [
                "Wait Time", "Agent Behavior", "Resolution", "Product Knowledge", "Billing",
                "Technical Issue", "Communication", "Process", "Follow-up", "Other",
            ]
        categories: list[str] = []
        category_confidence: list[float] = []
        feedback_values = analyzed["Verbatim Feedback"].fillna("").astype(str).tolist()
        for index, value in enumerate(feedback_values, start=1):
            category, confidence = _custom_category_for_text(value, custom_categories)
            categories.append(category)
            category_confidence.append(confidence)
            if index == 1 or index == len(feedback_values) or index % 10 == 0:
                pct = 86 + (index / max(len(feedback_values), 1)) * 8
                set_module_progress(pct, f"Categorizing row {index:,} of {len(feedback_values):,}.", index, len(feedback_values), index)
        analyzed["Custom Category"] = categories
        analyzed["Category Confidence"] = category_confidence
    if mode != "sentiment":
        if theme_engine in {"local", "rules", "local rules"}:
            set_module_progress(88, "Applying local theme defaults because Trained Theme Model was not selected.", total_rows, total_rows, total_rows)
            analyzed = _fill_owl_fallback_columns(analyzed, "User selected Local Rules for theme classification.")
        else:
            try:
                set_module_progress(88, "Loading Trained Theme Model model and tokenizer...", 0, total_rows, 0)

                def owl_progress(done: int, total: int, message: str | None = None) -> None:
                    pct = 88 + (min(max(done, 0), max(total, 1)) / max(total, 1)) * 7
                    label = message or f"Trained Theme Model is processing row {min(done, total):,} of {total:,}."
                    set_module_progress(pct, label, done, total, done)

                analyzed = add_theme_acpt_resolution_outputs(
                    analyzed,
                    feedback_col="Verbatim Feedback",
                    model_path=theme_model_path,
                    progress_callback=owl_progress,
                )
            except Exception as exc:
                set_module_progress(94, "Trained Theme Model unavailable. Filling theme fields with safe defaults...", total_rows, total_rows, total_rows)
                analyzed = _fill_owl_fallback_columns(analyzed, str(exc))

    def module_acpt_progress(done: int, total: int, message: str | None = None) -> None:
        pct = 94 + (min(max(done, 0), max(total, 1)) / max(total, 1)) * 2
        set_module_progress(pct, message or f"ACPT classification: {done:,}/{total:,} rows. Assigning Agent, Customer, Process, or Technology ownership.", done, total, done)

    trained_theme_supplied_acpt = (
        theme_engine not in {"local", "rules", "local rules"}
        and "ACPT Primary Category" in analyzed.columns
        and analyzed["ACPT Primary Category"].fillna("").astype(str).str.strip().ne("").any()
    )
    if trained_theme_supplied_acpt:
        set_module_progress(96, "Trained Theme Model supplied ACPT output. Building dashboard cards, word cloud, and export-ready tables...", total_rows)
    else:
        set_module_progress(94, "Starting ACPT verbatim classification: assigning Agent, Customer, Process, or Technology ownership.", 0, total_rows, 0)
        analyzed = _add_acpt_classification_outputs(analyzed, "Verbatim Feedback", module_acpt_progress)
        set_module_progress(96, "ACPT classification is complete. Building dashboard cards, word cloud, and export-ready tables...", total_rows)
    total = len(analyzed)
    sentiment = sentiment_summary(analyzed) if total else {"Positive": 0, "Neutral": 0, "Negative": 0}
    summary: dict[str, Any] = {"total": total, "sentiment": sentiment}
    if mode == "csat":
        counts = analyzed["CSAT Segment"].value_counts() if "CSAT Segment" in analyzed.columns else pd.Series(dtype=int)
        summary.update({
            "csat": round((int(counts.get("Satisfied", 0)) / max(total, 1)) * 100, 1),
            "satisfied": int(counts.get("Satisfied", 0)),
            "neutral": int(counts.get("Neutral", 0)),
            "dissatisfied": int(counts.get("Dissatisfied", 0)),
            "avgScore": round(float(pd.to_numeric(analyzed.get("CSAT Score", pd.Series(dtype=float)), errors="coerce").mean()), 2) if total else 0,
        })
    else:
        summary.update({
            "positive": sentiment.get("Positive", 0),
            "neutral": sentiment.get("Neutral", 0),
            "negative": sentiment.get("Negative", 0),
        })

    dimension_payload = []
    for column in dimensions:
        if column in analyzed.columns:
            dimension_payload.append({"name": column, "rows": _generic_dimension_rows(analyzed, column, mode)})
        elif column in df.columns:
            analyzed[column] = df[column].reset_index(drop=True).reindex(range(len(analyzed))).fillna("Unknown").astype(str).to_numpy()
            dimension_payload.append({"name": column, "rows": _generic_dimension_rows(analyzed, column, mode)})

    reason_cols = ["Custom Category"] if mode == "sentiment" else ["Owl Primary Driver", "Primary Reason", "Owl Secondary Driver", "Owl Issue Type"]
    theme_rows = []
    for column in reason_cols:
        if column in analyzed.columns:
            counts = analyzed[column].fillna("").astype(str)
            for label, count in counts[counts != ""].value_counts().head(12).items():
                theme_rows.append({"Field": column, "Theme": label, "Count": int(count)})

    if mode == "sentiment":
        visible_columns = [
            "Verbatim Feedback", "Sentiment", "Sentiment Score", "Custom Category", "Category Confidence",
            "Agent Name", "Manager/TL", "Feedback Date", "Analysis Source", "AI Rationale",
        ]
    else:
        visible_columns = [
            "Verbatim Feedback", "Custom Category", "Category Confidence", "CSAT Score", "CSAT Segment", "Sentiment", "Sentiment Score", "NPS Type", "NPS Score",
            "Agent Name", "Manager/TL", "LOB", "Feedback Date", "Owl Primary Driver", "Owl Secondary Driver",
            "Owl Issue Type", "Owl Customer Impact", "Analysis Source", "AI Rationale",
        ]
    result = {
        "ok": True,
        "mode": mode,
        "columns": list(df.columns),
        "guesses": _generic_guess_columns(list(df.columns), mode),
        "summary": summary,
        "dimensions": dimension_payload,
        "themes": theme_rows[:50],
        "insights": _build_verbatim_insights(analyzed) if mode == "sentiment" else [],
        "intelligenceCards": _build_holistic_sentiment_intelligence(analyzed) if mode == "sentiment" else [],
        "rows": _records_for_columns(analyzed, visible_columns, 1000),
    }
    set_module_progress(100, "Analysis complete.", total_rows, total_rows, total_rows)
    if analysis_id:
        with MODULE_PROGRESS_LOCK:
            if analysis_id in MODULE_PROGRESS:
                MODULE_PROGRESS[analysis_id]["complete"] = True
    return result


def _safe_records(df: pd.DataFrame, limit: int | None = 200) -> list[dict[str, Any]]:
    if df.empty:
        return []
    clean = df.copy() if limit is None else df.head(limit).copy()
    if "__row_id" not in clean.columns:
        clean.insert(0, "__row_id", clean.index.astype(str))
    clean = clean.where(pd.notna(clean), "")
    return clean.to_dict(orient="records")


def _records_for_columns(df: pd.DataFrame, columns: list[str], limit: int | None = 500) -> list[dict[str, Any]]:
    if df.empty:
        return []
    available = [column for column in columns if column in df.columns]
    if not available:
        return _safe_records(df, limit)
    source = df if limit is None else df.head(limit)
    clean = source[available].copy()
    if "__row_id" not in clean.columns:
        clean.insert(0, "__row_id", source.index.astype(str))
    clean = clean.where(pd.notna(clean), "")
    return clean.to_dict(orient="records")


def _stats_records(df: pd.DataFrame, limit: int | None = None) -> list[dict[str, Any]]:
    if df.empty:
        return []
    clean = df.copy() if limit is None else df.head(limit).copy()
    if "__row_id" not in clean.columns:
        clean.insert(0, "__row_id", clean.index.astype(str))
    clean = clean.where(pd.notna(clean), "")
    return clean.to_dict(orient="records")


def _statistics_source(name: str) -> pd.DataFrame:
    with STATE_LOCK:
        sources = {
            "Base Sheet": STATE.base_df.copy(),
            "Lookup Sheet": STATE.lookup_df.copy(),
            "Analyzed Feedback Rows": STATE.analyzed_df.copy(),
            "Theme Rows": STATE.analyzed_df.copy(),
        }
    frame = sources.get(name, pd.DataFrame())
    return _apply_date_filter(frame) if name in {"Analyzed Feedback Rows", "Theme Rows"} else frame


def _stat_number(value: Any) -> str:
    if value is None or pd.isna(value):
        return "0.0"
    return f"{float(value):.1f}"


def _stats_numeric_series(series: pd.Series) -> pd.Series:
    values = pd.to_numeric(series, errors="coerce").dropna()
    if values.empty:
        return values.astype(float)
    return values.astype(float)


def _numeric_pairs(frame: pd.DataFrame, first: str, second: str) -> pd.DataFrame:
    if first not in frame.columns or second not in frame.columns:
        return pd.DataFrame(columns=[first, second])
    pairs = pd.DataFrame({first: _stats_numeric_series(frame[first]), second: _stats_numeric_series(frame[second])})
    return pairs.dropna()


def _custom_statistics(payload: dict[str, Any]) -> dict[str, Any]:
    source = str(payload.get("source") or "Analyzed Feedback Rows")
    frame = _statistics_source(source)
    analysis_type = str(payload.get("type") or "desc")
    primary = str(payload.get("primary") or "")
    secondary = str(payload.get("secondary") or "")
    category = str(payload.get("category") or "")
    columns = [str(item) for item in payload.get("columns", []) if str(item) in frame.columns][:10]
    output: list[dict[str, Any]] = []

    numeric_columns = [str(column) for column in frame.columns if pd.to_numeric(frame[column], errors="coerce").notna().any()]
    selected = columns or numeric_columns[:8]
    if analysis_type in {"desc", "outlier", "volatility"}:
        for column in selected:
            values = _stats_numeric_series(frame[column])
            if values.empty:
                continue
            q1, median_value, q3 = values.quantile([0.25, 0.5, 0.75]).tolist()
            std = float(values.std(ddof=0))
            row: dict[str, Any] = {
                "Column": column, "Count": int(len(values)), "Mean": _stat_number(values.mean()),
                "Median": _stat_number(median_value), "Std Dev": _stat_number(std),
                "Min": _stat_number(values.min()), "Max": _stat_number(values.max()),
                "Q1": _stat_number(q1), "Q3": _stat_number(q3), "IQR": _stat_number(q3 - q1),
            }
            if analysis_type == "outlier":
                low, high = q1 - 1.5 * (q3 - q1), q3 + 1.5 * (q3 - q1)
                outliers = int(((values < low) | (values > high)).sum())
                row = {"Column": column, "Count": int(len(values)), "Lower Fence": _stat_number(low), "Upper Fence": _stat_number(high), "Outliers": outliers, "Outlier %": f"{outliers / max(len(values), 1) * 100:.1f}%"}
            elif analysis_type == "volatility":
                row["CV %"] = f"{std / max(abs(float(values.mean())), 1) * 100:.1f}%"
            output.append(row)
    elif analysis_type in {"mode", "pareto"} and category in frame.columns:
        values = frame[category].fillna("").astype(str).str.strip().replace("", "Blank")
        counts = values.value_counts(dropna=False).head(30)
        output = [{"Value": str(value), "Count": int(count), "Share": f"{count / max(len(frame), 1) * 100:.1f}%"} for value, count in counts.items()]
    elif analysis_type in {"pearson", "spearman", "scatter", "gap"}:
        pairs = _numeric_pairs(frame, primary, secondary)
        correlation_value = pairs[primary].corr(pairs[secondary], method="spearman" if analysis_type == "spearman" else "pearson") if len(pairs) > 1 else None
        output = [{"Metric A": primary, "Metric B": secondary, "Correlation": f"{correlation_value:.4f}" if correlation_value is not None and not pd.isna(correlation_value) else "n/a", "Sample": int(len(pairs)), "Mean A": _stat_number(pairs[primary].mean() if len(pairs) else None), "Mean B": _stat_number(pairs[secondary].mean() if len(pairs) else None), "Gap": _stat_number((pairs[primary] - pairs[secondary]).mean() if len(pairs) else None)}]
    elif analysis_type in {"benchmark", "segment"} and category in frame.columns and primary in frame.columns:
        working = pd.DataFrame({"Group": frame[category].fillna("Unknown").astype(str), "Value": pd.to_numeric(frame[primary], errors="coerce")}).dropna(subset=["Value"])
        overall = working["Value"].mean()
        for group, values in working.groupby("Group", sort=False)["Value"]:
            row = {"Group": str(group), "Responses": int(len(values)), "Average": _stat_number(values.mean()), "Median": _stat_number(values.median()), "Std Dev": _stat_number(values.std(ddof=0))}
            if analysis_type == "benchmark":
                row = {"Group": str(group), "Responses": int(len(values)), "Average": _stat_number(values.mean()), "Vs Benchmark": _stat_number(values.mean() - overall), "Benchmark": _stat_number(overall)}
            output.append(row)
        output.sort(key=lambda row: float(row.get("Average", 0)), reverse=True)
    elif analysis_type == "corrmatrix":
        numeric = frame[selected].apply(lambda column: _stats_numeric_series(column), axis=0) if selected else pd.DataFrame()
        matrix = numeric.corr()
        output = [{"Metric": column, **{other: (f"{matrix.loc[column, other]:.2f}" if not pd.isna(matrix.loc[column, other]) else "n/a") for other in selected}} for column in selected]
    elif analysis_type == "driver" and primary in frame.columns:
        for column in selected:
            if column == primary:
                continue
            pairs = _numeric_pairs(frame, primary, column)
            value = pairs[primary].corr(pairs[column]) if len(pairs) > 1 else None
            strength = "Not enough data" if value is None or pd.isna(value) else ("Strong" if abs(value) >= 0.7 else "Moderate" if abs(value) >= 0.4 else "Weak" if abs(value) >= 0.2 else "Very weak")
            output.append({"Predictor": column, "Target": primary, "Correlation": f"{value:.3f}" if value is not None and not pd.isna(value) else "n/a", "Sample": int(len(pairs)), "Strength": strength, "Direction": "n/a" if value is None or pd.isna(value) else ("Positive" if value >= 0 else "Negative")})
        output.sort(key=lambda row: abs(float(row["Correlation"])) if row["Correlation"] != "n/a" else -1, reverse=True)

    chart_rows: list[dict[str, Any]] = []
    if primary and secondary:
        pairs = _numeric_pairs(frame, primary, secondary)
        if len(pairs) > 500:
            pairs = pairs.sample(500, random_state=42)
        chart_rows = pairs.where(pd.notna(pairs), "").to_dict(orient="records")
    return {"ok": True, "source": source, "totalRows": int(len(frame)), "output": output, "chartRows": chart_rows}


def _paginated_statistics_rows(query: dict[str, list[str]]) -> dict[str, Any]:
    source = str(query.get("source", ["Analyzed Feedback Rows"])[0])
    frame = _statistics_source(source)
    search = str(query.get("search", [""])[0]).strip().lower()
    if search and not frame.empty:
        mask = frame.astype(str).apply(lambda column: column.str.lower().str.contains(search, regex=False, na=False)).any(axis=1)
        frame = frame[mask]
    page_size = max(25, min(200, int(query.get("pageSize", ["100"])[0] or 100)))
    page = max(1, int(query.get("page", ["1"])[0] or 1))
    start = (page - 1) * page_size
    return {"ok": True, "source": source, "page": page, "pageSize": page_size, "totalRows": int(len(frame)), "rows": _safe_records(frame.iloc[start:start + page_size], None)}


def _apply_date_filter(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty or "Feedback Date" not in df.columns:
        return df.copy()
    with STATE_LOCK:
        date_filter = dict(STATE.date_filter)
    mode = date_filter.get("mode", "All Time")
    if mode == "All Time":
        return df.copy()
    working = df.copy()
    dates = pd.to_datetime(working["Feedback Date"], errors="coerce")
    valid_dates = dates.dropna()
    if valid_dates.empty:
        return working.iloc[0:0].copy()
    latest = valid_dates.max().normalize()
    start = None
    end = latest
    if mode == "Today":
        start = latest
    elif mode == "Yesterday":
        start = latest - pd.Timedelta(days=1)
        end = start
    elif mode == "7D":
        start = latest - pd.Timedelta(days=6)
    elif mode == "This Month":
        start = latest.replace(day=1)
    elif mode == "Last 90D":
        start = latest - pd.Timedelta(days=89)
    elif mode == "Custom Range":
        start_text = date_filter.get("start") or ""
        end_text = date_filter.get("end") or ""
        start = pd.to_datetime(start_text, errors="coerce") if start_text else None
        end = pd.to_datetime(end_text, errors="coerce") if end_text else latest
    if start is None:
        return working.copy()
    start = pd.Timestamp(start).normalize()
    end = pd.Timestamp(end).normalize()
    mask = (dates.dt.normalize() >= start) & (dates.dt.normalize() <= end)
    return working[mask.fillna(False)].copy()


def _filter_explicit_date_range(df: pd.DataFrame, start_text: str, end_text: str) -> pd.DataFrame:
    if df.empty or "Feedback Date" not in df.columns:
        return df.copy()
    dates = pd.to_datetime(df["Feedback Date"], errors="coerce")
    valid_dates = dates.dropna()
    if valid_dates.empty:
        return df.iloc[0:0].copy()
    start = pd.to_datetime(start_text, errors="coerce") if start_text else valid_dates.min()
    end = pd.to_datetime(end_text, errors="coerce") if end_text else valid_dates.max()
    if pd.isna(start):
        start = valid_dates.min()
    if pd.isna(end):
        end = valid_dates.max()
    start = pd.Timestamp(start).normalize()
    end = pd.Timestamp(end).normalize()
    if start > end:
        start, end = end, start
    mask = (dates.dt.normalize() >= start) & (dates.dt.normalize() <= end)
    return df[mask.fillna(False)].copy()


def _range_metric_row(label: str, df: pd.DataFrame) -> dict[str, Any]:
    summary = nps_summary(df) if not df.empty else {"total": 0, "nps": 0, "promoters": 0, "passives": 0, "detractors": 0}
    counts = nps_composition_counts(df) if not df.empty else {"Promoter": 0, "Passive": 0, "Detractor": 0}
    sentiment = sentiment_summary(df) if not df.empty else {"Positive": 0, "Neutral": 0, "Negative": 0}
    return {
        "Range": label,
        "Responses": int(summary.get("total", 0)),
        "NPS": float(summary.get("nps", 0)),
        "Promoters": int(counts.get("Promoter", 0)),
        "Passives": int(counts.get("Passive", 0)),
        "Detractors": int(counts.get("Detractor", 0)),
        "Positive %": float(sentiment.get("Positive", 0)),
        "Neutral %": float(sentiment.get("Neutral", 0)),
        "Negative %": float(sentiment.get("Negative", 0)),
    }


def _comparison_payload(payload: dict[str, Any]) -> dict[str, Any]:
    with STATE_LOCK:
        raw_df = STATE.analyzed_df.copy()
    first_start = str(payload.get("currentStart") or "")
    first_end = str(payload.get("currentEnd") or "")
    second_start = str(payload.get("previousStart") or "")
    second_end = str(payload.get("previousEnd") or "")
    first = _filter_explicit_date_range(raw_df, first_start, first_end)
    second = _filter_explicit_date_range(raw_df, second_start, second_end)

    def parsed_date(value: str) -> pd.Timestamp:
        parsed = pd.to_datetime(value, errors="coerce")
        return parsed if not pd.isna(parsed) else pd.Timestamp.max

    def range_days(start: str, end: str) -> int:
        start_dt = parsed_date(start)
        end_dt = parsed_date(end)
        if start_dt == pd.Timestamp.max or end_dt == pd.Timestamp.max:
            return 999
        return abs((end_dt.normalize() - start_dt.normalize()).days) + 1

    weekly_labels = max(range_days(first_start, first_end), range_days(second_start, second_end)) <= 7
    older_label = "Previous Week" if weekly_labels else "Previous Period"
    newer_label = "This Week" if weekly_labels else "Current Period"
    if parsed_date(first_start) <= parsed_date(second_start):
        previous_row = _range_metric_row(older_label, first)
        current_row = _range_metric_row(newer_label, second)
    else:
        previous_row = _range_metric_row(older_label, second)
        current_row = _range_metric_row(newer_label, first)
    delta_row = {
        "Range": "Change",
        "Responses": current_row["Responses"] - previous_row["Responses"],
        "NPS": round(current_row["NPS"] - previous_row["NPS"], 1),
        "Promoters": current_row["Promoters"] - previous_row["Promoters"],
        "Passives": current_row["Passives"] - previous_row["Passives"],
        "Detractors": current_row["Detractors"] - previous_row["Detractors"],
        "Positive %": round(current_row["Positive %"] - previous_row["Positive %"], 1),
        "Neutral %": round(current_row["Neutral %"] - previous_row["Neutral %"], 1),
        "Negative %": round(current_row["Negative %"] - previous_row["Negative %"], 1),
    }

    def movement_rows(group_cols: list[str], limit: int = 50) -> list[dict[str, Any]]:
        usable_cols = [col for col in group_cols if col and col in raw_df.columns]
        if not usable_cols:
            return []

        def grouped(frame: pd.DataFrame, prefix: str) -> pd.DataFrame:
            if frame.empty:
                return pd.DataFrame(columns=usable_cols + [f"{prefix} Responses", f"{prefix} NPS"])
            work = frame.copy()
            result = (
                work.groupby(usable_cols)
                .agg(
                    Responses=(usable_cols[0], "size"),
                    Promoters=("NPS Type", lambda values: int((values == "Promoter").sum())),
                    Passives=("NPS Type", lambda values: int((values == "Passive").sum())),
                    Detractors=("NPS Type", lambda values: int((values == "Detractor").sum())),
                    Avg_Rating=("NPS Score", "mean"),
                )
                .reset_index()
            )
            result["NPS"] = ((result["Promoters"] / result["Responses"] * 100) - (result["Detractors"] / result["Responses"] * 100)).round(1)
            result["Avg_Rating"] = result["Avg_Rating"].round(2)
            return result.rename(columns={
                "Responses": f"{prefix} Responses",
                "NPS": f"{prefix} NPS",
                "Promoters": f"{prefix} Promoters",
                "Passives": f"{prefix} Passives",
                "Detractors": f"{prefix} Detractors",
                "Avg_Rating": f"{prefix} Avg Rating",
            })

        previous_group = grouped(previous_frame, "Previous")
        current_group = grouped(current_frame, "Current")
        merged = previous_group.merge(current_group, on=usable_cols, how="outer").fillna(0)
        merged["NPS Change"] = (pd.to_numeric(merged.get("Current NPS", 0), errors="coerce") - pd.to_numeric(merged.get("Previous NPS", 0), errors="coerce")).round(1)
        merged["Response Change"] = (pd.to_numeric(merged.get("Current Responses", 0), errors="coerce") - pd.to_numeric(merged.get("Previous Responses", 0), errors="coerce")).astype(int)
        sort_col = "NPS Change"
        merged["Abs Change"] = pd.to_numeric(merged[sort_col], errors="coerce").abs()
        ordered = merged.sort_values(["Abs Change", "Current Responses"], ascending=[False, False]).drop(columns=["Abs Change"])
        return _safe_records(ordered, limit)

    manager_col = _first_existing_column(raw_df, ["Manager/TL", "TL Name", "Team Manager", "Manager Name", "Manager", "Supervisor"])
    agent_cols = ["Agent Name"] + ([manager_col] if manager_col else [])

    def consistency_row(label: str, frame: pd.DataFrame) -> dict[str, Any]:
        calendar_settings = _state_calendar_settings()
        if frame.empty or "Feedback Date" not in frame.columns:
            weekly = pd.DataFrame()
        else:
            work = frame.copy()
            work["Week"] = week_period_start(work["Feedback Date"], calendar_settings["weekStart"])
            work = work.dropna(subset=["Week"])
            if work.empty:
                weekly = pd.DataFrame()
            else:
                weekly = (
                    work.groupby("Week")
                    .agg(
                        Responses=("Week", "size"),
                        Promoters=("NPS Type", lambda values: int((values == "Promoter").sum())),
                        Detractors=("NPS Type", lambda values: int((values == "Detractor").sum())),
                    )
                    .reset_index()
                )
                weekly["NPS"] = ((weekly["Promoters"] / weekly["Responses"] * 100) - (weekly["Detractors"] / weekly["Responses"] * 100)).round(1)
        nps_values = pd.to_numeric(weekly.get("NPS", pd.Series(dtype=float)), errors="coerce").dropna()
        response_values = pd.to_numeric(weekly.get("Responses", pd.Series(dtype=float)), errors="coerce").dropna()
        avg_responses = float(response_values.mean()) if not response_values.empty else 0.0
        response_std = float(response_values.std(ddof=0)) if len(response_values) > 1 else 0.0
        return {
            "Range": label,
            "Weeks": int(len(weekly)),
            "NPS Volatility": round(float(nps_values.max() - nps_values.min()), 1) if len(nps_values) > 1 else 0,
            "Average Weekly Responses": round(avg_responses, 1),
            "Response Variability %": round((response_std / avg_responses) * 100, 1) if avg_responses else 0,
        }

    previous_frame = first if parsed_date(first_start) <= parsed_date(second_start) else second
    current_frame = second if parsed_date(first_start) <= parsed_date(second_start) else first
    consistency_previous = consistency_row(older_label, previous_frame)
    consistency_current = consistency_row(newer_label, current_frame)
    consistency_change = {
        "Range": "Change",
        "Weeks": consistency_current["Weeks"] - consistency_previous["Weeks"],
        "NPS Volatility": round(consistency_current["NPS Volatility"] - consistency_previous["NPS Volatility"], 1),
        "Average Weekly Responses": round(consistency_current["Average Weekly Responses"] - consistency_previous["Average Weekly Responses"], 1),
        "Response Variability %": round(consistency_current["Response Variability %"] - consistency_previous["Response Variability %"], 1),
    }
    manager_rows = movement_rows([manager_col], 50) if manager_col else []
    if manager_col and manager_col != "Manager/TL":
        for row in manager_rows:
            if manager_col in row and "Manager/TL" not in row:
                row["Manager/TL"] = row[manager_col]

    agent_rows = movement_rows(agent_cols, 100)
    if manager_col and manager_col != "Manager/TL":
        for row in agent_rows:
            if manager_col in row and "Manager/TL" not in row:
                row["Manager/TL"] = row[manager_col]

    return {
        "ok": True,
        "rows": [current_row, previous_row, delta_row],
        "managerRows": manager_rows,
        "agentRows": agent_rows,
        "sentimentRows": [
            {key: current_row[key] for key in ["Range", "Positive %", "Neutral %", "Negative %"]},
            {key: previous_row[key] for key in ["Range", "Positive %", "Neutral %", "Negative %"]},
            {key: delta_row[key] for key in ["Range", "Positive %", "Neutral %", "Negative %"]},
        ],
        "consistencyRows": [consistency_current, consistency_previous, consistency_change],
    }


def _guess_columns(columns: list[str]) -> dict[str, str]:
    lowered = {column.lower(): column for column in columns}

    def pick(*needles: str) -> str:
        for needle in needles:
            for lower, original in lowered.items():
                if needle in lower:
                    return original
        return ""

    return {
        "feedback": pick("comment", "verbatim", "feedback"),
        "score": pick("nps", "csat", "satisfaction", "rating", "score"),
        "agent": pick("agent"),
        "manager": pick("manager/tl", "tl name", "team manager", "manager", "supervisor"),
        "lob": pick("lob", "line of business", "business unit", "stream", "process", "queue"),
        "date": pick("date", "response"),
        "wave": pick("wave", "batch", "cohort"),
        "tenure": pick("tenure", "tenurity", "tenure bucket", "tenure range"),
    }


def _merge_lookup(base: pd.DataFrame, lookup: pd.DataFrame, base_key: str, lookup_key: str) -> pd.DataFrame:
    if base.empty or lookup.empty or not base_key or not lookup_key:
        return base
    if base_key not in base.columns or lookup_key not in lookup.columns:
        return base
    deduped_lookup = lookup.drop_duplicates(subset=[lookup_key], keep="first")
    return base.merge(deduped_lookup, left_on=base_key, right_on=lookup_key, how="left", suffixes=("", "_Lookup"))


def _alerts(analyzed_df: pd.DataFrame, weekly_df: pd.DataFrame) -> list[dict[str, str]]:
    if analyzed_df.empty:
        return []
    summary = nps_summary(analyzed_df)
    total = int(summary.get("total", 0))
    nps = float(summary.get("nps", 0))
    detractor_pct = float(summary.get("detractors", 0))
    sentiment = sentiment_summary(analyzed_df)
    negative_pct = float(sentiment.get("Negative", 0))
    positive_pct = float(sentiment.get("Positive", 0))
    alerts: list[dict[str, str]] = []

    def add(title: str, status: str, tone: str, detail: str) -> None:
        alerts.append({"title": title, "status": status, "tone": tone, "detail": detail})

    add(
        "NPS Health Watch",
        "Needs Attention" if nps < 0 or detractor_pct >= 35 else "Clear",
        "negative" if nps < 0 or detractor_pct >= 35 else "clear",
        f"NPS {nps:.1f}; detractors {detractor_pct:.1f}% across {total:,} responses.",
    )
    add(
        "Negative Superiority Alert",
        "Needs Attention" if negative_pct > positive_pct else "Clear",
        "negative" if negative_pct > positive_pct else "clear",
        f"Negative {negative_pct:.1f}% vs positive {positive_pct:.1f}%.",
    )
    add(
        "Promoter Strength",
        "Positive Alert" if nps >= 50 or float(summary.get("promoters", 0)) >= 60 else "Clear",
        "positive" if nps >= 50 or float(summary.get("promoters", 0)) >= 60 else "clear",
        f"Promoters {summary.get('promoters', 0):.1f}%; NPS {nps:.1f}.",
    )
    if len(weekly_df) >= 2:
        previous = float(weekly_df.iloc[-2].get("NPS", 0) or 0)
        current = float(weekly_df.iloc[-1].get("NPS", 0) or 0)
        add(
            "Week-over-Week NPS Movement",
            "Needs Attention" if current < previous else "Positive Alert",
            "negative" if current < previous else "positive",
            f"Previous week {previous:.1f}; current week {current:.1f}; movement {current - previous:+.1f} pts.",
        )
    if "Silent Detractor Alert" in analyzed_df.columns:
        silent = int(pd.to_numeric(analyzed_df["Silent Detractor Alert"], errors="coerce").fillna(0).sum())
        add(
            "Silent Detractor Register",
            "Needs Attention" if silent else "Clear",
            "negative" if silent else "clear",
            f"{silent:,} passive/soft-text rows show hidden detractor risk.",
        )
    if {"NPS Type", "Sentiment"}.issubset(analyzed_df.columns):
        negative_promoters = int(((analyzed_df["NPS Type"] == "Promoter") & (analyzed_df["Sentiment"] == "Negative")).sum())
        positive_detractors = int(((analyzed_df["NPS Type"] == "Detractor") & (analyzed_df["Sentiment"] == "Positive")).sum())
        add(
            "Negative Promoter Mismatch",
            "Review" if negative_promoters else "Clear",
            "negative" if negative_promoters else "clear",
            f"{negative_promoters:,} promoters contain negative sentiment and may need review.",
        )
        add(
            "Positive Detractor Mismatch",
            "Review" if positive_detractors else "Clear",
            "positive" if positive_detractors else "clear",
            f"{positive_detractors:,} detractors contain positive sentiment; check scoring/context mismatch.",
        )
    if "Owl Resolution Status" in analyzed_df.columns:
        unresolved = int(analyzed_df["Owl Resolution Status"].astype(str).isin(["Unresolved", "Partially Resolved"]).sum())
        add(
            "Unresolved Theme Risk",
            "Needs Attention" if unresolved else "Clear",
            "negative" if unresolved else "clear",
            f"{unresolved:,} rows have unresolved or partially resolved Owl resolution status.",
        )
    return alerts


def _case_id_column(df: pd.DataFrame) -> str | None:
    return _first_existing_column(df, ["Case ID", "Case Id", "CaseID", "Interaction ID", "Contact ID", "Survey ID", "Ticket ID"])


def _sentiment_confidence(value: Any) -> str:
    try:
        return f"{abs(float(value)) * 100:.1f}%"
    except Exception:
        return ""


def _segment_sentiment_counts(df: pd.DataFrame) -> dict[str, dict[str, int]]:
    result: dict[str, dict[str, int]] = {}
    for segment in ["Promoter", "Passive", "Detractor"]:
        if df.empty or "NPS Type" not in df.columns or "Sentiment" not in df.columns:
            counts = pd.Series(dtype=int)
        else:
            counts = df.loc[df["NPS Type"] == segment, "Sentiment"].value_counts()
        result[segment] = {
            "Positive": int(counts.get("Positive", 0)),
            "Neutral": int(counts.get("Neutral", 0)),
            "Negative": int(counts.get("Negative", 0)),
        }
    return result


def _first_numeric_value(row: pd.Series, candidates: list[str]) -> float:
    for candidate in candidates:
        if candidate in row.index:
            value = pd.to_numeric(pd.Series([row.get(candidate)]), errors="coerce").iloc[0]
            if pd.notna(value):
                return float(value)
    return 0.0


def _executive_intelligence_cards(
    analyzed_df: pd.DataFrame,
    weekly_df: pd.DataFrame,
    reason_df: pd.DataFrame,
) -> list[dict[str, str]]:
    if analyzed_df.empty:
        return []
    cards: list[dict[str, str]] = []
    sentiment = sentiment_summary(analyzed_df)

    def add(label: str, value: str, note: str, tone: str = "neutral") -> None:
        cards.append({"label": label, "value": value, "note": note, "tone": tone})

    if not weekly_df.empty and "NPS" in weekly_df.columns:
        weekly = weekly_df.copy()
        weekly["NPS"] = pd.to_numeric(weekly["NPS"], errors="coerce")
        weekly = weekly.dropna(subset=["NPS"])
        if len(weekly) >= 2:
            previous = float(weekly.iloc[-2]["NPS"])
            current = float(weekly.iloc[-1]["NPS"])
            movement = current - previous
            add(
                "WoW NPS",
                f"{movement:+.1f} pts",
                f"Latest {current:.1f} vs prior {previous:.1f}.",
                "positive" if movement >= 0 else "negative",
            )
            current_volume = _first_numeric_value(weekly.iloc[-1], ["Responses", "Total Responses", "Survey Count"])
            previous_volume = _first_numeric_value(weekly.iloc[-2], ["Responses", "Total Responses", "Survey Count"])
            if current_volume or previous_volume:
                add(
                    "WoW Volume",
                    f"{current_volume - previous_volume:+.0f}",
                    f"Latest {current_volume:.0f} vs prior {previous_volume:.0f}.",
                    "positive" if current_volume >= previous_volume else "neutral",
                )
        recent = weekly.tail(4)
        if not recent.empty:
            add("4-Week Avg NPS", f"{recent['NPS'].mean():.1f}", "Rolling view across the latest available weeks.")
            best = weekly.loc[weekly["NPS"].idxmax()]
            low = weekly.loc[weekly["NPS"].idxmin()]
            add("Best Week", f"{float(best['NPS']):.1f}", f"{best.get('Week', 'Highest weekly NPS')}.", "positive")
            add("Lowest Week", f"{float(low['NPS']):.1f}", f"{low.get('Week', 'Lowest weekly NPS')}.", "negative")
            add("NPS Volatility", f"{weekly['NPS'].std(ddof=0):.1f}", "Lower is more consistent week to week.")

    add("Positive Sentiment", f"{float(sentiment.get('Positive', 0)):.1f}%", "Sparrow positive share.", "positive")
    add("Negative Sentiment", f"{float(sentiment.get('Negative', 0)):.1f}%", "Sparrow negative share.", "negative")

    if "NPS Type" in analyzed_df.columns:
        detractors = int((analyzed_df["NPS Type"] == "Detractor").sum())
        add("At-Risk Responses", f"{detractors:,}", "Detractor records in the selected range.", "negative" if detractors else "neutral")
    if "Silent Detractor Alert" in analyzed_df.columns:
        silent = int(pd.to_numeric(analyzed_df["Silent Detractor Alert"], errors="coerce").fillna(0).sum())
        add("Silent Detractors", f"{silent:,}", "Soft language with low-score risk signals.", "negative" if silent else "neutral")

    if not reason_df.empty:
        first = reason_df.iloc[0]
        label_col = next((column for column in reason_df.columns if column not in {"__row_id", "Responses", "Count", "NPS"}), None)
        value = str(first.get(label_col, "")) if label_col else ""
        volume = _first_numeric_value(first, ["Responses", "Count", "Total"])
        if value:
            note = f"{volume:.0f} mentions." if volume else "Top available driver."
            add("Top Driver", value[:28], note)

    return cards[:12]


def _fill_owl_fallback_columns(df: pd.DataFrame, reason: str) -> pd.DataFrame:
    working = df.copy()
    defaults = {
        "Owl Primary Driver": "Bucket Category",
        "Owl Issue Type": "Primary Reason",
    }
    for output_column in [
        "Owl Primary Driver",
        "Owl Secondary Driver",
        "Owl Tertiary Driver",
        "Owl People Sentiment",
        "Owl Process Sentiment",
        "Owl Tech Sentiment",
        "Owl Issue Type",
        "Owl Customer Impact",
        "Owl Resolution Status",
    ]:
        if output_column not in working.columns:
            working[output_column] = ""
    for output_column, source_column in defaults.items():
        if source_column in working.columns:
            working[output_column] = working[source_column].fillna("").astype(str)
    if "Bucket Category" in working.columns and "Sentiment" in working.columns:
        for category, output_column in [
            ("People", "Owl People Sentiment"),
            ("Process", "Owl Process Sentiment"),
            ("Technology", "Owl Tech Sentiment"),
            ("Tech", "Owl Tech Sentiment"),
        ]:
            mask = working["Bucket Category"].astype(str).str.lower().eq(category.lower())
            working.loc[mask, output_column] = working.loc[mask, "Sentiment"].astype(str)
    if "NPS Type" in working.columns:
        working["Owl Customer Impact"] = working["NPS Type"].map(
            {"Detractor": "High Risk", "Passive": "Moderate Risk", "Promoter": "Low Risk"}
        ).fillna("No Clear Impact")
    if "Silent Detractor Alert" in working.columns:
        silent = pd.to_numeric(working["Silent Detractor Alert"], errors="coerce").fillna(0).astype(bool)
        working.loc[silent, "Owl Customer Impact"] = "Hidden Churn Risk"
    working["Owl Resolution Status"] = "Review Required"
    working["Owl Analysis Source"] = f"Safe fallback: {reason}"
    return working


def _neutralize_blank_feedback_sentiment(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty or "Verbatim Feedback" not in df.columns or "Sentiment" not in df.columns:
        return df
    working = df.copy()
    blank_feedback = working["Verbatim Feedback"].fillna("").astype(str).str.strip().eq("")
    if not blank_feedback.any():
        return working
    working.loc[blank_feedback, "Sentiment"] = "Neutral"
    if "Sentiment Score" in working.columns:
        working.loc[blank_feedback, "Sentiment Score"] = 0.0
    if "Impact Score" in working.columns:
        working.loc[blank_feedback, "Impact Score"] = 0.0
    if "AI Rationale" in working.columns:
        working.loc[blank_feedback, "AI Rationale"] = "Blank comment treated as neutral sentiment."
    return working


ACPT_KEYWORDS: dict[str, tuple[str, ...]] = {
    "Agent": (
        "agent", "advisor", "representative", "rep", "staff", "associate", "executive",
        "rude", "polite", "helpful", "unhelpful", "knowledge", "explained", "listened",
        "empathy", "attitude", "tone", "courteous", "professional", "coaching", "callback",
        "did not call", "wrong information", "misinformed", "transferred me", "hold without",
    ),
    "Customer": (
        "customer", "i forgot", "i did not", "i didn't", "i was confused", "my mistake",
        "user error", "entered wrong", "provided wrong", "changed my mind", "not available",
        "entered the wrong", "wrong details", "wrong detail", "incorrect details",
        "missed appointment", "could not understand", "language issue", "asked repeatedly",
    ),
    "Process": (
        "process", "policy", "procedure", "approval", "documentation", "form", "paperwork",
        "waiting", "wait time", "queue", "delay", "delayed", "sla", "turnaround", "follow up",
        "follow-up", "refund", "billing", "invoice", "escalation", "resolution", "case",
        "ticket", "workflow", "verification", "compliance", "appointment", "schedule",
    ),
    "Technology": (
        "system", "website", "app", "application", "portal", "login", "password", "otp",
        "link", "page", "screen", "error", "bug", "crash", "slow", "server", "network",
        "software", "hardware", "device", "tool", "dashboard", "payment failed", "technical",
        "connectivity", "download", "upload", "browser", "email not received", "sms",
    ),
}


ACPT_CATEGORY_ORDER = ("Agent", "Customer", "Process", "Technology")


def _normalize_acpt_text(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "").strip().lower())


def _score_acpt_category(text: str, phrases: tuple[str, ...]) -> tuple[float, list[str]]:
    score = 0.0
    evidence: list[str] = []
    padded = f" {text} "
    for phrase in phrases:
        token = phrase.lower()
        if " " in token:
            if token in text:
                score += 2.4
                evidence.append(phrase)
        elif re.search(rf"\b{re.escape(token)}\b", padded):
            score += 1.0
            evidence.append(phrase)
    return score, evidence[:6]


def _classify_acpt_text(value: Any) -> dict[str, Any]:
    text = _normalize_acpt_text(value)
    if not text:
        return {
            "primary": "Unclassified",
            "secondary": "",
            "confidence": 0.0,
            "evidence": "Blank verbatim",
            "needsReview": "Yes",
        }
    scored: list[tuple[str, float, list[str]]] = []
    for category in ACPT_CATEGORY_ORDER:
        score, evidence = _score_acpt_category(text, ACPT_KEYWORDS[category])
        scored.append((category, score, evidence))
    scored.sort(key=lambda item: item[1], reverse=True)
    top_category, top_score, top_evidence = scored[0]
    second_category, second_score, _second_evidence = scored[1]
    total_score = sum(score for _category, score, _evidence in scored)
    if top_score <= 0:
        return {
            "primary": "Unclassified",
            "secondary": "",
            "confidence": 0.0,
            "evidence": "No ACPT keyword evidence found",
            "needsReview": "Yes",
        }
    confidence = min(0.98, max(0.35, (top_score / max(total_score, 1.0)) * (0.78 + min(top_score, 6.0) / 30)))
    close_second = second_score > 0 and (top_score - second_score) <= 1.2
    return {
        "primary": top_category,
        "secondary": second_category if close_second else "",
        "confidence": round(confidence, 3),
        "evidence": ", ".join(top_evidence) if top_evidence else "Keyword signal present",
        "needsReview": "Yes" if confidence < 0.58 or close_second else "No",
    }


def _add_acpt_classification_outputs(
    df: pd.DataFrame,
    feedback_col: str = "Verbatim Feedback",
    progress_callback: Any | None = None,
) -> pd.DataFrame:
    working = df.copy()
    if feedback_col not in working.columns:
        fallback = _first_existing_column(working, ["Verbatim Feedback", "Feedback", "Comment", "Comments", "Verbatim"])
        feedback_col = fallback if fallback else feedback_col
    values = working[feedback_col].fillna("").astype(str).tolist() if feedback_col in working.columns else [""] * len(working)
    total = len(values)
    primary: list[str] = []
    secondary: list[str] = []
    confidence: list[float] = []
    evidence: list[str] = []
    review: list[str] = []
    checkpoint = max(1, total // 25) if total else 1
    for index, value in enumerate(values, 1):
        result = _classify_acpt_text(value)
        primary.append(result["primary"])
        secondary.append(result["secondary"])
        confidence.append(result["confidence"])
        evidence.append(result["evidence"])
        review.append(result["needsReview"])
        if progress_callback and (index == total or index == 1 or index % checkpoint == 0):
            progress_callback(index, max(total, 1), f"ACPT classification: {index:,}/{max(total, 1):,} rows. Assigning Agent, Customer, Process, or Technology ownership.")
    working["ACPT Primary Category"] = primary
    working["ACPT Secondary Category"] = secondary
    working["ACPT Confidence"] = confidence
    working["ACPT Evidence"] = evidence
    working["ACPT Needs Review"] = review
    return working


def _read_worker_progress(progress_path: Path) -> dict[str, Any] | None:
    if not progress_path.exists():
        return None
    try:
        payload = json.loads(progress_path.read_text(encoding="utf-8"))
    except Exception:
        return None
    return payload if isinstance(payload, dict) else None


def _run_isolated_model_worker(
    *,
    task: str,
    df: pd.DataFrame,
    model_path: str,
    analysis_id: str,
    progress_start: float,
    progress_end: float,
    status: str,
    timeout_seconds: int,
    feedback_col: str = "Verbatim Feedback",
    score_col: str = "",
    agent_col: str = "",
    date_col: str = "",
) -> pd.DataFrame:
    worker = Path(__file__).resolve().parent / "model_worker.py"
    with tempfile.TemporaryDirectory(prefix=f"npshtml_{task}_") as temp_dir:
        temp_path = Path(temp_dir)
        input_path = temp_path / "input.pkl"
        output_path = temp_path / "output.pkl"
        progress_path = temp_path / "progress.json"
        with open(input_path, "wb") as handle:
            pickle.dump(df, handle, protocol=pickle.HIGHEST_PROTOCOL)

        portable_python = ROOT / "portable_python" / "python.exe"
        worker_python = portable_python if portable_python.exists() else Path(sys.executable)
        command = [
            str(worker_python),
            str(worker),
            "--task",
            task,
            "--input",
            str(input_path),
            "--output",
            str(output_path),
            "--feedback-col",
            feedback_col,
            "--score-col",
            score_col or "",
            "--agent-col",
            agent_col or "",
            "--date-col",
            date_col or "",
            "--model-path",
            model_path,
            "--progress",
            str(progress_path),
        ]
        env = dict(os.environ)
        package_path = str(ROOT / "python_packages")
        env["PYTHONPATH"] = package_path + (os.pathsep + env.get("PYTHONPATH", "") if env.get("PYTHONPATH") else "")
        env.setdefault("TOKENIZERS_PARALLELISM", "false")
        started = time.time()
        process = subprocess.Popen(
            command,
            cwd=str(ROOT),
            env=env,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
        )
        last_worker_progress: float | None = None
        while process.poll() is None:
            elapsed = time.time() - started
            if elapsed > timeout_seconds:
                process.kill()
                process.communicate(timeout=5)
                raise TimeoutError(f"{task.title()} model exceeded {timeout_seconds} seconds and was stopped.")
            worker_progress = _read_worker_progress(progress_path)
            if worker_progress:
                done = max(int(worker_progress.get("done", 0)), 0)
                total = max(int(worker_progress.get("total", 0)), 0)
                message = str(worker_progress.get("message") or "").strip()
                if total > 1:
                    ratio = min(done / total, 0.995)
                    last_worker_progress = ratio
                    _set_progress(
                        progress_start + (progress_end - progress_start) * ratio,
                        f"{status}: {done:,}/{total:,} rows ({int(elapsed)}s)...",
                        analysis_id,
                    )
                else:
                    ratio = last_worker_progress if last_worker_progress is not None else min(elapsed / max(timeout_seconds, 1), 0.95)
                    label = message or status
                    _set_progress(
                        progress_start + (progress_end - progress_start) * min(ratio, 0.95),
                        f"{label} ({int(elapsed)}s)...",
                        analysis_id,
                    )
            else:
                ratio = min(elapsed / max(timeout_seconds, 1), 0.95)
                _set_progress(progress_start + (progress_end - progress_start) * ratio, f"{status} ({int(elapsed)}s)...", analysis_id)
            time.sleep(2)

        stdout, stderr = process.communicate(timeout=10)
        if process.returncode != 0:
            detail = (stderr or stdout or f"{task.title()} worker exited with code {process.returncode}").strip()
            raise RuntimeError(detail[-1000:])
        if not output_path.exists():
            raise RuntimeError(f"{task.title()} worker finished without returning output.")
        with open(output_path, "rb") as handle:
            return pickle.load(handle)


def _dashboard_snapshot(df: pd.DataFrame) -> dict[str, Any]:
    if df.empty:
        return {
            "summary": {},
            "counts": {},
            "sentiment": {},
            "weekly": [],
            "reasons": [],
            "intelligence": [],
            "segmentSentiments": _segment_sentiment_counts(df),
            "insights": "",
        }
    is_csat = "CSAT Type" in df.columns
    summaries = _summaries_for_calendar(df, _state_calendar_settings())
    if is_csat:
        summaries = {key: _apply_csat_aliases(value) for key, value in summaries.items()}
    weekly_df = summaries.get("weekly", pd.DataFrame())
    reason_df = summaries.get("reasons", pd.DataFrame())
    try:
        insights = executive_snapshot_insights(df, reason_df)
    except Exception:
        insights = ""
    summary = nps_summary(df)
    counts = nps_composition_counts(df)
    if is_csat:
        summary, counts = _csat_summary_aliases(summary, counts)
    return {
        "summary": summary,
        "counts": counts,
        "sentiment": sentiment_summary(df),
        "weekly": _safe_records(weekly_df, 52),
        "reasons": _safe_records(reason_df, 20),
        "intelligence": _executive_intelligence_cards(df, weekly_df, reason_df),
        "segmentSentiments": _segment_sentiment_counts(df),
        "insights": insights,
    }


def _dashboard_snapshots_by(df: pd.DataFrame, column: str | None, all_label: str) -> dict[str, Any]:
    snapshots = {all_label: _dashboard_snapshot(df)}
    if df.empty or not column or column not in df.columns:
        return snapshots
    for value in sorted(df[column].dropna().astype(str).unique().tolist()):
        snapshots[value] = _dashboard_snapshot(df[df[column].astype(str) == value].copy())
    return snapshots


def _first_existing_column(df: pd.DataFrame, candidates: list[str]) -> str | None:
    lowered = {str(column).strip().lower(): column for column in df.columns if str(column).strip()}
    for candidate in candidates:
        needle = candidate.strip().lower()
        if needle and needle in lowered:
            return lowered[needle]
    for candidate in candidates:
        needle = candidate.strip().lower()
        if not needle:
            continue
        for lower, column in lowered.items():
            if needle in lower:
                return column
    return None


def _manager_summary(df: pd.DataFrame) -> pd.DataFrame:
    manager_col = _first_existing_column(df, ["Manager/TL", "TL Name", "Team Manager", "Manager Name", "Manager", "Supervisor"])
    if df.empty or manager_col is None or not str(manager_col).strip() or manager_col not in df.columns:
        return pd.DataFrame()
    working = df.copy()
    grouped = (
        working.groupby(manager_col)
        .agg(
            Responses=(manager_col, "size"),
            Promoters=("NPS Type", lambda values: int((values == "Promoter").sum())),
            Passives=("NPS Type", lambda values: int((values == "Passive").sum())),
            Detractors=("NPS Type", lambda values: int((values == "Detractor").sum())),
            Avg_Rating=("NPS Score", "mean"),
        )
        .reset_index()
        .rename(columns={manager_col: "Manager/TL"})
    )
    grouped["NPS"] = ((grouped["Promoters"] / grouped["Responses"] * 100) - (grouped["Detractors"] / grouped["Responses"] * 100)).round(2)
    grouped["Avg_Rating"] = grouped["Avg_Rating"].round(2)
    return grouped.sort_values(["NPS", "Responses"], ascending=[True, False])


def _agent_summary_with_manager(df: pd.DataFrame, agent_df: pd.DataFrame) -> pd.DataFrame:
    manager_col = _first_existing_column(df, ["Manager/TL", "TL Name", "Team Manager", "Manager Name", "Manager", "Supervisor"])
    if df.empty or agent_df.empty or manager_col is None or not str(manager_col).strip() or manager_col not in df.columns or "Agent Name" not in df.columns or "Agent Name" not in agent_df.columns:
        return agent_df
    mapping = (
        df[["Agent Name", manager_col]]
        .dropna(subset=["Agent Name"])
        .assign(**{"Agent Name": lambda frame: frame["Agent Name"].astype(str), manager_col: lambda frame: frame[manager_col].astype(str)})
        .groupby("Agent Name")[manager_col]
        .agg(lambda values: values.mode().iloc[0] if not values.mode().empty else values.iloc[0])
        .rename("Manager/TL")
        .reset_index()
    )
    enriched = agent_df.copy()
    enriched["Agent Name"] = enriched["Agent Name"].astype(str)
    if "Manager/TL" in enriched.columns:
        enriched = enriched.drop(columns=["Manager/TL"])
    insert_at = 1 if "Agent Name" in enriched.columns else 0
    enriched = enriched.merge(mapping, on="Agent Name", how="left")
    manager_series = enriched.pop("Manager/TL")
    enriched.insert(insert_at, "Manager/TL", manager_series.fillna(""))
    return enriched


def _quartile_summary(agent_df: pd.DataFrame) -> pd.DataFrame:
    if agent_df.empty or "Agent NPS" not in agent_df.columns:
        return pd.DataFrame()
    working = agent_df.copy().sort_values("Agent NPS", ascending=False).reset_index(drop=True)
    total = len(working)
    labels = []
    for index in range(total):
        rank_pct = (index + 1) / max(total, 1)
        if rank_pct <= 0.25:
            labels.append("Q1")
        elif rank_pct <= 0.50:
            labels.append("Q2")
        elif rank_pct <= 0.75:
            labels.append("Q3")
        else:
            labels.append("Q4")
    working["Quartile"] = labels
    columns = ["Quartile", "Agent Name", "Responses", "Agent NPS", "Promoters", "Detractors"]
    if "Average_Rating" in working.columns:
        working["Average Rating"] = pd.to_numeric(working["Average_Rating"], errors="coerce").round(2)
        columns.insert(4, "Average Rating")
    return working[columns]


REPORT_TAB_TITLES: dict[str, str] = {
    "executive": "Executive Dashboard",
    "alerts": "Alert Badges",
    "agent": "Agent Dashboard",
    "manager": "Manager Dashboard",
    "quartile": "Quartile Intelligence",
    "statistics": "Statistics",
    "feedback": "Sentiment Intelligence",
    "sentimentcompare": "Sentiment Comparison",
    "theme": "Theme Classification",
    "themebuilder": "Build Theme Classification",
    "themecompare": "Theme Comparison",
    "wave": "Wave Intelligence",
    "tenure": "Tenure Intelligence",
    "operations": "Operations",
    "rootcause": "Word Cloud Intelligence",
    "promoterdna": "Promoter DNA",
    "Satisfieddna": "Satisfied DNA",
    "satisfieddna": "Satisfied DNA",
    "gap": "Gap & Opportunity",
    "detail": "Detailed Insight",
    "churn": "Churn Risk Assessment",
    "analysis": "Analysis",
}


def _safe_columns(frame: pd.DataFrame, columns: list[str]) -> pd.DataFrame:
    if frame.empty:
        return pd.DataFrame()
    available = [column for column in columns if column in frame.columns]
    return frame[available].copy() if available else frame.copy()


def _report_tab_frames(selected: list[str]) -> dict[str, tuple[str, pd.DataFrame]]:
    summary = nps_summary(STATE.analyzed_df) if not STATE.analyzed_df.empty else {}
    sentiment = sentiment_summary(STATE.analyzed_df) if not STATE.analyzed_df.empty else {}
    executive_rows = [
        {"Metric": "Responses", "Value": summary.get("total_responses", len(STATE.analyzed_df))},
        {"Metric": "NPS", "Value": summary.get("nps", 0)},
        {"Metric": "Promoters", "Value": summary.get("promoters", 0)},
        {"Metric": "Passives", "Value": summary.get("passives", 0)},
        {"Metric": "Detractors", "Value": summary.get("detractors", 0)},
        {"Metric": "Positive Sentiment %", "Value": sentiment.get("Positive", 0)},
        {"Metric": "Neutral Sentiment %", "Value": sentiment.get("Neutral", 0)},
        {"Metric": "Negative Sentiment %", "Value": sentiment.get("Negative", 0)},
    ]
    promoter_df = STATE.analyzed_df[STATE.analyzed_df.get("NPS Type", pd.Series(dtype=str)).astype(str).str.lower().eq("promoter")].copy() if not STATE.analyzed_df.empty and "NPS Type" in STATE.analyzed_df.columns else pd.DataFrame()
    detractor_df = STATE.analyzed_df[STATE.analyzed_df.get("NPS Type", pd.Series(dtype=str)).astype(str).str.lower().eq("detractor")].copy() if not STATE.analyzed_df.empty and "NPS Type" in STATE.analyzed_df.columns else pd.DataFrame()
    satisfied_col = _first_existing_column(STATE.analyzed_df, ["CSAT Segment", "CSAT Type", "Satisfaction Segment"])
    satisfied_df = STATE.analyzed_df[STATE.analyzed_df[satisfied_col].astype(str).str.lower().eq("satisfied")].copy() if satisfied_col else promoter_df
    wave_col = _first_existing_column(STATE.analyzed_df, ["Batch#", "Batch", "Training Wave", "Cohort", "Wave"])
    tenure_col = _first_existing_column(STATE.analyzed_df, ["Tenure", "Tenure days", "Tenure Bucket", "Tenure Range", "Employee Tenure"])
    frame_map: dict[str, tuple[str, pd.DataFrame]] = {
        "executive": ("Executive", pd.DataFrame(executive_rows)),
        "alerts": ("Alert Badges", STATE.complaints_df),
        "agent": ("Agent Dashboard", STATE.agent_df),
        "manager": ("Manager Dashboard", STATE.manager_df),
        "quartile": ("Quartile Intelligence", _quartile_summary(STATE.agent_df)),
        "statistics": ("Statistics", STATE.weekly_df),
        "feedback": ("Sentiment Intelligence", _safe_columns(STATE.analyzed_df, ["Agent Name", "Manager/TL", "Case ID", "NPS Type", "NPS Score", "Sentiment", "Sentiment Confidence", "Verbatim Feedback"])),
        "sentimentcompare": ("Sentiment Comparison", _sentiment_movement_rows(STATE.weekly_df)),
        "theme": ("Theme Classification", _safe_columns(STATE.analyzed_df, ["Agent Name", "Manager/TL", "NPS Type", "NPS Score", "Theme", "Theme Confidence", "Primary Reason", "Verbatim Feedback"])),
        "wave": ("Wave Intelligence", _dimension_summary(STATE.analyzed_df, wave_col, "Wave")),
        "tenure": ("Tenure Intelligence", _dimension_summary(STATE.analyzed_df, tenure_col, "Tenure")),
        "operations": ("Operations", STATE.reason_df),
        "rootcause": ("Word Cloud Intel", STATE.reason_df),
        "promoterdna": ("Promoter DNA", _safe_columns(promoter_df, ["Agent Name", "Manager/TL", "NPS Score", "Sentiment", "Primary Reason", "Verbatim Feedback"])),
        "Satisfieddna": ("Satisfied DNA", _safe_columns(satisfied_df, ["Agent Name", "Manager/TL", "CSAT Score", "CSAT Segment", "Sentiment", "Primary Reason", "Verbatim Feedback"])),
        "satisfieddna": ("Satisfied DNA", _safe_columns(satisfied_df, ["Agent Name", "Manager/TL", "CSAT Score", "CSAT Segment", "Sentiment", "Primary Reason", "Verbatim Feedback"])),
        "gap": ("Gap Opportunity", STATE.passive_df),
        "detail": ("Detailed Insight", _safe_columns(STATE.analyzed_df, ["Agent Name", "Manager/TL", "NPS Type", "NPS Score", "Sentiment", "Primary Reason", "Verbatim Feedback"])),
        "churn": ("Churn Risk", _safe_columns(detractor_df, ["Agent Name", "Manager/TL", "NPS Score", "Sentiment", "Primary Reason", "Verbatim Feedback"])),
        "analysis": ("Movement Analysis", STATE.weekly_df),
    }
    return {key: frame_map[key] for key in selected if key in frame_map}


def _report_tab_frames_from_payload(selected: list[str], analysis: dict[str, Any]) -> dict[str, tuple[str, pd.DataFrame]]:
    summary = analysis.get("summary") or {}
    counts = analysis.get("counts") or {}
    sentiment = analysis.get("sentiment") or {}
    executive = pd.DataFrame(
        [
            {"Metric": "Responses", "Value": summary.get("total_responses") or summary.get("Responses") or 0},
            {"Metric": "NPS", "Value": summary.get("nps") or summary.get("NPS") or 0},
            {"Metric": "Promoters", "Value": counts.get("Promoters") or summary.get("promoters") or 0},
            {"Metric": "Passives", "Value": counts.get("Passives") or summary.get("passives") or 0},
            {"Metric": "Detractors", "Value": counts.get("Detractors") or summary.get("detractors") or 0},
            {"Metric": "Positive Sentiment %", "Value": sentiment.get("Positive", 0)},
            {"Metric": "Neutral Sentiment %", "Value": sentiment.get("Neutral", 0)},
            {"Metric": "Negative Sentiment %", "Value": sentiment.get("Negative", 0)},
        ]
    )

    def frame(key: str) -> pd.DataFrame:
        rows = analysis.get(key) or []
        return pd.DataFrame(rows if isinstance(rows, list) else [])

    frame_map: dict[str, tuple[str, pd.DataFrame]] = {
        "executive": ("Executive", executive),
        "alerts": ("Alert Badges", frame("alerts")),
        "agent": ("Agent Dashboard", frame("agents")),
        "manager": ("Manager Dashboard", frame("managers")),
        "quartile": ("Quartile Intelligence", frame("quartiles")),
        "statistics": ("Statistics", frame("weekly")),
        "feedback": ("Sentiment Intelligence", frame("feedbackTableRows") if analysis.get("feedbackTableRows") else frame("feedbackRows")),
        "sentimentcompare": ("Sentiment Comparison", frame("sentimentCompareRows") if analysis.get("sentimentCompareRows") else frame("sentimentMovement")),
        "theme": ("Theme Classification", frame("themeRows")),
        "themebuilder": ("Build Theme Classification", frame("themeBuilderRows")),
        "themecompare": ("Theme Comparison", frame("themeCompareRows")),
        "wave": ("Wave Intelligence", frame("wave")),
        "tenure": ("Tenure Intelligence", frame("tenure")),
        "operations": ("Operations", frame("operations")),
        "rootcause": ("Word Cloud Intel", frame("reasons")),
        "promoterdna": ("Promoter DNA", frame("feedbackRows")),
        "Satisfieddna": ("Satisfied DNA", frame("feedbackRows")),
        "satisfieddna": ("Satisfied DNA", frame("feedbackRows")),
        "gap": ("Gap Opportunity", frame("passives")),
        "detail": ("Detailed Insight", frame("feedbackRows")),
        "churn": ("Churn Risk", frame("churn")),
        "analysis": ("Movement Analysis", frame("weekly")),
    }
    return {key: frame_map[key] for key in selected if key in frame_map}


def _write_report_workbook(selected: list[str], analysis_payload: dict[str, Any] | None = None) -> bytes:
    frames = _report_tab_frames_from_payload(selected, analysis_payload or {}) if analysis_payload else _report_tab_frames(selected)
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        if not frames:
            pd.DataFrame({"Status": ["No matching data sheets were available for the selected tabs."]}).to_excel(
                writer,
                index=False,
                sheet_name="No Data",
            )
        for _key, (sheet_name, frame) in frames.items():
            (frame if not frame.empty else pd.DataFrame({"Status": ["No data available"]})).to_excel(
                writer,
                index=False,
                sheet_name=sheet_name[:31],
            )
    return output.getvalue()


def _add_ppt_textbox(slide, left, top, width, height, text: str, size: int = 11, bold: bool = False, color: tuple[int, int, int] = (12, 35, 64)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*color)
    return box


def _add_ppt_table(slide, frame: pd.DataFrame, left, top, width, height, max_rows: int = 8) -> None:
    if frame.empty:
        _add_ppt_textbox(slide, left, top, width, Inches(0.4), "No data available for this section.", 12)
        return
    working = frame.head(max_rows).copy()
    columns = [str(column)[:28] for column in list(working.columns)[:6]]
    working = working.iloc[:, : len(columns)]
    table_shape = slide.shapes.add_table(len(working) + 1, len(columns), left, top, width, height).table
    for index, column in enumerate(columns):
        cell = table_shape.cell(0, index)
        cell.text = column
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(237, 244, 248)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(8)
    for row_idx, (_, row) in enumerate(working.iterrows(), start=1):
        for col_idx, column in enumerate(working.columns):
            value = "" if pd.isna(row[column]) else str(row[column])
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = value[:120]
            cell.text_frame.paragraphs[0].font.size = Pt(7)


def _build_report_pptx(selected: list[str], analysis_payload: dict[str, Any] | None = None) -> bytes:
    frames = _report_tab_frames_from_payload(selected, analysis_payload or {}) if analysis_payload else _report_tab_frames(selected)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    navy = RGBColor(0, 61, 91)
    teal = RGBColor(0, 155, 155)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = navy
    _add_ppt_textbox(slide, Inches(0.55), Inches(0.55), Inches(6.5), Inches(0.35), "Feedback Intelligence Hub", 13, False, (255, 255, 255))
    _add_ppt_textbox(slide, Inches(0.55), Inches(1.4), Inches(10.8), Inches(0.8), "NPS Analyzer Executive Report", 30, True, (255, 255, 255))
    _add_ppt_textbox(slide, Inches(0.58), Inches(2.35), Inches(10), Inches(0.35), f"Selected sections: {', '.join(REPORT_TAB_TITLES.get(tab, tab) for tab in selected)}", 11, False, (214, 232, 240))

    slide = prs.slides.add_slide(blank)
    _add_ppt_textbox(slide, Inches(0.45), Inches(0.35), Inches(5.5), Inches(0.35), "Index & Reference Sheets", 20, True)
    index_frame = pd.DataFrame(
        [
            {"Ref": f"R{idx + 1:02d}", "Section": REPORT_TAB_TITLES.get(tab, tab), "Sheet": frames.get(tab, (tab, pd.DataFrame()))[0]}
            for idx, tab in enumerate(selected)
        ]
    )
    _add_ppt_table(slide, index_frame, Inches(0.55), Inches(1.0), Inches(12.1), Inches(5.6), max_rows=18)

    for idx, tab in enumerate(selected):
        title, frame = frames.get(tab, (REPORT_TAB_TITLES.get(tab, tab), pd.DataFrame()))
        slide = prs.slides.add_slide(blank)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
        accent.fill.solid()
        accent.fill.fore_color.rgb = teal
        accent.line.fill.background()
        _add_ppt_textbox(slide, Inches(0.45), Inches(0.42), Inches(1.0), Inches(0.3), f"R{idx + 1:02d}", 12, True, (0, 155, 155))
        _add_ppt_textbox(slide, Inches(1.2), Inches(0.35), Inches(6.5), Inches(0.45), title, 20, True)
        _add_ppt_textbox(slide, Inches(0.55), Inches(0.95), Inches(11.7), Inches(0.3), f"{len(frame):,} supporting rows available in the reference data.", 10, False, (82, 103, 123))
        _add_ppt_table(slide, frame, Inches(0.55), Inches(1.45), Inches(12.2), Inches(5.35), max_rows=10)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _build_boardroom_reference_pdf(payload: dict[str, Any]) -> bytes:
    """Build a concise, native PDF leadership pack without printing browser HTML."""
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase.pdfmetrics import stringWidth
    from reportlab.pdfgen.canvas import Canvas
    from reportlab.platypus import (
        KeepTogether,
        LongTable,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )
    from reportlab.graphics.shapes import Circle, Drawing, Line, PolyLine, Rect, String
    from xml.sax.saxutils import escape as xml_escape

    analysis = payload.get("analysis") if isinstance(payload.get("analysis"), dict) else {}
    selected = [str(item) for item in payload.get("tabs", []) if str(item).strip()]
    metric = str(payload.get("metric") or "NPS").strip().upper()
    options = payload.get("options") if isinstance(payload.get("options"), dict) else {}
    title = str(options.get("title") or f"{metric} Leadership Report").strip()
    account = str(options.get("accountName") or "Customer Intelligence Team").strip()
    prepared_for = str(options.get("preparedFor") or "Leadership Review").strip()
    footer_note = str(options.get("footerNote") or "Confidential - Krestrel Analysis Suite Review").strip()
    generated_at = _dt.datetime.now().strftime("%d %b %Y, %I:%M %p")
    frames = _report_tab_frames_from_payload(selected, analysis) if analysis else _report_tab_frames(selected)
    custom_dashboards = payload.get("customDashboards") if isinstance(payload.get("customDashboards"), list) else []

    output = BytesIO()
    page_size = landscape(A4)
    page_width, page_height = page_size
    margin_x = 15 * mm
    margin_top = 18 * mm
    margin_bottom = 15 * mm
    content_width = page_width - (2 * margin_x)
    navy = colors.HexColor("#003D5B")
    dark_teal = colors.HexColor("#004B57")
    teal = colors.HexColor("#009B9B")
    aqua = colors.HexColor("#2FE4D6")
    pale = colors.HexColor("#EEF7F8")
    pale_blue = colors.HexColor("#E8F1F6")
    ink = colors.HexColor("#0C2340")
    muted = colors.HexColor("#5D7185")
    line = colors.HexColor("#C8DCE4")
    positive = colors.HexColor("#0A8F74")
    negative = colors.HexColor("#C84B63")

    doc = SimpleDocTemplate(
        output,
        pagesize=page_size,
        leftMargin=margin_x,
        rightMargin=margin_x,
        topMargin=margin_top,
        bottomMargin=margin_bottom,
        title=title,
        author="Krestrel Analysis Suite",
        subject=f"{metric} board-room leadership report",
    )
    styles = getSampleStyleSheet()
    cover_eyebrow = ParagraphStyle("CoverEyebrow", parent=styles["Normal"], fontName="Helvetica-Bold", fontSize=10, leading=13, textColor=aqua, spaceAfter=10, uppercase=True)
    cover_title = ParagraphStyle("CoverTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=30, leading=34, textColor=colors.white, spaceAfter=12)
    cover_subtitle = ParagraphStyle("CoverSubtitle", parent=styles["Normal"], fontName="Helvetica", fontSize=12, leading=17, textColor=colors.HexColor("#D8F1F2"), spaceAfter=8)
    section_title = ParagraphStyle("SectionTitle", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=21, leading=25, textColor=ink, spaceAfter=5)
    section_intro = ParagraphStyle("SectionIntro", parent=styles["Normal"], fontName="Helvetica", fontSize=9.5, leading=14, textColor=muted, spaceAfter=10)
    small = ParagraphStyle("Small", parent=styles["Normal"], fontName="Helvetica", fontSize=8, leading=10.5, textColor=muted)
    table_head = ParagraphStyle("TableHead", parent=styles["Normal"], fontName="Helvetica-Bold", fontSize=7.3, leading=9, textColor=colors.white)
    table_cell = ParagraphStyle("TableCell", parent=styles["Normal"], fontName="Helvetica", fontSize=7.2, leading=9, textColor=ink)
    kpi_label = ParagraphStyle("KpiLabel", parent=styles["Normal"], fontName="Helvetica-Bold", fontSize=7.5, leading=9, textColor=teal, spaceAfter=5)
    kpi_value = ParagraphStyle("KpiValue", parent=styles["Normal"], fontName="Helvetica-Bold", fontSize=18, leading=21, textColor=ink)
    index_ref = ParagraphStyle("IndexRef", parent=styles["Normal"], fontName="Helvetica-Bold", fontSize=8, leading=10, textColor=teal)
    index_name = ParagraphStyle("IndexName", parent=styles["Normal"], fontName="Helvetica-Bold", fontSize=9, leading=11, textColor=ink)

    def clean(value: Any, limit: int = 180) -> str:
        if value is None or (isinstance(value, float) and math.isnan(value)):
            return ""
        if isinstance(value, (int, float)) and not isinstance(value, bool):
            numeric = float(value)
            text = f"{int(numeric):,}" if numeric.is_integer() else f"{numeric:,.2f}".rstrip("0").rstrip(".")
        else:
            text = re.sub(r"\s+", " ", str(value)).strip()
        if len(text) > limit:
            text = text[: max(0, limit - 3)].rstrip() + "..."
        return xml_escape(text)

    def first_value(*keys: str, default: Any = 0) -> Any:
        for source in (analysis.get("summary") or {}, analysis.get("counts") or {}, analysis.get("sentiment") or {}):
            for key in keys:
                value = source.get(key)
                if value not in (None, ""):
                    return value
        return default

    def number(value: Any, decimals: int = 1) -> str:
        try:
            numeric = float(value)
            if numeric.is_integer() and decimals == 0:
                return f"{int(numeric):,}"
            return f"{numeric:,.{decimals}f}"
        except Exception:
            return str(value or "Not available")

    if "executive" in frames:
        segment_labels = ("Satisfied", "Neutral", "Dissatisfied") if metric == "CSAT" else ("Promoters", "Passives", "Detractors")
        executive_rows = [
            {"Metric": "Responses", "Value": first_value("total_responses", "Responses", "total", default=0)},
            {"Metric": metric, "Value": first_value(metric.lower(), metric, "csat_score", "nps", "NPS", default=0)},
        ]
        for label in segment_labels:
            executive_rows.append({"Metric": label, "Value": first_value(label, label.lower(), f"{label.lower()}_count", default=0)})
        executive_rows.extend([
            {"Metric": "Positive Sentiment %", "Value": first_value("Positive", "positive", default=0)},
            {"Metric": "Neutral Sentiment %", "Value": first_value("Neutral", "neutral", default=0)},
            {"Metric": "Negative Sentiment %", "Value": first_value("Negative", "negative", default=0)},
        ])
        frames["executive"] = ("Executive", pd.DataFrame(executive_rows))

    def header_footer(canvas: Canvas, document) -> None:
        canvas.saveState()
        canvas.setFillColor(colors.white)
        canvas.rect(0, 0, page_width, page_height, fill=1, stroke=0)
        canvas.setFillColor(navy)
        canvas.rect(0, page_height - 10 * mm, page_width, 10 * mm, fill=1, stroke=0)
        canvas.setFillColor(aqua)
        canvas.rect(0, page_height - 10 * mm, 32 * mm, 1.2 * mm, fill=1, stroke=0)
        canvas.setFont("Helvetica-Bold", 8.5)
        canvas.setFillColor(colors.white)
        canvas.drawString(margin_x, page_height - 6.4 * mm, title[:80])
        canvas.setFont("Helvetica", 7.5)
        canvas.setFillColor(colors.HexColor("#D4ECEE"))
        canvas.drawRightString(page_width - margin_x, page_height - 6.4 * mm, prepared_for[:60])
        canvas.setStrokeColor(line)
        canvas.line(margin_x, 9 * mm, page_width - margin_x, 9 * mm)
        canvas.setFillColor(muted)
        canvas.setFont("Helvetica", 7)
        canvas.drawString(margin_x, 5.7 * mm, footer_note[:110])
        canvas.drawRightString(page_width - margin_x, 5.7 * mm, f"Page {document.page}")
        canvas.restoreState()

    def cover_page(canvas: Canvas, _document) -> None:
        canvas.saveState()
        canvas.setFillColor(dark_teal)
        canvas.rect(0, 0, page_width, page_height, fill=1, stroke=0)
        canvas.setFillColor(colors.HexColor("#075E68"))
        canvas.circle(page_width - 20 * mm, 18 * mm, 78 * mm, fill=1, stroke=0)
        canvas.setFillColor(colors.HexColor("#087C82"))
        canvas.circle(page_width - 2 * mm, page_height - 2 * mm, 42 * mm, fill=1, stroke=0)
        canvas.setFillColor(aqua)
        canvas.rect(margin_x, page_height - 20 * mm, 28 * mm, 1.5 * mm, fill=1, stroke=0)
        canvas.setFont("Helvetica", 7.5)
        canvas.setFillColor(colors.HexColor("#BDE4E6"))
        canvas.drawString(margin_x, 10 * mm, footer_note[:110])
        canvas.drawRightString(page_width - margin_x, 10 * mm, generated_at)
        canvas.restoreState()

    def section_header(ref: str, name: str, description: str) -> list[Any]:
        marker = Table([[Paragraph(clean(ref), index_ref), Paragraph(clean(name), section_title)]], colWidths=[18 * mm, content_width - 18 * mm])
        marker.setStyle(TableStyle([
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("BACKGROUND", (0, 0), (0, 0), pale),
            ("BOX", (0, 0), (0, 0), 0.8, colors.HexColor("#A8D6D5")),
            ("LEFTPADDING", (0, 0), (0, 0), 7),
            ("RIGHTPADDING", (0, 0), (0, 0), 7),
            ("TOPPADDING", (0, 0), (0, 0), 7),
            ("BOTTOMPADDING", (0, 0), (0, 0), 7),
            ("LEFTPADDING", (1, 0), (1, 0), 10),
            ("RIGHTPADDING", (1, 0), (1, 0), 0),
        ]))
        return [marker, Spacer(1, 3 * mm), Paragraph(clean(description, 450), section_intro)]

    def frame_table(frame: pd.DataFrame, max_rows: int = 12, max_cols: int = 7) -> list[Any]:
        if frame.empty:
            return [Paragraph("No populated rows were available for this section.", section_intro)]
        working = frame.copy().iloc[:max_rows, :max_cols]
        columns = [str(column) for column in working.columns]
        if not columns:
            return [Paragraph("No populated columns were available for this section.", section_intro)]
        data: list[list[Any]] = [[Paragraph(clean(column, 45), table_head) for column in columns]]
        for _, row in working.iterrows():
            data.append([Paragraph(clean(row[column]), table_cell) for column in working.columns])
        col_widths = [content_width / len(columns)] * len(columns)
        table = LongTable(data, colWidths=col_widths, repeatRows=1, hAlign="LEFT")
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), navy),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.35, line),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 5),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, pale_blue]),
        ]))
        result: list[Any] = [table]
        if len(frame) > max_rows or len(frame.columns) > max_cols:
            result = [Paragraph(
                f"Board-room view shows the first {min(len(frame), max_rows):,} of {len(frame):,} rows and {min(len(frame.columns), max_cols)} of {len(frame.columns)} columns. The Interactive HTML retains the detailed view.",
                small,
            ), Spacer(1, 2 * mm), table]
        return result

    def weekly_chart(rows: list[dict[str, Any]]) -> Drawing | None:
        if not rows:
            return None
        score_keys = [metric, "NPS", "CSAT", "Score", "Satisfaction"]
        score_key = next((key for key in score_keys if any(key in row for row in rows)), None)
        label_key = next((key for key in ("Week", "Month", "Period", "Date") if any(key in row for row in rows)), None)
        if not score_key or not label_key:
            return None
        points: list[tuple[str, float]] = []
        for row in rows[-12:]:
            try:
                points.append((str(row.get(label_key, ""))[:12], float(row.get(score_key))))
            except Exception:
                continue
        if len(points) < 2:
            return None
        width, height = content_width, 58 * mm
        drawing = Drawing(width, height)
        drawing.add(Rect(0, 0, width, height, rx=10, ry=10, fillColor=pale, strokeColor=line, strokeWidth=0.7))
        left, right, bottom, top = 38, width - 18, 30, height - 24
        values = [value for _, value in points]
        low, high = min(values), max(values)
        pad = max(5.0, (high - low) * 0.18)
        low, high = low - pad, high + pad
        if metric == "NPS":
            low, high = min(-100.0, low), max(100.0, high)
        for step in range(5):
            y = bottom + (top - bottom) * step / 4
            value = low + (high - low) * step / 4
            drawing.add(Line(left, y, right, y, strokeColor=colors.HexColor("#D4E4E8"), strokeWidth=0.5))
            drawing.add(String(6, y - 3, f"{value:.1f}", fontName="Helvetica", fontSize=6.5, fillColor=muted))
        poly_points: list[float] = []
        for index, (label, value) in enumerate(points):
            x = left + (right - left) * index / max(1, len(points) - 1)
            y = bottom + (top - bottom) * (value - low) / max(0.001, high - low)
            poly_points.extend([x, y])
            drawing.add(Circle(x, y, 3, fillColor=teal, strokeColor=colors.white, strokeWidth=1))
            drawing.add(String(x, y + 7, f"{value:.1f}", textAnchor="middle", fontName="Helvetica-Bold", fontSize=6.5, fillColor=ink))
            drawing.add(String(x, 12, label, textAnchor="middle", fontName="Helvetica", fontSize=5.5, fillColor=muted))
        drawing.add(PolyLine(poly_points, strokeColor=teal, strokeWidth=2.2))
        drawing.add(String(left, height - 14, f"{metric} trend - latest {len(points)} periods", fontName="Helvetica-Bold", fontSize=9, fillColor=ink))
        return drawing

    story: list[Any] = [
        Spacer(1, 27 * mm),
        Paragraph("KRESTREL ANALYSIS SUITE", cover_eyebrow),
        Paragraph(clean(title, 120), cover_title),
        Paragraph(f"Prepared for <b>{clean(prepared_for, 90)}</b>", cover_subtitle),
        Paragraph(f"{clean(account, 90)} | Generated {clean(generated_at)}", cover_subtitle),
        Spacer(1, 20 * mm),
        Paragraph("A concise leadership pack generated directly from the completed analysis. Detailed interactive exploration remains available in the companion HTML report.", cover_subtitle),
        PageBreak(),
    ]

    summary = analysis.get("summary") or {}
    total = first_value("total_responses", "Responses", "total", default=0)
    score = first_value(metric.lower(), metric, "csat_score", "nps", "NPS", "csat", "CSAT", default=0)
    promoter = first_value("promoters_pct", "satisfied_pct", "Promoters", "promoters", "Satisfied", "satisfied", default=0)
    detractor = first_value("detractors_pct", "dissatisfied_pct", "Detractors", "detractors", "Dissatisfied", "dissatisfied", default=0)
    story.extend(section_header("01", "Executive Snapshot", "The essential performance signals for leadership review."))
    kpis = [
        ("SURVEY VOLUME", number(total, 0)),
        (metric, number(score, 1)),
        ("PROMOTER / SATISFIED", number(promoter, 1)),
        ("DETRACTOR / DISSATISFIED", number(detractor, 1)),
    ]
    kpi_cells = []
    for label, value in kpis:
        card = Table([[Paragraph(clean(label), kpi_label)], [Paragraph(clean(value), kpi_value)]], colWidths=[content_width / 4 - 5 * mm])
        card.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), pale),
            ("BOX", (0, 0), (-1, -1), 0.7, line),
            ("LEFTPADDING", (0, 0), (-1, -1), 10),
            ("RIGHTPADDING", (0, 0), (-1, -1), 10),
            ("TOPPADDING", (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ]))
        kpi_cells.append(card)
    kpi_grid = Table([kpi_cells], colWidths=[content_width / 4] * 4)
    kpi_grid.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 3), ("RIGHTPADDING", (0, 0), (-1, -1), 3)]))
    story.extend([kpi_grid, Spacer(1, 6 * mm)])
    weekly_rows = analysis.get("weekly") if isinstance(analysis.get("weekly"), list) else []
    trend = weekly_chart(weekly_rows)
    if trend:
        story.extend([trend, Spacer(1, 4 * mm)])
    story.append(PageBreak())

    directory_rows = [[Paragraph("REF", table_head), Paragraph("SECTION", table_head), Paragraph("BOARD-ROOM CONTENT", table_head)]]
    directory_items: list[tuple[str, str, int]] = []
    directory_ref = 3
    for tab in selected:
        if tab == "customdashboards":
            continue
        title_item, frame = frames.get(tab, (REPORT_TAB_TITLES.get(tab, tab), pd.DataFrame()))
        if frame.empty:
            continue
        directory_items.append((f"{directory_ref:02d}", title_item, len(frame)))
        directory_ref += 1
    for dashboard in custom_dashboards:
        directory_items.append((f"{directory_ref:02d}", str(dashboard.get("title") or "Custom Dashboard"), len(dashboard.get("rows") or [])))
        directory_ref += 1
    for ref, name, row_count in directory_items:
        directory_rows.append([Paragraph(ref, index_ref), Paragraph(clean(name), index_name), Paragraph(f"{row_count:,} supporting rows available", table_cell)])
    story.extend(section_header("02", "Report Directory", "Selected sections included in this concise PDF pack."))
    directory = Table(directory_rows, colWidths=[20 * mm, 82 * mm, content_width - 102 * mm], repeatRows=1)
    directory.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), navy),
        ("GRID", (0, 0), (-1, -1), 0.35, line),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, pale_blue]),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 7),
        ("RIGHTPADDING", (0, 0), (-1, -1), 7),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.extend([directory, PageBreak()])

    section_number = 3
    for tab in selected:
        if tab == "customdashboards":
            continue
        name, frame = frames.get(tab, (REPORT_TAB_TITLES.get(tab, tab), pd.DataFrame()))
        if frame.empty:
            continue
        guide = REPORT_TAB_TITLES.get(tab, name)
        story.extend(section_header(f"{section_number:02d}", name, f"Leadership-ready reference view for {guide}."))
        story.extend(frame_table(frame))
        story.append(PageBreak())
        section_number += 1

    for dashboard in custom_dashboards:
        name = str(dashboard.get("title") or "Custom Dashboard")
        creator = str(dashboard.get("creator") or "Dashboard Maker")
        rows = dashboard.get("rows") if isinstance(dashboard.get("rows"), list) else []
        frame = pd.DataFrame(rows)
        story.extend(section_header(f"{section_number:02d}", name, f"Custom dashboard created in {creator}."))
        if not frame.empty:
            story.extend(frame_table(frame, max_rows=16, max_cols=8))
        else:
            text = clean(dashboard.get("text") or "No tabular values were available for this custom dashboard.", 900)
            story.append(Paragraph(text, section_intro))
        story.append(PageBreak())
        section_number += 1

    if story and isinstance(story[-1], PageBreak):
        story.pop()
    doc.build(story, onFirstPage=cover_page, onLaterPages=header_footer)
    return output.getvalue()


def _build_boardroom_pdf(payload: dict[str, Any]) -> bytes:
    """Create a narrative, data-led leadership storybook as a native PDF."""
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen.canvas import Canvas
    from reportlab.platypus import Paragraph
    from PIL import Image, ImageOps
    from xml.sax.saxutils import escape as xml_escape

    analysis = payload.get("analysis") if isinstance(payload.get("analysis"), dict) else {}
    options = payload.get("options") if isinstance(payload.get("options"), dict) else {}
    metric = str(payload.get("metric") or "NPS").strip().upper()
    title = str(options.get("title") or f"{metric} Leadership Report").strip()
    account = str(options.get("accountName") or "Customer Intelligence Team").strip()
    prepared_for = str(options.get("preparedFor") or "Leadership Review").strip()
    footer_note = str(options.get("footerNote") or "Confidential - Krestrel Analysis Suite Review").strip()
    home_image_raw = str(options.get("homeImage") or "")
    separator_images_raw = options.get("separatorImages") if isinstance(options.get("separatorImages"), list) else []
    app_url = str(options.get("appUrl") or "").strip()
    selected = [str(item) for item in payload.get("tabs", []) if str(item).strip()]
    custom_dashboards = payload.get("customDashboards") if isinstance(payload.get("customDashboards"), list) else []
    summary = analysis.get("summary") if isinstance(analysis.get("summary"), dict) else {}
    counts = analysis.get("counts") if isinstance(analysis.get("counts"), dict) else {}
    sentiment = analysis.get("sentiment") if isinstance(analysis.get("sentiment"), dict) else {}
    weekly = analysis.get("weekly") if isinstance(analysis.get("weekly"), list) else []
    agents = analysis.get("agents") if isinstance(analysis.get("agents"), list) else []
    managers = analysis.get("managers") if isinstance(analysis.get("managers"), list) else []
    quartiles = analysis.get("quartiles") if isinstance(analysis.get("quartiles"), list) else []
    quartile_rollup = analysis.get("quartileRollup") if isinstance(analysis.get("quartileRollup"), list) else []
    quartile_weekly = analysis.get("quartileWeekly") if isinstance(analysis.get("quartileWeekly"), list) else []
    wave_rows = analysis.get("wave") if isinstance(analysis.get("wave"), list) else []
    tenure_rows = analysis.get("tenure") if isinstance(analysis.get("tenure"), list) else []
    dynamic_dimensions = analysis.get("dynamicDimensions") if isinstance(analysis.get("dynamicDimensions"), list) else []
    sentiment_movement = analysis.get("sentimentMovement") if isinstance(analysis.get("sentimentMovement"), list) else []
    themes = analysis.get("themes") if isinstance(analysis.get("themes"), list) else []
    reasons = analysis.get("reasons") if isinstance(analysis.get("reasons"), list) else []
    feedback_rows = analysis.get("feedbackRows") if isinstance(analysis.get("feedbackRows"), list) else []
    alerts = analysis.get("alerts") if isinstance(analysis.get("alerts"), list) else []
    business_rules = analysis.get("businessRules") if isinstance(analysis.get("businessRules"), dict) else {}
    leadership_questions_raw = analysis.get("leadershipQuestions") if isinstance(analysis.get("leadershipQuestions"), list) else []
    overall_performance_read = analysis.get("overallPerformanceRead") if isinstance(analysis.get("overallPerformanceRead"), dict) else {}
    overall_performance_sections = overall_performance_read.get("sections") if isinstance(overall_performance_read.get("sections"), list) else []
    evidence_relationship = analysis.get("evidenceRelationship") if isinstance(analysis.get("evidenceRelationship"), dict) else {}

    output = BytesIO()
    page_width, page_height = A4
    canvas = Canvas(output, pagesize=A4, pageCompression=1)
    canvas.setTitle(title)
    canvas.setAuthor("Krestrel Analysis Suite")
    navy = colors.HexColor("#263F73")
    deep_teal = colors.HexColor("#1E646B")
    teal = colors.HexColor("#00A3A3")
    cyan = colors.HexColor("#2AB7D6")
    lime = colors.HexColor("#C8E52A")
    green = colors.HexColor("#33B986")
    amber = colors.HexColor("#F2A43A")
    red = colors.HexColor("#D85D72")
    pale = colors.HexColor("#F5F7F8")
    pale_teal = colors.HexColor("#EAF7F7")
    pale_blue = colors.HexColor("#EDF4FA")
    pale_red = colors.HexColor("#FCEFF2")
    ink = colors.HexColor("#203B70")
    muted = colors.HexColor("#536B87")
    line = colors.HexColor("#D7E0E8")
    white = colors.white
    margin = 18 * mm
    content_width = page_width - 2 * margin
    generated_at = _dt.datetime.now().strftime("%d %b %Y, %I:%M %p")

    def decode_pdf_image(value: Any) -> Any:
        text = str(value or "")
        if not text.startswith("data:image/") or "," not in text or len(text) > 18_000_000:
            return None
        try:
            raw = base64.b64decode(text.split(",", 1)[1], validate=False)
            if not raw or len(raw) > 13_000_000:
                return None
            image = Image.open(BytesIO(raw))
            image = ImageOps.exif_transpose(image).convert("RGB")
            if image.width < 16 or image.height < 16:
                return None
            return image
        except Exception:
            return None

    home_image = decode_pdf_image(home_image_raw)
    separator_images = [image for image in (decode_pdf_image(value) for value in separator_images_raw[:5]) if image is not None]
    if home_image is None and separator_images:
        home_image = separator_images[0]

    def prepared_image(image: Any, width: float, height: float) -> Any:
        if image is None:
            return None
        target_width = max(320, min(1400, int(width * 2.1)))
        target_height = max(220, min(1200, int(height * 2.1)))
        fitted = ImageOps.fit(image, (target_width, target_height), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))
        buffer = BytesIO()
        fitted.save(buffer, format="JPEG", quality=86, optimize=True, progressive=True)
        buffer.seek(0)
        return ImageReader(buffer)

    def draw_editorial_image(image: Any, x: float, y: float, width: float, height: float, label: str = "") -> None:
        reader = prepared_image(image, width, height)
        if reader is None:
            return
        canvas.saveState()
        clip = canvas.beginPath()
        clip.rect(x, y, width, height)
        canvas.clipPath(clip, stroke=0, fill=0)
        canvas.drawImage(reader, x, y, width=width, height=height, preserveAspectRatio=False, mask="auto")
        canvas.restoreState()
        canvas.setStrokeColor(line)
        canvas.setLineWidth(0.8)
        canvas.rect(x, y, width, height, fill=0, stroke=1)
        if label:
            canvas.setFillColor(muted)
            canvas.setFont("Helvetica-Bold", 6.2)
            canvas.drawRightString(x + width, y - 3.5 * mm, label.upper()[:54])

    def chapter_image(index: int) -> Any:
        if not separator_images:
            return None
        return separator_images[index % len(separator_images)]

    def draw_section_signature(index: int, label: str) -> None:
        image = chapter_image(index)
        if image is not None:
            draw_editorial_image(image, page_width - margin - 88 * mm, 20 * mm, 88 * mm, 46 * mm, label)

    body_style = ParagraphStyle("StoryBody", fontName="Helvetica", fontSize=10.2, leading=15.2, textColor=ink, alignment=TA_LEFT)
    body_small = ParagraphStyle("StorySmall", fontName="Helvetica", fontSize=8.3, leading=12.2, textColor=muted, alignment=TA_LEFT)
    card_title_style = ParagraphStyle("CardTitle", fontName="Helvetica-Bold", fontSize=11.2, leading=14, textColor=navy)
    card_body_style = ParagraphStyle("CardBody", fontName="Helvetica", fontSize=8.8, leading=12.8, textColor=ink)
    callout_style = ParagraphStyle("Callout", fontName="Helvetica", fontSize=10.2, leading=15, textColor=white)
    appendix_question_style = ParagraphStyle("AppendixQuestion", fontName="Helvetica-Bold", fontSize=7.6, leading=9.2, textColor=navy)
    appendix_answer_style = ParagraphStyle("AppendixAnswer", fontName="Helvetica", fontSize=6.8, leading=8.2, textColor=muted)

    def safe(value: Any, limit: int = 500) -> str:
        if value is None:
            return ""
        text = re.sub(r"\s+", " ", str(value)).strip()
        return xml_escape(text[:limit])

    def numeric(value: Any, default: float = 0.0) -> float:
        if isinstance(value, (int, float)) and not isinstance(value, bool):
            return float(value)
        match = re.search(r"-?\d+(?:\.\d+)?", str(value or "").replace(",", ""))
        return float(match.group(0)) if match else default

    def pick(sources: list[dict[str, Any]], keys: list[str], default: Any = 0) -> Any:
        for source in sources:
            for key in keys:
                if source.get(key) not in (None, ""):
                    return source.get(key)
        return default

    def fmt(value: Any, decimals: int = 1, suffix: str = "") -> str:
        number = numeric(value)
        if decimals == 0:
            return f"{int(round(number)):,}{suffix}"
        return f"{number:,.{decimals}f}{suffix}"

    def row_value(row: dict[str, Any], keys: list[str], default: Any = 0) -> Any:
        return pick([row], keys, default)

    total = numeric(pick([summary, counts], ["total_responses", "Responses", "total", "Total", "Survey Volume"], 0))
    score = numeric(pick([summary], [metric.lower(), metric, "csat_score", "nps", "NPS", "CSAT"], 0))
    segment_names = ["Satisfied", "Neutral", "Dissatisfied"] if metric == "CSAT" else ["Promoters", "Passives", "Detractors"]
    segment_counts = [numeric(pick([counts, summary], [name, name[:-1] if name.endswith("s") else name, name.lower()], 0)) for name in segment_names]
    if not total:
        total = sum(segment_counts)
    segment_shares = [(value / total * 100 if total else 0) for value in segment_counts]
    positive = numeric(pick([sentiment], ["Positive", "positive"], 0))
    neutral_sentiment = numeric(pick([sentiment], ["Neutral", "neutral"], 0))
    negative = numeric(pick([sentiment], ["Negative", "negative"], 0))
    if positive + neutral_sentiment + negative > 100.5 and total:
        positive, neutral_sentiment, negative = (value / total * 100 for value in (positive, neutral_sentiment, negative))

    score_keys = [metric, "NPS", "CSAT", "Score", "Agent NPS", "Agent CSAT"]
    period_key = next((key for key in ["Week", "Month", "Period", "Date"] if any(key in row for row in weekly)), "Week")
    score_key = next((key for key in score_keys if any(key in row for row in weekly)), metric)
    trend_points: list[tuple[str, float, float]] = []
    for row in weekly[-12:]:
        if not isinstance(row, dict):
            continue
        trend_points.append((str(row.get(period_key) or "")[:14], numeric(row.get(score_key)), numeric(row.get("Responses"))))
    latest = trend_points[-1][1] if trend_points else score
    previous = trend_points[-2][1] if len(trend_points) > 1 else latest
    movement = latest - previous
    trend_values = [point[1] for point in trend_points]
    average = sum(trend_values) / len(trend_values) if trend_values else score
    volatility = float(pd.Series(trend_values).std()) if len(trend_values) > 1 else 0
    best_period = max(trend_points, key=lambda point: point[1]) if trend_points else ("Not available", score, total)
    worst_period = min(trend_points, key=lambda point: point[1]) if trend_points else ("Not available", score, total)

    driver_rows = themes or reasons
    driver_data: list[tuple[str, float, str]] = []
    for row in driver_rows[:8]:
        if not isinstance(row, dict):
            continue
        label = str(row_value(row, ["Theme", "Primary Reason", "Bucket Category", "Reason", "Category", "Driver"], "Unclassified"))
        value = numeric(row_value(row, ["Share", "Responses", "Count", "Detractors", "Negative"], 0))
        issue = str(row_value(row, ["Top Issue", "Owl Issue Type", "Issue"], ""))
        driver_data.append((label, value, issue))
    driver_data = sorted(driver_data, key=lambda item: item[1], reverse=True)[:5]
    top_driver = driver_data[0][0] if driver_data else "No classified driver available"

    minimum_sample = int(numeric(pick([business_rules], ["minimumSample", "minSample", "minimum_sample"], 5), 5)) or 5

    def people_rank(rows: list[dict[str, Any]], entity_keys: list[str]) -> list[tuple[str, float, float]]:
        ranked: list[tuple[str, float, float]] = []
        for row in rows:
            if not isinstance(row, dict):
                continue
            name = str(row_value(row, entity_keys, "Unknown"))
            responses = numeric(row_value(row, ["Responses", "Survey Volume", "Count", "Surveys"], 0))
            entity_score = numeric(row_value(row, [metric, f"Agent {metric}", f"Manager {metric}", "NPS", "CSAT", "Score"], 0))
            if responses >= minimum_sample:
                ranked.append((name, entity_score, responses))
        return sorted(ranked, key=lambda item: item[1], reverse=True)

    ranked_agents = people_rank(agents, ["Agent Name", "Agent", "Name"])
    ranked_managers = people_rank(managers, ["Manager/TL", "Manager", "TL Name", "Name"])
    top_agents = ranked_agents[:3]
    bottom_agents = list(reversed(ranked_agents[-3:])) if ranked_agents else []

    def dimension_profile(name: str, rows: list[dict[str, Any]], label_keys: list[str]) -> dict[str, Any] | None:
        values: list[tuple[str, float]] = []
        for row in rows:
            if not isinstance(row, dict):
                continue
            score_value = None
            for key in [metric, f"Agent {metric}", f"Manager {metric}", "NPS", "CSAT", "Score", "Average"]:
                if row.get(key) not in (None, ""):
                    score_value = numeric(row.get(key))
                    break
            if score_value is None:
                continue
            label_value = next((str(row.get(key)) for key in label_keys if row.get(key) not in (None, "")), "Unknown")
            values.append((label_value, score_value))
        if len(values) < 3:
            return None
        series = pd.Series([value for _label, value in values], dtype="float64")
        q1_value = float(series.quantile(0.25))
        q3_value = float(series.quantile(0.75))
        return {
            "name": name,
            "iqr": q3_value - q1_value,
            "q1": q1_value,
            "q3": q3_value,
            "high": max(values, key=lambda item: item[1]),
            "low": min(values, key=lambda item: item[1]),
            "groups": len(values),
        }

    dimension_profiles: list[dict[str, Any]] = []
    for profile in [
        dimension_profile("Agent", agents, ["Agent Name", "Agent", "Name"]),
        dimension_profile("Manager / Team Leader", managers, ["Manager/TL", "Manager", "Name"]),
        dimension_profile("Wave", wave_rows, ["Wave", "Name"]),
        dimension_profile("Tenure", tenure_rows, ["Tenure", "Name"]),
    ]:
        if profile:
            dimension_profiles.append(profile)
    for dimension in dynamic_dimensions:
        if not isinstance(dimension, dict):
            continue
        profile = dimension_profile(str(dimension.get("name") or "Custom dimension"), dimension.get("rows") if isinstance(dimension.get("rows"), list) else [], [str(dimension.get("name") or ""), "Dimension", "Name"])
        if profile:
            dimension_profiles.append(profile)
    variance_dimension = max(dimension_profiles, key=lambda item: item["iqr"]) if dimension_profiles else None

    quartile_source = quartiles or quartile_rollup
    quartile_points: list[tuple[str, float, float]] = []
    for row in quartile_source:
        if not isinstance(row, dict):
            continue
        label = str(row_value(row, ["Quartile", "Segment", "Group"], "Unknown"))
        quartile_score = numeric(row_value(row, [metric, "NPS", "CSAT", "Average", "Score"], 0))
        quartile_volume = numeric(row_value(row, ["Rows", "Responses", "Game Changers", "Count"], 0))
        quartile_points.append((label, quartile_score, quartile_volume))
    best_quartile = max(quartile_points, key=lambda item: item[1]) if quartile_points else None
    weakest_quartile = min(quartile_points, key=lambda item: item[1]) if quartile_points else None
    quartile_gap = (best_quartile[1] - weakest_quartile[1]) if best_quartile and weakest_quartile else None

    projection_slope = 0.0
    if len(trend_values) >= 3:
        x_mean = (len(trend_values) - 1) / 2
        y_mean = sum(trend_values) / len(trend_values)
        denominator = sum((index - x_mean) ** 2 for index in range(len(trend_values)))
        if denominator:
            projection_slope = sum((index - x_mean) * (value - y_mean) for index, value in enumerate(trend_values)) / denominator
    score_floor, score_ceiling = (-100.0, 100.0) if metric == "NPS" else (0.0, 100.0)
    projected_4_weeks = max(score_floor, min(score_ceiling, latest + projection_slope * 4)) if len(trend_values) >= 3 else None
    projected_3_months = max(score_floor, min(score_ceiling, latest + projection_slope * 13)) if len(trend_values) >= 3 else None

    sentiment_points: list[tuple[float, float, float]] = []
    for row in (sentiment_movement or weekly)[-12:]:
        if not isinstance(row, dict):
            continue
        positive_value = numeric(row_value(row, ["Positive", "Positive %", "positive"], 0))
        neutral_value = numeric(row_value(row, ["Neutral", "Neutral %", "neutral"], 0))
        negative_value = numeric(row_value(row, ["Negative", "Negative %", "negative"], 0))
        sentiment_total = positive_value + neutral_value + negative_value
        if sentiment_total > 100.5:
            positive_value = positive_value / sentiment_total * 100
            neutral_value = neutral_value / sentiment_total * 100
            negative_value = negative_value / sentiment_total * 100
        if sentiment_total:
            sentiment_points.append((positive_value, neutral_value, negative_value))
    sentiment_positive_change = sentiment_points[-1][0] - sentiment_points[0][0] if len(sentiment_points) >= 2 else None
    sentiment_negative_change = sentiment_points[-1][2] - sentiment_points[0][2] if len(sentiment_points) >= 2 else None
    if sentiment_positive_change is None or sentiment_negative_change is None:
        sentiment_trend_label = "Not available"
    elif sentiment_positive_change > 1 and sentiment_negative_change < -1:
        sentiment_trend_label = "Improving"
    elif sentiment_positive_change < -1 and sentiment_negative_change > 1:
        sentiment_trend_label = "Deteriorating"
    else:
        sentiment_trend_label = "Mixed / stable"

    word_stop_words = {
        "a", "an", "and", "are", "as", "at", "be", "been", "but", "by", "can", "could", "did", "do", "does", "for", "from", "had", "has", "have", "he", "her", "his", "i", "in", "is", "it", "its", "me", "my", "no", "not", "of", "on", "or", "our", "she", "so", "that", "the", "their", "them", "there", "they", "this", "to", "too", "us", "was", "we", "were", "with", "you", "your",
        "agent", "customer", "service", "call", "chat", "issue", "problem", "please", "thanks", "thank",
    }
    word_counts: dict[str, int] = {}
    rows_with_words = 0
    for row in feedback_rows:
        if not isinstance(row, dict):
            continue
        text_value = " ".join(str(row.get(key) or "") for key in ["Verbatim Feedback", "Customer Comments", "Feedback", "Comments", "Comment", "Customer Feedback", "Text"])
        tokens = {
            token for token in re.sub(r"[^a-z0-9\s]", " ", text_value.lower()).split()
            if len(token) > 2 and token not in word_stop_words and not token.isdigit()
        }
        if not tokens:
            continue
        rows_with_words += 1
        for token in tokens:
            word_counts[token] = word_counts.get(token, 0) + 1
    top_meaningful_word = max(word_counts.items(), key=lambda item: (item[1], item[0])) if word_counts else None
    top_word_share = (top_meaningful_word[1] / rows_with_words * 100) if top_meaningful_word and rows_with_words else None

    low_positive_mismatch = 0
    negative_high_mismatch = 0
    for row in feedback_rows:
        if not isinstance(row, dict):
            continue
        segment = str(row_value(row, ["CSAT Type", "NPS Type"], "")).lower()
        tone = str(row.get("Sentiment") or "").lower()
        if ("detractor" in segment or "dissatisfied" in segment) and tone in {"positive", "neutral"}:
            low_positive_mismatch += 1
        if ("promoter" in segment or "satisfied" in segment or "passive" in segment or "neutral" in segment) and tone == "negative":
            negative_high_mismatch += 1

    def score_read(value: float) -> str:
        if metric == "CSAT":
            if value >= 90:
                return "very strong"
            if value >= 80:
                return "healthy"
            if value >= 70:
                return "mixed"
            return "under pressure"
        if value >= 50:
            return "excellent"
        if value >= 20:
            return "positive"
        if value >= 0:
            return "fragile"
        return "under pressure"

    leadership_questions: list[dict[str, Any]] = []
    for index, item in enumerate(leadership_questions_raw, start=1):
        if not isinstance(item, dict):
            continue
        question_text = str(item.get("question") or "").strip()
        answer_text = str(item.get("text") or item.get("answer") or "").strip()
        if not question_text and not answer_text:
            continue
        status_text = str(item.get("status") or "Monitor").strip()
        category_text = str(item.get("category") or item.get("area") or ("Score" if index <= 50 else "Voice of Customer")).strip()
        leadership_questions.append({
            "number": item.get("number") or index,
            "category": category_text,
            "question": question_text or "Leadership evidence check",
            "answer": answer_text or "No answer was returned by the completed analysis.",
            "status": status_text,
            "method": str(item.get("method") or item.get("statistics") or "").strip(),
            "evidence": item.get("evidence") if isinstance(item.get("evidence"), list) else [],
        })

    def question_tone(item: dict[str, Any]) -> str:
        status_value = str(item.get("status") or "").lower()
        answer_value = str(item.get("answer") or "").lower()
        combined = f"{status_value} {answer_value}"
        if any(term in combined for term in ["no evidence", "not available", "insufficient", "not returned"]):
            return "unavailable"
        if any(term in status_value for term in ["action", "investigate", "attention", "review", "risk"]):
            return "attention"
        if any(term in status_value for term in ["strong", "positive", "no action", "healthy"]):
            return "strength"
        return "monitor"

    question_tones = [question_tone(item) for item in leadership_questions]
    question_counts = {tone: question_tones.count(tone) for tone in ["strength", "attention", "monitor", "unavailable"]}
    question_categories: dict[str, int] = {}
    for item in leadership_questions:
        category = str(item.get("category") or "Other")
        question_categories[category] = question_categories.get(category, 0) + 1
    question_category_rows = sorted(question_categories.items(), key=lambda pair: pair[1], reverse=True)[:4]
    question_strengths = [item for item in leadership_questions if question_tone(item) == "strength"][:4]
    question_attention = [item for item in leadership_questions if question_tone(item) == "attention"][:4]
    overall_highlights: list[str] = []
    for section in overall_performance_sections:
        if not isinstance(section, dict):
            continue
        section_heading = str(section.get("heading") or "Performance insight").strip()
        section_items = section.get("items") if isinstance(section.get("items"), list) else []
        if section_items:
            overall_highlights.append(f"{section_heading}: {str(section_items[0]).strip()}")
    overall_highlights = overall_highlights[:4]

    def paragraph(html: str, x: float, top: float, width: float, style=body_style, max_height: float = 120 * mm) -> float:
        item = Paragraph(html, style)
        _, height = item.wrap(width, max_height)
        item.drawOn(canvas, x, top - height)
        return top - height

    def page_background() -> None:
        canvas.setFillColor(colors.HexColor("#FBFBFC"))
        canvas.rect(0, 0, page_width, page_height, fill=1, stroke=0)
        canvas.setFillColor(colors.HexColor("#DFF5F3"))
        canvas.circle(page_width + 18 * mm, -12 * mm, 58 * mm, fill=1, stroke=0)
        canvas.setFillColor(colors.HexColor("#E9F4FB"))
        canvas.circle(page_width - 20 * mm, -40 * mm, 62 * mm, fill=1, stroke=0)
        canvas.setStrokeColor(lime)
        canvas.setLineWidth(2.2)
        canvas.line(page_width - 70 * mm, 0, page_width, 76 * mm)

    def page_header(page_number: int) -> None:
        page_background()
        canvas.setFont("Helvetica", 7.2)
        canvas.setFillColor(navy)
        canvas.drawRightString(page_width - margin, page_height - 12 * mm, f"{title}   |   {page_number}")
        canvas.setStrokeColor(line)
        canvas.line(page_width - margin - 1 * mm, page_height - 15 * mm, page_width - margin, page_height - 15 * mm)
        canvas.setFont("Helvetica", 6.8)
        canvas.setFillColor(muted)
        canvas.drawString(margin, 8 * mm, footer_note[:95])

    def heading(title_text: str, subtitle: str = "") -> float:
        top = page_height - 28 * mm
        canvas.setFillColor(deep_teal)
        canvas.setFont("Helvetica-Bold", 27)
        for index, line_text in enumerate(title_text.split("\n")):
            canvas.drawString(margin, top - index * 10 * mm, line_text)
        top -= max(1, len(title_text.split("\n"))) * 10 * mm + 2 * mm
        if subtitle:
            top = paragraph(safe(subtitle), margin, top, content_width, body_style)
        return top

    def rounded_card(x: float, y: float, width: float, height: float, fill=pale, stroke=line) -> None:
        canvas.setFillColor(colors.HexColor("#E3E6E9"))
        canvas.roundRect(x + 1.2 * mm, y - 1.2 * mm, width, height, 4 * mm, fill=1, stroke=0)
        canvas.setFillColor(fill)
        canvas.setStrokeColor(stroke)
        canvas.setLineWidth(0.6)
        canvas.roundRect(x, y, width, height, 4 * mm, fill=1, stroke=1)

    def kpi_card(x: float, y: float, width: float, label: str, value: str, note: str, accent=teal) -> None:
        height = 30 * mm
        rounded_card(x, y, width, height, white)
        canvas.setFillColor(accent)
        canvas.rect(x, y + height - 2 * mm, width, 2 * mm, fill=1, stroke=0)
        canvas.setFillColor(accent)
        canvas.setFont("Helvetica-Bold", 7.3)
        canvas.drawString(x + 5 * mm, y + height - 8 * mm, label[:30].upper())
        canvas.setFillColor(navy)
        value_text = str(value)
        value_font_size = 18 if len(value_text) <= 12 else 14 if len(value_text) <= 20 else 11
        canvas.setFont("Helvetica-Bold", value_font_size)
        canvas.drawString(x + 5 * mm, y + 10.5 * mm, value_text)
        paragraph(safe(note), x + 5 * mm, y + 8.5 * mm, width - 10 * mm, body_small, 8 * mm)

    def callout(x: float, y: float, width: float, height: float, title_text: str, body_html: str, fill=deep_teal) -> None:
        rounded_card(x, y, width, height, fill, fill)
        canvas.setFillColor(white)
        canvas.setFont("Helvetica-Bold", 11)
        canvas.drawString(x + 6 * mm, y + height - 10 * mm, title_text)
        paragraph(body_html, x + 6 * mm, y + height - 15 * mm, width - 12 * mm, callout_style, height - 20 * mm)

    def intelligence_card(x: float, y: float, width: float, height: float, kicker: str, headline: str, body_html: str, accent=teal, fill=white) -> None:
        rounded_card(x, y, width, height, fill)
        canvas.setFillColor(accent)
        canvas.rect(x, y + height - 2.2 * mm, width, 2.2 * mm, fill=1, stroke=0)
        canvas.setFillColor(accent)
        canvas.setFont("Helvetica-Bold", 7.2)
        canvas.drawString(x + 6 * mm, y + height - 9 * mm, kicker.upper()[:42])
        headline_style = ParagraphStyle("IntelligenceHeadline", fontName="Helvetica-Bold", fontSize=13, leading=15.5, textColor=navy)
        paragraph(safe(headline, 80), x + 6 * mm, y + height - 13 * mm, width - 12 * mm, headline_style, 18 * mm)
        paragraph(body_html, x + 6 * mm, y + height - 29 * mm, width - 12 * mm, card_body_style, height - 34 * mm)

    def horizontal_bars(x: float, top: float, width: float, data: list[tuple[str, float]], palette: list[Any], percent: bool = False) -> float:
        maximum = max([value for _, value in data] + [1])
        for index, (label, value) in enumerate(data):
            y = top - index * 15 * mm
            canvas.setFillColor(ink)
            canvas.setFont("Helvetica-Bold", 8.3)
            canvas.drawString(x, y, label[:32])
            canvas.setFillColor(colors.HexColor("#E4EAEE"))
            canvas.roundRect(x, y - 7 * mm, width, 4 * mm, 2 * mm, fill=1, stroke=0)
            bar_width = width * max(0, value) / maximum
            canvas.setFillColor(palette[index % len(palette)])
            canvas.roundRect(x, y - 7 * mm, max(1.5 * mm, bar_width), 4 * mm, 2 * mm, fill=1, stroke=0)
            canvas.setFillColor(muted)
            canvas.setFont("Helvetica-Bold", 7.5)
            display = f"{value:.1f}%" if percent else f"{value:,.0f}"
            canvas.drawRightString(x + width, y, display)
        return top - len(data) * 15 * mm

    def line_chart(x: float, y: float, width: float, height: float, points: list[tuple[str, float, float]]) -> None:
        rounded_card(x, y, width, height, white)
        if len(points) < 2:
            paragraph("Not enough period-level data is available to draw a trend.", x + 8 * mm, y + height - 12 * mm, width - 16 * mm, body_style)
            return
        chart_left, chart_right = x + 12 * mm, x + width - 7 * mm
        chart_bottom, chart_top = y + 17 * mm, y + height - 13 * mm
        values = [item[1] for item in points]
        low, high = min(values), max(values)
        pad = max(4, (high - low) * 0.18)
        low, high = low - pad, high + pad
        if metric == "NPS":
            low, high = min(low, -100), max(high, 100)
        for step in range(5):
            line_y = chart_bottom + (chart_top - chart_bottom) * step / 4
            canvas.setStrokeColor(colors.HexColor("#E1E7EC"))
            canvas.setLineWidth(0.5)
            canvas.line(chart_left, line_y, chart_right, line_y)
            canvas.setFillColor(muted)
            canvas.setFont("Helvetica", 6.5)
            canvas.drawRightString(chart_left - 2 * mm, line_y - 1.5 * mm, f"{low + (high-low)*step/4:.0f}")
        coords: list[tuple[float, float]] = []
        for index, (_label, value, _volume) in enumerate(points):
            point_x = chart_left + (chart_right - chart_left) * index / max(1, len(points) - 1)
            point_y = chart_bottom + (chart_top - chart_bottom) * (value - low) / max(0.001, high - low)
            coords.append((point_x, point_y))
        canvas.setStrokeColor(teal)
        canvas.setLineWidth(2.1)
        path = canvas.beginPath()
        path.moveTo(*coords[0])
        for point_x, point_y in coords[1:]:
            path.lineTo(point_x, point_y)
        canvas.drawPath(path, stroke=1, fill=0)
        for index, ((label, value, _volume), (point_x, point_y)) in enumerate(zip(points, coords)):
            canvas.setFillColor(teal)
            canvas.circle(point_x, point_y, 2.2, fill=1, stroke=0)
            canvas.setFillColor(ink)
            canvas.setFont("Helvetica-Bold", 6.3)
            canvas.drawCentredString(point_x, point_y + 3 * mm, f"{value:.1f}")
            if index % 2 == 0 or len(points) <= 8:
                canvas.setFillColor(muted)
                canvas.setFont("Helvetica", 5.7)
                canvas.drawCentredString(point_x, y + 7 * mm, label[:10])

    def end_page() -> None:
        canvas.showPage()

    custom_page_count = min(3, len(custom_dashboards))
    signal_horizon_page = 9
    evidence_relationship_page = 10
    question_summary_page = 11 if leadership_questions else 0
    custom_start_page = 12 if leadership_questions else 11
    question_appendix_page_count = math.ceil(len(leadership_questions) / 8) if leadership_questions else 0
    chapters = [
        ("Overall Performance Lens", 3),
        ("Performance trajectory", 4),
        ("Customer composition and sentiment", 5),
        ("Drivers of experience", 6),
        ("People performance and coaching", 7),
        ("Risk and score-sentiment alignment", 8),
        ("Signal Horizon", signal_horizon_page),
        ("Evidence & Relationship Intelligence", evidence_relationship_page),
    ]
    if leadership_questions:
        chapters.append(("100-question leadership results", question_summary_page))
    for index in range(custom_page_count):
        chapters.append((f"Custom dashboard insight {index + 1}", custom_start_page + index))
    action_page = custom_start_page + custom_page_count
    methodology_page = action_page + 1
    appendix_start_page = methodology_page + 1
    conclusion_page = methodology_page + question_appendix_page_count + 1
    chapters.extend([
        ("Recommended leadership actions", action_page),
        ("Methodology and evidence", methodology_page),
    ])
    if leadership_questions:
        chapters.append(("Detailed 100-question appendix", appendix_start_page))
    chapters.append(("Conclusion and next steps", conclusion_page))

    # Cover
    canvas.setFillColor(colors.HexColor("#FAFAFB"))
    canvas.rect(0, 0, page_width, page_height, fill=1, stroke=0)
    canvas.setFillColor(colors.HexColor("#DFF5F3"))
    canvas.circle(page_width + 10 * mm, 25 * mm, 88 * mm, fill=1, stroke=0)
    canvas.setFillColor(colors.HexColor("#54C8DB"))
    canvas.circle(page_width + 22 * mm, -18 * mm, 70 * mm, fill=1, stroke=0)
    canvas.setStrokeColor(lime)
    canvas.setLineWidth(3)
    canvas.line(30 * mm, 34 * mm, page_width, 128 * mm)
    if home_image is not None:
        draw_editorial_image(home_image, page_width - 106 * mm, 60 * mm, 88 * mm, 110 * mm, "Customer experience evidence")
    canvas.setFillColor(teal)
    canvas.rect(margin, page_height - 36 * mm, 9 * mm, 9 * mm, fill=1, stroke=0)
    canvas.setFillColor(cyan)
    canvas.rect(margin + 10 * mm, page_height - 36 * mm, 9 * mm, 9 * mm, fill=1, stroke=0)
    canvas.setFillColor(navy)
    canvas.setFont("Helvetica-Bold", 12)
    canvas.drawString(margin + 23 * mm, page_height - 33 * mm, "Krestrel Analysis Suite")
    canvas.setFillColor(navy)
    canvas.setFont("Helvetica-Bold", 37)
    canvas.drawString(margin, page_height - 82 * mm, metric)
    canvas.setFont("Helvetica-Bold", 27)
    canvas.drawString(margin, page_height - 96 * mm, "Customer experience leadership read")
    canvas.setFont("Helvetica", 14)
    canvas.drawString(margin, page_height - 111 * mm, "From performance signals to practical action")
    canvas.setFillColor(deep_teal)
    canvas.setFont("Helvetica-Bold", 9)
    canvas.drawString(margin, 57 * mm, prepared_for.upper()[:70])
    canvas.setFillColor(muted)
    canvas.setFont("Helvetica", 8.5)
    canvas.drawString(margin, 49 * mm, f"{account}  |  {generated_at}")
    canvas.drawString(margin, 15 * mm, footer_note[:100])
    end_page()

    # Contents
    page_header(2)
    top = heading("Contents", "A concise data story designed for leadership discussion, not dashboard replication.")
    top -= 7 * mm
    contents_step = min(20 * mm, 198 * mm / max(1, len(chapters)))
    for index, (chapter, chapter_page) in enumerate(chapters, start=1):
        y = top - (index - 1) * contents_step
        canvas.setFillColor(teal)
        canvas.setFont("Helvetica-Bold", 8)
        canvas.drawString(margin, y, f"{index:02d}")
        canvas.setFillColor(ink)
        canvas.setFont("Helvetica-Bold", 11)
        canvas.drawString(margin + 14 * mm, y, chapter)
        canvas.setStrokeColor(line)
        canvas.line(margin + 14 * mm, y - 3 * mm, page_width - margin - 12 * mm, y - 3 * mm)
        canvas.setFillColor(muted)
        canvas.setFont("Helvetica", 9)
        canvas.drawRightString(page_width - margin, y, str(chapter_page))
    end_page()

    # Overall Performance Lens executive summary
    page_header(3)
    top = heading("Overall Performance Lens", f"Executive summary: the current {metric} picture is {score_read(score)}. The evidence combines {int(total):,} customer surveys with sentiment, driver and people-performance signals.")
    card_width = (content_width - 8 * mm) / 3
    card_y = top - 39 * mm
    kpi_card(margin, card_y, card_width, metric, fmt(score), f"Overall score across the selected reporting period.", teal)
    kpi_card(margin + card_width + 4 * mm, card_y, card_width, "Survey volume", fmt(total, 0), "Customer responses supporting this leadership read.", cyan)
    kpi_card(margin + 2 * (card_width + 4 * mm), card_y, card_width, "Latest movement", f"{movement:+.1f} pts", "Latest reporting period compared with the prior period.", green if movement >= 0 else red)
    callout(margin, card_y - 55 * mm, content_width, 43 * mm, "The leadership read",
            f"Overall performance is <b>{score_read(score)}</b> at <b>{score:.1f}</b>. The latest period moved <b>{movement:+.1f} points</b>. "
            f"<b>{segment_names[-1]}</b> represent <b>{segment_shares[-1]:.1f}%</b> of classified responses, while negative sentiment is <b>{negative:.1f}%</b>. "
            f"The leading classified experience driver is <b>{safe(top_driver)}</b>.")
    canvas.setFillColor(navy)
    canvas.setFont("Helvetica-Bold", 12)
    canvas.drawString(margin, card_y - 68 * mm, "What leadership should discuss")
    priorities = overall_highlights[:3] or [
        f"Whether the {movement:+.1f}-point latest movement represents a sustained change or short-term variation.",
        f"Why {top_driver} is the leading classified experience driver and which part is operationally controllable.",
        f"Which teams or agents need coaching, monitoring or recognition based on sufficient survey evidence.",
    ]
    for index, item in enumerate(priorities, start=1):
        y = card_y - (75 + (index - 1) * 17) * mm
        canvas.setFillColor(teal)
        canvas.circle(margin + 3 * mm, y + 2 * mm, 3 * mm, fill=1, stroke=0)
        canvas.setFillColor(white)
        canvas.setFont("Helvetica-Bold", 7)
        canvas.drawCentredString(margin + 3 * mm, y + 0.5 * mm, str(index))
        paragraph(safe(item, 180), margin + 10 * mm, y + 5 * mm, content_width - 10 * mm, body_small, 15 * mm)
    draw_section_signature(0, "Leadership overview")
    end_page()

    # Performance trajectory
    page_header(4)
    top = heading("Performance trajectory", f"Performance averaged {average:.1f} across the available periods. The latest result is {latest:.1f}, compared with {previous:.1f} previously.")
    chart_y = top - 92 * mm
    line_chart(margin, chart_y, content_width, 83 * mm, trend_points)
    insight_y = chart_y - 52 * mm
    box_width = (content_width - 8 * mm) / 3
    insights = [
        ("Momentum", f"{movement:+.1f} pts", "Improvement in the latest period." if movement >= 0 else "A decline requiring recovery attention.", green if movement >= 0 else red),
        ("Range", f"{best_period[1] - worst_period[1]:.1f} pts", f"Best: {best_period[0]}; lowest: {worst_period[0]}.", cyan),
        ("Volatility", f"{volatility:.1f}", "Lower values indicate more consistent customer outcomes.", amber),
    ]
    for index, (label, value, note, accent) in enumerate(insights):
        kpi_card(margin + index * (box_width + 4 * mm), insight_y, box_width, label, value, note, accent)
    draw_section_signature(1, "Performance movement")
    end_page()

    # Composition and sentiment
    page_header(5)
    top = heading("Customer composition and sentiment", "Score categories show what customers selected; sentiment shows how customers expressed the experience in their own words.")
    column_width = (content_width - 12 * mm) / 2
    rounded_card(margin, top - 100 * mm, column_width, 90 * mm, white)
    rounded_card(margin + column_width + 12 * mm, top - 100 * mm, column_width, 90 * mm, white)
    canvas.setFillColor(navy)
    canvas.setFont("Helvetica-Bold", 12)
    canvas.drawString(margin + 7 * mm, top - 20 * mm, "Score composition")
    canvas.drawString(margin + column_width + 19 * mm, top - 20 * mm, "Verbatim sentiment")
    horizontal_bars(margin + 7 * mm, top - 33 * mm, column_width - 14 * mm, list(zip(segment_names, segment_shares)), [green, amber, red], True)
    horizontal_bars(margin + column_width + 19 * mm, top - 33 * mm, column_width - 14 * mm, [("Positive", positive), ("Neutral", neutral_sentiment), ("Negative", negative)], [green, amber, red], True)
    callout(margin, top - 150 * mm, content_width, 38 * mm, "Interpretation",
            f"The score mix contains <b>{segment_shares[-1]:.1f}% {safe(segment_names[-1].lower())}</b>, while negative sentiment is <b>{negative:.1f}%</b>. "
            f"The difference between score and sentiment is useful: score identifies the outcome; language helps explain the cause and whether it is agent, process, policy, product or technology related.", navy)
    draw_section_signature(2, "Customer voice")
    end_page()

    # Drivers
    page_header(6)
    top = heading("Drivers of experience", "The most frequent classified topics indicate where customer experience pressure is concentrated.")
    left_width = 92 * mm
    if driver_data:
        horizontal_bars(margin, top - 15 * mm, left_width, [(label, value) for label, value, _issue in driver_data], [teal, cyan, green, amber, red], False)
    else:
        paragraph("No classified theme or reason rows were returned in this run.", margin, top - 15 * mm, left_width, body_style)
    narrative_x = margin + left_width + 14 * mm
    rounded_card(narrative_x, top - 105 * mm, content_width - left_width - 14 * mm, 94 * mm, pale_teal)
    canvas.setFillColor(deep_teal)
    canvas.setFont("Helvetica-Bold", 12)
    canvas.drawString(narrative_x + 7 * mm, top - 22 * mm, "What the data says")
    if driver_data:
        driver_sentence = ", ".join(item[0] for item in driver_data[:3])
        driver_body = f"The leading themes are <b>{safe(driver_sentence)}</b>. <b>{safe(top_driver)}</b> has the largest available evidence base in this run. "
        if driver_data[0][2]:
            driver_body += f"Its most common associated issue is <b>{safe(driver_data[0][2])}</b>. "
        driver_body += "Leadership should separate structural causes from agent-controllable causes before assigning individual coaching actions."
    else:
        driver_body = "Theme evidence is not available. Use the Interactive HTML to confirm mapping and theme classification coverage before drawing driver conclusions."
    paragraph(driver_body, narrative_x + 7 * mm, top - 30 * mm, content_width - left_width - 28 * mm, card_body_style, 70 * mm)
    callout(margin, top - 153 * mm, content_width, 36 * mm, "Decision lens",
            f"Prioritise <b>{safe(top_driver)}</b> when it combines high response volume, negative sentiment and a controllable root cause. Validate representative interactions before changing policy, process or coaching plans.")
    draw_section_signature(3, "Experience drivers")
    end_page()

    # People performance
    page_header(7)
    top = heading("People performance and coaching", f"Rankings include only entities with at least {minimum_sample} responses, reducing the risk of conclusions based on very small samples.")
    panel_width = (content_width - 10 * mm) / 2
    panel_y = top - 103 * mm
    for panel_index, (panel_title, records, fill_color) in enumerate([("Practices to investigate", top_agents, pale_teal), ("Coaching priority", bottom_agents, pale_red)]):
        x = margin + panel_index * (panel_width + 10 * mm)
        rounded_card(x, panel_y, panel_width, 92 * mm, fill_color)
        canvas.setFillColor(deep_teal if panel_index == 0 else red)
        canvas.setFont("Helvetica-Bold", 12)
        canvas.drawString(x + 7 * mm, top - 23 * mm, panel_title)
        if not records:
            paragraph("No eligible agent ranking is available.", x + 7 * mm, top - 34 * mm, panel_width - 14 * mm, card_body_style)
        for index, (name, entity_score, responses) in enumerate(records, start=1):
            y = top - (37 + (index - 1) * 19) * mm
            canvas.setFillColor(teal if panel_index == 0 else red)
            canvas.circle(x + 8 * mm, y + 2 * mm, 3.2 * mm, fill=1, stroke=0)
            canvas.setFillColor(white)
            canvas.setFont("Helvetica-Bold", 7)
            canvas.drawCentredString(x + 8 * mm, y + 0.5 * mm, str(index))
            canvas.setFillColor(ink)
            canvas.setFont("Helvetica-Bold", 10)
            canvas.drawString(x + 15 * mm, y + 3 * mm, name[:28])
            canvas.setFont("Helvetica", 8)
            canvas.setFillColor(muted)
            canvas.drawString(x + 15 * mm, y - 2 * mm, f"{metric} {entity_score:.1f}  |  {responses:,.0f} responses")
    manager_note = f"{len(ranked_managers)} managers/team leaders and {len(ranked_agents)} agents meet the minimum sample requirement. "
    manager_note += "Use high performers as investigation candidates, not automatic best-practice examples; interaction-level validation is still required."
    callout(margin, panel_y - 48 * mm, content_width, 36 * mm, "How to use this page", manager_note, navy)
    draw_section_signature(4, "People and coaching")
    end_page()

    # Risk and alignment
    page_header(8)
    top = heading("Risk and score-sentiment alignment", "Misaligned score and verbatim signals identify interactions that deserve review before coaching or process action.")
    box_width = (content_width - 8 * mm) / 2
    kpi_card(margin, top - 42 * mm, box_width, "Low score / non-negative text", fmt(low_positive_mismatch, 0), "Review for policy, product, process or technology causes.", amber)
    kpi_card(margin + box_width + 8 * mm, top - 42 * mm, box_width, "Negative text / acceptable score", fmt(negative_high_mismatch, 0), "Hidden dissatisfaction may not be visible in the score alone.", red)
    rounded_card(margin, top - 135 * mm, content_width, 78 * mm, white)
    canvas.setFillColor(navy)
    canvas.setFont("Helvetica-Bold", 12)
    canvas.drawString(margin + 7 * mm, top - 69 * mm, "Risk signals to review")
    risk_items = []
    for alert in alerts[:4]:
        if isinstance(alert, dict):
            risk_items.append(str(row_value(alert, ["detail", "title", "message", "Alert"], "")))
    if not risk_items:
        risk_items = [
            f"Review {low_positive_mismatch:,} low-score interactions whose language is neutral or positive.",
            f"Review {negative_high_mismatch:,} acceptable-score interactions containing negative sentiment.",
            f"Validate whether {top_driver} is structurally driven or individually controllable.",
        ]
    for index, item in enumerate(risk_items[:4], start=1):
        y = top - (82 + (index - 1) * 14) * mm
        canvas.setFillColor(red if index == 1 else amber)
        canvas.rect(margin + 8 * mm, y - 1 * mm, 3 * mm, 3 * mm, fill=1, stroke=0)
        paragraph(safe(item), margin + 15 * mm, y + 4 * mm, content_width - 24 * mm, body_small, 12 * mm)
    callout(margin, top - 182 * mm, content_width, 34 * mm, "Control principle",
            "Do not convert every detractor or dissatisfied response into an agent coaching action. First determine whether the evidence points to people, process, policy, product or technology.")
    end_page()

    # Signal Horizon - variance, voice, outlook and sentiment
    page_header(signal_horizon_page)
    top = heading("Signal Horizon", "Variance, voice and forward view: five decision signals retrieved or derived from the completed analysis.")
    upper_width = (content_width - 8 * mm) / 2
    upper_y = top - 68 * mm
    if variance_dimension:
        variance_headline = f"{variance_dimension['name']}: {variance_dimension['iqr']:.1f} pts IQR"
        variance_body = (
            f"The widest interquartile variance is across <b>{safe(variance_dimension['name'])}</b>. The middle 50% of group scores spans "
            f"<b>{variance_dimension['q1']:.1f} to {variance_dimension['q3']:.1f}</b>. Highest: <b>{safe(variance_dimension['high'][0])} {variance_dimension['high'][1]:.1f}</b>; "
            f"lowest: <b>{safe(variance_dimension['low'][0])} {variance_dimension['low'][1]:.1f}</b>."
        )
    else:
        variance_headline = "Dimension IQR not available"
        variance_body = "At least three scored groups are required within a dimension before an interquartile variance is reported."
    intelligence_card(margin, upper_y, upper_width, 56 * mm, "Variance Radar", variance_headline, variance_body, cyan, pale_blue)

    if best_quartile and weakest_quartile:
        quartile_headline = f"{best_quartile[0]} leads; {weakest_quartile[0]} is the recovery quartile"
        quartile_body = (
            f"Best score: <b>{best_quartile[1]:.1f}</b>. Weakest score: <b>{weakest_quartile[1]:.1f}</b>. "
            f"The quartile gap is <b>{quartile_gap:.1f} points</b>. Use the top quartile as an investigation pool and the weakest quartile as a validation-led coaching focus."
        )
    else:
        quartile_headline = "Quartile evidence not available"
        quartile_body = "Quartile callouts require populated Quartile Intelligence rows from the completed analysis."
    intelligence_card(margin + upper_width + 8 * mm, upper_y, upper_width, 56 * mm, "Quartile Pulse", quartile_headline, quartile_body, green, pale_teal)

    lower_gap = 5 * mm
    lower_width = (content_width - 2 * lower_gap) / 3
    lower_y = upper_y - 79 * mm
    if top_meaningful_word:
        word_headline = f'"{top_meaningful_word[0]}"'
        word_body = f"Leading meaningful word after stop-word removal, appearing in <b>{top_meaningful_word[1]:,}</b> feedback records (<b>{top_word_share:.1f}%</b> of records with usable words)."
    else:
        word_headline = "No meaningful word available"
        word_body = "No usable verbatim tokens remained after blank text, short tokens, numbers and stop words were excluded."
    intelligence_card(margin, lower_y, lower_width, 68 * mm, "Voiceprint", word_headline, word_body, amber, white)

    if projected_4_weeks is not None and projected_3_months is not None:
        forecast_headline = f"4 weeks: {projected_4_weeks:.1f} | 3 months: {projected_3_months:.1f}"
        forecast_body = f"A straight-line continuation of the available weekly trend ({projection_slope:+.2f} points per week). This is a directional scenario, not a causal or model-based forecast."
    else:
        forecast_headline = "Trend projection unavailable"
        forecast_body = "At least three valid weekly score points are required to extend the same observed trend into future weeks and months."
    intelligence_card(margin + lower_width + lower_gap, lower_y, lower_width, 68 * mm, "Forward Curve", forecast_headline, forecast_body, teal, white)

    if sentiment_points:
        latest_positive, _latest_neutral, latest_negative = sentiment_points[-1]
        sentiment_headline = sentiment_trend_label
        if sentiment_positive_change is not None and sentiment_negative_change is not None:
            sentiment_body = f"Latest mix: <b>{latest_positive:.1f}% positive</b> and <b>{latest_negative:.1f}% negative</b>. Across the available trend, positive moved <b>{sentiment_positive_change:+.1f} pts</b> and negative moved <b>{sentiment_negative_change:+.1f} pts</b>."
        else:
            sentiment_body = f"Latest available sentiment is <b>{latest_positive:.1f}% positive</b> and <b>{latest_negative:.1f}% negative</b>; a trend comparison needs at least two periods."
    else:
        sentiment_headline = "Sentiment trend unavailable"
        sentiment_body = "Weekly or period-level positive, neutral and negative sentiment values were not returned in this analysis."
    intelligence_card(margin + 2 * (lower_width + lower_gap), lower_y, lower_width, 68 * mm, "Sentiment Current", sentiment_headline, sentiment_body, red if sentiment_trend_label == "Deteriorating" else green, white)

    callout(margin, 31 * mm, content_width, 30 * mm, "Interpretation guardrail",
            "IQR compares scored groups within each available dimension. Word intelligence excludes configured stop words. Forward values extend the same linear trend and should be used for scenario discussion, not as guaranteed outcomes.", navy)
    end_page()

    # Evidence & Relationship Intelligence
    page_header(evidence_relationship_page)
    evidence_dimensions = evidence_relationship.get("dimensions") if isinstance(evidence_relationship.get("dimensions"), list) else []
    evidence_dimensions = [item for item in evidence_dimensions if isinstance(item, dict)]
    evidence_dimensions.sort(key=lambda item: numeric(item.get("Effect Size"), -1), reverse=True)
    evidence_strongest = evidence_relationship.get("strongestDimension") if isinstance(evidence_relationship.get("strongestDimension"), dict) else (evidence_dimensions[0] if evidence_dimensions else {})
    evidence_usable = int(numeric(evidence_relationship.get("usableResponses"), total))
    evidence_score = numeric(evidence_relationship.get("score"), score)
    evidence_low = evidence_relationship.get("confidenceLow")
    evidence_high = evidence_relationship.get("confidenceHigh")
    evidence_margin = evidence_relationship.get("marginOfError")
    evidence_rating = str(evidence_relationship.get("evidenceRating") or "Insufficient")
    top = heading("Evidence & Relationship Intelligence", "Statistical confidence and relationship strength across the dimensions selected during Setup. Associations support investigation; they do not establish causation.")
    evidence_card_width = (content_width - 8 * mm) / 3
    evidence_card_y = top - 39 * mm
    kpi_card(margin, evidence_card_y, evidence_card_width, f"{metric} with 95% CI", fmt(evidence_score), f"{fmt(evidence_low)} to {fmt(evidence_high)}; margin +/- {fmt(evidence_margin)} pts.", teal)
    kpi_card(margin + evidence_card_width + 4 * mm, evidence_card_y, evidence_card_width, "Usable evidence", fmt(evidence_usable, 0), f"Evidence rating: {evidence_rating}.", cyan)
    kpi_card(margin + 2 * (evidence_card_width + 4 * mm), evidence_card_y, evidence_card_width, "Strongest dimension", str(evidence_strongest.get("Dimension") or "Not available")[:22], f"{evidence_strongest.get('Relationship') or 'Unavailable'} relationship; effect {fmt(evidence_strongest.get('Effect Size'), 4)}.", green)

    ranking_y = evidence_card_y - 95 * mm
    ranking_width = 92 * mm
    rounded_card(margin, ranking_y, ranking_width, 82 * mm, white)
    canvas.setFillColor(navy)
    canvas.setFont("Helvetica-Bold", 11.5)
    canvas.drawString(margin + 7 * mm, ranking_y + 70 * mm, "Selected-dimension relationship ranking")
    ranking_rows = [(str(item.get("Dimension") or "Unknown"), max(0.0, numeric(item.get("Effect Size")) * 100)) for item in evidence_dimensions[:6] if item.get("Effect Size") is not None]
    if ranking_rows:
        horizontal_bars(margin + 7 * mm, ranking_y + 59 * mm, ranking_width - 14 * mm, ranking_rows, [teal, cyan, green, amber, red, navy], False)
    else:
        paragraph("No selected dimension currently has enough eligible groups for an effect-size ranking.", margin + 7 * mm, ranking_y + 56 * mm, ranking_width - 14 * mm, body_small, 38 * mm)

    relationship_x = margin + ranking_width + 10 * mm
    relationship_width = content_width - ranking_width - 10 * mm
    rounded_card(relationship_x, ranking_y, relationship_width, 82 * mm, pale_teal)
    canvas.setFillColor(deep_teal)
    canvas.setFont("Helvetica-Bold", 11.5)
    canvas.drawString(relationship_x + 7 * mm, ranking_y + 70 * mm, "What leadership should take from this")
    if evidence_strongest:
        relationship_body = (
            f"<b>{safe(evidence_strongest.get('Dimension'))}</b> has the strongest measurable association with {metric}. "
            f"Its relationship is <b>{safe(evidence_strongest.get('Relationship') or 'unavailable')}</b> with effect size <b>{fmt(evidence_strongest.get('Effect Size'), 4)}</b>. "
            f"The eligible group score spread is <b>{fmt(evidence_strongest.get('Score Spread'))} points</b>, from <b>{safe(evidence_strongest.get('Lowest Value'))}</b> to <b>{safe(evidence_strongest.get('Highest Value'))}</b>. "
            "Use this as a prioritisation signal, then validate volume, operational context and representative interactions before assigning action."
        )
    else:
        relationship_body = "Selected dimensions were retained, but the current run does not provide enough eligible group variation for a reliable statistical relationship ranking."
    paragraph(relationship_body, relationship_x + 7 * mm, ranking_y + 61 * mm, relationship_width - 14 * mm, card_body_style, 51 * mm)

    detail_y = ranking_y - 58 * mm
    rounded_card(margin, detail_y, content_width, 47 * mm, white)
    canvas.setFillColor(navy)
    canvas.setFont("Helvetica-Bold", 10.5)
    canvas.drawString(margin + 7 * mm, detail_y + 36 * mm, "Evidence safeguards")
    safeguards = [
        f"All {len(evidence_relationship.get('selectedDimensions') or [])} Setup-selected dimensions remain visible, including fields with insufficient or missing evidence.",
        f"Group comparisons require at least {int(numeric(evidence_relationship.get('minimumSample'), 5))} responses per value.",
        "Pearson correlation and simple regression are reported only for genuinely numeric dimensions; categorical dimensions use effect size.",
        "Statistical significance and association do not prove causation.",
    ]
    for index, item in enumerate(safeguards):
        y = detail_y + (27 - index * 7.5) * mm
        canvas.setFillColor(teal)
        canvas.circle(margin + 8 * mm, y + 1 * mm, 1.5 * mm, fill=1, stroke=0)
        paragraph(item, margin + 12 * mm, y + 3 * mm, content_width - 19 * mm, body_small, 7 * mm)
    end_page()

    # Existing 100-question leadership results
    if leadership_questions:
        page_header(question_summary_page)
        top = heading("100-question leadership results", "A decision-oriented summary of the existing score and sentiment question framework. No questions are recalculated for this PDF.")
        result_card_width = (content_width - 8 * mm) / 3
        result_card_y = top - 39 * mm
        usable_count = len(leadership_questions) - question_counts["unavailable"]
        kpi_card(margin, result_card_y, result_card_width, "Questions returned", fmt(len(leadership_questions), 0), f"{usable_count} contain usable answer text.", teal)
        kpi_card(margin + result_card_width + 4 * mm, result_card_y, result_card_width, "Action / review", fmt(question_counts["attention"], 0), "Questions whose existing status calls for investigation or review.", red)
        kpi_card(margin + 2 * (result_card_width + 4 * mm), result_card_y, result_card_width, "Unavailable", fmt(question_counts["unavailable"], 0), "Questions without sufficient returned evidence.", amber)

        category_panel_y = result_card_y - 82 * mm
        category_panel_width = 96 * mm
        rounded_card(margin, category_panel_y, category_panel_width, 70 * mm, white)
        canvas.setFillColor(navy)
        canvas.setFont("Helvetica-Bold", 11.5)
        canvas.drawString(margin + 7 * mm, category_panel_y + 59 * mm, "Coverage by question area")
        if question_category_rows:
            horizontal_bars(margin + 7 * mm, category_panel_y + 49 * mm, category_panel_width - 14 * mm, [(label, float(value)) for label, value in question_category_rows], [teal, cyan, green, amber, red, navy], False)

        narrative_x = margin + category_panel_width + 10 * mm
        narrative_width = content_width - category_panel_width - 10 * mm
        rounded_card(narrative_x, category_panel_y, narrative_width, 70 * mm, pale_teal)
        canvas.setFillColor(deep_teal)
        canvas.setFont("Helvetica-Bold", 11.5)
        canvas.drawString(narrative_x + 7 * mm, category_panel_y + 59 * mm, "How leadership should read it")
        result_read = (
            f"The framework returned <b>{len(leadership_questions)}</b> answers: <b>{question_counts['attention']}</b> require action or review, "
            f"<b>{question_counts['monitor']}</b> should be monitored, <b>{question_counts['strength']}</b> indicate positive or no-action evidence, and "
            f"<b>{question_counts['unavailable']}</b> lack sufficient evidence. The appendix preserves the question-level answer and status."
        )
        paragraph(result_read, narrative_x + 7 * mm, category_panel_y + 50 * mm, narrative_width - 14 * mm, card_body_style, 46 * mm)

        focus_y = category_panel_y - 68 * mm
        focus_width = (content_width - 10 * mm) / 2
        for panel_index, (panel_title, items, panel_fill, accent) in enumerate([
            ("Evidence to sustain", question_strengths or [item for item in leadership_questions if question_tone(item) == "monitor"][:4], pale_teal, green),
            ("Questions needing attention", question_attention, pale_red, red),
        ]):
            x = margin + panel_index * (focus_width + 10 * mm)
            rounded_card(x, focus_y, focus_width, 56 * mm, panel_fill)
            canvas.setFillColor(accent)
            canvas.setFont("Helvetica-Bold", 10.5)
            canvas.drawString(x + 6 * mm, focus_y + 45 * mm, panel_title)
            if not items:
                paragraph("No questions carry this status in the current results.", x + 6 * mm, focus_y + 37 * mm, focus_width - 12 * mm, body_small, 25 * mm)
            for item_index, item in enumerate(items[:4]):
                y = focus_y + (34 - item_index * 9.5) * mm
                canvas.setFillColor(accent)
                canvas.circle(x + 7 * mm, y + 1 * mm, 1.5 * mm, fill=1, stroke=0)
                paragraph(safe(item.get("question"), 90), x + 11 * mm, y + 3 * mm, focus_width - 17 * mm, body_small, 8 * mm)
        end_page()

    # Custom dashboard narrative pages
    for custom_index, dashboard in enumerate(custom_dashboards[:3], start=1):
        page_number = custom_start_page + custom_index - 1
        page_header(page_number)
        dashboard_title = str(dashboard.get("title") or f"Custom dashboard {custom_index}")
        creator = str(dashboard.get("creator") or "Dashboard Maker")
        top = heading(dashboard_title, f"A user-configured view from {creator}, interpreted as part of the leadership story.")
        rows = dashboard.get("rows") if isinstance(dashboard.get("rows"), list) else []
        label_column = next(iter(rows[0].keys()), "Period") if rows else "Period"
        numeric_column = None
        if rows:
            for key in rows[0].keys():
                if key == label_column:
                    continue
                if any(re.search(r"-?\d", str(row.get(key) or "")) for row in rows):
                    numeric_column = key
                    break
        custom_points = []
        if numeric_column:
            for row in rows[:12]:
                custom_points.append((str(row.get(label_column) or "")[:12], numeric(row.get(numeric_column)), 0))
        line_chart(margin, top - 98 * mm, content_width, 88 * mm, custom_points)
        if custom_points:
            values = [point[1] for point in custom_points]
            custom_movement = values[-1] - values[-2] if len(values) > 1 else 0
            custom_read = f"The latest <b>{safe(numeric_column)}</b> value is <b>{values[-1]:.1f}</b>. It moved <b>{custom_movement:+.1f}</b> from the previous point. The displayed range is <b>{min(values):.1f} to {max(values):.1f}</b>."
        else:
            custom_read = safe(dashboard.get("text") or "No numeric series was available for charting in this custom dashboard.", 900)
        callout(margin, top - 145 * mm, content_width, 34 * mm, "Leadership interpretation", custom_read, navy)
        end_page()

    # Recommended actions
    page_header(action_page)
    top = heading("Recommended leadership actions", "Actions are derived from the available score, sentiment, driver, people and alignment evidence. Validate owners and targets during the review.")
    actions = [
        ("1", "Stabilise the latest movement", f"Investigate the {movement:+.1f}-point latest change and confirm whether the same direction appears across teams and dimensions.", "Next 1-2 weeks", teal),
        ("2", f"Address {top_driver}", "Review representative interactions, separate structural and controllable causes, and assign the correct operational owner.", "Next 2-4 weeks", cyan),
        ("3", "Target coaching with evidence", f"Prioritise eligible lower performers while validating the interaction evidence and preserving recognition for repeatable practices.", "Next 4 weeks", amber),
        ("4", "Close alignment blind spots", f"Review {low_positive_mismatch + negative_high_mismatch:,} score-sentiment mismatch cases and feed confirmed learning into QA and process governance.", "Monthly governance", red),
    ]
    action_height = 43 * mm
    for index, (number_text, action_title, action_body, timing, accent) in enumerate(actions):
        y = top - (index + 1) * action_height + 2 * mm
        rounded_card(margin, y, content_width, action_height - 7 * mm, white)
        canvas.setFillColor(accent)
        canvas.circle(margin + 10 * mm, y + 18 * mm, 6 * mm, fill=1, stroke=0)
        canvas.setFillColor(white)
        canvas.setFont("Helvetica-Bold", 10)
        canvas.drawCentredString(margin + 10 * mm, y + 15.5 * mm, number_text)
        canvas.setFillColor(navy)
        canvas.setFont("Helvetica-Bold", 11)
        canvas.drawString(margin + 21 * mm, y + 24 * mm, action_title[:72])
        paragraph(safe(action_body), margin + 21 * mm, y + 19 * mm, content_width - 58 * mm, card_body_style, 19 * mm)
        canvas.setFillColor(accent)
        canvas.setFont("Helvetica-Bold", 7.5)
        canvas.drawRightString(page_width - margin - 7 * mm, y + 18 * mm, timing.upper())
    draw_section_signature(0, "Leadership action")
    end_page()

    # Methodology
    page_header(methodology_page)
    top = heading("Methodology and evidence", "This report retrieves and interprets existing analysis outputs. It does not rerun or alter the underlying NPS or CSAT analysis.")
    evidence_items = [
        ("Evidence base", f"{int(total):,} analyzed customer surveys; {len(trend_points)} period-level trend points; {len(ranked_agents)} eligible agents; {len(ranked_managers)} eligible managers/team leaders."),
        ("Minimum sample", f"People rankings use a minimum sample of {minimum_sample} responses. Entities below this threshold are excluded from ranking."),
        ("Interpretation", "Score indicates the recorded customer outcome. Sentiment and classified drivers provide context but should be validated against representative interactions."),
        ("Detail retained", "The Interactive HTML remains the detailed evidence environment. This PDF summarises decision-relevant signals and retains the existing 100-question answers in its appendix."),
    ]
    for index, (label, description) in enumerate(evidence_items):
        y = top - (index + 1) * 39 * mm + 8 * mm
        rounded_card(margin, y, content_width, 31 * mm, white)
        canvas.setFillColor(teal)
        canvas.setFont("Helvetica-Bold", 9)
        canvas.drawString(margin + 7 * mm, y + 20 * mm, label.upper())
        paragraph(safe(description), margin + 7 * mm, y + 15 * mm, content_width - 14 * mm, card_body_style, 12 * mm)
    selected_names = [REPORT_TAB_TITLES.get(tab, tab) for tab in selected if tab != "customdashboards"]
    callout(margin, 30 * mm, content_width, 34 * mm, "Included report scope",
            safe(", ".join(selected_names[:16]) or "Executive narrative based on the completed analysis.", 850), navy)

    if leadership_questions:
        end_page()
        for appendix_index in range(question_appendix_page_count):
            start_index = appendix_index * 8
            page_questions = leadership_questions[start_index:start_index + 8]
            end_index = start_index + len(page_questions)
            appendix_page_number = appendix_start_page + appendix_index
            page_header(appendix_page_number)
            top = heading("Detailed question appendix", f"Questions {start_index + 1}-{end_index} of {len(leadership_questions)}. Answers and evidence statuses are retrieved directly from the existing Results output.")
            row_height = 24 * mm
            row_step = 26 * mm
            for row_index, item in enumerate(page_questions):
                y = top - (row_index + 1) * row_step + 2 * mm
                tone = question_tone(item)
                accent = red if tone == "attention" else green if tone == "strength" else amber if tone == "unavailable" else teal
                rounded_card(margin, y, content_width, row_height, white)
                canvas.setFillColor(accent)
                canvas.rect(margin, y, 2.2 * mm, row_height, fill=1, stroke=0)
                canvas.setFillColor(accent)
                canvas.setFont("Helvetica-Bold", 7)
                number_text = str(item.get("number") or start_index + row_index + 1)
                category_text = str(item.get("category") or "Leadership result")
                canvas.drawString(margin + 6 * mm, y + 18.5 * mm, f"Q{number_text}  |  {category_text[:34].upper()}")
                canvas.drawRightString(page_width - margin - 6 * mm, y + 18.5 * mm, str(item.get("status") or "Monitor")[:30].upper())
                paragraph(safe(item.get("question"), 110), margin + 6 * mm, y + 16 * mm, content_width - 12 * mm, appendix_question_style, 7 * mm)
                paragraph(safe(item.get("answer"), 190), margin + 6 * mm, y + 8.5 * mm, content_width - 12 * mm, appendix_answer_style, 7 * mm)
            if appendix_index < question_appendix_page_count - 1:
                end_page()

    # Conclusion is intentionally rendered last, after methodology and appendices.
    end_page()
    page_header(conclusion_page)
    top = heading("Conclusion and next steps", "A concise close for the leadership discussion, with the detailed evidence retained in the application.")
    conclusion_panel_width = 105 * mm
    conclusion_panel_y = top - 91 * mm
    rounded_card(margin, conclusion_panel_y, conclusion_panel_width, 78 * mm, pale_teal)
    canvas.setFillColor(deep_teal)
    canvas.setFont("Helvetica-Bold", 13)
    canvas.drawString(margin + 8 * mm, conclusion_panel_y + 64 * mm, "Overall conclusion")
    conclusion_text = (
        f"The completed analysis places {metric} at <b>{score:.1f}</b>, a position that is <b>{score_read(score)}</b>. "
        f"The latest period moved <b>{movement:+.1f} points</b>, while <b>{top_driver}</b> is the leading classified experience driver. "
        f"Negative sentiment is <b>{negative:.1f}%</b> and the {segment_names[-1].lower()} population represents <b>{segment_shares[-1]:.1f}%</b> of responses. "
        "Leadership should confirm the movement across teams, validate representative interactions, and assign structural and coaching actions to the correct owners."
    )
    paragraph(conclusion_text, margin + 8 * mm, conclusion_panel_y + 56 * mm, conclusion_panel_width - 16 * mm, card_body_style, 48 * mm)
    conclusion_image = chapter_image(2)
    if conclusion_image is None:
        conclusion_image = home_image
    if conclusion_image is not None:
        draw_editorial_image(conclusion_image, margin + conclusion_panel_width + 10 * mm, conclusion_panel_y, content_width - conclusion_panel_width - 10 * mm, 78 * mm, "Leadership close")

    conclusion_card_width = (content_width - 8 * mm) / 3
    conclusion_card_y = conclusion_panel_y - 42 * mm
    kpi_card(margin, conclusion_card_y, conclusion_card_width, "Evidence base", fmt(total, 0), "Analyzed customer responses supporting the report.", teal)
    kpi_card(margin + conclusion_card_width + 4 * mm, conclusion_card_y, conclusion_card_width, "Priority driver", str(top_driver)[:18], "Validate controllability before assigning action.", amber)
    kpi_card(margin + 2 * (conclusion_card_width + 4 * mm), conclusion_card_y, conclusion_card_width, "Questions retained", fmt(len(leadership_questions), 0), "Detailed Results evidence included in the appendix.", cyan)

    app_reference = "Krestrel Analysis Suite - Interactive Board Room HTML"
    if re.match(r"^https?://", app_url, flags=re.IGNORECASE):
        app_reference += f"<br/><font size='8'>{safe(app_url, 130)}</font>"
    callout(margin, 35 * mm, content_width, 47 * mm, "Continue the investigation in the app",
            f"Use <b>{app_reference}</b> for interactive filters, detailed tables, record-level evidence, role-based reads, custom dashboards, and the complete 100-question Results trail. This PDF is the leadership narrative; the application remains the source for drill-down analysis.", navy)
    canvas.save()
    return output.getvalue()


def _quartile_rollup(agent_df: pd.DataFrame) -> list[dict[str, Any]]:
    frame = _quartile_summary(agent_df)
    if frame.empty:
        return []
    rollup = (
        frame.groupby("Quartile")
        .agg(Game_Changers=("Agent Name", "size"), Average=("Agent NPS", "mean"))
        .reindex(["Q1", "Q2", "Q3", "Q4"])
        .fillna(0)
        .reset_index()
    )
    rollup["Game_Changers"] = rollup["Game_Changers"].astype(int)
    rollup["Average"] = rollup["Average"].round(1)
    return rollup.to_dict(orient="records")


def _quartile_weekly_trend(analyzed_df: pd.DataFrame) -> list[dict[str, Any]]:
    if analyzed_df.empty or "Feedback Date" not in analyzed_df.columns or "Agent Name" not in analyzed_df.columns:
        return []
    working = analyzed_df.dropna(subset=["Feedback Date"]).copy()
    if working.empty:
        return []
    working["Week"] = week_period_start(working["Feedback Date"], _state_calendar_settings()["weekStart"])
    rows: list[dict[str, Any]] = []
    for week, week_df in working.dropna(subset=["Week"]).groupby("Week"):
        agent_frame = (
            week_df.groupby("Agent Name")
            .agg(
                Responses=("Agent Name", "size"),
                Promoters=("NPS Type", lambda values: int((values == "Promoter").sum())),
                Detractors=("NPS Type", lambda values: int((values == "Detractor").sum())),
            )
            .reset_index()
        )
        if agent_frame.empty:
            continue
        agent_frame["Agent NPS"] = ((agent_frame["Promoters"] / agent_frame["Responses"] * 100) - (agent_frame["Detractors"] / agent_frame["Responses"] * 100)).round(1)
        quartiles = _quartile_summary(agent_frame)
        if quartiles.empty:
            continue
        record: dict[str, Any] = {"Week": week}
        for quartile in ["Q1", "Q2", "Q3", "Q4"]:
            subset = quartiles[quartiles["Quartile"] == quartile]
            record[quartile] = round(float(subset["Agent NPS"].mean()), 1) if not subset.empty else 0
        rows.append(record)
    return rows


def _theme_summary(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return pd.DataFrame()
    driver_col = _first_existing_column(df, ["Owl Primary Driver", "Primary Driver", "Bucket Category"])
    issue_col = _first_existing_column(df, ["Owl Issue Type", "Issue Type", "Primary Reason"])
    if driver_col is None:
        return pd.DataFrame()
    grouped = (
        df.groupby(driver_col)
        .agg(
            Responses=(driver_col, "size"),
            Detractors=("NPS Type", lambda values: int((values == "Detractor").sum())),
            Negative=("Sentiment", lambda values: int((values == "Negative").sum())),
        )
        .reset_index()
        .rename(columns={driver_col: "Theme"})
        .sort_values(["Detractors", "Negative", "Responses"], ascending=[False, False, False])
    )
    grouped["Share"] = (grouped["Responses"] / max(len(df), 1) * 100).round(1)
    if issue_col:
        top_issue = df.groupby([driver_col, issue_col]).size().reset_index(name="Count").sort_values("Count", ascending=False)
        top_issue = top_issue.drop_duplicates(driver_col).set_index(driver_col)[issue_col].to_dict()
        grouped["Top Issue"] = grouped["Theme"].map(top_issue).fillna("")
    return grouped


def _analysis_summary_rows(analyzed_df: pd.DataFrame, weekly_df: pd.DataFrame) -> list[dict[str, Any]]:
    if analyzed_df.empty:
        return []
    summary = nps_summary(analyzed_df)
    rows = [
        {"View": "Current NPS", "Value": f"{float(summary.get('nps', 0)):.1f}", "Interpretation": "Selected date range NPS."},
        {"View": "Responses", "Value": f"{int(summary.get('total', 0)):,}", "Interpretation": "Survey volume in the selected range."},
        {"View": "Promoters", "Value": f"{float(summary.get('promoters', 0)):.1f}%", "Interpretation": "Promoter share in the selected range."},
        {"View": "Detractors", "Value": f"{float(summary.get('detractors', 0)):.1f}%", "Interpretation": "Detractor share in the selected range."},
    ]
    if not weekly_df.empty and len(weekly_df) >= 2:
        previous = float(weekly_df.iloc[-2].get("NPS", 0) or 0)
        current = float(weekly_df.iloc[-1].get("NPS", 0) or 0)
        rows.append(
            {
                "View": "Latest Week Movement",
                "Value": f"{current - previous:+.1f} pts",
                "Interpretation": f"Latest week NPS {current:.1f} versus previous week {previous:.1f}.",
            }
        )
    return rows


def _sentiment_movement_rows(weekly_df: pd.DataFrame) -> list[dict[str, Any]]:
    if weekly_df.empty:
        return []
    available = [column for column in ["Week", "Positive", "Neutral", "Negative", "NPS", "Responses"] if column in weekly_df.columns]
    return _safe_records(weekly_df[available], 52)


def _consistency_rows(weekly_df: pd.DataFrame) -> list[dict[str, Any]]:
    if weekly_df.empty or "NPS" not in weekly_df.columns:
        return []
    nps_values = pd.to_numeric(weekly_df["NPS"], errors="coerce").dropna()
    if nps_values.empty:
        return []
    return [
        {"Metric": "Average Weekly NPS", "Value": f"{nps_values.mean():.1f}", "Leadership Read": "Central weekly performance level."},
        {"Metric": "Best Week", "Value": f"{nps_values.max():.1f}", "Leadership Read": "Peak weekly NPS in the selected range."},
        {"Metric": "Lowest Week", "Value": f"{nps_values.min():.1f}", "Leadership Read": "Lowest weekly NPS requiring review."},
        {"Metric": "Volatility", "Value": f"{nps_values.std(ddof=0):.1f}", "Leadership Read": "Higher volatility means less predictable customer experience."},
    ]


def _formula_rows() -> list[dict[str, Any]]:
    return [
        {"Formula": "NPS", "Definition": "% Promoters minus % Detractors."},
        {"Formula": "Promoter", "Definition": "NPS score 9 or 10."},
        {"Formula": "Passive", "Definition": "NPS score 7 or 8."},
        {"Formula": "Detractor", "Definition": "NPS score 0 to 6."},
        {"Formula": "Quartile", "Definition": "Agents ranked by Agent NPS within the selected date range."},
    ]


def _operations_summary(df: pd.DataFrame) -> list[dict[str, Any]]:
    if df.empty:
        return []
    rows: list[dict[str, Any]] = []
    rows.append({"Metric": "Responses in Scope", "Value": f"{len(df):,}", "Signal": "Volume available for leadership read."})
    silent = int(pd.to_numeric(df.get("Silent Detractor Alert", pd.Series(dtype=int)), errors="coerce").fillna(0).sum())
    rows.append({"Metric": "Silent Detractors", "Value": f"{silent:,}", "Signal": "Soft comments paired with low NPS scores."})
    if "Feedback Date" in df.columns:
        dated = pd.to_datetime(df["Feedback Date"], errors="coerce").dropna()
        if not dated.empty:
            rows.append({"Metric": "Date Coverage", "Value": f"{dated.min().date()} to {dated.max().date()}", "Signal": "Current analysis period."})
    if "Agent Name" in df.columns:
        agents = df["Agent Name"].nunique()
        rows.append({"Metric": "Agents Covered", "Value": f"{agents:,}", "Signal": "Unique agents represented in the file."})
    return rows


def _churn_summary_rows(df: pd.DataFrame) -> list[dict[str, Any]]:
    if df.empty:
        return []
    high_risk = df[(df["NPS Type"] == "Detractor") | (pd.to_numeric(df.get("Silent Detractor Alert", pd.Series(dtype=int)), errors="coerce").fillna(0) > 0)]
    return [
        {"Risk Level": "High", "Responses": int(len(high_risk)), "Definition": "Detractors and silent detractor signals."},
        {"Risk Level": "Medium", "Responses": int((df["NPS Type"] == "Passive").sum()), "Definition": "Passive responses that may convert either way."},
        {"Risk Level": "Low", "Responses": int((df["NPS Type"] == "Promoter").sum()), "Definition": "Promoter responses with lower churn concern."},
    ]


def _tenure_bucket(value: Any) -> str:
    if value is None or pd.isna(value):
        return "Unknown"
    text = str(value).strip()
    if not text or text.lower() in {"nan", "none", "-", "unknown"}:
        return "Unknown"
    normalized = text.lower().replace(" ", "")
    known = {
        "0-30": "0-30",
        "31-60": "31-60",
        "61-90": "61-90",
        "90-180": "91-180",
        "91-180": "91-180",
        "181-365": "181-365",
        "180-365": "181-365",
        ">1year": ">1 year",
        "1+year": ">1 year",
        ">2years": ">2 years",
        "2+years": ">2 years",
        ">3years": ">3 years",
        "3+years": ">3 years",
    }
    if normalized in known:
        return known[normalized]
    numeric = pd.to_numeric(pd.Series([text]), errors="coerce").iloc[0]
    if pd.isna(numeric):
        extracted = "".join(ch if ch.isdigit() or ch in ".-" else " " for ch in text).split()
        numeric = pd.to_numeric(pd.Series([extracted[0] if extracted else None]), errors="coerce").iloc[0]
    if pd.isna(numeric):
        return text
    days = float(numeric)
    if days <= 30:
        return "0-30"
    if days <= 60:
        return "31-60"
    if days <= 90:
        return "61-90"
    if days <= 180:
        return "91-180"
    if days <= 365:
        return "181-365"
    if days <= 730:
        return ">1 year"
    if days <= 1095:
        return ">2 years"
    return ">3 years"


def _dimension_summary(df: pd.DataFrame, column: str | None, label: str) -> pd.DataFrame:
    if df.empty or not column or column not in df.columns:
        return pd.DataFrame()
    working = df.copy()
    if label == "Tenure":
        working[column] = working[column].apply(_tenure_bucket)
    else:
        working[column] = working[column].fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
    grouped = (
        working.groupby(column)
        .agg(
            Responses=(column, "size"),
            Promoters=("NPS Type", lambda values: int((values == "Promoter").sum())),
            Passives=("NPS Type", lambda values: int((values == "Passive").sum())),
            Detractors=("NPS Type", lambda values: int((values == "Detractor").sum())),
            Avg_Rating=("NPS Score", "mean"),
            Positive=("Sentiment", lambda values: int((values == "Positive").sum())),
            Neutral=("Sentiment", lambda values: int((values == "Neutral").sum())),
            Negative=("Sentiment", lambda values: int((values == "Negative").sum())),
            Silent_Detractors=("Silent Detractor Alert", lambda values: int(pd.to_numeric(values, errors="coerce").fillna(0).sum())),
        )
        .reset_index()
        .rename(columns={column: label})
    )
    grouped["NPS"] = ((grouped["Promoters"] / grouped["Responses"] * 100) - (grouped["Detractors"] / grouped["Responses"] * 100)).round(2)
    grouped["Promoter %"] = (grouped["Promoters"] / grouped["Responses"] * 100).round(2)
    grouped["Negative %"] = (grouped["Negative"] / grouped["Responses"] * 100).round(2)
    grouped["Avg_Rating"] = grouped["Avg_Rating"].round(2)
    drivers = []
    if "Primary Reason" in working.columns:
        for value in grouped[label].tolist():
            subset = working[working[column].astype(str).str.strip().replace("", "Unknown") == str(value)]
            drivers.append(str(subset["Primary Reason"].mode().iloc[0]) if not subset.empty and not subset["Primary Reason"].mode().empty else "")
        grouped["Top Driver"] = drivers
    return grouped.sort_values(["NPS", "Responses"], ascending=[True, False])


def _dimension_cards(df: pd.DataFrame, summary_df: pd.DataFrame, label: str) -> list[dict[str, Any]]:
    if df.empty or summary_df.empty:
        return []
    best = summary_df.sort_values(["NPS", "Responses"], ascending=[False, False]).head(1)
    risk = summary_df.sort_values(["NPS", "Responses"], ascending=[True, False]).head(1)
    high_negative = summary_df.sort_values(["Negative %", "Responses"], ascending=[False, False]).head(1)
    covered_label = "Waves Covered" if label == "Wave" else "Tenure Bands Covered"
    return [
        {"Metric": covered_label, "Value": int(summary_df[label].nunique()), "Comment": "Unique segments found in the selected data."},
        {"Metric": "Highest NPS", "Value": f"{best.iloc[0]['NPS']:.1f}" if not best.empty else "0.0", "Comment": str(best.iloc[0][label]) if not best.empty else "No segment available."},
        {"Metric": "Lowest NPS", "Value": f"{risk.iloc[0]['NPS']:.1f}" if not risk.empty else "0.0", "Comment": str(risk.iloc[0][label]) if not risk.empty else "No segment available."},
        {"Metric": "Highest Negative Share", "Value": f"{high_negative.iloc[0]['Negative %']:.1f}%" if not high_negative.empty else "0.0%", "Comment": str(high_negative.iloc[0][label]) if not high_negative.empty else "No segment available."},
    ]


def _dynamic_dimension_payload(df: pd.DataFrame, selected_columns: list[str]) -> list[dict[str, Any]]:
    payload: list[dict[str, Any]] = []
    seen: set[str] = set()
    for column in selected_columns:
        label = str(column or "").strip()
        if not label or label in seen or label not in df.columns:
            continue
        seen.add(label)
        summary_df = _dimension_summary(df, label, label)
        payload.append(
            {
                "name": label,
                "rows": _safe_records(summary_df, 100),
                "cards": _dimension_cards(df, summary_df, label),
            }
        )
    return payload


def _selected_dimension_columns(df: pd.DataFrame, dynamic_dimensions: list[str], mapping: dict[str, Any]) -> list[str]:
    candidates = list(dynamic_dimensions or [])
    for mapping_key, canonical in (("agent", "Agent Name"), ("manager", "Manager/TL"), ("wave", "Wave"), ("tenure", "Tenure")):
        if str(mapping.get(mapping_key) or "").strip():
            candidates.append(canonical)
    selected: list[str] = []
    seen: set[str] = set()
    for candidate in candidates:
        label = str(candidate or "").strip()
        if not label or label in seen:
            continue
        seen.add(label)
        selected.append(label)
    return selected


def _evidence_relationship_payload(
    df: pd.DataFrame,
    selected_dimensions: list[str],
    is_csat: bool,
    minimum_sample: int = 5,
) -> dict[str, Any]:
    metric = "CSAT" if is_csat else "NPS"
    empty = {
        "metric": metric,
        "minimumSample": minimum_sample,
        "selectedDimensions": list(selected_dimensions or []),
        "usableResponses": 0,
        "score": None,
        "confidenceLow": None,
        "confidenceHigh": None,
        "marginOfError": None,
        "evidenceRating": "Insufficient",
        "dimensions": [],
        "strongestDimension": None,
    }
    if df.empty:
        return empty

    type_column = "CSAT Type" if is_csat and "CSAT Type" in df.columns else "NPS Type"
    if type_column not in df.columns:
        return empty
    labels = df[type_column].fillna("").astype(str).str.strip().str.lower()
    if is_csat:
        outcome = labels.map({"satisfied": 100.0, "neutral": 0.0, "dissatisfied": 0.0, "promoter": 100.0, "passive": 0.0, "detractor": 0.0})
    else:
        outcome = labels.map({"promoter": 100.0, "passive": 0.0, "detractor": -100.0})
    usable = outcome.dropna()
    if usable.empty:
        return empty
    score = float(usable.mean())
    standard_error = float(usable.std(ddof=1) / math.sqrt(len(usable))) if len(usable) > 1 else 0.0
    margin = 1.96 * standard_error
    lower_bound = 0.0 if is_csat else -100.0
    confidence_low = max(lower_bound, score - margin)
    confidence_high = min(100.0, score + margin)
    evidence_rating = "Strong" if len(usable) >= 1000 else "Moderate" if len(usable) >= 300 else "Directional" if len(usable) >= 50 else "Insufficient"

    rows: list[dict[str, Any]] = []
    total_variation = float(((usable - score) ** 2).sum())
    for dimension in selected_dimensions:
        if dimension not in df.columns:
            rows.append({
                "Dimension": dimension,
                "Status": "Not returned in analysis output",
                "Populated": 0,
                "Missing %": 100.0,
                "Values": 0,
                "Eligible Values": 0,
                "Effect Size": None,
                "Relationship": "Unavailable",
                "Score Spread": None,
                "Highest Value": "Not available",
                "Highest Score": None,
                "Lowest Value": "Not available",
                "Lowest Score": None,
                "Pearson r": None,
                "Regression Slope": None,
                "R Squared": None,
                "Numeric": False,
            })
            continue

        raw_dimension = df.loc[usable.index, dimension]
        text_dimension = raw_dimension.fillna("").astype(str).str.strip()
        populated_mask = text_dimension.ne("") & text_dimension.str.lower().ne("unknown")
        populated_count = int(populated_mask.sum())
        missing_pct = round((1 - populated_count / max(len(usable), 1)) * 100, 2)
        value_count = int(text_dimension[populated_mask].nunique())
        grouped_scores: list[dict[str, Any]] = []
        if populated_count:
            grouped_frame = pd.DataFrame({"Dimension": text_dimension[populated_mask], "Outcome": usable.loc[text_dimension[populated_mask].index]})
            for value, group in grouped_frame.groupby("Dimension", dropna=False):
                volume = int(len(group))
                if volume < minimum_sample:
                    continue
                grouped_scores.append({"value": str(value), "score": float(group["Outcome"].mean()), "volume": volume})
        grouped_scores.sort(key=lambda item: (item["score"], item["volume"]), reverse=True)
        eligible_count = len(grouped_scores)
        highest = grouped_scores[0] if grouped_scores else None
        lowest = grouped_scores[-1] if grouped_scores else None
        spread = float(highest["score"] - lowest["score"]) if highest and lowest and eligible_count > 1 else None

        eligible_values = {item["value"] for item in grouped_scores}
        eligible_mask = populated_mask & text_dimension.isin(eligible_values)
        between_variation = 0.0
        if eligible_mask.any() and total_variation > 0:
            eligible_frame = pd.DataFrame({"Dimension": text_dimension[eligible_mask], "Outcome": usable.loc[text_dimension[eligible_mask].index]})
            eligible_mean = float(eligible_frame["Outcome"].mean())
            eligible_total_variation = float(((eligible_frame["Outcome"] - eligible_mean) ** 2).sum())
            if eligible_total_variation > 0:
                for _value, group in eligible_frame.groupby("Dimension"):
                    between_variation += len(group) * (float(group["Outcome"].mean()) - eligible_mean) ** 2
                effect_size = max(0.0, min(1.0, between_variation / eligible_total_variation))
            else:
                effect_size = 0.0
        else:
            effect_size = None
        relationship = "Unavailable" if effect_size is None else "Strong" if effect_size >= 0.14 else "Moderate" if effect_size >= 0.06 else "Weak" if effect_size >= 0.01 else "Minimal"

        numeric_values = pd.to_numeric(raw_dimension, errors="coerce")
        numeric_quality = float(numeric_values.notna().sum() / max(populated_count, 1)) if populated_count else 0.0
        numeric_mask = numeric_values.notna() & outcome.loc[numeric_values.index].notna()
        is_numeric = numeric_quality >= 0.8 and int(numeric_values[numeric_mask].nunique()) >= 3
        pearson = slope = r_squared = None
        if is_numeric and int(numeric_mask.sum()) >= max(minimum_sample * 2, 10):
            x = numeric_values[numeric_mask].astype(float)
            y = outcome.loc[x.index].astype(float)
            x_variance = float(((x - x.mean()) ** 2).sum())
            if x_variance > 0:
                pearson_value = x.corr(y)
                covariance = float(((x - x.mean()) * (y - y.mean())).sum())
                slope_value = covariance / x_variance
                if pd.notna(pearson_value):
                    pearson = float(pearson_value)
                    slope = float(slope_value)
                    r_squared = float(pearson_value ** 2)

        status = "Ready"
        if populated_count == 0:
            status = "No populated values"
        elif value_count < 2:
            status = "No measurable variation"
        elif is_numeric and pearson is not None:
            status = "Ready - numeric relationship"
        elif eligible_count < 2:
            status = "Insufficient sample across values"
        elif value_count > 250:
            status = "High-cardinality; summarized with safeguards"
        rows.append({
            "Dimension": dimension,
            "Status": status,
            "Populated": populated_count,
            "Missing %": missing_pct,
            "Values": value_count,
            "Eligible Values": eligible_count,
            "Effect Size": round(effect_size, 4) if effect_size is not None else None,
            "Relationship": relationship,
            "Score Spread": round(spread, 2) if spread is not None else None,
            "Highest Value": highest["value"] if highest else "Not available",
            "Highest Score": round(highest["score"], 2) if highest else None,
            "Lowest Value": lowest["value"] if lowest else "Not available",
            "Lowest Score": round(lowest["score"], 2) if lowest else None,
            "Pearson r": round(pearson, 4) if pearson is not None else None,
            "Regression Slope": round(slope, 4) if slope is not None else None,
            "R Squared": round(r_squared, 4) if r_squared is not None else None,
            "Numeric": is_numeric,
        })

    rows.sort(key=lambda item: (item.get("Effect Size") is not None, item.get("Effect Size") or -1, item.get("Populated") or 0), reverse=True)
    strongest = next((row for row in rows if row.get("Effect Size") is not None), None)
    return {
        "metric": metric,
        "minimumSample": minimum_sample,
        "selectedDimensions": list(selected_dimensions or []),
        "usableResponses": int(len(usable)),
        "score": round(score, 2),
        "confidenceLow": round(confidence_low, 2),
        "confidenceHigh": round(confidence_high, 2),
        "marginOfError": round(margin, 2),
        "evidenceRating": evidence_rating,
        "dimensions": rows,
        "strongestDimension": strongest,
    }


def _numeric_quality(series: pd.Series) -> tuple[pd.Series, float]:
    populated = series[(series.notna()) & (series.astype(str).str.strip() != "")]
    numeric = pd.to_numeric(populated, errors="coerce")
    quality = 0.0 if populated.empty else float(numeric.notna().sum() / len(populated))
    return pd.to_numeric(series, errors="coerce"), quality


def _unique_column_name(df: pd.DataFrame, preferred: str) -> str:
    label = preferred.strip() or "Metric"
    if label not in df.columns:
        return label
    index = 2
    while f"{label} {index}" in df.columns:
        index += 1
    return f"{label} {index}"


def _add_analysis_ready_numeric_fields(analyzed: pd.DataFrame, source: pd.DataFrame, mapping: dict[str, Any]) -> pd.DataFrame:
    if analyzed.empty or source.empty:
        return analyzed
    working = analyzed.copy()
    source_aligned = source.reset_index(drop=True).reindex(range(len(working))) if len(source) == len(working) else source.reindex(working.index)
    mapped_columns = {str(value or "").strip() for value in mapping.values() if str(value or "").strip()}
    candidates = list(dict.fromkeys([*mapped_columns, *source.columns]))
    for column in candidates:
        if column not in source_aligned.columns or str(column).startswith("__"):
            continue
        numeric_values, quality = _numeric_quality(source_aligned[column])
        if quality < 0.85:
            continue
        if column in working.columns:
            _, existing_quality = _numeric_quality(working[column])
            if existing_quality >= 0.85:
                continue
            target = _unique_column_name(working, f"{column} Numeric")
        else:
            target = column
        working[target] = numeric_values.to_numpy()
    return working


def _apply_optional_dimensions(analyzed: pd.DataFrame, source: pd.DataFrame, mapping: dict[str, Any]) -> pd.DataFrame:
    working = _add_analysis_ready_numeric_fields(analyzed, source, mapping)
    manager_column = str(mapping.get("manager") or "").strip()
    if manager_column and manager_column in source.columns:
        values = source[manager_column]
        if len(values) == len(working):
            working["Manager/TL"] = values.reset_index(drop=True).reindex(range(len(working))).fillna("Unknown").astype(str).to_numpy()
        else:
            working["Manager/TL"] = values.reindex(working.index).fillna("Unknown").astype(str)
    for key, target in (("wave", "Wave"), ("tenure", "Tenure")):
        column = str(mapping.get(key) or "").strip()
        if column and column in source.columns:
            values = source[column].apply(_tenure_bucket) if target == "Tenure" else source[column]
            if len(values) == len(working):
                working[target] = values.reset_index(drop=True).reindex(range(len(working))).fillna("Unknown").astype(str).to_numpy()
            else:
                working[target] = values.reindex(working.index).fillna("Unknown").astype(str)
        elif target not in working.columns:
            fallback = _first_existing_column(working, ["Batch#", "Batch", "Training Wave", "Cohort", "Wave"] if target == "Wave" else ["Tenure", "Tenure days", "Tenure Bucket", "Tenure Range", "Employee Tenure"])
            if fallback and fallback in working.columns:
                working[target] = working[fallback].apply(_tenure_bucket) if target == "Tenure" else working[fallback].fillna("Unknown").astype(str)
    return working


def _apply_dynamic_dimensions(analyzed: pd.DataFrame, source: pd.DataFrame, selected_columns: list[str]) -> pd.DataFrame:
    if analyzed.empty or source.empty or not selected_columns:
        return analyzed
    working = analyzed.copy()
    source_aligned = source.reset_index(drop=True).reindex(range(len(working))) if len(source) == len(working) else source.reindex(working.index)
    seen: set[str] = set()
    for column in selected_columns:
        label = str(column or "").strip()
        if not label or label in seen or label not in source_aligned.columns:
            continue
        seen.add(label)
        values = source_aligned[label].fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
        if len(values) == len(working):
            working[label] = values.to_numpy()
    return working


def _analysis_payload() -> dict[str, Any]:
    with STATE_LOCK:
        raw_analyzed_df = STATE.analyzed_df.copy()
        date_filter = dict(STATE.date_filter)
        status = STATE.status
        progress = STATE.progress
        analysis_running = STATE.analysis_running
        analysis_error = STATE.analysis_error
        analysis_id = STATE.analysis_id
        dynamic_dimensions = list(STATE.dynamic_dimensions)
        analysis_engines = dict(STATE.analysis_engines)
        model_paths = dict(STATE.model_paths)
        calendar_settings = dict(STATE.calendar_settings)
        mapping_config = dict(STATE.last_run_config.get("mapping", {}))
    analyzed_df = _apply_date_filter(raw_analyzed_df)
    analyzed_df = _apply_reporting_calendar(analyzed_df, calendar_settings)
    if analyzed_df.empty:
        weekly_df = pd.DataFrame()
        agent_df = pd.DataFrame()
        manager_df = pd.DataFrame()
        reason_df = pd.DataFrame()
        complaints_df = pd.DataFrame()
        passive_df = pd.DataFrame()
    else:
        summaries = _summaries_for_calendar(analyzed_df, calendar_settings)
        if "CSAT Type" in analyzed_df.columns:
            summaries = {key: _apply_csat_aliases(value) for key, value in summaries.items()}
            analyzed_df = _apply_csat_aliases(analyzed_df)
        weekly_df = summaries.get("weekly", pd.DataFrame())
        agent_df = summaries.get("agent", pd.DataFrame())
        agent_df = _agent_summary_with_manager(analyzed_df, agent_df)
        manager_df = _manager_summary(analyzed_df)
        reason_df = summaries.get("reasons", pd.DataFrame())
        complaints_df = summaries.get("complaints", pd.DataFrame())
        passive_df = summaries.get("passive", pd.DataFrame())
    summary = nps_summary(analyzed_df) if not analyzed_df.empty else {}
    counts = nps_composition_counts(analyzed_df) if not analyzed_df.empty else {}
    is_csat = "CSAT Type" in analyzed_df.columns
    if is_csat:
        summary, counts = _csat_summary_aliases(summary, counts)
    selected_dimension_columns = _selected_dimension_columns(analyzed_df, dynamic_dimensions, mapping_config)
    evidence_relationship = _evidence_relationship_payload(analyzed_df, selected_dimension_columns, is_csat)
    sentiment = sentiment_summary(analyzed_df) if not analyzed_df.empty else {}
    try:
        insights = executive_snapshot_insights(analyzed_df, reason_df) if not analyzed_df.empty else ""
    except Exception:
        insights = "Executive insights could not be generated for the current filtered view, but the calculated tables remain available."
    case_id = _case_id_column(analyzed_df)
    manager_col = _first_existing_column(analyzed_df, ["Manager/TL", "TL Name", "Team Manager", "Manager Name", "Manager", "Supervisor"])
    wave_col = _first_existing_column(analyzed_df, ["Batch#", "Batch", "Training Wave", "Cohort", "Wave"])
    tenure_col = _first_existing_column(analyzed_df, ["Tenure", "Tenure days", "Tenure Bucket", "Tenure Range", "Employee Tenure"])
    wave_df = _dimension_summary(analyzed_df, wave_col, "Wave")
    tenure_df = _dimension_summary(analyzed_df, tenure_col, "Tenure")
    if is_csat:
        weekly_df = _apply_csat_aliases(weekly_df)
        agent_df = _apply_csat_aliases(agent_df)
        manager_df = _apply_csat_aliases(manager_df)
        wave_df = _apply_csat_aliases(wave_df)
        tenure_df = _apply_csat_aliases(tenure_df)
    feedback_df = analyzed_df.copy()
    if not feedback_df.empty:
        feedback_df["Case ID"] = feedback_df[case_id].astype(str) if case_id else ""
        feedback_df["Manager/TL"] = feedback_df[manager_col].astype(str) if manager_col else ""
        feedback_df["Wave"] = feedback_df[wave_col].astype(str) if wave_col else ""
        feedback_df["Tenure"] = feedback_df[tenure_col].apply(_tenure_bucket).astype(str) if tenure_col else ""
        feedback_df["Sentiment Confidence"] = feedback_df.get("Sentiment Score", pd.Series(dtype=float)).apply(_sentiment_confidence)
    feedback_columns = [
        "Verbatim Feedback",
        "Bucket Category",
        "Sentiment",
        "Sentiment Score",
        "Primary Reason",
        "NPS Score",
        "CSAT Score",
        "NPS Type",
        "CSAT Type",
        "Agent Name",
        "Manager/TL",
        "Wave",
        "Tenure",
        "Feedback Date",
        "Silent Detractor Alert",
        "Silent Dissatisfied Alert",
        "Case ID",
        "Sentiment Confidence",
        "Owl Primary Driver",
        "Owl Secondary Driver",
        "Owl Tertiary Driver",
        "Owl People Sentiment",
        "Owl Process Sentiment",
        "Owl Tech Sentiment",
        "Owl Issue Type",
        "Owl Customer Impact",
        "Owl Resolution Status",
        "ACPT Primary Category",
        "ACPT Secondary Category",
        "ACPT Confidence",
        "ACPT Evidence",
        "ACPT Needs Review",
        "Analysis Source",
        "AI Rationale",
        "Manual Override Notes",
    ]
    feedback_table_columns = [
        "Agent Name",
        "Manager/TL",
        "Wave",
        "Tenure",
        "Case ID",
        "NPS Type",
        "CSAT Type",
        "NPS Score",
        "CSAT Score",
        "Sentiment",
        "Sentiment Score",
        "Sentiment Confidence",
        "ACPT Primary Category",
        "ACPT Confidence",
        "ACPT Needs Review",
        "Analysis Source",
        "Manual Override Notes",
        "Verbatim Feedback",
    ]
    theme_detail_columns = [
        "Agent Name",
        "Manager/TL",
        "Wave",
        "Tenure",
        "Case ID",
        "NPS Type",
        "CSAT Type",
        "NPS Score",
        "CSAT Score",
        "Owl Primary Driver",
        "Owl Secondary Driver",
        "Owl Tertiary Driver",
        "Owl People Sentiment",
        "Owl Process Sentiment",
        "Owl Tech Sentiment",
        "Owl Issue Type",
        "Owl Customer Impact",
        "Owl Resolution Status",
        "Sentiment",
        "Sentiment Score",
        "Sentiment Confidence",
        "Verbatim Feedback",
        "Owl Analysis Source",
        "Analysis Source",
        "Manual Override Notes",
    ]
    population_daily: list[dict[str, Any]] = []
    if not analyzed_df.empty and "Feedback Date" in analyzed_df.columns:
        population_work = analyzed_df.copy()
        population_work["__Population Date"] = pd.to_datetime(population_work["Feedback Date"], errors="coerce").dt.normalize()
        population_work = population_work.dropna(subset=["__Population Date"])
        for period, period_df in population_work.groupby("__Population Date", sort=True):
            period_counts = nps_composition_counts(period_df)
            period_summary = nps_summary(period_df)
            if is_csat:
                period_summary, period_counts = _csat_summary_aliases(period_summary, period_counts)
            period_sentiment = sentiment_summary(period_df)
            satisfied = int(period_counts.get("Satisfied", period_counts.get("Promoter", 0)) or 0)
            neutral = int(period_counts.get("Neutral", period_counts.get("Passive", 0)) or 0)
            dissatisfied = int(period_counts.get("Dissatisfied", period_counts.get("Detractor", 0)) or 0)
            promoter = int(period_counts.get("Promoter", period_counts.get("Satisfied", 0)) or 0)
            passive = int(period_counts.get("Passive", period_counts.get("Neutral", 0)) or 0)
            detractor = int(period_counts.get("Detractor", period_counts.get("Dissatisfied", 0)) or 0)
            period_score = float(period_summary.get("CSAT", period_summary.get("nps", period_summary.get("NPS", 0))) or 0)
            population_daily.append(
                {
                    "Period": pd.Timestamp(period).strftime("%Y-%m-%d"),
                    "Key": pd.Timestamp(period).strftime("%Y-%m-%d"),
                    "NPS": period_score,
                    "CSAT": period_score,
                    "Responses": satisfied + neutral + dissatisfied,
                    "Promoter": promoter,
                    "Passive": passive,
                    "Detractor": detractor,
                    "Satisfied": satisfied,
                    "Neutral": neutral,
                    "Dissatisfied": dissatisfied,
                    "Positive": float(period_sentiment.get("Positive", 0) or 0),
                    "Negative": float(period_sentiment.get("Negative", 0) or 0),
                }
            )
    classified_population = (
        int(counts.get("Satisfied", counts.get("Promoter", 0)) or 0)
        + int(counts.get("Neutral", counts.get("Passive", 0)) or 0)
        + int(counts.get("Dissatisfied", counts.get("Detractor", 0)) or 0)
    )
    return {
        "status": status,
        "progress": progress,
        "running": analysis_running,
        "error": analysis_error,
        "analysisId": analysis_id,
        "analysisEngines": analysis_engines,
        "modelPaths": model_paths,
        "calendar": calendar_settings,
        "mapping": mapping_config,
        "selectedDimensionColumns": selected_dimension_columns,
        "businessRules": dict(STATE.last_run_config.get("businessRules", {})),
        "timings": {
            "totalSeconds": round(max(0.0, STATE.analysis_completed_at - STATE.analysis_started_at), 2)
            if STATE.analysis_started_at and STATE.analysis_completed_at else None,
        },
        "completedAt": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(STATE.analysis_completed_at))
            if STATE.analysis_completed_at else "",
        "dateFilter": date_filter,
        "summary": summary,
        "counts": counts,
        "sentiment": sentiment,
        "segmentSentiments": _segment_sentiment_counts(analyzed_df),
        "population": {
            "rows": int(len(analyzed_df)),
            "classifiedRows": classified_population,
            "evidenceRowsReturned": int(min(len(feedback_df), 5000)),
            "evidenceLimit": 5000,
            "daily": population_daily,
        },
        "weekly": _safe_records(weekly_df, 52),
        "agents": _safe_records(agent_df, 5000),
        "managers": _safe_records(manager_df, 1000),
        "reasons": _safe_records(reason_df, 20),
        "complaints": _safe_records(complaints_df, 50),
        "passives": _safe_records(passive_df, 50),
        "quartiles": _safe_records(_quartile_summary(agent_df), 100),
        "quartileRollup": _quartile_rollup(agent_df),
        "quartileWeekly": _safe_records(pd.DataFrame(_quartile_weekly_trend(analyzed_df)), 52),
        "wave": _safe_records(wave_df, 100),
        "waveCards": _dimension_cards(analyzed_df, wave_df, "Wave"),
        "tenure": _safe_records(tenure_df, 100),
        "tenureCards": _dimension_cards(analyzed_df, tenure_df, "Tenure"),
        "dynamicDimensions": _dynamic_dimension_payload(analyzed_df, dynamic_dimensions),
        "evidenceRelationship": evidence_relationship,
        "themes": _safe_records(_theme_summary(analyzed_df), 100),
        "themeRows": _records_for_columns(feedback_df, theme_detail_columns, 5000),
        "operations": _operations_summary(analyzed_df),
        "churn": _churn_summary_rows(analyzed_df),
        "alerts": _alerts(analyzed_df, weekly_df),
        "feedbackRows": _records_for_columns(feedback_df, feedback_columns, 5000),
        "feedbackTableRows": _records_for_columns(feedback_df, feedback_table_columns, 5000),
        "executiveDashboard": _dashboard_snapshot(analyzed_df),
        "agentDashboards": _dashboard_snapshots_by(analyzed_df, "Agent Name", "All Agents"),
        "managerDashboards": _dashboard_snapshots_by(analyzed_df, manager_col, "All Managers"),
        "analysisSummary": _analysis_summary_rows(analyzed_df, weekly_df),
        "sentimentMovement": _sentiment_movement_rows(weekly_df),
        "consistency": _consistency_rows(weekly_df),
        "formulas": _formula_rows(),
        "preview": _safe_records(analyzed_df, 100),
        "insights": insights,
    }


def _set_progress(progress: float, status: str, analysis_id: str | None = None) -> None:
    rounded_progress = round(float(progress), 2)
    stage, _component = _analysis_log_component(status)
    rows_done = 0
    rows_total = 0
    if analysis_id:
        context = ANALYSIS_RUN_LOGS.get(analysis_id) or {}
        rows_done, rows_total = _parse_progress_rows(status, int(context.get("rows", 0) or 0))
    with STATE_LOCK:
        if analysis_id is not None and STATE.analysis_id != analysis_id:
            return
        STATE.progress = rounded_progress
        STATE.status = status
        STATE.analysis_stage = stage
        if rows_done:
            STATE.analysis_rows_processed = rows_done
        if rows_total:
            STATE.analysis_total_rows = rows_total
    if analysis_id:
        _analysis_log_progress(analysis_id, rounded_progress, status)


def _normalize_csat_category(value: object) -> str:
    text = str(value or "").strip().lower()
    if not text:
        return ""
    if text in {"satisfied", "satisfaction", "sat", "positive", "promoter", "good", "yes", "happy"}:
        return "Satisfied"
    if text in {"neutral", "neither", "passive", "average", "ok", "okay", "mixed"}:
        return "Neutral"
    if text in {"dissatisfied", "dissat", "unsatisfied", "negative", "detractor", "bad", "no", "unhappy"}:
        return "Dissatisfied"
    if "dissatisfied" in text or "unsatisfied" in text or "detractor" in text:
        return "Dissatisfied"
    if "satisfied" in text or "promoter" in text:
        return "Satisfied"
    if "neutral" in text or "passive" in text:
        return "Neutral"
    return ""


def _csat_category_to_nps_type(category: str) -> str:
    return {
        "Satisfied": "Promoter",
        "Neutral": "Passive",
        "Dissatisfied": "Detractor",
    }.get(category, "Unknown")


def _csat_band_config(raw: Any) -> dict[str, float | str]:
    config = raw if isinstance(raw, dict) else {}
    scale = str(config.get("scale") or "5").strip().lower()
    presets = {
        "5": {"satisfiedMin": 4.0, "neutralMin": 3.0},
        "10": {"satisfiedMin": 9.0, "neutralMin": 7.0},
        "100": {"satisfiedMin": 80.0, "neutralMin": 60.0},
        "custom": {"satisfiedMin": 4.0, "neutralMin": 3.0},
    }
    preset = presets.get(scale, presets["5"])
    try:
        satisfied_min = float(config.get("satisfiedMin", preset["satisfiedMin"]))
    except (TypeError, ValueError):
        satisfied_min = preset["satisfiedMin"]
    try:
        neutral_min = float(config.get("neutralMin", preset["neutralMin"]))
    except (TypeError, ValueError):
        neutral_min = preset["neutralMin"]
    return {"scale": scale, "satisfiedMin": satisfied_min, "neutralMin": neutral_min}


def _classify_csat(score: object, band_config: Any = None) -> tuple[str, str]:
    try:
        value = float(score)
    except (TypeError, ValueError):
        return "Unknown", "Unknown"
    if pd.isna(value):
        return "Unknown", "Unknown"
    config = _csat_band_config(band_config)
    if value >= float(config["satisfiedMin"]):
        return "Satisfied", "Promoter"
    if value >= float(config["neutralMin"]):
        return "Neutral", "Passive"
    return "Dissatisfied", "Detractor"


def _apply_csat_scoring(analyzed: pd.DataFrame, score_col: str | None, satisfaction_col: str | None = None, band_config: Any = None) -> pd.DataFrame:
    working = analyzed.copy()
    source_scores = pd.to_numeric(working[score_col], errors="coerce") if score_col and score_col in working.columns else pd.to_numeric(working.get("NPS Score"), errors="coerce")
    if satisfaction_col and satisfaction_col in working.columns:
        explicit_types = working[satisfaction_col].apply(_normalize_csat_category)
    else:
        explicit_types = pd.Series([""] * len(working), index=working.index)
    classifications = source_scores.apply(lambda score: _classify_csat(score, band_config))
    working["CSAT Score"] = source_scores
    working["CSAT Type"] = [
        explicit if explicit else scored[0]
        for explicit, scored in zip(explicit_types, classifications)
    ]
    working["NPS Type"] = working["CSAT Type"].apply(_csat_category_to_nps_type)
    working["Silent Detractor Alert"] = (working["CSAT Type"] == "Neutral") & (working["Sentiment"] == "Negative")
    working["Silent Dissatisfied Alert"] = working["Silent Detractor Alert"]
    working["Analysis Mode"] = "CSAT"
    return working


def _csat_summary_aliases(summary: dict[str, Any], counts: dict[str, Any]) -> tuple[dict[str, Any], dict[str, Any]]:
    adjusted_summary = dict(summary)
    adjusted_counts = dict(counts)
    satisfied_pct = float(adjusted_summary.get("promoters", 0.0) or 0.0)
    neutral_pct = float(adjusted_summary.get("passives", 0.0) or 0.0)
    dissatisfied_pct = float(adjusted_summary.get("detractors", 0.0) or 0.0)
    adjusted_summary.update(
        {
            "csat": satisfied_pct,
            "nps": satisfied_pct,
            "CSAT": satisfied_pct,
            "Satisfied": satisfied_pct,
            "Neutral": neutral_pct,
            "Dissatisfied": dissatisfied_pct,
            "Satisfieds": satisfied_pct,
            "Neutrals": neutral_pct,
            "Dissatisfieds": dissatisfied_pct,
        }
    )
    adjusted_counts.update(
        {
            "Satisfied": int(adjusted_counts.get("Promoter", 0) or 0),
            "Neutral": int(adjusted_counts.get("Passive", 0) or 0),
            "Dissatisfied": int(adjusted_counts.get("Detractor", 0) or 0),
        }
    )
    return adjusted_summary, adjusted_counts


def _apply_csat_aliases(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df
    working = df.copy()
    rename_map = {
        "Promoters": "Satisfieds",
        "Passives": "Neutrals",
        "Detractors": "Dissatisfieds",
        "Promoter": "Satisfied_Count",
        "Passive": "Neutral_Count",
        "Detractor": "Dissatisfied_Count",
        "Promoter_Count": "Satisfied_Count",
        "Passive_Count": "Neutral_Count",
        "Detractor_Count": "Dissatisfied_Count",
        "Promoter %": "Satisfied %",
        "Agent NPS": "Agent CSAT",
        "Manager NPS": "Manager CSAT",
    }
    for old, new in rename_map.items():
        if old in working.columns and new not in working.columns:
            working[new] = working[old]
    if "NPS Type" in working.columns and "CSAT Type" not in working.columns:
        working["CSAT Type"] = working["NPS Type"].map({"Promoter": "Satisfied", "Passive": "Neutral", "Detractor": "Dissatisfied"}).fillna(working["NPS Type"])
    if "NPS Score" in working.columns and "CSAT Score" not in working.columns:
        working["CSAT Score"] = working["NPS Score"]
    if {"Responses", "Satisfieds"}.issubset(working.columns):
        responses = pd.to_numeric(working["Responses"], errors="coerce").replace(0, pd.NA)
        satisfied = pd.to_numeric(working["Satisfieds"], errors="coerce").fillna(0)
        working["CSAT"] = (satisfied / responses * 100).fillna(0).round(2)
        working["NPS"] = working["CSAT"]
    if {"Responses", "Satisfied_Count"}.issubset(working.columns):
        responses = pd.to_numeric(working["Responses"], errors="coerce").replace(0, pd.NA)
        satisfied = pd.to_numeric(working["Satisfied_Count"], errors="coerce").fillna(0)
        working["CSAT"] = (satisfied / responses * 100).fillna(0).round(2)
        working["NPS"] = working["CSAT"]
    if "CSAT" in working.columns:
        if "Agent Name" in working.columns:
            working["Agent CSAT"] = working["CSAT"]
            if "Agent NPS" in working.columns:
                working["Agent NPS"] = working["CSAT"]
        if "Manager/TL" in working.columns:
            working["Manager CSAT"] = working["CSAT"]
            if "Manager NPS" in working.columns:
                working["Manager NPS"] = working["CSAT"]
    if "Satisfied_Count" in working.columns and "Satisfied" not in working.columns:
        working["Satisfied"] = working["Satisfied_Count"]
    if "Neutral_Count" in working.columns and "Neutral" not in working.columns:
        working["Neutral"] = working["Neutral_Count"]
    if "Dissatisfied_Count" in working.columns and "Dissatisfied" not in working.columns:
        working["Dissatisfied"] = working["Dissatisfied_Count"]
    if "Silent Detractor Alert" in working.columns and "Silent Dissatisfied Alert" not in working.columns:
        working["Silent Dissatisfied Alert"] = working["Silent Detractor Alert"]
    return working


def _run_analysis_job(payload: dict[str, Any], analysis_id: str) -> None:
    try:
        analysis_mode = str(payload.get("mode") or "nps").strip().lower()
        calendar_settings = _calendar_settings(payload)
        mapping = payload.get("mapping", {})
        feedback = mapping.get("feedback")
        score = mapping.get("score")
        satisfaction = mapping.get("satisfaction")
        agent = mapping.get("agent")
        date = mapping.get("date")
        csat_bands = payload.get("csatBands") if isinstance(payload.get("csatBands"), dict) else {}
        engines = payload.get("engines") if isinstance(payload.get("engines"), dict) else {}
        sentiment_engine = str(engines.get("sentiment") or "sparrow").strip().lower()
        theme_engine = str(engines.get("theme") or "local").strip().lower()
        model_paths = payload.get("modelPaths") if isinstance(payload.get("modelPaths"), dict) else {}
        sparrow_model_path = str(model_paths.get("sparrow") or _default_model_path("sparrow")).strip()
        theme_model_path = str(model_paths.get("theme") or model_paths.get("owl") or _default_model_path("theme")).strip()
        dynamic_dimensions = [str(item).strip() for item in payload.get("dynamicDimensions", []) if str(item).strip()]
        _set_progress(3, "Validating selected columns and files...", analysis_id)
        merged = _merge_lookup(
            STATE.base_df,
            STATE.lookup_df,
            payload.get("baseKey", ""),
            payload.get("lookupKey", ""),
        )
        _set_progress(8, f"Preparing local analysis for {len(merged):,} rows: cleaning feedback and checking selected columns...", analysis_id)
        def fallback_progress(done: int, total: int, message: str | None = None) -> None:
            total = max(total, 1)
            pct = 18 + (done / total) * 28
            if analysis_mode == "csat":
                total_records = max(len(merged), 1)
                estimated_row = min(total_records, max(0, round((done / total) * total_records)))
                detail = message or "Running safe local sentiment and theme rules."
                _set_progress(pct, f"Local Rules: {estimated_row:,}/{total_records:,} rows. {detail}", analysis_id)
            else:
                _set_progress(pct, message or f"Local fallback rules: step {done}/{total}...", analysis_id)

        if analysis_mode == "csat":
            _set_progress(18, f"Local Rules: 0/{len(merged):,} rows. Running safe local sentiment rules...", analysis_id)
        else:
            _set_progress(18, f"Running safe local sentiment rules for {len(merged):,} rows...", analysis_id)
        analyzed = build_analysis(
            merged,
            feedback,
            score if score else None,
            agent if agent else None,
            date if date else None,
            progress_callback=fallback_progress,
        )
        analyzed = _apply_optional_dimensions(analyzed, merged, mapping)
        analyzed["Analysis Source"] = "Local Fallback: safe rules completed before AI model attempt."

        if sentiment_engine in {"local", "rules", "local rules"}:
            analyzed["Analysis Source"] = "Local Rules"
            analyzed["AI Rationale"] = "User selected Local Rules for sentiment analysis."
            if analysis_mode == "csat":
                _set_progress(60, f"Local Rules: {len(merged):,}/{len(merged):,} rows. Local sentiment rules completed. Sparrow was not used for this run.", analysis_id)
            else:
                _set_progress(60, "Local sentiment rules completed. Sparrow was not used for this run.", analysis_id)
        else:
            try:
                if not _sparrow_runtime_enabled():
                    raise RuntimeError(
                        "Sparrow neural model loading is disabled by environment setting."
                    )
                if _isolated_worker_enabled():
                    _set_progress(48, "Attempting Sparrow AI sentiment in isolated worker...", analysis_id)
                    analyzed = _run_isolated_model_worker(
                        task="sparrow",
                        df=merged,
                        model_path=sparrow_model_path,
                        analysis_id=analysis_id,
                        progress_start=48,
                        progress_end=60,
                        status="Sparrow AI sentiment is running in an isolated worker",
                        timeout_seconds=900,
                        feedback_col=feedback,
                        score_col=score if score else "",
                        agent_col=agent if agent else "",
                        date_col=date if date else "",
                    )
                else:
                    def sparrow_progress(done: int, total: int, message: str | None = None) -> None:
                        total = max(total, 1)
                        pct = 48 + (done / total) * 12
                        if total > 1:
                            label = f"Sparrow AI sentiment: {done:,}/{total:,} rows ({(done / total) * 100:.1f}%). On track."
                        else:
                            label = message or "Preparing Sparrow AI sentiment..."
                        _set_progress(pct, label, analysis_id)

                    _set_progress(48, "Starting Sparrow AI sentiment. Local model retrieval is in progress...", analysis_id)
                    analyzed = build_analysis_with_local_model(
                        merged,
                        feedback,
                        score if score else None,
                        agent if agent else None,
                        date if date else None,
                        model_path=sparrow_model_path,
                        progress_callback=sparrow_progress,
                    )
                analyzed = _apply_optional_dimensions(analyzed, merged, mapping)
                _set_progress(60, "Sparrow AI sentiment completed successfully.", analysis_id)
            except Exception as exc:
                _set_progress(60, "Sparrow AI unavailable. Continuing with safe local sentiment results.", analysis_id)
                analyzed = _apply_optional_dimensions(analyzed, merged, mapping)
                analyzed["Analysis Source"] = f"Local Fallback: Sparrow AI unavailable: {exc}"
        if theme_engine in {"local", "rules", "local rules"}:
            _set_progress(62, f"Starting local theme rules for {len(analyzed):,} rows. Trained Theme Model is not selected for this run.", analysis_id)
            analyzed = _fill_owl_fallback_columns(analyzed, "User selected Local Rules for theme classification.")
            if analysis_mode == "csat":
                _set_progress(86, f"Local Rules: {len(analyzed):,}/{len(analyzed):,} rows. Local theme rules completed. Trained Theme Model was not used for this run.", analysis_id)
            else:
                _set_progress(86, "Local theme rules completed. Trained Theme Model was not used for this run.", analysis_id)
        else:
            try:
                _set_progress(62, "Trained Theme Model selected. Checking local model files before theme classification begins...", analysis_id)
                if not _theme_model_runtime_enabled():
                    raise RuntimeError(
                        ""
                        "Using safe theme fallback fields instead."
                    )

                if False and _isolated_worker_enabled():
                    analyzed = _run_isolated_model_worker(
                        task="theme",
                        df=analyzed,
                        feedback_col="Verbatim Feedback",
                        model_path=theme_model_path,
                        analysis_id=analysis_id,
                        progress_start=62,
                        progress_end=86,
                        status="Trained Theme Model is running in an isolated worker",
                        timeout_seconds=900,
                    )
                else:
                    def owl_progress(done: int, total: int, message: str | None = None) -> None:
                        total = max(total, 1)
                        pct = 62 + (done / total) * 24
                        label = f"Trained Theme Model: {done:,}/{total:,} rows ({(done / total) * 100:.1f}%). On track."
                        _set_progress(pct, message or label, analysis_id)

                    analyzed = add_theme_acpt_resolution_outputs(
                        analyzed,
                        feedback_col="Verbatim Feedback",
                        model_path=theme_model_path,
                        progress_callback=owl_progress,
                    )
                _set_progress(86, "Trained Theme Model completed successfully.", analysis_id)
            except Exception as exc:
                _set_progress(86, "Trained Theme Model unavailable. Filling theme fields with safe defaults...", analysis_id)
                analyzed = _fill_owl_fallback_columns(analyzed, str(exc))
        analyzed = _apply_optional_dimensions(analyzed, merged, mapping)
        analyzed = _apply_dynamic_dimensions(analyzed, merged, dynamic_dimensions)
        analyzed = _neutralize_blank_feedback_sentiment(analyzed)
        analyzed = _apply_reporting_calendar(analyzed, calendar_settings)
        if analysis_mode == "csat":
            analyzed = _apply_csat_scoring(analyzed, score if score else None, satisfaction if satisfaction else None, csat_bands)
        def acpt_progress(done: int, total: int, message: str | None = None) -> None:
            total = max(total, 1)
            pct = 86 + (done / total) * 3
            _set_progress(pct, message or f"ACPT classification: {done:,}/{total:,} rows. Assigning Agent, Customer, Process, or Technology ownership.", analysis_id)

        trained_theme_supplied_acpt = (
            theme_engine not in {"local", "rules", "local rules"}
            and "ACPT Primary Category" in analyzed.columns
            and analyzed["ACPT Primary Category"].fillna("").astype(str).str.strip().ne("").any()
        )
        if trained_theme_supplied_acpt:
            _set_progress(89, f"Trained Theme Model supplied ACPT output for {len(analyzed):,} rows. Building weekly, agent, driver, and alert summaries next.", analysis_id)
        else:
            _set_progress(86, f"Starting ACPT verbatim classification for {len(analyzed):,} rows. I am identifying whether each comment points to Agent, Customer, Process, or Technology.", analysis_id)
            analyzed = _add_acpt_classification_outputs(analyzed, "Verbatim Feedback", acpt_progress)
            _set_progress(89, f"ACPT classification completed for {len(analyzed):,} rows. Building weekly, agent, driver, and alert summaries next.", analysis_id)
        _set_progress(91, "Building weekly, agent, driver, and alert summaries. Final dashboard assembly is on track...", analysis_id)
        summaries = _summaries_for_calendar(analyzed, calendar_settings)
        if analysis_mode == "csat":
            summaries = {key: _apply_csat_aliases(value) for key, value in summaries.items()}
            analyzed = _apply_csat_aliases(analyzed)
        _set_progress(95, "Preparing dashboard response and export-ready data. Almost ready...", analysis_id)
        _set_progress(97, f"Profiling {len(analyzed):,} analyzed rows for the Column Explorer. This is the final full-data scan...", analysis_id)
        analyzed_column_profile = _column_profile(analyzed)
        _set_progress(99, "Column profiling complete. Publishing the final dashboard state...", analysis_id)
        with STATE_LOCK:
            if STATE.analysis_id != analysis_id:
                return
            STATE.analyzed_df = analyzed
            STATE.weekly_df = summaries.get("weekly", pd.DataFrame())
            STATE.agent_df = _agent_summary_with_manager(analyzed, summaries.get("agent", pd.DataFrame()))
            STATE.manager_df = _manager_summary(analyzed)
            STATE.reason_df = summaries.get("reasons", pd.DataFrame())
            STATE.complaints_df = summaries.get("complaints", pd.DataFrame())
            STATE.passive_df = summaries.get("passive", pd.DataFrame())
            STATE.analyzed_column_profile = analyzed_column_profile
            STATE.dynamic_dimensions = dynamic_dimensions
            STATE.analysis_engines = {"sentiment": sentiment_engine, "theme": theme_engine}
            STATE.model_paths = {"sparrow": sparrow_model_path, "theme": theme_model_path}
            STATE.calendar_settings = calendar_settings
            STATE.progress = 100
            STATE.status = "Analysis complete"
            _write_audit_record("SYSTEM", "ANALYSIS_COMPLETED", "Feedback analysis completed.", {"analysisId": analysis_id, "rows": int(len(analyzed))})
            STATE.analysis_error = ""
            STATE.analysis_running = False
            STATE.analysis_completed_at = time.time()
        _finalize_analysis_log(analysis_id, analyzed)
    except Exception as exc:
        traceback.print_exc()
        failure_details = _analysis_log_failure(analysis_id, exc)
        _write_audit_record(
            "SYSTEM",
            "ANALYSIS_FAILED",
            "Feedback analysis failed.",
            {
                "analysisId": analysis_id,
                "error": f"{type(exc).__name__}: {exc}",
                **failure_details,
            },
            level="ERROR",
        )
        _finalize_analysis_log(analysis_id, None, exc)
        with STATE_LOCK:
            if STATE.analysis_id != analysis_id:
                return
            STATE.progress = 100
            STATE.status = f"Failed: {exc}"
            STATE.analysis_error = str(exc)
            STATE.analysis_running = False


def _apply_feedback_override(payload: dict[str, Any]) -> dict[str, Any]:
    row_id = str(payload.get("rowId", "")).strip()
    updates = payload.get("updates") or {}
    allowed_columns = {
        "Sentiment",
        "Bucket Category",
        "Primary Reason",
        "Owl Primary Driver",
        "Owl Secondary Driver",
        "Owl Tertiary Driver",
        "Owl People Sentiment",
        "Owl Process Sentiment",
        "Owl Tech Sentiment",
        "Owl Issue Type",
        "Owl Customer Impact",
        "Owl Resolution Status",
        "Manual Override Notes",
    }
    with STATE_LOCK:
        if STATE.analyzed_df.empty:
            return {"ok": False, "error": "Run analysis before overriding feedback."}
        matching_index = None
        for index in STATE.analyzed_df.index:
            if str(index) == row_id:
                matching_index = index
                break
        if matching_index is None:
            return {"ok": False, "error": "Feedback row was not found in the current analysis."}
        for column, value in updates.items():
            if column not in allowed_columns:
                continue
            if column not in STATE.analyzed_df.columns:
                STATE.analyzed_df[column] = ""
            STATE.analyzed_df.at[matching_index, column] = str(value or "").strip()
        sentiment = str(STATE.analyzed_df.at[matching_index, "Sentiment"] or "")
        if sentiment == "Positive":
            STATE.analyzed_df.at[matching_index, "Sentiment Score"] = 1.0
        elif sentiment == "Negative":
            STATE.analyzed_df.at[matching_index, "Sentiment Score"] = -1.0
        elif sentiment == "Neutral":
            STATE.analyzed_df.at[matching_index, "Sentiment Score"] = 0.0
        STATE.analyzed_df.at[matching_index, "Analysis Source"] = "Manual Override"
        if "AI Rationale" not in STATE.analyzed_df.columns:
            STATE.analyzed_df["AI Rationale"] = ""
        STATE.analyzed_df.at[matching_index, "AI Rationale"] = "User-reviewed feedback intelligence override."
        summaries = _summaries_for_calendar(STATE.analyzed_df, STATE.calendar_settings)
        STATE.weekly_df = summaries.get("weekly", pd.DataFrame())
        STATE.agent_df = _agent_summary_with_manager(STATE.analyzed_df, summaries.get("agent", pd.DataFrame()))
        STATE.manager_df = _manager_summary(STATE.analyzed_df)
        STATE.reason_df = summaries.get("reasons", pd.DataFrame())
        STATE.complaints_df = summaries.get("complaints", pd.DataFrame())
        STATE.passive_df = summaries.get("passive", pd.DataFrame())
    return {"ok": True, "analysis": _analysis_payload()}


LEADERSHIP_MIN_SAMPLE = 10


def _leadership_question_framework() -> list[dict[str, Any]]:
    path = APPS / "csat-analyzer-next" / "question_framework.json"
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
        return payload if isinstance(payload, list) else []
    except Exception:
        return []


def _nps_leadership_question_framework() -> list[dict[str, Any]]:
    path = APPS / "nps-analyzer" / "question_framework.json"
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
        return payload if isinstance(payload, list) else []
    except Exception:
        return []


def _chi_square_survival_approx(statistic: float, degrees_freedom: int) -> float:
    if statistic <= 0 or degrees_freedom <= 0:
        return 1.0
    # Wilson-Hilferty normal approximation keeps the local engine dependency-light.
    transformed = ((statistic / degrees_freedom) ** (1 / 3) - (1 - 2 / (9 * degrees_freedom))) / math.sqrt(2 / (9 * degrees_freedom))
    return 0.5 * math.erfc(transformed / math.sqrt(2))


def _first_existing_column(df: pd.DataFrame, candidates: list[str]) -> str:
    lowered = {str(column).strip().lower(): column for column in df.columns if str(column).strip()}
    for candidate in candidates:
        needle = candidate.strip().lower()
        if needle and needle in lowered:
            return lowered[needle]
    for candidate in candidates:
        needle = candidate.strip().lower()
        if not needle:
            continue
        for lower, column in lowered.items():
            if needle in lower:
                return column
    return ""


def _safe_float(value: Any, default: float = 0.0) -> float:
    try:
        result = float(value)
        return result if math.isfinite(result) else default
    except (TypeError, ValueError):
        return default


def _wilson_interval(successes: float, total: float, z: float = 1.96) -> tuple[float, float]:
    if total <= 0:
        return (0.0, 0.0)
    proportion = successes / total
    denominator = 1 + (z * z / total)
    center = (proportion + z * z / (2 * total)) / denominator
    margin = z * math.sqrt((proportion * (1 - proportion) / total) + (z * z / (4 * total * total))) / denominator
    return (max(0.0, (center - margin) * 100), min(100.0, (center + margin) * 100))


def _two_proportion_test(success_a: float, total_a: float, success_b: float, total_b: float) -> dict[str, float]:
    if min(total_a, total_b) <= 0:
        return {"z": 0.0, "pValue": 1.0, "effectSize": 0.0}
    p_a, p_b = success_a / total_a, success_b / total_b
    pooled = (success_a + success_b) / (total_a + total_b)
    standard_error = math.sqrt(max(0.0, pooled * (1 - pooled) * ((1 / total_a) + (1 / total_b))))
    z_score = (p_a - p_b) / standard_error if standard_error else 0.0
    p_value = math.erfc(abs(z_score) / math.sqrt(2))
    effect_size = 2 * (math.asin(math.sqrt(max(0.0, min(1.0, p_a)))) - math.asin(math.sqrt(max(0.0, min(1.0, p_b)))))
    return {"z": z_score, "pValue": p_value, "effectSize": effect_size}


def _linear_trend(values: list[float]) -> dict[str, float]:
    clean = [float(value) for value in values if math.isfinite(float(value))]
    count = len(clean)
    if count < 2:
        return {"slope": 0.0, "rSquared": 0.0}
    x_mean = (count - 1) / 2
    y_mean = sum(clean) / count
    denominator = sum((index - x_mean) ** 2 for index in range(count))
    slope = sum((index - x_mean) * (value - y_mean) for index, value in enumerate(clean)) / denominator if denominator else 0.0
    fitted = [y_mean + slope * (index - x_mean) for index in range(count)]
    total_ss = sum((value - y_mean) ** 2 for value in clean)
    residual_ss = sum((value - prediction) ** 2 for value, prediction in zip(clean, fitted))
    return {"slope": slope, "rSquared": max(0.0, 1 - residual_ss / total_ss) if total_ss else 0.0}


def _leadership_statistics(target: float = 85.0, minimum_sample: int = LEADERSHIP_MIN_SAMPLE, satisfied_min: float = 4.0, neutral_min: float = 3.0) -> dict[str, Any]:
    minimum_sample = max(2, min(500, int(minimum_sample or LEADERSHIP_MIN_SAMPLE)))
    with STATE_LOCK:
        df = STATE.analyzed_df.copy()
        calendar_settings = dict(STATE.calendar_settings)
    df = _apply_date_filter(df)
    if df.empty:
        return {"ok": False, "error": "No analyzed rows are available for the selected date range."}

    csat_type = _first_existing_column(df, ["CSAT Type", "Satisfaction Level", "SatisfactionLevel", "NPS Type"])
    if not csat_type:
        return {"ok": False, "error": "The analyzed data does not contain a satisfaction classification."}
    normalized = df[csat_type].fillna("").astype(str).str.strip().str.lower()
    recognized = normalized.str.contains(r"^satisfied$|^promoter$|^positive$|^neutral$|^passive$|^dissatisfied$|^detractor$|^negative$", regex=True)
    if not bool(recognized.any()):
        score_col = _first_existing_column(df, ["CSAT Score", "Score", "CSAT", "Rating"])
        scores = pd.to_numeric(df[score_col], errors="coerce") if score_col else pd.Series([math.nan] * len(df), index=df.index)
        if not bool(scores.notna().any()) and "Score" in df.columns:
            scores = pd.to_numeric(df["Score"], errors="coerce")
        satisfied = (scores >= float(satisfied_min)).astype(float)
    else:
        satisfied = normalized.str.contains(r"^satisfied$|^promoter$|^positive$", regex=True).astype(float)
    total = int(len(satisfied))
    successes = float(satisfied.sum())
    overall = successes / total * 100 if total else 0.0
    overall_ci = _wilson_interval(successes, total)
    agent_col = _first_existing_column(df, ["Agent Name", "AgentName", "Agent"])
    manager_col = _first_existing_column(df, ["Manager/TL", "ManagerName", "Manager", "Team"])
    date_col = _first_existing_column(df, ["Feedback Date", "CallDateTime", "Date", "Response Date"])
    working = df.copy()
    working["__satisfied"] = satisfied

    def entity_rows(column: str) -> list[dict[str, Any]]:
        if not column:
            return []
        grouped = working.assign(__entity=working[column].fillna("Not available").astype(str)).groupby("__entity", dropna=False)["__satisfied"].agg(["sum", "count"])
        rows: list[dict[str, Any]] = []
        for name, row in grouped.iterrows():
            count = int(row["count"])
            sat = float(row["sum"])
            mean = sat / count * 100 if count else 0.0
            low, high = _wilson_interval(sat, count)
            rows.append({"name": str(name), "mean": mean, "n": count, "ciLow": low, "ciHigh": high, "gap": mean - overall, "targetGap": mean - target, "reliable": count >= minimum_sample})
        return sorted(rows, key=lambda item: item["mean"], reverse=True)

    agents = entity_rows(agent_col)
    managers = entity_rows(manager_col)
    reliable_agents = [row for row in agents if row["reliable"]]
    reliable_managers = [row for row in managers if row["reliable"]]

    periods: list[dict[str, Any]] = []
    if date_col:
        dates = pd.to_datetime(working[date_col], errors="coerce")
        period_frame = pd.DataFrame({"date": dates, "satisfied": satisfied}).dropna(subset=["date"])
        if not period_frame.empty:
            period_frame["period"] = week_period_start(period_frame["date"], calendar_settings["weekStart"])
            grouped = period_frame.groupby("period")["satisfied"].agg(["sum", "count"]).reset_index()
            for _, row in grouped.iterrows():
                count = int(row["count"])
                sat = float(row["sum"])
                low, high = _wilson_interval(sat, count)
                periods.append({"period": str(row["period"].date()), "mean": sat / count * 100 if count else 0.0, "successes": sat, "n": count, "ciLow": low, "ciHigh": high})
    trend = _linear_trend([row["mean"] for row in periods])
    latest_test = {"z": 0.0, "pValue": 1.0, "effectSize": 0.0}
    if len(periods) >= 2:
        latest_test = _two_proportion_test(periods[-1]["successes"], periods[-1]["n"], periods[-2]["successes"], periods[-2]["n"])
    means = [row["mean"] for row in periods]
    mean_period = sum(means) / len(means) if means else overall
    period_sd = math.sqrt(sum((value - mean_period) ** 2 for value in means) / max(1, len(means) - 1)) if len(means) > 1 else 0.0
    sorted_means = sorted(means)
    q1 = float(pd.Series(sorted_means).quantile(0.25)) if sorted_means else 0.0
    q3 = float(pd.Series(sorted_means).quantile(0.75)) if sorted_means else 0.0
    iqr = q3 - q1
    outliers = [row for row in periods if (period_sd and abs((row["mean"] - mean_period) / period_sd) >= 2) or row["mean"] < q1 - 1.5 * iqr or row["mean"] > q3 + 1.5 * iqr]

    def fmt_row(row: dict[str, Any] | None) -> str:
        if not row:
            return "not available"
        return f"{row['name']} {row['mean']:.1f}% (n={row['n']}, 95% CI {row['ciLow']:.1f}-{row['ciHigh']:.1f})"

    best_agent = reliable_agents[0] if reliable_agents else None
    worst_agent = reliable_agents[-1] if reliable_agents else None
    best_manager = reliable_managers[0] if reliable_managers else None
    worst_manager = reliable_managers[-1] if reliable_managers else None
    latest = periods[-1] if periods else None
    previous = periods[-2] if len(periods) >= 2 else None
    movement = (latest["mean"] - previous["mean"]) if latest and previous else 0.0
    significance = "statistically significant" if latest_test["pValue"] < 0.05 else "directional, not statistically significant"
    effect_label = "negligible" if abs(latest_test["effectSize"]) < 0.2 else "small" if abs(latest_test["effectSize"]) < 0.5 else "material"
    below_agents = [row for row in reliable_agents if row["mean"] < target]
    above_agents = [row for row in reliable_agents if row["mean"] >= target]
    excluded_agents = len(agents) - len(reliable_agents)
    excluded_managers = len(managers) - len(reliable_managers)
    trend_direction = "improving" if trend["slope"] > 0.5 else "declining" if trend["slope"] < -0.5 else "stable"

    def entity_trends(column: str) -> list[dict[str, Any]]:
        if not column or not date_col:
            return []
        frame = pd.DataFrame({"entity": working[column].fillna("Not available").astype(str), "date": pd.to_datetime(working[date_col], errors="coerce"), "satisfied": satisfied}).dropna(subset=["date"])
        if frame.empty:
            return []
        frame["period"] = week_period_start(frame["date"], calendar_settings["weekStart"])
        grouped = frame.groupby(["entity", "period"])["satisfied"].agg(["sum", "count"]).reset_index()
        results: list[dict[str, Any]] = []
        for entity, entity_rows in grouped.groupby("entity"):
            entity_rows = entity_rows.sort_values("period")
            period_values = [float(row["sum"]) / int(row["count"]) * 100 for _, row in entity_rows.iterrows() if int(row["count"]) > 0]
            total_n = int(entity_rows["count"].sum())
            if len(period_values) < 2:
                continue
            entity_trend = _linear_trend(period_values)
            recent = float(period_values[-1])
            historical = float(sum(period_values[:-1]) / len(period_values[:-1]))
            positive_changes = sum(1 for prior, current in zip(period_values, period_values[1:]) if current > prior)
            entity_sd = float(pd.Series(period_values).std(ddof=1)) if len(period_values) > 1 else 0.0
            results.append({"name": str(entity), "n": total_n, "periods": len(period_values), "recentMean": recent, "historicalMean": historical, "change": recent - historical, "slope": entity_trend["slope"], "rSquared": entity_trend["rSquared"], "sd": entity_sd, "positiveChanges": positive_changes, "reliable": total_n >= minimum_sample})
        return results

    manager_trends = entity_trends(manager_col)
    agent_trends = entity_trends(agent_col)

    team_rows: list[dict[str, Any]] = []
    if manager_col and agent_col:
        team_frame = pd.DataFrame({"manager": working[manager_col].fillna("Not available").astype(str), "agent": working[agent_col].fillna("Not available").astype(str), "satisfied": satisfied})
        agent_team = team_frame.groupby(["manager", "agent"])["satisfied"].agg(["sum", "count"]).reset_index()
        agent_team["mean"] = agent_team["sum"] / agent_team["count"] * 100
        agent_team = agent_team[agent_team["count"] >= minimum_sample]
        for manager, rows in agent_team.groupby("manager"):
            values = [float(value) for value in rows["mean"].tolist()]
            if not values:
                continue
            team_mean = sum(values) / len(values)
            team_sd = float(pd.Series(values).std(ddof=1)) if len(values) > 1 else 0.0
            above_count = sum(1 for value in values if value >= target)
            team_rows.append({"name": str(manager), "agents": len(values), "mean": team_mean, "sd": team_sd, "cv": team_sd / team_mean * 100 if team_mean else 0.0, "range": max(values) - min(values), "aboveTarget": above_count, "aboveTargetPct": above_count / len(values) * 100})
        team_rows.sort(key=lambda row: row["sd"])

    def omnibus(rows: list[dict[str, Any]]) -> dict[str, Any]:
        reliable = [row for row in rows if row.get("reliable") and row.get("n", 0) > 0]
        if len(reliable) < 2:
            return {"groups": len(reliable), "statistic": 0.0, "degreesFreedom": 0, "pValue": 1.0, "significant": False}
        pooled = sum(row["mean"] / 100 * row["n"] for row in reliable) / sum(row["n"] for row in reliable)
        statistic = 0.0
        for row in reliable:
            observed_success = row["mean"] / 100 * row["n"]
            expected_success = row["n"] * pooled
            expected_failure = row["n"] * (1 - pooled)
            if expected_success > 0:
                statistic += (observed_success - expected_success) ** 2 / expected_success
            observed_failure = row["n"] - observed_success
            if expected_failure > 0:
                statistic += (observed_failure - expected_failure) ** 2 / expected_failure
        degrees = len(reliable) - 1
        p_value = _chi_square_survival_approx(statistic, degrees)
        return {"groups": len(reliable), "statistic": statistic, "degreesFreedom": degrees, "pValue": p_value, "significant": p_value < 0.05}

    manager_omnibus = omnibus(reliable_managers)
    agent_omnibus = omnibus(reliable_agents)
    manager_confidence = sorted(reliable_managers, key=lambda row: row["ciHigh"] - row["ciLow"])
    low_volume_agents = [row for row in agents if not row["reliable"]]
    improving_managers = sorted([row for row in manager_trends if row["reliable"] and row["slope"] > 0], key=lambda row: (row["slope"], row["rSquared"]), reverse=True)
    sustained_agents = sorted([row for row in agent_trends if row["reliable"] and row["slope"] > 0 and row["positiveChanges"] >= 2], key=lambda row: (row["positiveChanges"], row["slope"]), reverse=True)
    recent_agent_outliers = sorted([row for row in agent_trends if row["reliable"] and abs(row["change"]) >= max(5.0, period_sd * 1.5)], key=lambda row: abs(row["change"]), reverse=True)
    near_target_agents = sorted(reliable_agents, key=lambda row: abs(row["mean"] - target))[:10]
    confidence_agents = sorted(reliable_agents, key=lambda row: (row["ciLow"], row["mean"]), reverse=True)
    composite_managers: list[dict[str, Any]] = []
    team_by_name = {row["name"]: row for row in team_rows}
    for manager in reliable_managers:
        team = team_by_name.get(manager["name"], {})
        stability_penalty = float(team.get("sd", 0.0))
        composite_managers.append({**manager, "teamSd": stability_penalty, "compositeScore": manager["mean"] - stability_penalty})
    composite_managers.sort(key=lambda row: row["compositeScore"], reverse=True)

    def answer(question: str, text: str, method: str, status: str = "Actionable", evidence: list[dict[str, Any]] | None = None) -> dict[str, Any]:
        return {"question": question, "text": text, "method": method, "status": status, "evidence": evidence or []}

    questions = [
        answer("CSAT trend over time", f"CSAT is {trend_direction}; weekly slope {trend['slope']:+.2f} pts, R-squared {trend['rSquared']:.2f}. Latest movement {movement:+.1f} pts is {significance} (p={latest_test['pValue']:.3f}, Cohen's h={latest_test['effectSize']:.2f}, {effect_label} effect).", "Weekly aggregation + linear regression + two-proportion z-test", "Actionable" if latest_test["pValue"] < 0.05 else "Monitor", periods[-6:]),
        answer("Overall CSAT for the period", f"Overall CSAT is {overall:.1f}% from n={total:,}; 95% Wilson CI {overall_ci[0]:.1f}-{overall_ci[1]:.1f}. Target gap {overall-target:+.1f} pts.", "Satisfied responses / valid responses; Wilson confidence interval"),
        answer("Highest and lowest managers", f"Best reliable manager: {fmt_row(best_manager)}. Weakest: {fmt_row(worst_manager)}. {excluded_managers} manager(s) excluded below n={minimum_sample}.", "Rank only entities meeting minimum sample; report uncertainty"),
        answer("Highest and lowest agents", f"Best reliable agent: {fmt_row(best_agent)}. Weakest: {fmt_row(worst_agent)}. {excluded_agents} agent(s) excluded below n={minimum_sample}.", "Rank only entities meeting minimum sample; report uncertainty"),
        answer("Manager improvement over time", f"Current aggregate movement is {movement:+.1f} pts ({significance}); use manager drill-down only where equal periods each meet n={minimum_sample}.", "Equal-period comparison + significance + effect size", "Monitor"),
        answer("Manager decline over time", f"A decline is actionable only when direction, p<0.05, practical effect, and minimum volume agree. Current aggregate p={latest_test['pValue']:.3f}, effect={latest_test['effectSize']:.2f}.", "Equal-period decline test with practical significance", "Monitor"),
        answer("Agent improvement over time", f"Improvement requires repeated period evidence; isolated gains are directional. Current overall slope is {trend['slope']:+.2f} with R-squared {trend['rSquared']:.2f}.", "Entity-period regression and equal-period testing", "Monitor"),
        answer("Agent decline over time", f"Coaching should require sufficient volume plus sustained negative slope, not one low observation. Minimum sample is n={minimum_sample} per comparison.", "Entity-period regression, effect size and reliability guardrail", "Monitor"),
        answer("Consistent manager performance", f"Most consistent manager: {min(manager_trends,key=lambda row:row['sd'])['name'] if manager_trends else 'not available'} with weekly SD {min(manager_trends,key=lambda row:row['sd'])['sd'] if manager_trends else 0:.2f} pts. Performance level and stability are reported separately.", "Standard deviation, coefficient of variation and IQR", "Monitor"),
        answer("Volatile manager performance", f"Most volatile manager: {max(manager_trends,key=lambda row:row['sd'])['name'] if manager_trends else 'not available'} with weekly SD {max(manager_trends,key=lambda row:row['sd'])['sd'] if manager_trends else 0:.2f} pts.", "SD, range, z-score and IQR outliers", "Review required" if manager_trends and max(row['sd'] for row in manager_trends) >= 10 else "Monitor"),
        answer("Consistent agent performance", f"Most consistent reliable agent: {min(agent_trends,key=lambda row:row['sd'])['name'] if agent_trends else 'not available'} with weekly SD {min(agent_trends,key=lambda row:row['sd'])['sd'] if agent_trends else 0:.2f} pts.", "Agent-period SD, CV, rolling SD and IQR", "Monitor"),
        answer("Volatile agent performance", f"Most volatile reliable agent: {max(agent_trends,key=lambda row:row['sd'])['name'] if agent_trends else 'not available'} with weekly SD {max(agent_trends,key=lambda row:row['sd'])['sd'] if agent_trends else 0:.2f} pts.", "Agent-period SD, range, CV and outlier count", "Review required" if agent_trends and max(row['sd'] for row in agent_trends) >= 10 else "Monitor"),
        answer("Highest and lowest week", f"Best week: {max(periods,key=lambda r:r['mean'])['period'] if periods else 'not available'} at {max(means) if means else 0:.1f}%. Weakest: {min(periods,key=lambda r:r['mean'])['period'] if periods else 'not available'} at {min(means) if means else 0:.1f}%.", "Period ranking with volume and Wilson intervals", evidence=periods),
        answer("Unusual spikes or drops", f"Detected {len(outliers)} unusual period(s) using |z|>=2 or the 1.5xIQR rule.", "Z-score plus IQR outlier detection", "Review required" if outliers else "No action required", outliers),
        answer("Managers above the organization average", f"{sum(1 for row in reliable_managers if row['mean']>overall)} of {len(reliable_managers)} reliable managers are above the organizational CSAT of {overall:.1f}%.", "Manager mean vs organization mean with confidence intervals", evidence=reliable_managers),
        answer("Agents above the organization average", f"{sum(1 for row in reliable_agents if row['mean']>overall)} of {len(reliable_agents)} reliable agents are above the organizational CSAT of {overall:.1f}%.", "Agent mean vs organization mean with confidence intervals"),
        answer("CSAT distribution by manager", f"Reliable manager spread is {(best_manager['mean']-worst_manager['mean']) if best_manager and worst_manager else 0:.1f} pts across {len(reliable_managers)} managers.", "Mean, median, quartiles, IQR and uncertainty", evidence=reliable_managers),
        answer("CSAT distribution by agent", f"Reliable agent spread is {(best_agent['mean']-worst_agent['mean']) if best_agent and worst_agent else 0:.1f} pts across {len(reliable_agents)} agents.", "Mean, median, quartiles, IQR and uncertainty"),
        answer("Manager performance gap", f"Best-to-worst reliable manager gap is {(best_manager['mean']-worst_manager['mean']) if best_manager and worst_manager else 0:.1f} pts.", "Highest minus lowest reliable manager result"),
        answer("Balanced team performance", "A balanced team requires low within-team agent SD/IQR as well as adequate sample sizes; a high average alone is not sufficient.", "Within-manager agent SD, CV, range and IQR", "Monitor"),
        answer("Top 10 agents", "Top reliable agents: " + "; ".join(fmt_row(row) for row in reliable_agents[:10]), "Volume-filtered rank with Wilson intervals", evidence=reliable_agents[:10]),
        answer("Bottom 10 agents", "Bottom reliable agents: " + "; ".join(fmt_row(row) for row in reliable_agents[-10:]), "Volume-filtered rank with Wilson intervals", evidence=reliable_agents[-10:]),
        answer("Agents improving month over month", "Agent improvement requires at least two comparable monthly periods and a positive regression slope with adequate R-squared.", "Monthly regression, R-squared, percent change and significance", "Monitor"),
        answer("Agents needing coaching", f"{len(below_agents)} reliable agent(s) are below the {target:.1f}% target; prioritize only after combining gap, trend, variability, and confidence interval.", "Multi-signal coaching rule: target gap + trend + volatility + sample", evidence=below_agents[:10]),
        answer("Managers needing attention", f"Weakest reliable manager is {fmt_row(worst_manager)}; confirm team trend and within-team spread before action.", "Team mean + trend + spread + uncertainty", evidence=reliable_managers[-5:]),
        answer("Agents above target", f"{len(above_agents)} of {len(reliable_agents)} reliable agents ({(len(above_agents)/len(reliable_agents)*100) if reliable_agents else 0:.1f}%) are at or above {target:.1f}%.", "Count and percentage after minimum-sample filtering"),
        answer("Agents below target", f"{len(below_agents)} of {len(reliable_agents)} reliable agents ({(len(below_agents)/len(reliable_agents)*100) if reliable_agents else 0:.1f}%) are below {target:.1f}%.", "Count and percentage after minimum-sample filtering", evidence=below_agents),
        answer("Manager with most high-performing agents", "This requires agent-to-manager membership plus reliable agent samples; the final ranking should report team size and a proportion confidence interval.", "Proportion above target by manager with Wilson interval", "Monitor"),
        answer("Leadership highlights and concerns", f"CSAT {overall:.1f}% (95% CI {overall_ci[0]:.1f}-{overall_ci[1]:.1f}, n={total:,}), {trend_direction} trend, {len(outliers)} outlier period(s), and {len(below_agents)} reliable agents below target. Findings are associations, not causal conclusions.", "Integrated statistically guarded summary"),
    ]
    summary_answer = questions.pop()
    questions.append(answer("Changes compared with the previous period", f"Latest CSAT changed {movement:+.1f} pts versus the previous equal week; the change is {significance} with p={latest_test['pValue']:.3f} and Cohen's h={latest_test['effectSize']:.2f}.", "Equal-period two-proportion comparison + effect size", "Actionable" if latest_test["pValue"] < 0.05 and abs(latest_test["effectSize"]) >= 0.2 else "Monitor", periods[-2:]))
    questions.append(summary_answer)
    questions.extend([
        answer("Manager estimate confidence", f"Highest-confidence manager estimate: {fmt_row(manager_confidence[0] if manager_confidence else None)}. Confidence is based on sample size and CI width, not score level.", "Survey count + Wilson 95% CI width", "Monitor", manager_confidence[:10]),
        answer("Low-volume agent reliability", f"{len(low_volume_agents)} agent(s) are below the n={minimum_sample} reliability threshold and should not receive definitive ranks.", "Minimum sample threshold + confidence interval width", "Review required" if low_volume_agents else "No action required", low_volume_agents[:20]),
        answer("Managers improving despite fluctuations", f"{len(improving_managers)} reliable manager(s) have a positive weekly regression slope. Strongest signal: {improving_managers[0]['name'] if improving_managers else 'none'}.", "Rolling period means + regression slope + R-squared + SD", "Monitor", improving_managers[:10]),
        answer("Agents with sustained improvement", f"{len(sustained_agents)} reliable agent(s) show a positive slope with at least two positive period movements.", "Consecutive positive changes + regression slope + R-squared", "Monitor", sustained_agents[:10]),
        answer("Managers stable despite volume changes", f"Most stable available manager trend: {min(manager_trends,key=lambda row:row['sd'])['name'] if manager_trends else 'not available'}, based on weekly CSAT dispersion.", "Period means + rolling SD + control-limit interpretation", "Monitor", sorted(manager_trends,key=lambda row:row['sd'])[:10]),
        answer("Recently emerging agent outliers", f"{len(recent_agent_outliers)} reliable agent(s) have a recent-vs-historical movement outside the practical outlier screen.", "Recent mean vs historical mean + standardized movement + sample check", "Review required" if recent_agent_outliers else "No action required", recent_agent_outliers[:10]),
        answer("Most consistent teams", f"Most consistent team: {team_rows[0]['name'] if team_rows else 'not available'} with within-team agent SD {team_rows[0]['sd'] if team_rows else 0:.1f} pts.", "Within-team agent mean SD + CV + range", "Monitor", team_rows[:10]),
        answer("Teams with unusually high variation", f"Highest-variation team: {team_rows[-1]['name'] if team_rows else 'not available'} with within-team SD {team_rows[-1]['sd'] if team_rows else 0:.1f} pts.", "Within-team SD + CV + range + benchmark comparison", "Review required" if team_rows and team_rows[-1]["sd"] >= 10 else "Monitor", list(reversed(team_rows[-10:]))),
        answer("Largest manager performance shifts", f"Largest available manager shift: {max(manager_trends,key=lambda row:abs(row['change']))['name'] if manager_trends else 'not available'} at {max(manager_trends,key=lambda row:abs(row['change']))['change'] if manager_trends else 0:+.1f} pts.", "Recent mean vs historical mean + absolute movement + trend strength", "Review required", sorted(manager_trends,key=lambda row:abs(row['change']),reverse=True)[:10]),
        answer("Agents closest to target", f"Closest reliable agent to the {target:.1f}% target: {near_target_agents[0]['name'] if near_target_agents else 'not available'} at {near_target_agents[0]['mean'] if near_target_agents else 0:.1f}%.", "Absolute target gap + confidence interval + sample threshold", "Monitor", near_target_agents),
        answer("Managers above average throughout", f"{sum(1 for row in manager_trends if row['reliable'] and row['historicalMean']>overall and row['recentMean']>overall)} manager(s) remained above the organizational average in both historical and recent views.", "Entity period means vs organizational mean", "Monitor", manager_trends[:10]),
        answer("Agents above average throughout", f"{sum(1 for row in agent_trends if row['reliable'] and row['historicalMean']>overall and row['recentMean']>overall)} agent(s) remained above the organizational average in both historical and recent views.", "Agent period means vs organizational mean + sample filtering", "Monitor", agent_trends[:10]),
        answer("Statistical differences between managers", f"Across {manager_omnibus['groups']} reliable managers, chi-square={manager_omnibus['statistic']:.2f}, df={manager_omnibus['degreesFreedom']}, p={manager_omnibus['pValue']:.3f}; differences are {'statistically significant' if manager_omnibus['significant'] else 'not statistically distinguishable from random variation'}.", "Multi-group satisfied/not-satisfied chi-square screen", "Actionable" if manager_omnibus["significant"] else "No action required", reliable_managers),
        answer("Statistical differences between agents", f"Across {agent_omnibus['groups']} reliable agents, chi-square={agent_omnibus['statistic']:.2f}, df={agent_omnibus['degreesFreedom']}, p={agent_omnibus['pValue']:.3f}; differences are {'statistically significant' if agent_omnibus['significant'] else 'not statistically distinguishable from random variation'}.", "Multi-group chi-square screen with minimum-sample filtering", "Actionable" if agent_omnibus["significant"] else "No action required", reliable_agents[:20]),
        answer("Most stable time periods", f"Lowest expected score dispersion occurred in {min(periods,key=lambda row:row['mean']*(100-row['mean']))['period'] if periods else 'not available'}; interpret stability separately from performance level.", "Period dispersion + sample count + control limits", "Monitor", periods),
        answer("Most volatile time periods", f"Overall weekly performance SD is {period_sd:.2f} pts and {len(outliers)} period(s) were outside the z-score/IQR screen.", "Period SD + range + rolling SD + outlier count", "Review required" if outliers else "Monitor", outliers or periods[-6:]),
        answer("Sustained improvement or short-term fluctuation", f"The long-run slope is {trend['slope']:+.2f} pts with R-squared {trend['rSquared']:.2f}; latest movement is {movement:+.1f} pts. This is {'a sustained signal' if trend['rSquared']>=0.5 and trend['slope']*movement>0 else 'more consistent with short-term fluctuation'}.", "Recent change vs long-run regression + R-squared", "Monitor", periods),
        answer("Confidence-adjusted agent ranking", f"Top confidence-adjusted agent: {fmt_row(confidence_agents[0] if confidence_agents else None)}, ranked by lower 95% confidence bound.", "Lower Wilson confidence bound + sample threshold", "Monitor", confidence_agents[:10]),
        answer("Managers combining performance and consistency", f"Strongest combined manager: {composite_managers[0]['name'] if composite_managers else 'not available'} with composite score {composite_managers[0]['compositeScore'] if composite_managers else 0:.1f} (mean minus within-team SD penalty).", "Reliable manager mean + within-team SD stability penalty", "Monitor", composite_managers[:10]),
        answer("Top leadership priorities", f"Priorities: {'close the ' + format(target-overall, '.1f') + '-pt shortfall' if overall < target else 'protect the ' + format(overall-target, '.1f') + '-pt target surplus'}, review {len(below_agents)} reliable below-target agents, examine {len(recent_agent_outliers)} recent agent outlier(s), and validate {len(outliers)} unusual time period(s). These are focus areas, not causal findings.", "Reliability-weighted target gap + trend + volatility + outlier synthesis", "Actionable"),
    ])
    framework = _leadership_question_framework()
    for index, question in enumerate(questions):
        if index >= len(framework):
            break
        specification = framework[index]
        question.update({"number": specification.get("number", index + 1), "question": specification.get("question", question["question"]), "logic": specification.get("logic", ""), "statistics": specification.get("statistics", ""), "guardrail": specification.get("guardrail", "")})
    overall_evidence = [{"CSAT": overall, "Target": target, "Target Gap": overall - target, "Responses": total, "CI Low": overall_ci[0], "CI High": overall_ci[1], "Trend": trend_direction, "Slope": trend["slope"], "R Squared": trend["rSquared"]}]
    manager_period_questions = {5, 6, 9, 10, 15, 33, 35, 39, 41}
    agent_period_questions = {7, 8, 11, 12, 16, 23, 24, 34, 36, 42}
    team_questions = {19, 20, 28, 37, 38, 49}
    manager_rank_questions = {3, 17, 25, 31, 43}
    agent_rank_questions = {4, 18, 21, 22, 26, 27, 32, 40, 44, 48}
    period_questions = {1, 13, 14, 29, 30, 45, 46, 47, 50}
    for question in questions:
        if question.get("evidence"):
            continue
        number = int(question.get("number") or 0)
        if number in manager_period_questions:
            question["evidence"] = sorted(manager_trends, key=lambda row: row.get("sd", 0))[:20]
        elif number in agent_period_questions:
            question["evidence"] = sorted(agent_trends, key=lambda row: row.get("sd", 0))[:20]
        elif number in team_questions:
            question["evidence"] = team_rows[:20]
        elif number in manager_rank_questions:
            question["evidence"] = reliable_managers[:20]
        elif number in agent_rank_questions:
            question["evidence"] = reliable_agents[:20]
        elif number in period_questions:
            question["evidence"] = periods[-20:] or overall_evidence
        else:
            question["evidence"] = overall_evidence
    return {"ok": True, "minimumSample": minimum_sample, "target": target, "overall": {"csat": overall, "n": total, "ciLow": overall_ci[0], "ciHigh": overall_ci[1]}, "questions": questions}


def _nps_standard_error(promoters: float, detractors: float, total: float) -> float:
    if total <= 0:
        return 0.0
    promoter_rate = promoters / total
    detractor_rate = detractors / total
    variance = max(0.0, (promoter_rate + detractor_rate - (promoter_rate - detractor_rate) ** 2) / total)
    return math.sqrt(variance) * 100


def _nps_interval(promoters: float, detractors: float, total: float, z: float = 1.96) -> tuple[float, float]:
    if total <= 0:
        return (0.0, 0.0)
    nps_value = ((promoters - detractors) / total) * 100
    margin = z * _nps_standard_error(promoters, detractors, total)
    return (max(-100.0, nps_value - margin), min(100.0, nps_value + margin))


def _nps_change_test(current: dict[str, Any] | None, previous: dict[str, Any] | None) -> dict[str, float]:
    if not current or not previous:
        return {"difference": 0.0, "z": 0.0, "pValue": 1.0, "effectSize": 0.0}
    current_n = float(current.get("n", 0) or 0)
    previous_n = float(previous.get("n", 0) or 0)
    if current_n <= 0 or previous_n <= 0:
        return {"difference": 0.0, "z": 0.0, "pValue": 1.0, "effectSize": 0.0}
    difference = float(current.get("nps", 0) or 0) - float(previous.get("nps", 0) or 0)
    standard_error = math.sqrt((_nps_standard_error(current.get("promoters", 0), current.get("detractors", 0), current_n) ** 2) + (_nps_standard_error(previous.get("promoters", 0), previous.get("detractors", 0), previous_n) ** 2))
    z_score = difference / standard_error if standard_error else 0.0
    return {"difference": difference, "z": z_score, "pValue": math.erfc(abs(z_score) / math.sqrt(2)), "effectSize": difference}


def _nps_leadership_statistics(target: float = 0.0, minimum_sample: int = LEADERSHIP_MIN_SAMPLE, promoter_min: float = 9.0, passive_min: float = 7.0) -> dict[str, Any]:
    minimum_sample = max(2, min(500, int(minimum_sample or LEADERSHIP_MIN_SAMPLE)))
    with STATE_LOCK:
        df = STATE.analyzed_df.copy()
        calendar_settings = dict(STATE.calendar_settings)
    df = _apply_date_filter(df)
    if df.empty:
        return {"ok": False, "error": "No analyzed rows are available for the selected date range."}

    nps_type_col = _first_existing_column(df, ["NPS Type", "NPS Category", "NPS Segment"])
    score_col = _first_existing_column(df, ["NPS Score", "NPS", "Score", "Rating"])
    normalized = df[nps_type_col].fillna("").astype(str).str.strip().str.lower() if nps_type_col else pd.Series([""] * len(df), index=df.index)
    scores = pd.to_numeric(df[score_col], errors="coerce") if score_col else pd.Series([math.nan] * len(df), index=df.index)

    promoter = normalized.str.contains(r"^promoter$|^positive$|^satisfied$", regex=True)
    passive = normalized.str.contains(r"^passive$|^neutral$", regex=True)
    detractor = normalized.str.contains(r"^detractor$|^negative$|^dissatisfied$", regex=True)
    known = promoter | passive | detractor
    promoter = promoter | (~known & scores.ge(float(promoter_min)))
    passive = passive | (~known & scores.lt(float(promoter_min)) & scores.ge(float(passive_min)))
    detractor = detractor | (~known & scores.lt(float(passive_min)))
    valid = promoter | passive | detractor
    if not bool(valid.any()):
        return {"ok": False, "error": "The analyzed data does not contain a usable NPS score or NPS classification."}

    working = df.loc[valid].copy()
    promoter = promoter.loc[valid].astype(float)
    passive = passive.loc[valid].astype(float)
    detractor = detractor.loc[valid].astype(float)
    working["__promoter"] = promoter
    working["__passive"] = passive
    working["__detractor"] = detractor
    total = int(len(working))
    promoters = float(promoter.sum())
    passives = float(passive.sum())
    detractors = float(detractor.sum())
    promoter_pct = promoters / total * 100 if total else 0.0
    passive_pct = passives / total * 100 if total else 0.0
    detractor_pct = detractors / total * 100 if total else 0.0
    overall = promoter_pct - detractor_pct
    overall_ci = _nps_interval(promoters, detractors, total)
    overall_se = _nps_standard_error(promoters, detractors, total)

    agent_col = _first_existing_column(working, ["Agent Name", "AgentName", "Agent"])
    manager_col = _first_existing_column(working, ["Manager/TL", "ManagerName", "Manager", "Team"])
    date_col = _first_existing_column(working, ["Feedback Date", "CallDateTime", "Date", "Response Date"])

    def entity_rows(column: str) -> list[dict[str, Any]]:
        if not column:
            return []
        rows: list[dict[str, Any]] = []
        for name, group in working.assign(__entity=working[column].fillna("Not available").astype(str)).groupby("__entity", dropna=False):
            count = int(len(group))
            pro = float(group["__promoter"].sum())
            pas = float(group["__passive"].sum())
            det = float(group["__detractor"].sum())
            nps_value = (pro - det) / count * 100 if count else 0.0
            low, high = _nps_interval(pro, det, count)
            rows.append({
                "name": str(name),
                "nps": nps_value,
                "n": count,
                "promoters": pro,
                "passives": pas,
                "detractors": det,
                "promoterPct": pro / count * 100 if count else 0.0,
                "passivePct": pas / count * 100 if count else 0.0,
                "detractorPct": det / count * 100 if count else 0.0,
                "ciLow": low,
                "ciHigh": high,
                "ciWidth": high - low,
                "gap": nps_value - overall,
                "targetGap": nps_value - target,
                "reliable": count >= minimum_sample,
            })
        rows.sort(key=lambda item: item["nps"], reverse=True)
        total_rows = len(rows)
        for index, row in enumerate(rows, 1):
            row["percentileRank"] = 100.0 if total_rows <= 1 else (total_rows - index) / (total_rows - 1) * 100
            row["reliabilityScore"] = min(100.0, (row["n"] / minimum_sample) * 50 + max(0.0, 50 - row["ciWidth"] / 2))
        return rows

    agents = entity_rows(agent_col)
    managers = entity_rows(manager_col)
    reliable_agents = [row for row in agents if row["reliable"]]
    reliable_managers = [row for row in managers if row["reliable"]]

    periods: list[dict[str, Any]] = []
    if date_col:
        frame = working[["__promoter", "__passive", "__detractor"]].copy()
        frame["date"] = pd.to_datetime(working[date_col], errors="coerce")
        frame = frame.dropna(subset=["date"])
        if not frame.empty:
            frame["period"] = week_period_start(frame["date"], calendar_settings["weekStart"])
            grouped = frame.groupby("period")[["__promoter", "__passive", "__detractor"]].sum()
            counts = frame.groupby("period").size()
            for period, row in grouped.iterrows():
                count = int(counts.loc[period])
                pro = float(row["__promoter"])
                pas = float(row["__passive"])
                det = float(row["__detractor"])
                low, high = _nps_interval(pro, det, count)
                periods.append({
                    "period": str(period.date()),
                    "nps": (pro - det) / count * 100 if count else 0.0,
                    "n": count,
                    "promoters": pro,
                    "passives": pas,
                    "detractors": det,
                    "promoterPct": pro / count * 100 if count else 0.0,
                    "passivePct": pas / count * 100 if count else 0.0,
                    "detractorPct": det / count * 100 if count else 0.0,
                    "ciLow": low,
                    "ciHigh": high,
                })
            periods.sort(key=lambda row: row["period"])

    nps_values = [row["nps"] for row in periods]
    trend = _linear_trend(nps_values)
    promoter_trend = _linear_trend([row["promoterPct"] for row in periods])
    passive_trend = _linear_trend([row["passivePct"] for row in periods])
    detractor_trend = _linear_trend([row["detractorPct"] for row in periods])
    latest = periods[-1] if periods else None
    previous = periods[-2] if len(periods) >= 2 else None
    latest_test = _nps_change_test(latest, previous)
    mean_period = sum(nps_values) / len(nps_values) if nps_values else overall
    period_sd = float(pd.Series(nps_values).std(ddof=1)) if len(nps_values) > 1 else 0.0
    q1 = float(pd.Series(nps_values).quantile(0.25)) if nps_values else 0.0
    q3 = float(pd.Series(nps_values).quantile(0.75)) if nps_values else 0.0
    iqr = q3 - q1
    outliers = [row for row in periods if (period_sd and abs((row["nps"] - mean_period) / period_sd) >= 2) or row["nps"] < q1 - 1.5 * iqr or row["nps"] > q3 + 1.5 * iqr]
    trend_direction = "improving" if trend["slope"] > 0.5 else "declining" if trend["slope"] < -0.5 else "stable"
    significance = "statistically significant" if latest_test["pValue"] < 0.05 else "directional, not statistically significant"
    practical = abs(latest_test["difference"]) >= 5.0

    def fmt_entity(row: dict[str, Any] | None) -> str:
        if not row:
            return "not available"
        return f"{row['name']} NPS {row['nps']:.1f} (n={row['n']}, 95% CI {row['ciLow']:.1f}-{row['ciHigh']:.1f})"

    def compare_rows(rows: list[dict[str, Any]], column: str) -> list[dict[str, Any]]:
        if not column or not date_col:
            return []
        frame = working[[column, "__promoter", "__passive", "__detractor"]].copy()
        frame["date"] = pd.to_datetime(working[date_col], errors="coerce")
        frame = frame.dropna(subset=["date"])
        if frame.empty:
            return []
        cutoff = frame["date"].median()
        frame["window"] = frame["date"].apply(lambda value: "Current" if value >= cutoff else "Previous")
        results: list[dict[str, Any]] = []
        for entity, entity_frame in frame.assign(__entity=frame[column].fillna("Not available").astype(str)).groupby("__entity", dropna=False):
            windows: dict[str, dict[str, Any]] = {}
            for window, group in entity_frame.groupby("window"):
                count = int(len(group))
                pro = float(group["__promoter"].sum())
                pas = float(group["__passive"].sum())
                det = float(group["__detractor"].sum())
                windows[window] = {"n": count, "promoters": pro, "passives": pas, "detractors": det, "nps": (pro - det) / count * 100 if count else 0.0, "detractorPct": det / count * 100 if count else 0.0}
            current = windows.get("Current")
            prior = windows.get("Previous")
            if not current or not prior:
                continue
            test = _nps_change_test(current, prior)
            period_values = []
            dated = entity_frame.copy()
            dated["period"] = week_period_start(dated["date"], calendar_settings["weekStart"])
            for _, group in dated.groupby("period"):
                count = int(len(group))
                if count:
                    period_values.append((float(group["__promoter"].sum()) - float(group["__detractor"].sum())) / count * 100)
            entity_trend = _linear_trend(period_values)
            entity_sd = float(pd.Series(period_values).std(ddof=1)) if len(period_values) > 1 else 0.0
            positive_changes = sum(1 for prior_value, current_value in zip(period_values, period_values[1:]) if current_value > prior_value)
            results.append({
                "name": str(entity),
                "n": int(entity_frame.shape[0]),
                "previousNps": prior["nps"],
                "currentNps": current["nps"],
                "change": test["difference"],
                "pValue": test["pValue"],
                "effectSize": test["effectSize"],
                "currentN": current["n"],
                "previousN": prior["n"],
                "detractorPctChange": current["detractorPct"] - prior["detractorPct"],
                "periods": len(period_values),
                "slope": entity_trend["slope"],
                "rSquared": entity_trend["rSquared"],
                "sd": entity_sd,
                "positiveChanges": positive_changes,
                "reliable": current["n"] >= minimum_sample and prior["n"] >= minimum_sample,
            })
        return results

    manager_changes = compare_rows(managers, manager_col)
    agent_changes = compare_rows(agents, agent_col)
    manager_improvers = sorted([row for row in manager_changes if row["reliable"] and row["change"] > 0], key=lambda row: (row["change"], -row["pValue"]), reverse=True)
    manager_decliners = sorted([row for row in manager_changes if row["reliable"] and row["change"] < 0], key=lambda row: row["change"])
    agent_improvers = sorted([row for row in agent_changes if row["reliable"] and row["change"] > 0], key=lambda row: (row["change"], -row["pValue"]), reverse=True)
    agent_decliners = sorted([row for row in agent_changes if row["reliable"] and row["change"] < 0], key=lambda row: row["change"])
    manager_stability = sorted([row for row in manager_changes if row["periods"] >= 2], key=lambda row: row["sd"])
    agent_stability = sorted([row for row in agent_changes if row["periods"] >= 2], key=lambda row: row["sd"])
    manager_outperformers = [row for row in reliable_managers if row["nps"] > overall]
    agent_outperformers = [row for row in reliable_agents if row["nps"] > overall]
    above_agents = [row for row in reliable_agents if row["nps"] >= target]
    below_agents = [row for row in reliable_agents if row["nps"] < target]
    high_det_managers = sorted(reliable_managers, key=lambda row: row["detractorPct"], reverse=True)
    high_det_agents = sorted(reliable_agents, key=lambda row: row["detractorPct"], reverse=True)
    high_pro_managers = sorted(reliable_managers, key=lambda row: row["promoterPct"], reverse=True)
    high_pro_agents = sorted(reliable_agents, key=lambda row: row["promoterPct"], reverse=True)
    high_nps_low_reliability_managers = sorted([row for row in managers if row["nps"] >= overall and not row["reliable"]], key=lambda row: row["nps"], reverse=True)
    high_nps_low_reliability_agents = sorted([row for row in agents if row["nps"] >= overall and not row["reliable"]], key=lambda row: row["nps"], reverse=True)
    low_nps_high_reliability_managers = sorted([row for row in reliable_managers if row["nps"] < overall], key=lambda row: row["nps"])
    low_nps_high_reliability_agents = sorted([row for row in reliable_agents if row["nps"] < overall], key=lambda row: row["nps"])
    confidence_managers = sorted(reliable_managers, key=lambda row: (row["ciWidth"], -row["nps"]))
    confidence_agents = sorted(reliable_agents, key=lambda row: (row["ciWidth"], -row["nps"]))

    team_rows: list[dict[str, Any]] = []
    if manager_col and agent_col:
        for manager, manager_frame in working.assign(__manager=working[manager_col].fillna("Not available").astype(str), __agent=working[agent_col].fillna("Not available").astype(str)).groupby("__manager"):
            agent_values = []
            for agent, group in manager_frame.groupby("__agent"):
                count = int(len(group))
                if count < minimum_sample:
                    continue
                agent_values.append({"agent": str(agent), "nps": (float(group["__promoter"].sum()) - float(group["__detractor"].sum())) / count * 100, "n": count})
            values = [row["nps"] for row in agent_values]
            if not values:
                continue
            above_count = sum(1 for value in values if value >= target)
            team_rows.append({"name": str(manager), "agents": len(values), "mean": sum(values) / len(values), "sd": float(pd.Series(values).std(ddof=1)) if len(values) > 1 else 0.0, "range": max(values) - min(values), "iqr": float(pd.Series(values).quantile(0.75) - pd.Series(values).quantile(0.25)) if len(values) > 1 else 0.0, "aboveTarget": above_count, "aboveTargetPct": above_count / len(values) * 100})
        team_rows.sort(key=lambda row: row["sd"])

    def summary_stats(rows: list[dict[str, Any]], key: str = "nps") -> dict[str, Any]:
        values = [float(row[key]) for row in rows if math.isfinite(float(row.get(key, 0)))]
        if not values:
            return {"count": 0, "mean": 0.0, "median": 0.0, "q1": 0.0, "q3": 0.0, "iqr": 0.0, "sd": 0.0, "range": 0.0}
        series = pd.Series(values)
        return {"count": len(values), "mean": float(series.mean()), "median": float(series.median()), "q1": float(series.quantile(0.25)), "q3": float(series.quantile(0.75)), "iqr": float(series.quantile(0.75) - series.quantile(0.25)), "sd": float(series.std(ddof=1)) if len(values) > 1 else 0.0, "range": max(values) - min(values)}

    manager_distribution = summary_stats(reliable_managers)
    agent_distribution = summary_stats(reliable_agents)
    strongest_changes = sorted(manager_improvers + agent_improvers + manager_decliners + agent_decliners, key=lambda row: (row["pValue"], -abs(row["change"])))[:20]
    highest_confidence_findings = confidence_agents[:5] + confidence_managers[:5]
    best_period = max(periods, key=lambda row: row["nps"]) if periods else None
    worst_period = min(periods, key=lambda row: row["nps"]) if periods else None
    volatility_direction = "more stable" if len(nps_values) >= 4 and float(pd.Series(nps_values[-2:]).std(ddof=1)) < float(pd.Series(nps_values[:-2]).std(ddof=1)) else "more volatile" if len(nps_values) >= 4 else "not enough periods"

    def answer(question: str, text: str, method: str, status: str = "Actionable", evidence: list[dict[str, Any]] | None = None) -> dict[str, Any]:
        return {"question": question, "text": text, "method": method, "status": status, "evidence": evidence or []}

    questions = [
        answer("Overall NPS", f"Overall NPS is {overall:.1f} from n={total:,}; promoters {promoter_pct:.1f}%, passives {passive_pct:.1f}%, detractors {detractor_pct:.1f}%. 95% CI {overall_ci[0]:.1f} to {overall_ci[1]:.1f}; standard error {overall_se:.2f}.", "Promoter % minus detractor % with NPS standard error and 95% CI"),
        answer("NPS trend over time", f"NPS is {trend_direction}; weekly slope {trend['slope']:+.2f} pts, R-squared {trend['rSquared']:.2f}. Latest movement {latest_test['difference']:+.1f} pts is {significance}.", "Weekly NPS aggregation + linear regression + NPS change z-test", "Actionable" if latest_test["pValue"] < 0.05 and practical else "Monitor", periods[-12:]),
        answer("NPS change vs previous period", f"Latest NPS changed {latest_test['difference']:+.1f} pts versus the previous equal period; p={latest_test['pValue']:.3f}. Practical significance threshold is 5 pts.", "NPS difference z-test using promoter/detractor variance", "Actionable" if latest_test["pValue"] < 0.05 and practical else "Monitor", periods[-2:]),
        answer("Promoter, passive and detractor distribution", f"Distribution: promoters {int(promoters):,} ({promoter_pct:.1f}%), passives {int(passives):,} ({passive_pct:.1f}%), detractors {int(detractors):,} ({detractor_pct:.1f}%) across {total:,} surveys.", "NPS category counts and Wilson confidence intervals"),
        answer("Managers with highest NPS", "Top reliable managers: " + "; ".join(fmt_entity(row) for row in reliable_managers[:5]), "Minimum-sample manager ranking with NPS CI", evidence=reliable_managers[:10]),
        answer("Managers with lowest NPS", "Lowest reliable managers: " + "; ".join(fmt_entity(row) for row in reliable_managers[-5:]), "Minimum-sample manager ranking with NPS CI", "Review required" if reliable_managers else "No evidence", reliable_managers[-10:]),
        answer("Agents with highest NPS", "Top reliable agents: " + "; ".join(fmt_entity(row) for row in reliable_agents[:10]), "Minimum-sample agent ranking with NPS CI", evidence=reliable_agents[:10]),
        answer("Agents with lowest NPS", "Lowest reliable agents: " + "; ".join(fmt_entity(row) for row in reliable_agents[-10:]), "Minimum-sample agent ranking with NPS CI", "Review required" if reliable_agents else "No evidence", reliable_agents[-10:]),
        answer("Managers improved the most", f"{len(manager_improvers)} reliable manager(s) improved across equal windows. Strongest: {manager_improvers[0]['name'] if manager_improvers else 'not available'}.", "Equal-window manager NPS change + z-test", "Actionable" if manager_improvers and manager_improvers[0]["pValue"] < 0.05 else "Monitor", manager_improvers[:10]),
        answer("Managers declined the most", f"{len(manager_decliners)} reliable manager(s) declined across equal windows. Largest decline: {manager_decliners[0]['name'] if manager_decliners else 'not available'}.", "Equal-window manager NPS change + detractor movement", "Review required" if manager_decliners else "No action required", manager_decliners[:10]),
        answer("Agents improved the most", f"{len(agent_improvers)} reliable agent(s) improved across equal windows. Strongest: {agent_improvers[0]['name'] if agent_improvers else 'not available'}.", "Equal-window agent NPS change + z-test", "Actionable" if agent_improvers and agent_improvers[0]["pValue"] < 0.05 else "Monitor", agent_improvers[:10]),
        answer("Agents declined the most", f"{len(agent_decliners)} reliable agent(s) declined across equal windows. Largest decline: {agent_decliners[0]['name'] if agent_decliners else 'not available'}.", "Equal-window agent NPS change + detractor movement", "Review required" if agent_decliners else "No action required", agent_decliners[:10]),
        answer("Managers with most consistent NPS", f"Most consistent manager: {manager_stability[0]['name'] if manager_stability else 'not available'} with period SD {manager_stability[0]['sd'] if manager_stability else 0:.1f} pts.", "Period NPS SD, rolling SD, CV and IQR", "Monitor", manager_stability[:10]),
        answer("Managers with most volatile NPS", f"Most volatile manager: {manager_stability[-1]['name'] if manager_stability else 'not available'} with period SD {manager_stability[-1]['sd'] if manager_stability else 0:.1f} pts.", "Period NPS SD, range and outlier count", "Review required" if manager_stability and manager_stability[-1]["sd"] >= 10 else "Monitor", list(reversed(manager_stability[-10:]))),
        answer("Agents with most consistent NPS", f"Most consistent agent: {agent_stability[0]['name'] if agent_stability else 'not available'} with period SD {agent_stability[0]['sd'] if agent_stability else 0:.1f} pts.", "Agent period NPS SD and IQR", "Monitor", agent_stability[:10]),
        answer("Agents with most volatile NPS", f"Most volatile agent: {agent_stability[-1]['name'] if agent_stability else 'not available'} with period SD {agent_stability[-1]['sd'] if agent_stability else 0:.1f} pts.", "Agent period NPS SD, range and control-limit screen", "Review required" if agent_stability and agent_stability[-1]["sd"] >= 10 else "Monitor", list(reversed(agent_stability[-10:]))),
        answer("Highest NPS week/month", f"Highest period: {best_period['period'] if best_period else 'not available'} at NPS {best_period['nps'] if best_period else 0:.1f}.", "Period NPS rank with CI", "Monitor", periods),
        answer("Lowest NPS week/month", f"Lowest period: {worst_period['period'] if worst_period else 'not available'} at NPS {worst_period['nps'] if worst_period else 0:.1f}.", "Period NPS rank with CI", "Review required" if worst_period else "No evidence", periods),
        answer("Unusual NPS spikes or drops", f"Detected {len(outliers)} unusual NPS period(s) using z-score/IQR screening.", "Z-score plus IQR outlier detection", "Review required" if outliers else "No action required", outliers),
        answer("Managers outperforming organization", f"{len(manager_outperformers)} of {len(reliable_managers)} reliable managers are above organization NPS {overall:.1f}.", "Manager NPS minus organization NPS with CI", "Monitor", manager_outperformers[:20]),
        answer("Agents outperforming organization", f"{len(agent_outperformers)} of {len(reliable_agents)} reliable agents are above organization NPS {overall:.1f}.", "Agent NPS minus organization NPS with CI", "Monitor", agent_outperformers[:20]),
        answer("NPS distribution across managers", f"Reliable manager NPS distribution: median {manager_distribution['median']:.1f}, IQR {manager_distribution['iqr']:.1f}, SD {manager_distribution['sd']:.1f}, range {manager_distribution['range']:.1f}.", "Mean, median, quartiles, IQR, SD and range", "Monitor", reliable_managers),
        answer("NPS distribution across agents", f"Reliable agent NPS distribution: median {agent_distribution['median']:.1f}, IQR {agent_distribution['iqr']:.1f}, SD {agent_distribution['sd']:.1f}, range {agent_distribution['range']:.1f}.", "Mean, median, quartiles, IQR, SD and percentiles", "Monitor", reliable_agents[:50]),
        answer("Managers with most balanced teams", f"Most balanced team: {team_rows[0]['name'] if team_rows else 'not available'} with within-team SD {team_rows[0]['sd'] if team_rows else 0:.1f} pts.", "Within-team agent NPS SD, range and IQR", "Monitor", team_rows[:10]),
        answer("Largest within-team performance gap", f"Largest team gap: {team_rows[-1]['name'] if team_rows else 'not available'} with range {team_rows[-1]['range'] if team_rows else 0:.1f} pts.", "Highest minus lowest reliable agent NPS within manager", "Review required" if team_rows else "No evidence", list(reversed(team_rows[-10:]))),
        answer("Top 10 agents by NPS", "Top 10 reliable agents: " + "; ".join(fmt_entity(row) for row in reliable_agents[:10]), "Confidence-adjusted reliable ranking", "Monitor", reliable_agents[:10]),
        answer("Bottom 10 agents by NPS", "Bottom 10 reliable agents: " + "; ".join(fmt_entity(row) for row in reliable_agents[-10:]), "Low-end reliable ranking with CI width", "Review required" if reliable_agents else "No evidence", reliable_agents[-10:]),
        answer("Agents improving month over month", f"{sum(1 for row in agent_changes if row['reliable'] and row['slope'] > 0 and row['positiveChanges'] >= 2)} reliable agent(s) show sustained positive NPS slope.", "Monthly/period slope + consecutive improvement count", "Monitor", sorted(agent_changes, key=lambda row: row["slope"], reverse=True)[:10]),
        answer("Managers improving month over month", f"{sum(1 for row in manager_changes if row['reliable'] and row['slope'] > 0)} reliable manager(s) show positive NPS trend slope.", "Monthly/period slope + R-squared", "Monitor", sorted(manager_changes, key=lambda row: row["slope"], reverse=True)[:10]),
        answer("Agents needing attention based on NPS", f"{len(below_agents)} reliable agent(s) are below target {target:.1f}; highest detractor concentration is {high_det_agents[0]['name'] if high_det_agents else 'not available'}.", "Low NPS + high detractor rate + trend + sample validation", "Review required" if below_agents else "No action required", below_agents[:20]),
        answer("Managers needing attention based on NPS", f"{len([row for row in reliable_managers if row['nps'] < target])} reliable manager(s) are below target {target:.1f}; highest detractor concentration is {high_det_managers[0]['name'] if high_det_managers else 'not available'}.", "Low NPS + detractor rate + team variability + CI", "Review required", [row for row in reliable_managers if row["nps"] < target][:20]),
        answer("Agents above NPS target", f"{len(above_agents)} of {len(reliable_agents)} reliable agents ({(len(above_agents)/len(reliable_agents)*100) if reliable_agents else 0:.1f}%) are at or above target {target:.1f}.", "Reliable agent target comparison", "Monitor", above_agents[:20]),
        answer("Agents below NPS target", f"{len(below_agents)} of {len(reliable_agents)} reliable agents ({(len(below_agents)/len(reliable_agents)*100) if reliable_agents else 0:.1f}%) are below target {target:.1f}.", "Reliable agent target comparison", "Review required" if below_agents else "No action required", below_agents[:20]),
        answer("Manager with highest percentage of agents above target", f"Best team target coverage: {max(team_rows,key=lambda row:row['aboveTargetPct'])['name'] if team_rows else 'not available'} at {max(team_rows,key=lambda row:row['aboveTargetPct'])['aboveTargetPct'] if team_rows else 0:.1f}% of reliable agents above target.", "Within-manager proportion of reliable agents above target", "Monitor", sorted(team_rows, key=lambda row: row["aboveTargetPct"], reverse=True)[:10]),
        answer("Manager with highest detractor concentration", f"Highest reliable manager detractor concentration: {fmt_entity(high_det_managers[0] if high_det_managers else None)} with detractors {high_det_managers[0]['detractorPct'] if high_det_managers else 0:.1f}%.", "Detractor rate by manager with CI and organization comparison", "Review required" if high_det_managers else "No evidence", high_det_managers[:10]),
        answer("Agent with highest detractor concentration", f"Highest reliable agent detractor concentration: {fmt_entity(high_det_agents[0] if high_det_agents else None)} with detractors {high_det_agents[0]['detractorPct'] if high_det_agents else 0:.1f}%.", "Detractor rate by agent with sample validation", "Review required" if high_det_agents else "No evidence", high_det_agents[:10]),
        answer("Manager with highest promoter concentration", f"Highest reliable manager promoter concentration: {fmt_entity(high_pro_managers[0] if high_pro_managers else None)} with promoters {high_pro_managers[0]['promoterPct'] if high_pro_managers else 0:.1f}%.", "Promoter rate by manager with organization comparison", "Monitor", high_pro_managers[:10]),
        answer("Agent with highest promoter concentration", f"Highest reliable agent promoter concentration: {fmt_entity(high_pro_agents[0] if high_pro_agents else None)} with promoters {high_pro_agents[0]['promoterPct'] if high_pro_agents else 0:.1f}%.", "Promoter rate by agent with sample validation", "Monitor", high_pro_agents[:10]),
        answer("Previous reporting period changes", f"Latest period changed {latest_test['difference']:+.1f} NPS pts; promoter change {((latest['promoterPct']-previous['promoterPct']) if latest and previous else 0):+.1f} pts, passive change {((latest['passivePct']-previous['passivePct']) if latest and previous else 0):+.1f} pts, detractor change {((latest['detractorPct']-previous['detractorPct']) if latest and previous else 0):+.1f} pts.", "Comparable period NPS and category movement", "Actionable" if latest_test["pValue"] < 0.05 and practical else "Monitor", periods[-2:]),
        answer("Strongest statistically significant NPS changes", f"{sum(1 for row in strongest_changes if row['pValue'] < 0.05 and abs(row['change']) >= 5)} entity change(s) meet p<0.05 and 5-pt practical movement.", "Sort by p-value and NPS effect size", "Actionable" if strongest_changes else "No evidence", strongest_changes),
        answer("Managers with high NPS but low reliability", f"{len(high_nps_low_reliability_managers)} manager(s) have above-average NPS but low reliability due to volume/CI width.", "High NPS + low sample or wide CI", "Monitor", high_nps_low_reliability_managers[:20]),
        answer("Agents with high NPS but low reliability", f"{len(high_nps_low_reliability_agents)} agent(s) have above-average NPS but low reliability due to volume/CI width.", "High NPS + low sample or wide CI", "Monitor", high_nps_low_reliability_agents[:20]),
        answer("Managers with low NPS and high reliability", f"{len(low_nps_high_reliability_managers)} reliable manager(s) are below organization NPS. Lowest: {fmt_entity(low_nps_high_reliability_managers[0] if low_nps_high_reliability_managers else None)}.", "Low NPS + adequate sample + CI reliability", "Review required" if low_nps_high_reliability_managers else "No action required", low_nps_high_reliability_managers[:20]),
        answer("Agents with low NPS and high reliability", f"{len(low_nps_high_reliability_agents)} reliable agent(s) are below organization NPS. Lowest: {fmt_entity(low_nps_high_reliability_agents[0] if low_nps_high_reliability_agents else None)}.", "Low NPS + adequate sample + CI reliability", "Review required" if low_nps_high_reliability_agents else "No action required", low_nps_high_reliability_agents[:20]),
        answer("Organizational NPS stability", f"Organizational NPS is {volatility_direction}; period SD is {period_sd:.1f} pts across {len(periods)} period(s).", "Rolling SD, period SD and control-limit screen", "Monitor", periods),
        answer("Passive responses increasing or decreasing", f"Passive trend slope is {passive_trend['slope']:+.2f} pts per period; latest passive movement is {((latest['passivePct']-previous['passivePct']) if latest and previous else 0):+.1f} pts.", "Passive percentage trend slope and period change", "Monitor", periods[-12:]),
        answer("Detractors increasing or decreasing", f"Detractor trend slope is {detractor_trend['slope']:+.2f} pts per period; latest detractor movement is {((latest['detractorPct']-previous['detractorPct']) if latest and previous else 0):+.1f} pts.", "Detractor percentage trend and z-test context", "Review required" if detractor_trend["slope"] > 0.5 else "Monitor", periods[-12:]),
        answer("Promoters increasing or decreasing", f"Promoter trend slope is {promoter_trend['slope']:+.2f} pts per period; latest promoter movement is {((latest['promoterPct']-previous['promoterPct']) if latest and previous else 0):+.1f} pts.", "Promoter percentage trend and z-test context", "Monitor", periods[-12:]),
        answer("Highest-confidence NPS findings", f"Highest-confidence read combines {len(confidence_agents)} reliable agents and {len(confidence_managers)} reliable managers with narrow CI width. Best agent confidence: {fmt_entity(confidence_agents[0] if confidence_agents else None)}.", "Sample size + CI width + reliability score", "Monitor", highest_confidence_findings),
        answer("Key NPS insights", f"NPS {overall:.1f} (95% CI {overall_ci[0]:.1f} to {overall_ci[1]:.1f}), {trend_direction} trend, promoter/passive/detractor split {promoter_pct:.1f}%/{passive_pct:.1f}%/{detractor_pct:.1f}%, {len(outliers)} outlier period(s), {len(below_agents)} reliable agents below target.", "Statistically guarded executive NPS summary", "Actionable", [{"NPS": overall, "Target": target, "Promoter %": promoter_pct, "Passive %": passive_pct, "Detractor %": detractor_pct, "Responses": total, "CI Low": overall_ci[0], "CI High": overall_ci[1], "Trend": trend_direction}]),
    ]

    framework = _nps_leadership_question_framework()
    for index, question in enumerate(questions):
        if index >= len(framework):
            break
        specification = framework[index]
        question.update({"number": specification.get("number", index + 1), "question": specification.get("question", question["question"]), "logic": specification.get("logic", ""), "statistics": specification.get("statistics", ""), "guardrail": specification.get("guardrail", "")})

    return {
        "ok": True,
        "mode": "nps",
        "minimumSample": minimum_sample,
        "target": target,
        "overall": {
            "nps": overall,
            "n": total,
            "promoterPct": promoter_pct,
            "passivePct": passive_pct,
            "detractorPct": detractor_pct,
            "ciLow": overall_ci[0],
            "ciHigh": overall_ci[1],
            "standardError": overall_se,
        },
        "questions": questions,
    }


def _leadership_results_workbook(payload: dict[str, Any]) -> bytes:
    mode = str(payload.get("mode") or "csat").lower()
    metric_label = "NPS" if mode == "nps" else "CSAT"
    promoter_label = "Promoter Minimum" if mode == "nps" else "Satisfied Minimum"
    passive_label = "Passive Minimum" if mode == "nps" else "Neutral Minimum"
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    target = _safe_float(payload.get("target"), 85.0)
    minimum_sample = int(_safe_float(payload.get("minimumSample"), LEADERSHIP_MIN_SAMPLE))
    satisfied_min = _safe_float(payload.get("satisfiedMin"), 4.0)
    neutral_min = _safe_float(payload.get("neutralMin"), 3.0)
    if isinstance(payload.get("calendar"), dict) or isinstance(payload.get("rules"), dict):
        calendar_source = payload.get("calendar") if isinstance(payload.get("calendar"), dict) else payload.get("rules")
        with STATE_LOCK:
            STATE.calendar_settings = _calendar_settings({"calendar": calendar_source})
    if mode == "nps":
        result = _nps_leadership_statistics(target, minimum_sample, _safe_float(payload.get("promoterMin"), satisfied_min), _safe_float(payload.get("passiveMin"), neutral_min))
    else:
        result = _leadership_statistics(target, minimum_sample, satisfied_min, neutral_min)
    if not result.get("ok"):
        raise ValueError(result.get("error") or "Leadership results are not available.")

    with STATE_LOCK:
        analyzed_df = STATE.analyzed_df.copy()
        weekly_df = STATE.weekly_df.copy()
        agent_df = STATE.agent_df.copy()
        manager_df = STATE.manager_df.copy()
        reason_df = STATE.reason_df.copy()
        complaints_df = STATE.complaints_df.copy()
        passive_df = STATE.passive_df.copy()
        base_profile = dict(STATE.base_column_profile or {})
        analyzed_profile = dict(STATE.analyzed_column_profile or {})
        files = dict(STATE.files)
        file_sizes = dict(STATE.file_sizes)
        analysis_engines = dict(STATE.analysis_engines)
        model_paths = dict(STATE.model_paths)

    workbook = Workbook()
    results_sheet = workbook.active
    results_sheet.title = "Results"
    header_fill = PatternFill("solid", fgColor="0B7180")
    header_font = Font(color="FFFFFF", bold=True)
    sub_fill = PatternFill("solid", fgColor="E8F4F5")
    thin = Side(style="thin", color="D4E5E7")
    border = Border(bottom=thin)

    def clean(value: Any) -> Any:
        if value is None:
            return ""
        try:
            if pd.isna(value):
                return ""
        except Exception:
            pass
        if isinstance(value, (dict, list, tuple)):
            value = json.dumps(value, ensure_ascii=False)
        if isinstance(value, str):
            value = "".join(ch for ch in value if ch in "\t\n\r" or ord(ch) >= 32)
            return value[:32000]
        return value

    def style_sheet(sheet) -> None:
        if sheet.max_row >= 1:
            for cell in sheet[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(vertical="center", wrap_text=True)
            sheet.freeze_panes = "A2"
            sheet.auto_filter.ref = sheet.dimensions
        for row in sheet.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
                cell.border = border

    def fit_columns(sheet, widths: list[int] | None = None) -> None:
        for index, column_cells in enumerate(sheet.columns, 1):
            letter = column_cells[0].column_letter
            if widths and index <= len(widths):
                sheet.column_dimensions[letter].width = widths[index - 1]
                continue
            max_length = max((len(str(cell.value or "")) for cell in column_cells[:100]), default=12)
            sheet.column_dimensions[letter].width = min(max(max_length + 2, 12), 46)

    def write_rows(sheet_name: str, rows: list[dict[str, Any]], columns: list[str] | None = None, widths: list[int] | None = None):
        sheet = workbook.create_sheet(sheet_name[:31])
        if not rows:
            sheet.append(["Status"])
            sheet.append(["No data available"])
            style_sheet(sheet)
            fit_columns(sheet, widths)
            return sheet
        if columns is None:
            columns = []
            for row in rows:
                for key in row.keys():
                    if key not in columns:
                        columns.append(str(key))
        sheet.append(columns)
        for row in rows:
            sheet.append([clean(row.get(column)) for column in columns])
        style_sheet(sheet)
        fit_columns(sheet, widths)
        return sheet

    def write_frame(sheet_name: str, frame: pd.DataFrame, columns: list[str] | None = None):
        if frame is None or frame.empty:
            write_rows(sheet_name, [])
            return
        working = frame.copy()
        if columns:
            available = [column for column in columns if column in working.columns]
            if available:
                working = working[available]
        rows = working.where(pd.notna(working), "").to_dict(orient="records")
        write_rows(sheet_name, rows)

    def answer_rows(items: list[dict[str, Any]], area: str, offset: int = 0) -> list[dict[str, Any]]:
        output = []
        for index, item in enumerate(items or [], 1):
            status = item.get("status", "")
            output.append({
                "#": item.get("number") or index + offset,
                "Area": area,
                "Question": item.get("question", ""),
                "Answer": item.get("text", ""),
                "Evidence Status": status,
                "Recommended Action": recommended_action(status),
                "Statistical Method": item.get("method", ""),
                "Logic": item.get("logic", ""),
                "Statistics Used": item.get("statistics", ""),
                "Interpretation Guardrail": item.get("guardrail", ""),
                "Outcome": item.get("outcome", ""),
            })
        return output

    def evidence_rows(items: list[dict[str, Any]], area: str, offset: int = 0) -> list[dict[str, Any]]:
        rows = []
        for index, item in enumerate(items or [], 1):
            number = item.get("number") or index + offset
            evidence = item.get("evidence") if isinstance(item.get("evidence"), list) else []
            if not evidence:
                rows.append({"Question #": number, "Area": area, "Question": item.get("question", ""), "Artifact Row": 0, "Data Point": "Evidence", "Value": "No row-level artifact available"})
                continue
            for artifact_index, artifact in enumerate(evidence, 1):
                if isinstance(artifact, dict):
                    for key, value in artifact.items():
                        rows.append({"Question #": number, "Area": area, "Question": item.get("question", ""), "Artifact Row": artifact_index, "Data Point": key, "Value": clean(value)})
                else:
                    rows.append({"Question #": number, "Area": area, "Question": item.get("question", ""), "Artifact Row": artifact_index, "Data Point": "Value", "Value": clean(artifact)})
        return rows

    def recommended_action(status: str) -> str:
        value = str(status or "").lower()
        if "no action" in value:
            return "No action required"
        if "action" in value:
            return "Investigate"
        if "review" in value:
            return "Review"
        if "monitor" in value:
            return "Monitor"
        if "no evidence" in value:
            return "No evidence"
        return "Validate evidence"

    score_answers = payload.get("scoreAnswers") if isinstance(payload.get("scoreAnswers"), list) else result["questions"]
    sentiment_answers = payload.get("sentimentAnswers") if isinstance(payload.get("sentimentAnswers"), list) else []
    analysis_payload = payload.get("analysis") if isinstance(payload.get("analysis"), dict) else {}
    base_payload = payload.get("base") if isinstance(payload.get("base"), dict) else {}
    rules_payload = payload.get("rules") if isinstance(payload.get("rules"), dict) else {}

    result_headers = ["#", "Question", "Answer", "Evidence Status", "Recommended Action", "Statistical Method", "Logic", "Statistics Used", "Interpretation Guardrail"]
    results_sheet.append(result_headers)

    for item in score_answers:
        results_sheet.append([
            item.get("number"), item.get("question"), item.get("text"), item.get("status"), recommended_action(item.get("status", "")),
            item.get("method"), item.get("logic"), item.get("statistics"), item.get("guardrail"),
        ])
    for index, item in enumerate(sentiment_answers, 1):
        number = item.get("number") or len(score_answers) + index
        results_sheet.append([
            number, item.get("question"), item.get("text"), item.get("status"), recommended_action(item.get("status", "")),
            item.get("method"), item.get("logic"), item.get("statistics"), item.get("guardrail"),
        ])
    style_sheet(results_sheet)
    fit_columns(results_sheet, [6, 46, 62, 18, 22, 42, 58, 48, 52])

    combined_answers = answer_rows(score_answers, "Score") + answer_rows(sentiment_answers, "Sentiment", len(score_answers))
    write_rows("All Question Answers", combined_answers, ["#", "Area", "Question", "Answer", "Evidence Status", "Recommended Action", "Statistical Method", "Logic", "Statistics Used", "Interpretation Guardrail", "Outcome"], [7, 14, 48, 64, 20, 22, 42, 56, 50, 56, 14])
    write_rows("Score Calculations", answer_rows(score_answers, "Score"), ["#", "Area", "Question", "Answer", "Evidence Status", "Recommended Action", "Statistical Method", "Logic", "Statistics Used", "Interpretation Guardrail", "Outcome"], [7, 14, 48, 64, 20, 22, 42, 56, 50, 56, 14])
    write_rows("Sentiment Calculations", answer_rows(sentiment_answers, "Sentiment", len(score_answers)), ["#", "Area", "Question", "Answer", "Evidence Status", "Recommended Action", "Statistical Method", "Logic", "Statistics Used", "Interpretation Guardrail", "Outcome"], [7, 14, 48, 64, 20, 22, 42, 56, 50, 56, 14])
    write_rows("Score Evidence", evidence_rows(score_answers, "Score"), ["Question #", "Area", "Question", "Artifact Row", "Data Point", "Value"], [12, 14, 52, 14, 30, 48])
    write_rows("Sentiment Evidence", evidence_rows(sentiment_answers, "Sentiment", len(score_answers)), ["Question #", "Area", "Question", "Artifact Row", "Data Point", "Value"], [12, 14, 52, 14, 30, 48])

    sentiment_columns = [
        "Case ID", "Agent Name", "Manager/TL", "Feedback Date", "Verbatim Feedback", "CSAT Score", "CSAT Segment", "CSAT Type",
        "Sentiment", "Sentiment Score", "Sentiment Confidence", "AI Rationale", "Analysis Source",
        "ACPT Primary Category", "ACPT Secondary Category", "ACPT Confidence", "ACPT Evidence", "ACPT Needs Review",
        "Custom Category", "Category Confidence", "Owl Primary Driver", "Owl Secondary Driver", "Owl Issue Type",
        "Owl Resolution Status", "Owl People Sentiment", "Owl Process Sentiment", "Owl Tech Sentiment",
    ]
    write_frame("Verbatim Sentiment", analyzed_df, sentiment_columns)
    write_frame("ACPT Verbatim Classification", analyzed_df, [
        "Case ID", "Agent Name", "Manager/TL", "Feedback Date", "Verbatim Feedback",
        "Sentiment", "CSAT Score", "CSAT Segment",
        "ACPT Primary Category", "ACPT Secondary Category", "ACPT Confidence", "ACPT Evidence", "ACPT Needs Review",
    ])
    if "ACPT Primary Category" in analyzed_df.columns:
        acpt_summary = (
            analyzed_df.groupby("ACPT Primary Category", dropna=False)
            .agg(
                Responses=("ACPT Primary Category", "size"),
                Average_Confidence=("ACPT Confidence", "mean"),
                Needs_Review=("ACPT Needs Review", lambda values: int((values.astype(str).str.lower() == "yes").sum())),
            )
            .reset_index()
            .rename(columns={"ACPT Primary Category": "ACPT Category", "Average_Confidence": "Average Confidence", "Needs_Review": "Needs Review"})
            .sort_values("Responses", ascending=False)
        )
        write_frame("ACPT Summary", acpt_summary)
    write_frame("Analyzed Data Full", analyzed_df)
    write_frame("Weekly Trend", weekly_df)
    write_frame("Agent Scorecard", agent_df)
    write_frame("Manager Scorecard", manager_df)
    write_frame("Root Causes", reason_df)
    write_frame("Complaints", complaints_df)
    write_frame("Passive Insights", passive_df)

    def payload_rows(key: str) -> list[dict[str, Any]]:
        rows = analysis_payload.get(key)
        return rows if isinstance(rows, list) else []

    write_rows("Themes Summary", payload_rows("themes"))
    write_rows("Theme Detail Rows", payload_rows("themeRows"))
    write_rows("Feedback Table Rows", payload_rows("feedbackTableRows") or payload_rows("feedbackRows"))
    write_rows("Analysis Summary", payload_rows("analysisSummary"))
    write_rows("Sentiment Movement", payload_rows("sentimentMovement"))

    sentiment_engine_raw = str(analysis_engines.get("sentiment", rules_payload.get("sparrow", "")) or "").strip().lower()
    theme_engine_raw = str(analysis_engines.get("theme", rules_payload.get("theme", "")) or "").strip().lower()
    sentiment_is_sparrow = sentiment_engine_raw in {"sparrow", "model", "trained"}
    theme_is_owl = theme_engine_raw in {"owl", "theme", "model", "trained"}
    sparrow_model_path = model_paths.get("sparrow") or rules_payload.get("sparrowPath", "")
    owl_model_path = (
        model_paths.get("theme")
        or model_paths.get("owl")
        or model_paths.get("owl")
        or rules_payload.get("themePath", "")
        or rules_payload.get("owlPath", "")
        or rules_payload.get("owlPath", "")
    )
    sentiment_engine_label = "Sparrow Model" if sentiment_is_sparrow else "Local Rules"
    theme_engine_label = "Owl Model" if theme_is_owl else "Local Rules"
    dynamic_dimension_list = payload.get("dynamicDimensions") or analysis_payload.get("dynamicDimensions") or []
    if isinstance(dynamic_dimension_list, list):
        dynamic_dimension_value = ", ".join(str(item) for item in dynamic_dimension_list if item)
    else:
        dynamic_dimension_value = str(dynamic_dimension_list or "")
    output_rows_count = len(analyzed_df)

    setup_rows = [
        {"Analysis Setting": "Generated", "Value": time.strftime("%Y-%m-%d %H:%M:%S")},
        {"Analysis Setting": "Base File", "Value": base_payload.get("fileName") or files.get("base", "")},
        {"Analysis Setting": "Base File Size", "Value": base_payload.get("fileSize") or file_sizes.get("base", "")},
        {"Analysis Setting": "Rows Processed", "Value": base_payload.get("rows") or len(analyzed_df)},
        {"Analysis Setting": "Column Count", "Value": len(base_payload.get("columns") or []) or len(analyzed_df.columns)},
        {"Analysis Setting": "Processing Time", "Value": base_payload.get("processingTime", "")},
        {"Analysis Setting": f"{metric_label} Target", "Value": target},
        {"Analysis Setting": promoter_label, "Value": payload.get("promoterMin", satisfied_min)},
        {"Analysis Setting": passive_label, "Value": payload.get("passiveMin", neutral_min)},
        {"Analysis Setting": "Minimum Sample for Ranking", "Value": minimum_sample},
        {"Analysis Setting": "Sentiment Engine", "Value": sentiment_engine_label},
        {"Analysis Setting": "Theme Engine", "Value": theme_engine_label},
        {"Analysis Setting": "Sparrow Model Path", "Value": sparrow_model_path},
        {"Analysis Setting": "Owl Model Path", "Value": owl_model_path},
        {"Analysis Setting": "Question Framework", "Value": "100 leadership questions: 50 score questions plus 50 sentiment questions"},
    ]
    write_rows("Analysis Setup", setup_rows, ["Analysis Setting", "Value"], [34, 78])

    run_detail_rows = [
        {"Area": "Run", "Detail": "Generated At", "Value": time.strftime("%Y-%m-%d %H:%M:%S")},
        {"Area": "Run", "Detail": "Stream", "Value": metric_label},
        {"Area": "Input", "Detail": "Base File", "Value": base_payload.get("fileName") or files.get("base", "")},
        {"Area": "Input", "Detail": "Base File Size", "Value": base_payload.get("fileSize") or file_sizes.get("base", "")},
        {"Area": "Input", "Detail": "Lookup File", "Value": files.get("lookup", "")},
        {"Area": "Input", "Detail": "Lookup File Size", "Value": file_sizes.get("lookup", "")},
        {"Area": "Data", "Detail": "Rows Processed", "Value": base_payload.get("rows") or len(analyzed_df)},
        {"Area": "Data", "Detail": "Output Rows", "Value": output_rows_count},
        {"Area": "Data", "Detail": "Column Count", "Value": len(base_payload.get("columns") or []) or len(analyzed_df.columns)},
        {"Area": "Rules", "Detail": f"{metric_label} Target", "Value": target},
        {"Area": "Rules", "Detail": promoter_label, "Value": payload.get("promoterMin", satisfied_min)},
        {"Area": "Rules", "Detail": passive_label, "Value": payload.get("passiveMin", neutral_min)},
        {"Area": "Rules", "Detail": "Minimum Sample for Ranking", "Value": minimum_sample},
        {"Area": "Calendar", "Detail": "Week Start", "Value": rules_payload.get("weekStart", "") or payload.get("weekStart", "")},
        {"Area": "Calendar", "Detail": "Fiscal Start", "Value": rules_payload.get("fiscalStart", "") or payload.get("fiscalStart", "")},
        {"Area": "Setup", "Detail": "Dynamic Dimensions", "Value": dynamic_dimension_value},
        {"Area": "Intelligence", "Detail": "Sentiment Engine", "Value": sentiment_engine_label},
        {"Area": "Intelligence", "Detail": "Theme Engine", "Value": theme_engine_label},
        {"Area": "Intelligence", "Detail": "Sparrow Model Path", "Value": sparrow_model_path},
        {"Area": "Intelligence", "Detail": "Owl Model Path", "Value": owl_model_path},
        {"Area": "Output", "Detail": "Export Scope", "Value": "Full analyzed output with every available row and generated column."},
        {"Area": "Output", "Detail": "Question Framework", "Value": "100 leadership questions: 50 score questions plus 50 sentiment questions"},
    ]
    write_rows("Run Details", run_detail_rows, ["Area", "Detail", "Value"], [24, 34, 92])

    model_rows = [
        {"Section": "Sentiment Analysis", "Metric": "AI Model Used", "Value": "Sparrow Sentiment" if sentiment_is_sparrow else "Local Rules"},
        {"Section": "Sentiment Analysis", "Metric": "Model Path", "Value": sparrow_model_path},
        {"Section": "Sentiment Analysis", "Metric": "Model Type", "Value": "Local Fine-Tuned Sparrow Sentiment Model" if sentiment_is_sparrow else "Local rule-based engine"},
        {"Section": "Sentiment Analysis", "Metric": "Model Size", "Value": "475.5 MB" if sentiment_is_sparrow else "Not applicable"},
        {"Section": "Sentiment Analysis", "Metric": "Number of Output Classes", "Value": "3 (Positive, Neutral, Negative)"},
        {"Section": "Theme Classification", "Metric": "AI Model Used", "Value": "Owl Theme/ACPT/Resolution Model" if theme_is_owl else "Local Rules"},
        {"Section": "Theme Classification", "Metric": "Model Path", "Value": owl_model_path},
        {"Section": "Theme Classification", "Metric": "Model Type", "Value": "Local Fine-Tuned Owl Theme Model" if theme_is_owl else "Local rule-based engine"},
    ]
    write_rows("Model Information", model_rows, ["Section", "Metric", "Value"], [24, 32, 78])

    profile_rows = []
    merged_profile = dict(base_profile)
    for column, profile in (payload.get("base", {}).get("columnStats") or {}).items():
        if column not in merged_profile:
            merged_profile[column] = profile
    for source, profile in [("Base", merged_profile), ("Analyzed", analyzed_profile)]:
        for column, details in profile.items():
            row = {"Source": source, "Column": column}
            if isinstance(details, dict):
                row.update(details)
            else:
                row["Profile"] = details
            profile_rows.append(row)
    write_rows("Column Profiles", profile_rows)

    for sheet in workbook.worksheets:
        sheet.sheet_view.showGridLines = False

    stream = BytesIO()
    workbook.save(stream)
    return stream.getvalue()


class NPSHandler(BaseHTTPRequestHandler):
    server_version = "NPSHTML/0.1"

    def log_message(self, format: str, *args: Any) -> None:
        return

    def _serve_file(self, path: Path) -> None:
        if not path.exists() or path.is_dir():
            self.send_error(404)
            return
        content = path.read_bytes()
        mime = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        self.send_response(200)
        self.send_header("Content-Type", mime)
        if path.suffix.lower() in {".html", ".js", ".css"}:
            self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
            self.send_header("Pragma", "no-cache")
            self.send_header("Expires", "0")
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)

    def do_GET(self) -> None:
        parsed_path = urlparse(self.path)
        if parsed_path.path in {"", "/", "/index.html"}:
            self._serve_file(FRONTEND / "index.html")
            return
        if parsed_path.path == "/healthz":
            _json_response(self, {"ok": True, "status": "ready"})
            return
        if parsed_path.path == "/api/upload/progress":
            upload_id = str(parse_qs(parsed_path.query).get("id", [""])[0])
            with UPLOAD_PROGRESS_LOCK:
                progress_payload = dict(UPLOAD_PROGRESS.get(upload_id, {}))
            _json_response(self, progress_payload or {"ok": True, "uploadId": upload_id, "percent": 60, "stage": "Transfer", "message": "Waiting for the local server to begin workbook processing.", "complete": False})
            return
        if parsed_path.path == "/api/analysis/progress":
            with STATE_LOCK:
                _json_response(
                    self,
                    {
                        "ok": True,
                        "analysisId": STATE.analysis_id,
                        "progress": STATE.progress,
                        "status": STATE.status,
                        "running": STATE.analysis_running,
                        "error": STATE.analysis_error,
                        "analysisEngines": dict(STATE.analysis_engines),
                        "stage": STATE.analysis_stage,
                        "rowsProcessed": STATE.analysis_rows_processed,
                        "totalRows": STATE.analysis_total_rows,
                        "logFile": STATE.analysis_log_file,
                    },
                )
            return
        if parsed_path.path == "/api/auth/status":
            session = _session_from_handler(self)
            _json_response(self, {"ok": True, "authenticated": bool(session), "user": session or None})
            return
        if parsed_path.path == "/api/auth/users":
            allowed, session = _require_admin(self)
            if not allowed:
                return
            store = _load_user_store()
            _json_response(self, {"ok": True, "users": [_public_user(user) for user in store.get("users", [])]})
            return
        if parsed_path.path == "/api/audit/logs":
            allowed, session = _require_admin(self)
            if not allowed:
                return
            query = parse_qs(parsed_path.query)
            limit = max(25, min(2000, int(query.get("limit", ["300"])[0] or 300)))
            _audit_from_handler(self, session.get("username", "admin"), "AUDIT_LOG_VIEW", "Admin reviewed audit logs.", {"limit": limit})
            _json_response(self, {"ok": True, "entries": _audit_log_entries(limit), "auditFolder": str(AUDIT_DIR), "userLogFolder": str(USER_LOG_DIR)})
            return
        if parsed_path.path == "/api/status":
            _json_response(
                self,
                {
                    "ok": True,
                    "files": STATE.files,
                    "model_status": _model_status(),
                    "base_columns": list(STATE.base_df.columns),
                    "lookup_columns": list(STATE.lookup_df.columns),
                    "base_column_stats": STATE.base_column_profile or _column_profile(STATE.base_df),
                    "lookup_column_stats": STATE.lookup_column_profile or _column_profile(STATE.lookup_df),
                    "dynamic_dimensions": list(STATE.dynamic_dimensions),
                    "guesses": _guess_columns(list(STATE.base_df.columns)),
                    "analysis": _analysis_payload(),
                    "sparrow_training": _training_snapshot(),
                },
            )
            return
        if parsed_path.path == "/api/statistics-data":
            with STATE_LOCK:
                base_df = STATE.base_df.copy()
                lookup_df = STATE.lookup_df.copy()
                analyzed_df = STATE.analyzed_df.copy()
            _json_response(
                self,
                {
                    "ok": True,
                    "baseRows": _stats_records(base_df, 200),
                    "lookupRows": _stats_records(lookup_df, 200),
                    "analyzedRows": _stats_records(analyzed_df, 200),
                    "rowCounts": {
                        "baseRows": int(len(base_df)),
                        "lookupRows": int(len(lookup_df)),
                        "analyzedRows": int(len(analyzed_df)),
                    },
                    "sampleLimit": 200,
                    "columnProfiles": {
                        "base": STATE.base_column_profile or _column_profile(base_df),
                        "lookup": STATE.lookup_column_profile or _column_profile(lookup_df),
                        "analyzed": STATE.analyzed_column_profile or _column_profile(analyzed_df),
                    },
                    "baseColumns": list(base_df.columns),
                    "lookupColumns": list(lookup_df.columns),
                    "analyzedColumns": list(analyzed_df.columns),
                },
            )
            return
        if parsed_path.path == "/api/statistics-rows":
            _json_response(self, _paginated_statistics_rows(parse_qs(parsed_path.query)))
            return
        if parsed_path.path == "/api/module/progress":
            analysis_id = parse_qs(parsed_path.query).get("id", [""])[0]
            with MODULE_PROGRESS_LOCK:
                payload = dict(MODULE_PROGRESS.get(analysis_id, {}))
            if not payload:
                payload = {
                    "ok": True,
                    "id": analysis_id,
                    "percent": 0,
                    "message": "Waiting for analysis to start...",
                    "done": 0,
                    "total": 0,
                    "currentRow": 0,
                    "complete": False,
                }
            _json_response(self, payload)
            return
        if parsed_path.path == "/api/sparrow-training/status":
            _json_response(self, _training_snapshot())
            return
        if parsed_path.path == "/api/model-defaults":
            _json_response(
                self,
                {
                    "ok": True,
                    "sparrow": _validate_model_path("sparrow"),
                    "theme": _validate_model_path("theme"),
                },
            )
            return
        if parsed_path.path == "/api/ollama/models":
            models = _ollama_models()
            _json_response(
                self,
                {
                    "ok": True,
                    "models": models,
                    "qwen7b": _resolve_ollama_model("qwen2.5:7b-instruct"),
                    "ollamaReachable": bool(models),
                },
            )
            return
        if parsed_path.path == "/api/export/excel":
            session = _session_from_handler(self)
            _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "EXPORT_EXCEL", "User exported Excel workbook.")
            self._export_excel()
            return
        if parsed_path.path == "/api/export/raw-csv":
            session = _session_from_handler(self)
            _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "EXPORT_RAW_CSV", "User exported complete analyzed raw data as CSV.")
            self._export_raw_csv()
            return
        self._serve_static()

    def do_POST(self) -> None:
        try:
            if _handle_theme_training_api(self):
                return
            if self.path == "/api/auth/login":
                payload = _read_json(self)
                username = str(payload.get("username") or "").strip()
                password = str(payload.get("password") or "")
                store = _load_user_store()
                user = next((item for item in store.get("users", []) if str(item.get("username", "")).lower() == username.lower()), None)
                if not user or not bool(user.get("active", True)) or not _verify_password(password, user):
                    _audit_from_handler(self, username or "unknown", "LOGIN_FAILED", "Login failed.", {"username": username}, "WARN")
                    _json_response(self, {"ok": False, "error": "Invalid username or password."}, 401)
                    return
                token = secrets.token_urlsafe(32)
                user["lastLoginAt"] = _audit_timestamp()
                user["updatedAt"] = _audit_timestamp()
                _save_user_store(store)
                public = _public_user(user)
                with SESSION_LOCK:
                    SESSIONS[token] = {**public, "token": token, "loginAt": _audit_timestamp()}
                _audit_from_handler(self, public.get("username", ""), "LOGIN_SUCCESS", "User logged in.", {"role": public.get("role")})
                _json_response(self, {"ok": True, "token": token, "user": public})
                return
            if self.path == "/api/auth/logout":
                session = _session_from_handler(self)
                token = self.headers.get("X-Session-Token", "")
                if token:
                    with SESSION_LOCK:
                        SESSIONS.pop(token, None)
                if session:
                    _audit_from_handler(self, session.get("username", ""), "LOGOUT", "User logged out.")
                _json_response(self, {"ok": True})
                return
            if self.path == "/api/auth/users":
                allowed, session = _require_admin(self)
                if not allowed:
                    return
                payload = _read_json(self)
                action = str(payload.get("action") or "").strip().lower()
                store = _load_user_store()
                users = store.setdefault("users", [])
                username = str(payload.get("username") or "").strip()
                actor = session.get("username", "admin") if session else "admin"
                if action == "create":
                    password = str(payload.get("password") or "").strip()
                    if not username or not password:
                        _json_response(self, {"ok": False, "error": "Username and password are required."}, 400)
                        return
                    if any(str(item.get("username", "")).lower() == username.lower() for item in users):
                        _json_response(self, {"ok": False, "error": "User already exists."}, 409)
                        return
                    salt = secrets.token_hex(16)
                    users.append({
                        "username": username,
                        "displayName": str(payload.get("displayName") or username).strip(),
                        "role": "admin" if str(payload.get("role") or "user").lower() == "admin" else "user",
                        "active": bool(payload.get("active", True)),
                        "salt": salt,
                        "passwordHash": _hash_password(password, salt),
                        "iterations": PASSWORD_ITERATIONS,
                        "createdAt": _audit_timestamp(),
                        "updatedAt": _audit_timestamp(),
                        "mustChangePassword": bool(payload.get("mustChangePassword", False)),
                    })
                    _audit_from_handler(self, actor, "USER_CREATED", f"Admin created user {username}.", {"username": username})
                elif action in {"update", "reset", "delete"}:
                    user = next((item for item in users if str(item.get("username", "")).lower() == username.lower()), None)
                    if not user:
                        _json_response(self, {"ok": False, "error": "User not found."}, 404)
                        return
                    if action == "delete":
                        if username.lower() == str(actor).lower():
                            _json_response(self, {"ok": False, "error": "You cannot delete the signed-in admin user."}, 400)
                            return
                        store["users"] = [item for item in users if str(item.get("username", "")).lower() != username.lower()]
                        _audit_from_handler(self, actor, "USER_DELETED", f"Admin deleted user {username}.", {"username": username})
                    elif action == "reset":
                        password = str(payload.get("password") or "").strip()
                        if not password:
                            _json_response(self, {"ok": False, "error": "New password is required."}, 400)
                            return
                        salt = secrets.token_hex(16)
                        user["salt"] = salt
                        user["passwordHash"] = _hash_password(password, salt)
                        user["iterations"] = PASSWORD_ITERATIONS
                        user["mustChangePassword"] = bool(payload.get("mustChangePassword", True))
                        user["updatedAt"] = _audit_timestamp()
                        _audit_from_handler(self, actor, "PASSWORD_RESET", f"Admin reset password for {username}.", {"username": username})
                    else:
                        if "displayName" in payload:
                            user["displayName"] = str(payload.get("displayName") or username).strip()
                        if "role" in payload:
                            user["role"] = "admin" if str(payload.get("role") or "user").lower() == "admin" else "user"
                        if "active" in payload:
                            user["active"] = bool(payload.get("active"))
                        user["updatedAt"] = _audit_timestamp()
                        _audit_from_handler(self, actor, "USER_UPDATED", f"Admin updated user {username}.", {"username": username})
                else:
                    _json_response(self, {"ok": False, "error": "Unknown user action."}, 400)
                    return
                _save_user_store(store)
                _json_response(self, {"ok": True, "users": [_public_user(user) for user in _load_user_store().get("users", [])]})
                return
            if self.path == "/api/audit/event":
                payload = _read_json(self)
                session = _session_from_handler(self)
                username = session.get("username", "unknown") if session else str(payload.get("username") or "unknown")
                event = str(payload.get("event") or "USER_EVENT")[:80]
                action = str(payload.get("action") or "User activity")[:240]
                details = payload.get("details") if isinstance(payload.get("details"), dict) else {}
                _audit_from_handler(self, username, event, action, details)
                _json_response(self, {"ok": True})
                return
            if self.path == "/api/leadership-statistics":
                payload = _read_json(self)
                if isinstance(payload.get("calendar"), dict):
                    with STATE_LOCK:
                        STATE.calendar_settings = _calendar_settings(payload)
                target = _safe_float(payload.get("target"), 85.0)
                minimum_sample = int(_safe_float(payload.get("minimumSample"), LEADERSHIP_MIN_SAMPLE))
                satisfied_min = _safe_float(payload.get("satisfiedMin"), 4.0)
                neutral_min = _safe_float(payload.get("neutralMin"), 3.0)
                if str(payload.get("mode") or "").strip().lower() == "nps":
                    promoter_min = _safe_float(payload.get("promoterMin"), satisfied_min)
                    passive_min = _safe_float(payload.get("passiveMin"), neutral_min)
                    _json_response(self, _nps_leadership_statistics(target, minimum_sample, promoter_min, passive_min))
                else:
                    _json_response(self, _leadership_statistics(target, minimum_sample, satisfied_min, neutral_min))
                return
            if self.path == "/api/export/leadership-results":
                payload = _read_json(self)
                content = _leadership_results_workbook(payload)
                filename = "NPS_Leadership_Results.xlsx" if str(payload.get("mode") or "").lower() == "nps" else "CSAT_Leadership_Results.xlsx"
                self.send_response(200)
                self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                self.send_header("Content-Disposition", f"attachment; filename={filename}")
                self.send_header("Content-Length", str(len(content)))
                self.end_headers()
                self.wfile.write(content)
                return
            if self.path == "/api/analysis-file-preview":
                result = _analysis_file_preview(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 400)
                return
            if self.path == "/api/export-analysis-project":
                payload = _read_json(self)
                with STATE_LOCK:
                    analyzed_df = STATE.analyzed_df.copy()
                    base_columns = list(STATE.base_df.columns) if not STATE.base_df.empty else []
                    lookup_columns = list(STATE.lookup_df.columns) if not STATE.lookup_df.empty else []
                    files = dict(STATE.files)
                    file_sizes = dict(STATE.file_sizes)
                    base_profile = dict(STATE.base_column_profile)
                    lookup_profile = dict(STATE.lookup_column_profile)
                    analyzed_profile = dict(STATE.analyzed_column_profile or _column_profile(STATE.analyzed_df))
                    dynamic_dimensions = list(STATE.dynamic_dimensions)
                    date_filter = dict(STATE.date_filter)
                    calendar_settings = dict(STATE.calendar_settings)
                analysis_payload = _analysis_payload()
                project = {
                    "schema": "cx-suite-analysis-project",
                    "schemaVersion": 1,
                    "exportedAt": _dt.datetime.utcnow().isoformat(timespec="seconds") + "Z",
                    "appMode": str(payload.get("appMode") or payload.get("mode") or ""),
                    "preferredName": str(payload.get("preferredName") or ""),
                    "files": files,
                    "fileSizes": file_sizes,
                    "baseColumns": base_columns,
                    "lookupColumns": lookup_columns,
                    "baseColumnStats": base_profile,
                    "lookupColumnStats": lookup_profile,
                    "analyzedColumnStats": analyzed_profile,
                    "dynamicDimensions": dynamic_dimensions,
                    "dateFilter": date_filter,
                    "calendar": calendar_settings,
                    "setup": payload.get("setup") if isinstance(payload.get("setup"), dict) else {},
                    "analysis": analysis_payload,
                    "fullAnalyzedRows": _safe_records(analyzed_df, max(len(analyzed_df), 1)),
                }
                _json_response(self, {"ok": True, "project": project})
                return
            if self.path == "/api/import-analysis-project":
                payload = _read_json(self)
                project = payload.get("project") if isinstance(payload.get("project"), dict) else payload
                analysis = project.get("analysis") if isinstance(project.get("analysis"), dict) else {}
                rows = project.get("fullAnalyzedRows") or analysis.get("fullAnalyzedRows") or analysis.get("feedbackRows") or analysis.get("feedbackTableRows") or analysis.get("preview") or []
                if not isinstance(rows, list) or not rows:
                    _json_response(self, {"ok": False, "error": "The selected analysis file does not contain restorable analysis rows."}, 400)
                    return
                restored_df = pd.DataFrame(rows)
                calendar_settings = _calendar_settings({"calendar": project.get("calendar") if isinstance(project.get("calendar"), dict) else analysis.get("calendar") if isinstance(analysis.get("calendar"), dict) else {}})
                summaries = _summaries_for_calendar(restored_df, calendar_settings)
                with STATE_LOCK:
                    STATE.analyzed_df = restored_df
                    STATE.weekly_df = summaries.get("weekly", pd.DataFrame())
                    STATE.agent_df = _agent_summary_with_manager(restored_df, summaries.get("agent", pd.DataFrame()))
                    STATE.manager_df = _manager_summary(restored_df)
                    STATE.reason_df = summaries.get("reasons", pd.DataFrame())
                    STATE.analyzed_column_profile = _column_profile(restored_df)
                    STATE.base_column_profile = project.get("baseColumnStats") if isinstance(project.get("baseColumnStats"), dict) else {}
                    STATE.lookup_column_profile = project.get("lookupColumnStats") if isinstance(project.get("lookupColumnStats"), dict) else {}
                    STATE.files = project.get("files") if isinstance(project.get("files"), dict) else {"analysis": project.get("preferredName") or "Imported analysis"}
                    STATE.file_sizes = project.get("fileSizes") if isinstance(project.get("fileSizes"), dict) else {}
                    STATE.dynamic_dimensions = project.get("dynamicDimensions") if isinstance(project.get("dynamicDimensions"), list) else []
                    STATE.last_run_config = {
                        "mapping": dict(analysis.get("mapping", {})) if isinstance(analysis.get("mapping"), dict) else {},
                        "businessRules": dict(analysis.get("businessRules", {})) if isinstance(analysis.get("businessRules"), dict) else {},
                    }
                    STATE.date_filter = project.get("dateFilter") if isinstance(project.get("dateFilter"), dict) else analysis.get("dateFilter") if isinstance(analysis.get("dateFilter"), dict) else {"mode": "All Time", "start": "", "end": ""}
                    STATE.calendar_settings = calendar_settings
                    STATE.status = "Imported analysis ready"
                    STATE.progress = 100.0
                    STATE.analysis_running = False
                    STATE.analysis_error = ""
                _json_response(self, {"ok": True, "analysis": _analysis_payload(), "rows": len(restored_df)})
                return
            if self.path == "/api/signout":
                _json_response(self, {"ok": True, "message": "Server stopped. You can safely close this window."})
                def _stop_server() -> None:
                    global KESTRELIQ_PROCESS, AUDIT_PROCESS
                    time.sleep(0.5)
                    with STATE_LOCK:
                        STATE.analysis_running = False
                        STATE.analysis_error = ""
                        STATE.status = "Signed out"
                    stopped_server_processes = _terminate_other_cx_suite_servers()
                    stopped_browser_processes = _terminate_cx_suite_browser_processes(int(getattr(self.server, "server_address", ["", 0])[1] or 0))
                    _write_audit_record(
                        "SYSTEM",
                        "SIGN_OUT_STOP",
                        "User signed out; stopping all local CX Suite servers, matching local browser windows, and related run processes.",
                        {
                            "launcherProcessId": SERVER_LAUNCHER_PROCESS_ID,
                            "stoppedServerProcessIds": stopped_server_processes,
                            "stoppedBrowserProcessIds": stopped_browser_processes,
                        },
                    )
                    _terminate_popen(KESTRELIQ_PROCESS)
                    _terminate_popen(AUDIT_PROCESS)
                    KESTRELIQ_PROCESS = None
                    AUDIT_PROCESS = None
                    self.server.shutdown()
                    time.sleep(0.5)
                    _terminate_process_tree(SERVER_LAUNCHER_PROCESS_ID)
                threading.Thread(target=_stop_server, daemon=True, name="signout-shutdown").start()
                return
            if self.path == "/api/upload/sheet-preview":
                result = _excel_sheet_preview(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 400)
                return
            if self.path == "/api/uploaded-file-preview":
                result = _uploaded_file_preview(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 400)
                return
            if self.path == "/api/upload":
                payload = _read_json(self)
                kind = payload.get("kind")
                upload_id = str(payload.get("uploadId") or "")
                _set_upload_progress(upload_id, 62, "Received", "The local server received the full file. Preparing workbook decoding.")
                df, sheet_names, selected_sheet = _decode_excel(payload, upload_id)
                if sheet_names and not selected_sheet:
                    _json_response(self, {"ok": True, "kind": kind, "needsSheetSelection": True, "sheetNames": sheet_names, "sheetCount": len(sheet_names)})
                    return
                _set_upload_progress(upload_id, 88, "Profile", f"Scanning {len(df):,} rows from sheet '{selected_sheet}' to identify column types, blanks, and unique values for the Explorer.")
                column_profile = _column_profile(df)
                if kind == "base":
                    STATE.base_df = df
                    STATE.base_column_profile = column_profile
                    STATE.analyzed_column_profile = {}
                    STATE.analyzed_df = pd.DataFrame()
                    STATE.weekly_df = pd.DataFrame()
                    STATE.agent_df = pd.DataFrame()
                    STATE.manager_df = pd.DataFrame()
                    STATE.reason_df = pd.DataFrame()
                    STATE.complaints_df = pd.DataFrame()
                    STATE.passive_df = pd.DataFrame()
                    STATE.status = "Ready"
                    STATE.progress = 0.0
                    STATE.analysis_running = False
                    STATE.analysis_error = ""
                    STATE.analysis_id = ""
                    STATE.date_filter = {"mode": "All Time", "start": "", "end": ""}
                    STATE.dynamic_dimensions = []
                elif kind == "lookup":
                    STATE.lookup_df = df
                    STATE.lookup_column_profile = column_profile
                else:
                    _json_response(self, {"ok": False, "error": "Unknown upload kind."}, 400)
                    return
                _set_upload_progress(upload_id, 96, "Finalize", "Column profiling is complete. Saving mappings and preparing the 200-row browser preview.")
                STATE.files[str(kind)] = payload.get("name", "Uploaded file")
                STATE.files[f"{kind}_sheet"] = selected_sheet
                data_text = str(payload.get("data") or "")
                encoded_text = data_text.partition(",")[2] if "," in data_text else data_text
                if encoded_text:
                    padding = len(encoded_text) - len(encoded_text.rstrip("="))
                    STATE.file_sizes[str(kind)] = max(0, (len(encoded_text) * 3 // 4) - padding)
                session = _session_from_handler(self)
                _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "FILE_UPLOAD", f"Uploaded {kind} file.", {"kind": kind, "name": payload.get("name", "Uploaded file"), "rows": len(df), "columns": len(df.columns)})
                _json_response(
                    self,
                    {
                        "ok": True,
                        "kind": kind,
                        "rows": len(df),
                        "columns": list(df.columns),
                        "sheetNames": sheet_names,
                        "sheetCount": len(sheet_names) or 1,
                        "selectedSheet": selected_sheet,
                        "columnStats": column_profile,
                        "guesses": _guess_columns(list(STATE.base_df.columns)),
                    },
                )
                _set_upload_progress(upload_id, 100, "Complete", f"Upload complete. {len(df):,} rows are ready.", True)
                return
            if self.path == "/api/analyze":
                self._analyze(_read_json(self))
                return
            if self.path == "/api/custom-statistics":
                _json_response(self, _custom_statistics(_read_json(self)))
                return
            if self.path == "/api/weekly-trend-dashboard":
                result = _weekly_trend_dashboard(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 400)
                return
            if self.path == "/api/weekly-trend-matrix":
                result = _weekly_trend_matrix(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 400)
                return
            if self.path == "/api/module/inspect":
                payload = _read_json(self)
                df = _decode_tabular(payload)
                mode = str(payload.get("mode") or "csat").strip().lower()
                _json_response(
                    self,
                    {
                        "ok": True,
                        "rows": len(df),
                        "columns": list(df.columns),
                        "guesses": _generic_guess_columns(list(df.columns), mode),
                        "columnStats": _column_profile(df),
                    },
                )
                return
            if self.path == "/api/module/analyze":
                _json_response(self, _generic_analysis_payload(_read_json(self)))
                return
            if self.path == "/api/date-filter":
                payload = _read_json(self)
                with STATE_LOCK:
                    STATE.date_filter = {
                        "mode": str(payload.get("mode") or "All Time"),
                        "start": str(payload.get("start") or ""),
                        "end": str(payload.get("end") or ""),
                    }
                _json_response(self, {"ok": True, "analysis": _analysis_payload()})
                return
            if self.path == "/api/compare-ranges":
                _json_response(self, _comparison_payload(_read_json(self)))
                return
            if self.path == "/api/feedback-override":
                override_payload = _read_json(self)
                result = _apply_feedback_override(override_payload)
                session = _session_from_handler(self)
                _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "FEEDBACK_OVERRIDE", "User applied manual feedback override.", {"rowId": override_payload.get("rowId")})
                _json_response(self, result, 200 if result.get("ok") else 400)
                return
            if self.path == "/api/model/validate":
                payload = _read_json(self)
                kind = str(payload.get("kind") or "").strip().lower()
                if kind not in {"sparrow", "theme", "owl"}:
                    _json_response(self, {"ok": False, "error": "Unknown model kind."}, 400)
                    return
                result = _validate_model_path(kind, str(payload.get("path") or ""))
                _json_response(self, result, 200 if result.get("ok") else 400)
                return
            if self.path == "/api/model/browse":
                payload = _read_json(self)
                result = _choose_folder(str(payload.get("initialPath") or ""))
                _json_response(self, result, 200 if result.get("ok") else 500)
                return
            if self.path == "/api/model/list":
                payload = _read_json(self)
                kind = str(payload.get("kind") or "").strip().lower()
                if kind not in {"theme", "owl", ""}:
                    _json_response(self, {"ok": False, "error": "Unknown model kind."}, 400)
                    return
                result = _list_owl_models()
                _json_response(self, result, 200 if result.get("ok") else 500)
                return
            if self.path == "/api/train/open":
                payload = _read_json(self)
                kind = str(payload.get("kind") or "").strip().lower()
                result = _launch_training_tool(kind)
                _json_response(self, result, 200 if result.get("ok") else 404)
                return
            if self.path == "/api/custom-category/preview":
                result = _ollama_category_preview(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 500)
                return
            if self.path == "/api/theme-builder/discover":
                result = _theme_builder_discover(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 500)
                return
            if self.path == "/api/theme-builder/classify":
                result = _theme_builder_classify(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 500)
                return
            if self.path == "/api/executive-lens/benchmark":
                result = _executive_lens_benchmark(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 500)
                return
            if self.path == "/api/sparrow-training/start":
                result = _start_sparrow_training(_read_json(self))
                _json_response(self, result, 200 if result.get("ok") else 409)
                return
            if self.path == "/api/sparrow-training/status":
                _json_response(self, _training_snapshot())
                return
            if self.path == "/api/export/selected":
                export_payload = _read_json(self)
                session = _session_from_handler(self)
                _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "EXPORT_SELECTED", "User exported selected data.", {"selected": export_payload.get("selected")})
                self._export_selected(export_payload)
                return
            if self.path == "/api/export/report-workbook":
                report_payload = _read_json(self)
                session = _session_from_handler(self)
                _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "EXPORT_REPORT_WORKBOOK", "User exported board-room workbook.")
                self._export_report_workbook(report_payload)
                return
            if self.path == "/api/export/report-pptx":
                ppt_payload = _read_json(self)
                session = _session_from_handler(self)
                _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "EXPORT_REPORT_PPTX", "User exported board-room PPT/PDF pack.")
                self._export_report_pptx(ppt_payload)
                return
            if self.path == "/api/export/boardroom-pdf":
                pdf_payload = _read_json(self)
                session = _session_from_handler(self)
                _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "EXPORT_BOARDROOM_PDF", "User exported a native board-room PDF pack.")
                self._export_boardroom_pdf(pdf_payload)
                return
            _json_response(self, {"ok": False, "error": "Unknown endpoint."}, 404)
        except Exception as exc:
            traceback.print_exc()
            STATE.status = f"Failed: {exc}"
            _json_response(self, {"ok": False, "error": str(exc)}, 500)

    def _serve_static(self) -> None:
        requested = unquote(self.path.split("?", 1)[0].lstrip("/"))
        if not requested:
            requested = "index.html"
        if requested == "apps" or requested.startswith("apps/"):
            parts = requested.split("/")
            if len(parts) < 2 or not parts[1]:
                self.send_error(404)
                return
            candidate_app = (APPS / parts[1]).resolve()
            frontend_root = (candidate_app / "frontend").resolve()
            app_root = frontend_root if frontend_root.exists() else candidate_app
            relative = "/".join(parts[2:]) or "index.html"
            path = (app_root / relative).resolve()
            allowed_root = app_root
        else:
            path = (FRONTEND / requested).resolve()
            allowed_root = FRONTEND
        if allowed_root not in path.parents and path != allowed_root:
            self.send_error(403)
            return
        if path.is_dir():
            path = path / "index.html"
        if not path.exists() or path.is_dir():
            self.send_error(404)
            return
        content = path.read_bytes()
        mime = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        self.send_response(200)
        self.send_header("Content-Type", mime)
        if path.suffix.lower() in {".html", ".js", ".css"}:
            self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
            self.send_header("Pragma", "no-cache")
            self.send_header("Expires", "0")
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)

    def _analyze(self, payload: dict[str, Any]) -> None:
        if STATE.base_df.empty:
            _json_response(self, {"ok": False, "error": "Upload a base file first."}, 400)
            return
        mapping = payload.get("mapping", {})
        feedback = mapping.get("feedback")
        score = mapping.get("score")
        agent = mapping.get("agent")
        date = mapping.get("date")
        if not feedback or feedback not in STATE.base_df.columns:
            _json_response(self, {"ok": False, "error": "Select a valid verbatim/comment column."}, 400)
            return
        with STATE_LOCK:
            if STATE.analysis_running:
                _json_response(self, {"ok": False, "error": "Analysis is already running."}, 409)
                return
            analysis_id = uuid.uuid4().hex
            STATE.analysis_running = True
            STATE.analysis_error = ""
            STATE.analysis_id = analysis_id
            STATE.analysis_log_file = ""
            STATE.analysis_stage = "Starting analysis"
            STATE.analysis_rows_processed = 0
            STATE.analysis_total_rows = int(len(STATE.base_df))
            STATE.analysis_started_at = time.time()
            STATE.analysis_completed_at = 0.0
            bands = payload.get("csatBands") if isinstance(payload.get("csatBands"), dict) else {}
            calendar = payload.get("calendar") if isinstance(payload.get("calendar"), dict) else {}
            STATE.last_run_config = {
                "mapping": dict(payload.get("mapping", {})),
                "businessRules": {
                    "target": payload.get("target"),
                    "scoreScale": bands.get("scale") or payload.get("scoreScale"),
                    "satisfiedStartsAt": bands.get("satisfiedMin"),
                    "neutralStartsAt": bands.get("neutralMin"),
                    "promoterStartsAt": payload.get("promoterStartsAt"),
                    "passiveStartsAt": payload.get("passiveStartsAt"),
                    "weekStart": calendar.get("weekStartDay") or calendar.get("weekStart"),
                    "fiscalStart": calendar.get("fiscalStartMonth") or calendar.get("fiscalYearStartMonth"),
                },
            }
            STATE.analyzed_df = pd.DataFrame()
            STATE.weekly_df = pd.DataFrame()
            STATE.agent_df = pd.DataFrame()
            STATE.manager_df = pd.DataFrame()
            STATE.reason_df = pd.DataFrame()
            STATE.complaints_df = pd.DataFrame()
            STATE.passive_df = pd.DataFrame()
            STATE.analyzed_column_profile = {}
            STATE.progress = 1
            STATE.status = "Analysis started"
        try:
            log_path = _start_analysis_log(payload, analysis_id)
            with STATE_LOCK:
                if STATE.analysis_id == analysis_id:
                    STATE.analysis_log_file = str(log_path)
        except Exception as exc:
            traceback.print_exc()
            _json_response(self, {"ok": False, "error": f"Could not create the required analysis log: {exc}"}, 500)
            with STATE_LOCK:
                STATE.analysis_running = False
                STATE.analysis_error = str(exc)
            return
        thread = threading.Thread(target=_run_analysis_job, args=(payload, analysis_id), daemon=True)
        thread.start()
        session = _session_from_handler(self)
        _audit_from_handler(self, session.get("username", "unknown") if session else "unknown", "ANALYSIS_STARTED", "Feedback analysis started.", {"analysisId": analysis_id, "rows": int(len(STATE.base_df)), "sentimentEngine": payload.get("sentimentEngine"), "themeEngine": payload.get("themeEngine")})
        _json_response(self, {"ok": True, "started": True, "analysisId": analysis_id, "logFile": str(log_path), "analysis": _analysis_payload()})

    def _export_excel(self) -> None:
        try:
            if STATE.analyzed_df.empty:
                _json_response(self, {"ok": False, "error": "Run analysis first."}, 400)
                return
            content = export_workbook(
                STATE.analyzed_df,
                STATE.weekly_df,
                STATE.agent_df,
                STATE.complaints_df,
                STATE.reason_df,
                STATE.passive_df,
            )
        except Exception as exc:
            traceback.print_exc()
            _json_response(self, {"ok": False, "error": f"Excel export failed: {exc}"}, 500)
            return
        self.send_response(200)
        self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        self.send_header("Content-Disposition", "attachment; filename=NPSHTML_Analysis_Output.xlsx")
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)

    def _export_raw_csv(self) -> None:
        try:
            if STATE.analyzed_df.empty:
                _json_response(self, {"ok": False, "error": "Run analysis first."}, 400)
                return
            content = STATE.analyzed_df.to_csv(index=False).encode("utf-8-sig")
        except Exception as exc:
            traceback.print_exc()
            _json_response(self, {"ok": False, "error": f"CSV export failed: {exc}"}, 500)
            return
        self.send_response(200)
        self.send_header("Content-Type", "text/csv; charset=utf-8")
        self.send_header("Content-Disposition", "attachment; filename=Analysis_Raw_Data.csv")
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)

    def _export_selected(self, payload: dict[str, Any]) -> None:
        try:
            if STATE.analyzed_df.empty:
                _json_response(self, {"ok": False, "error": "Run analysis first."}, 400)
                return
            selected = [str(item) for item in payload.get("sheets", []) if str(item).strip()]
            if not selected:
                _json_response(self, {"ok": False, "error": "Select at least one output."}, 400)
                return
            sheet_map: dict[str, tuple[str, pd.DataFrame]] = {
                "analyzed": ("Analyzed Feedback", STATE.analyzed_df),
                "weekly": ("Weekly Trend", STATE.weekly_df),
                "agents": ("Agent Scorecard", STATE.agent_df),
                "managers": ("Manager Scorecard", STATE.manager_df),
                "complaints": ("Agent Complaints", STATE.complaints_df),
                "reasons": ("Root Causes", STATE.reason_df),
                "passives": ("Passive Insights", STATE.passive_df),
                "quartiles": ("Quartiles", _quartile_summary(STATE.agent_df)),
            }
            output = BytesIO()
            with pd.ExcelWriter(output, engine="openpyxl") as writer:
                for key in selected:
                    if key not in sheet_map:
                        continue
                    sheet_name, frame = sheet_map[key]
                    frame.to_excel(writer, index=False, sheet_name=sheet_name[:31])
            content = output.getvalue()
        except Exception as exc:
            traceback.print_exc()
            _json_response(self, {"ok": False, "error": f"Selected export failed: {exc}"}, 500)
            return
        self.send_response(200)
        self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        self.send_header("Content-Disposition", "attachment; filename=NPSHTML_Selected_Outputs.xlsx")
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)

    def _export_report_workbook(self, payload: dict[str, Any]) -> None:
        try:
            analysis_payload = payload.get("analysis") if isinstance(payload.get("analysis"), dict) else {}
            if STATE.analyzed_df.empty and not analysis_payload:
                _json_response(self, {"ok": False, "error": "Run analysis first."}, 400)
                return
            selected = [str(item) for item in payload.get("tabs", []) if str(item).strip()]
            if not selected:
                _json_response(self, {"ok": False, "error": "Select at least one tab."}, 400)
                return
            content = _write_report_workbook(selected, analysis_payload)
        except Exception as exc:
            traceback.print_exc()
            _json_response(self, {"ok": False, "error": f"Report workbook failed: {exc}"}, 500)
            return
        self.send_response(200)
        self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        self.send_header("Content-Disposition", "attachment; filename=NPS_Selected_Report_Data.xlsx")
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)

    def _export_report_pptx(self, payload: dict[str, Any]) -> None:
        try:
            analysis_payload = payload.get("analysis") if isinstance(payload.get("analysis"), dict) else {}
            if STATE.analyzed_df.empty and not analysis_payload:
                _json_response(self, {"ok": False, "error": "Run analysis first."}, 400)
                return
            selected = [str(item) for item in payload.get("tabs", []) if str(item).strip()]
            if not selected:
                _json_response(self, {"ok": False, "error": "Select at least one tab."}, 400)
                return
            content = _build_report_pptx(selected, analysis_payload)
        except Exception as exc:
            traceback.print_exc()
            _json_response(self, {"ok": False, "error": f"PPT export failed: {exc}"}, 500)
            return
        self.send_response(200)
        self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        self.send_header("Content-Disposition", "attachment; filename=NPS_Analyzer_Selected_Report.pptx")
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)

    def _export_boardroom_pdf(self, payload: dict[str, Any]) -> None:
        try:
            analysis_payload = payload.get("analysis") if isinstance(payload.get("analysis"), dict) else {}
            if STATE.analyzed_df.empty and not analysis_payload:
                _json_response(self, {"ok": False, "error": "Run analysis first."}, 400)
                return
            selected = [str(item) for item in payload.get("tabs", []) if str(item).strip()]
            custom_dashboards = payload.get("customDashboards") if isinstance(payload.get("customDashboards"), list) else []
            if not selected and not custom_dashboards:
                _json_response(self, {"ok": False, "error": "Select at least one report section or add a custom dashboard."}, 400)
                return
            content = _build_boardroom_pdf(payload)
            metric = re.sub(r"[^A-Za-z0-9_-]+", "_", str(payload.get("metric") or "NPS").upper()).strip("_") or "NPS"
            filename = f"{metric}_Board_Room_Leadership_Report.pdf"
        except Exception as exc:
            traceback.print_exc()
            _json_response(self, {"ok": False, "error": f"Board-room PDF export failed: {exc}"}, 500)
            return
        self.send_response(200)
        self.send_header("Content-Type", "application/pdf")
        self.send_header("Content-Disposition", f"attachment; filename={filename}")
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)


def _process_is_running(process_id: int) -> bool:
    if process_id <= 0:
        return False
    if os.name == "nt":
        try:
            import ctypes

            process_query_limited_information = 0x1000
            handle = ctypes.windll.kernel32.OpenProcess(process_query_limited_information, False, process_id)
            if handle:
                ctypes.windll.kernel32.CloseHandle(handle)
                return True
            return ctypes.windll.kernel32.GetLastError() == 5
        except Exception:
            return True
    try:
        os.kill(process_id, 0)
        return True
    except OSError:
        return False


def _terminate_popen(process: subprocess.Popen | None, timeout: float = 2.0) -> None:
    if not process:
        return
    try:
        if process.poll() is not None:
            return
        process.terminate()
        try:
            process.wait(timeout=timeout)
        except subprocess.TimeoutExpired:
            process.kill()
    except Exception:
        pass


def _terminate_process_tree(process_id: int | None) -> None:
    if not process_id or process_id <= 0 or process_id == os.getpid():
        return
    try:
        if os.name == "nt":
            subprocess.Popen(
                ["taskkill", "/PID", str(process_id), "/T", "/F"],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
            )
        else:
            os.kill(process_id, 15)
    except Exception:
        pass


def _cx_suite_server_process_ids() -> set[int]:
    current = os.getpid()
    matches: set[int] = set()
    if os.name != "nt":
        return matches
    try:
        command = [
            "powershell",
            "-NoProfile",
            "-Command",
            (
                "Get-CimInstance Win32_Process | "
                "Where-Object { $_.CommandLine -match 'backend[/\\\\]server\\.py' } | "
                "ForEach-Object { \"$($_.ProcessId),$($_.ParentProcessId)\" }"
            ),
        ]
        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=6,
            creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
        )
        for line in (result.stdout or "").splitlines():
            parts = [part.strip() for part in line.split(",") if part.strip()]
            if not parts:
                continue
            try:
                process_id = int(parts[0])
            except ValueError:
                continue
            if process_id and process_id != current:
                matches.add(process_id)
    except Exception:
        pass
    return matches


def _terminate_other_cx_suite_servers() -> list[int]:
    process_ids = sorted(_cx_suite_server_process_ids())
    for process_id in process_ids:
        _terminate_process_tree(process_id)
    return process_ids


def _terminate_cx_suite_browser_processes(port: int | None) -> list[int]:
    if os.name != "nt" or not port:
        return []
    current = os.getpid()
    matches: set[int] = set()
    try:
        command = [
            "powershell",
            "-NoProfile",
            "-Command",
            (
                f"$patterns=@('127\\.0\\.0\\.1:{int(port)}','localhost:{int(port)}'); "
                "Get-CimInstance Win32_Process | "
                "Where-Object { $_.Name -match '^(msedge|chrome|firefox|brave|opera)\\.exe$' -and $_.CommandLine } | "
                "Where-Object { $cmd=$_.CommandLine; $patterns | Where-Object { $cmd -match $_ } } | "
                "ForEach-Object { $_.ProcessId }"
            ),
        ]
        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=6,
            creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
        )
        for line in (result.stdout or "").splitlines():
            try:
                process_id = int(str(line).strip())
            except ValueError:
                continue
            if process_id and process_id != current:
                matches.add(process_id)
    except Exception:
        pass
    for process_id in sorted(matches):
        _terminate_process_tree(process_id)
    return sorted(matches)


def _watch_launcher_process(
    server: ThreadingHTTPServer,
    launcher_process_id: int,
    shutdown_file: Path | None,
) -> None:
    if os.environ.get("CX_TOOLKIT_KEEP_SERVER") == "1":
        return
    time.sleep(5)
    while True:
        if shutdown_file and shutdown_file.exists():
            _write_audit_record(
                "SYSTEM",
                "SERVER_STOP_REQUESTED",
                "The launcher requested a verified clean shutdown.",
                {
                    "launcherProcessId": launcher_process_id,
                    "shutdownFile": str(shutdown_file),
                    "analysisRunning": STATE.analysis_running,
                },
            )
            server.shutdown()
            return
        if not _process_is_running(launcher_process_id):
            _write_audit_record(
                "SYSTEM",
                "LAUNCHER_CLOSED",
                "The launcher process closed; stopping its local server.",
                {"launcherProcessId": launcher_process_id, "analysisRunning": STATE.analysis_running},
            )
            server.shutdown()
            return
        time.sleep(2)


def main() -> None:
    global SERVER_LAUNCHER_PROCESS_ID
    OUTPUTS.mkdir(exist_ok=True)
    session_log = _start_server_session_log()
    _ensure_security_files()
    launcher_process_id = os.getppid()
    SERVER_LAUNCHER_PROCESS_ID = launcher_process_id
    _write_audit_record("SYSTEM", "SERVER_STARTED", "Feedback Intelligence Suite server started.", {"root": str(ROOT), "sessionId": SERVER_SESSION_ID, "sessionLog": str(session_log), "launcherProcessId": launcher_process_id})
    render_port = os.environ.get("PORT", "").strip()
    port = int(render_port or os.environ.get("CX_TOOLKIT_PORT", "8765"))
    host = "0.0.0.0" if render_port else os.environ.get("CX_TOOLKIT_HOST", "127.0.0.1")
    server = ThreadingHTTPServer((host, port), NPSHandler)
    shutdown_file_value = os.environ.get("CX_TOOLKIT_SHUTDOWN_FILE", "").strip()
    shutdown_file = Path(shutdown_file_value) if shutdown_file_value else None
    launcher_watchdog = threading.Thread(
        target=_watch_launcher_process,
        args=(server, launcher_process_id, shutdown_file),
        daemon=True,
        name="launcher-watchdog",
    )
    launcher_watchdog.start()
    display_host = "127.0.0.1" if host == "0.0.0.0" else host
    url = f"http://{display_host}:{port}/index.html"
    print(f"Feedback Intelligence Suite running at {url}")
    if not render_port and os.environ.get("CX_TOOLKIT_NO_BROWSER") != "1":
        try:
            webbrowser.open(url)
        except Exception:
            pass
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        _write_audit_record("SYSTEM", "SERVER_STOP_REQUESTED", "Server stop requested from the console.")
    except Exception as exc:
        _write_audit_record("SYSTEM", "SERVER_CRASHED", "Server stopped because of an unhandled error.", {"error": f"{type(exc).__name__}: {exc}", "traceback": traceback.format_exc()}, "ERROR")
        raise
    finally:
        elapsed = max(0.0, time.time() - SERVER_SESSION_STARTED)
        _write_audit_record("SYSTEM", "SERVER_STOPPED", "Feedback Intelligence Suite server stopped.", {"sessionId": SERVER_SESSION_ID, "durationSeconds": round(elapsed, 2), "activeUsers": len(SESSIONS), "analysisRunning": STATE.analysis_running})
        server.server_close()


if __name__ == "__main__":
    main()
















